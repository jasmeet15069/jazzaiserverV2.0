"""
╔══════════════════════════════════════════════════════════════════════════════╗
║                        JAZZ AI  —  v14.0  Production                        ║
║                                                                              ║
║  Architecture : FastAPI · SQLite/WAL · ChromaDB · APScheduler               ║
║  Auth         : JWT + Refresh · API-keys (scoped) · Rate-limiting            ║
║  Chat         : Session · streaming SSE · branching · edit · regenerate      ║
║  Context      : Sliding-window + rolling summaries + memory + RAG            ║
║  RAG          : Chunk → embed → hybrid retrieval → re-rank                  ║
║  Memory       : Auto-extract · tiered priority · confidence scoring          ║
║  Agents       : Think → Act → Observe · tool registry · cron jobs           ║
║  OAuth        : Full PKCE OAuth2 for 25+ connectors (FIXED)                 ║
║  Smart Conn.  : Gmail · Calendar · Slack · Notion · GitHub · Zapier · more  ║
║  MCP          : Full MCP server CRUD · tool discovery · JSON-RPC 2.0        ║
║  LiveKit      : Voice rooms · STT (Whisper) · TTS (PlayAI) · multilingual   ║
║  Computer     : Shell · browser · file ops · code execution                  ║
║  Admin        : Multi-model compare · Impersonation · GDPR · audit log      ║
╚══════════════════════════════════════════════════════════════════════════════╝
"""

from __future__ import annotations

# ─── stdlib ───────────────────────────────────────────────────────────────────
import ast, asyncio, base64, csv, hashlib, hmac, html as html_lib, io, json, logging, math, os, platform, socket
import smtplib, ssl
import re, secrets, shlex, shutil, subprocess, sys, tempfile, time, traceback, uuid
import urllib.parse, urllib.request, urllib.error
import zipfile
from collections import defaultdict
from concurrent.futures import ThreadPoolExecutor
from contextlib import asynccontextmanager
from datetime import datetime, timedelta, timezone
from email.message import EmailMessage
from pathlib import Path
from typing import Any, AsyncGenerator, Dict, List, Optional, Set, Tuple

# ─── third-party ──────────────────────────────────────────────────────────────
import aiosqlite
from apscheduler.schedulers.background import BackgroundScheduler
from apscheduler.triggers.cron import CronTrigger
from apscheduler.triggers.interval import IntervalTrigger
from cryptography.fernet import Fernet
from dotenv import dotenv_values, find_dotenv, load_dotenv
from fastapi import (BackgroundTasks, Depends, FastAPI, File, HTTPException,
                     Request, UploadFile, status, WebSocket, WebSocketDisconnect)
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import (FileResponse, HTMLResponse, JSONResponse,
                                RedirectResponse, StreamingResponse)
from fastapi.security import HTTPAuthorizationCredentials, HTTPBearer
from fastapi.staticfiles import StaticFiles
from jose import JWTError, jwt
from openai import OpenAI
from passlib.context import CryptContext
from pydantic import BaseModel, Field, field_validator

try:
    from huggingface_hub import InferenceClient
    HAS_HF_INFERENCE_CLIENT = True
except ImportError:
    InferenceClient = None
    HAS_HF_INFERENCE_CLIENT = False

try:
    import fitz; HAS_FITZ = True
except ImportError:
    HAS_FITZ = False

try:
    import openpyxl; HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from docx import Document as DocxDocument; HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    Presentation = None
    HAS_PPTX = False

try:
    import chromadb
    from chromadb.utils import embedding_functions
    HAS_CHROMA = True
except ImportError:
    HAS_CHROMA = False

try:
    import aiohttp; HAS_AIOHTTP = True
except ImportError:
    HAS_AIOHTTP = False

try:
    from livekit import api as livekit_api; HAS_LIVEKIT = True
except ImportError:
    HAS_LIVEKIT = False

_DOTENV_PATH = find_dotenv(usecwd=True)
load_dotenv(_DOTENV_PATH or None)

# ══════════════════════════════════════════════════════════════════════════════
# §1  LOGGING
# ══════════════════════════════════════════════════════════════════════════════

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s — %(message)s",
    handlers=[
        logging.FileHandler("jazz.log", encoding="utf-8"),
        logging.StreamHandler(),
    ],
)
logger = logging.getLogger("jazz")

# ══════════════════════════════════════════════════════════════════════════════
# §2  CONFIGURATION
# ══════════════════════════════════════════════════════════════════════════════

JWT_SECRET      = os.getenv("JWT_SECRET_KEY") or secrets.token_urlsafe(64)
JWT_ALGORITHM   = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES  = 60 * 24 * 7
REFRESH_TOKEN_EXPIRE_MINUTES = 60 * 24 * 30
try:
    BCRYPT_ROUNDS = max(8, min(14, int(os.getenv("BCRYPT_ROUNDS", "10") or "10")))
except Exception:
    BCRYPT_ROUNDS = 10

_FERNET_KEY_RAW = os.getenv("FERNET_KEY", "") or Fernet.generate_key().decode()
_fernet = Fernet(
    _FERNET_KEY_RAW.encode() if isinstance(_FERNET_KEY_RAW, str) else _FERNET_KEY_RAW
)

def _encrypt(data: dict) -> str:
    return _fernet.encrypt(json.dumps(data).encode()).decode()

def _decrypt(token: str) -> dict:
    return json.loads(_fernet.decrypt(token.encode()).decode())

ADMIN_EMAIL    = os.getenv("ADMIN_EMAIL", "jasmeet.15069@gmail.com")
ADMIN_PASSWORD = os.getenv("ADMIN_PASSWORD", "Acx@POWER@12345jassi789")
GROQ_API_KEY   = os.getenv("GROQ_API_KEY", "")
HF_TOKEN       = os.getenv("HF_TOKEN", "")
HF_TOKEN_BACKUP = os.getenv("HF_TOKEN_BACKUP", "")
HF_TOKEN_BACKUPS = os.getenv("HF_TOKEN_BACKUPS", "")
HF_TOKENS      = os.getenv("HF_TOKENS", "")
HUGGINGFACE_API_KEY = os.getenv("HUGGINGFACE_API_KEY", "")
HUGGINGFACEHUB_API_TOKEN = os.getenv("HUGGINGFACEHUB_API_TOKEN", "")
NVIDIA_API_KEY = os.getenv("NVIDIA_API_KEY", "")

# ── OAuth Provider Credentials ────────────────────────────────────────────────
GOOGLE_CLIENT_ID        = os.getenv("GOOGLE_CLIENT_ID", "")
GOOGLE_CLIENT_SECRET    = os.getenv("GOOGLE_CLIENT_SECRET", "")
GITHUB_CLIENT_ID        = os.getenv("GITHUB_CLIENT_ID", "")
GITHUB_CLIENT_SECRET    = os.getenv("GITHUB_CLIENT_SECRET", "")
GITLAB_CLIENT_ID        = os.getenv("GITLAB_CLIENT_ID", "")
GITLAB_CLIENT_SECRET    = os.getenv("GITLAB_CLIENT_SECRET", "")
NOTION_CLIENT_ID        = os.getenv("NOTION_CLIENT_ID", "")
NOTION_CLIENT_SECRET    = os.getenv("NOTION_CLIENT_SECRET", "")
SLACK_CLIENT_ID         = os.getenv("SLACK_CLIENT_ID", "")
SLACK_CLIENT_SECRET     = os.getenv("SLACK_CLIENT_SECRET", "")
SLACK_BOT_TOKEN         = os.getenv("SLACK_BOT_TOKEN", "")
ATLASSIAN_CLIENT_ID     = os.getenv("ATLASSIAN_CLIENT_ID", "")
ATLASSIAN_CLIENT_SECRET = os.getenv("ATLASSIAN_CLIENT_SECRET", "")
MICROSOFT_CLIENT_ID     = os.getenv("MICROSOFT_CLIENT_ID", "")
MICROSOFT_CLIENT_SECRET = os.getenv("MICROSOFT_CLIENT_SECRET", "")
AIRTABLE_CLIENT_ID      = os.getenv("AIRTABLE_CLIENT_ID", "")
AIRTABLE_CLIENT_SECRET  = os.getenv("AIRTABLE_CLIENT_SECRET", "")
DROPBOX_CLIENT_ID       = os.getenv("DROPBOX_CLIENT_ID", "")
DROPBOX_CLIENT_SECRET   = os.getenv("DROPBOX_CLIENT_SECRET", "")
DISCORD_CLIENT_ID       = os.getenv("DISCORD_CLIENT_ID", "")
DISCORD_CLIENT_SECRET   = os.getenv("DISCORD_CLIENT_SECRET", "")
ZOOM_CLIENT_ID          = os.getenv("ZOOM_CLIENT_ID", "")
ZOOM_CLIENT_SECRET      = os.getenv("ZOOM_CLIENT_SECRET", "")
ASANA_CLIENT_ID         = os.getenv("ASANA_CLIENT_ID", "")
ASANA_CLIENT_SECRET     = os.getenv("ASANA_CLIENT_SECRET", "")
HUBSPOT_CLIENT_ID       = os.getenv("HUBSPOT_CLIENT_ID", "")
HUBSPOT_CLIENT_SECRET   = os.getenv("HUBSPOT_CLIENT_SECRET", "")
LINEAR_CLIENT_ID        = os.getenv("LINEAR_CLIENT_ID", "")
LINEAR_CLIENT_SECRET    = os.getenv("LINEAR_CLIENT_SECRET", "")
SALESFORCE_CLIENT_ID    = os.getenv("SALESFORCE_CLIENT_ID", "")
SALESFORCE_CLIENT_SECRET= os.getenv("SALESFORCE_CLIENT_SECRET", "")
LINKEDIN_CLIENT_ID      = os.getenv("LINKEDIN_CLIENT_ID", "")
LINKEDIN_CLIENT_SECRET  = os.getenv("LINKEDIN_CLIENT_SECRET", "")
BOX_CLIENT_ID           = os.getenv("BOX_CLIENT_ID", "")
BOX_CLIENT_SECRET       = os.getenv("BOX_CLIENT_SECRET", "")
CLICKUP_CLIENT_ID       = os.getenv("CLICKUP_CLIENT_ID", "")
CLICKUP_CLIENT_SECRET   = os.getenv("CLICKUP_CLIENT_SECRET", "")
STRIPE_SECRET_KEY       = os.getenv("STRIPE_SECRET_KEY", "")
RAZORPAY_KEY_ID         = os.getenv("RAZORPAY_KEY_ID", "")
RAZORPAY_KEY_SECRET     = os.getenv("RAZORPAY_KEY_SECRET", "")
RAZORPAY_CURRENCY       = os.getenv("RAZORPAY_CURRENCY", "INR").upper()
RAZORPAY_PRO_MONTHLY    = int(os.getenv("RAZORPAY_PRO_MONTHLY_AMOUNT", "20000") or "20000")
RAZORPAY_PREMIUM_MONTHLY= int(os.getenv("RAZORPAY_PREMIUM_MONTHLY_AMOUNT", "99900") or "99900")
RAZORPAY_ENT_MONTHLY    = int(os.getenv("RAZORPAY_ENTERPRISE_MONTHLY_AMOUNT", "499900") or "499900")

APP_BASE_URL       = os.getenv("APP_BASE_URL", "https://imperceptibly-hymnlike-leesa.ngrok-free.dev")
AUTH_LINK_BASE_URL = os.getenv("AUTH_LINK_BASE_URL") or os.getenv("PUBLIC_APP_URL", "https://www.jazzai.online")
LIVEKIT_URL        = os.getenv("LIVEKIT_URL", "")
LIVEKIT_API_KEY    = os.getenv("LIVEKIT_API_KEY", "")
LIVEKIT_API_SECRET = os.getenv("LIVEKIT_API_SECRET", "")
TTS_VOICE          = os.getenv("TTS_VOICE", "Fritz-PlayAI")
LIVY_URL           = os.getenv("LIVY_URL", "http://localhost:8998")
LIVY_USER          = os.getenv("LIVY_USER", "")
LIVY_PASSWORD      = os.getenv("LIVY_PASSWORD", "")

_RUNTIME_ENV_KEYS = [
    "ADMIN_EMAIL", "ADMIN_PASSWORD", "GROQ_API_KEY", "HF_TOKEN", "HF_TOKEN_BACKUP",
    "HF_TOKEN_BACKUPS", "HF_TOKENS", "HUGGINGFACE_API_KEY", "HUGGINGFACEHUB_API_TOKEN",
    "NVIDIA_API_KEY",
    "GOOGLE_CLIENT_ID", "GOOGLE_CLIENT_SECRET",
    "GITHUB_CLIENT_ID", "GITHUB_CLIENT_SECRET",
    "GITLAB_CLIENT_ID", "GITLAB_CLIENT_SECRET",
    "NOTION_CLIENT_ID", "NOTION_CLIENT_SECRET",
    "SLACK_CLIENT_ID", "SLACK_CLIENT_SECRET", "SLACK_BOT_TOKEN",
    "ATLASSIAN_CLIENT_ID", "ATLASSIAN_CLIENT_SECRET",
    "MICROSOFT_CLIENT_ID", "MICROSOFT_CLIENT_SECRET",
    "AIRTABLE_CLIENT_ID", "AIRTABLE_CLIENT_SECRET",
    "DISCORD_CLIENT_ID", "DISCORD_CLIENT_SECRET",
    "ZOOM_CLIENT_ID", "ZOOM_CLIENT_SECRET",
    "ASANA_CLIENT_ID", "ASANA_CLIENT_SECRET",
    "LINEAR_CLIENT_ID", "LINEAR_CLIENT_SECRET",
    "LINKEDIN_CLIENT_ID", "LINKEDIN_CLIENT_SECRET",
    "HUBSPOT_CLIENT_ID", "HUBSPOT_CLIENT_SECRET",
    "SALESFORCE_CLIENT_ID", "SALESFORCE_CLIENT_SECRET",
    "DROPBOX_CLIENT_ID", "DROPBOX_CLIENT_SECRET",
    "BOX_CLIENT_ID", "BOX_CLIENT_SECRET",
    "CLICKUP_CLIENT_ID", "CLICKUP_CLIENT_SECRET",
    "STRIPE_SECRET_KEY", "RAZORPAY_KEY_ID", "RAZORPAY_KEY_SECRET", "RAZORPAY_CURRENCY",
    "RAZORPAY_PRO_MONTHLY_AMOUNT", "RAZORPAY_PREMIUM_MONTHLY_AMOUNT", "RAZORPAY_ENTERPRISE_MONTHLY_AMOUNT",
    "APP_BASE_URL", "AUTH_LINK_BASE_URL", "PUBLIC_APP_URL",
    "LIVEKIT_URL", "LIVEKIT_API_KEY", "LIVEKIT_API_SECRET",
    "TTS_VOICE", "LIVY_URL", "LIVY_USER", "LIVY_PASSWORD",
    "NEXT_PUBLIC_SUPABASE_URL", "NEXT_PUBLIC_SUPABASE_PUBLISHABLE_KEY",
    "SUPABASE_URL", "SUPABASE_PUBLISHABLE_KEY",
    "SUPABASE_OAUTH_CLIENT_ID", "SUPABASE_OAUTH_CLIENT_SECRET",
    "SUPABASE_OAUTH_SCOPES", "SUPABASE_OAUTH_AUTH_METHOD",
    "REQUIRE_EMAIL_VERIFICATION", "EMAIL_VERIFICATION_EXPIRE_HOURS",
    "EMAIL_VERIFICATION_RETURN_LINK",
    "PASSWORD_RESET_EXPIRE_MINUTES", "PASSWORD_RESET_RETURN_LINK",
    "BCRYPT_ROUNDS",
    "EMAIL_FROM", "EMAIL_FROM_NAME", "RESEND_API_KEY",
    "SMTP_HOST", "SMTP_PORT", "SMTP_USERNAME", "SMTP_PASSWORD",
    "SMTP_FROM", "SMTP_USE_TLS", "SMTP_USE_SSL",
]
_RUNTIME_ENV_DEFAULTS = {
    "ADMIN_EMAIL": "jasmeet.15069@gmail.com",
    "ADMIN_PASSWORD": "Acx@POWER@12345jassi789",
    "APP_BASE_URL": "https://imperceptibly-hymnlike-leesa.ngrok-free.dev",
    "AUTH_LINK_BASE_URL": "https://www.jazzai.online",
    "PUBLIC_APP_URL": "https://www.jazzai.online",
    "LIVY_URL": "http://localhost:8998",
    "TTS_VOICE": "Fritz-PlayAI",
    "REQUIRE_EMAIL_VERIFICATION": "1",
    "EMAIL_VERIFICATION_EXPIRE_HOURS": "24",
    "EMAIL_VERIFICATION_RETURN_LINK": "0",
    "PASSWORD_RESET_EXPIRE_MINUTES": "30",
    "PASSWORD_RESET_RETURN_LINK": "0",
    "BCRYPT_ROUNDS": "10",
    "EMAIL_FROM": "",
    "EMAIL_FROM_NAME": "JAZZ AI",
    "RESEND_API_KEY": "",
    "SUPABASE_OAUTH_CLIENT_ID": "",
    "SUPABASE_OAUTH_CLIENT_SECRET": "",
    "SUPABASE_OAUTH_SCOPES": "openid email profile",
    "SUPABASE_OAUTH_AUTH_METHOD": "auto",
    "SMTP_HOST": "",
    "SMTP_PORT": "587",
    "SMTP_USERNAME": "",
    "SMTP_PASSWORD": "",
    "SMTP_FROM": "",
    "SMTP_USE_TLS": "1",
    "SMTP_USE_SSL": "0",
}

def _env_file_path() -> Path:
    configured = os.getenv("ENV_FILE_PATH") or os.getenv("ENV_PATH")
    if configured:
        return Path(configured).expanduser().resolve()
    discovered = find_dotenv(usecwd=True)
    if discovered:
        return Path(discovered).resolve()
    return (Path(__file__).resolve().parent / ".env").resolve()

def _reload_runtime_env() -> None:
    path = _env_file_path()
    if path.exists():
        load_dotenv(path, override=True)
    for key in _RUNTIME_ENV_KEYS:
        if key in globals():
            globals()[key] = os.getenv(key, _RUNTIME_ENV_DEFAULTS.get(key, ""))

DB_PATH       = str(Path(os.getenv("DB_PATH", "./jazz_v14.db")).resolve())
WORKSPACE_DIR = Path("workspace"); WORKSPACE_DIR.mkdir(exist_ok=True)
SITES_DIR     = WORKSPACE_DIR / "sites";   SITES_DIR.mkdir(exist_ok=True)
UPLOADS_DIR   = WORKSPACE_DIR / "uploads"; UPLOADS_DIR.mkdir(exist_ok=True)
BACKUPS_DIR   = WORKSPACE_DIR / "backups"; BACKUPS_DIR.mkdir(exist_ok=True)

CHUNK_SIZE         = 500
CHUNK_OVERLAP      = 50
TOP_K_RETRIEVAL    = 6
RERANK_TOP_K       = 3
MAX_CONTEXT_TOKENS = int(os.getenv("MAX_CONTEXT_TOKENS") or "258000")
CONTEXT_OUTPUT_RESERVE_TOKENS = int(os.getenv("CONTEXT_OUTPUT_RESERVE_TOKENS") or "4096")
CONTEXT_AUTO_COMPACT_PCT = float(os.getenv("CONTEXT_AUTO_COMPACT_PCT") or "0.90")
CONTEXT_COMPACT_KEEP_RECENT = int(os.getenv("CONTEXT_COMPACT_KEEP_RECENT") or "10")
SUMMARY_TRIGGER    = 20
MAX_FILE_SIZE      = 50 * 1024 * 1024
MAX_IMAGE_SIZE     = 20 * 1024 * 1024
CODE_EXEC_TIMEOUT  = 30
CODE_EXEC_MAX_OUT  = 100_000
MAX_MEMORIES       = 200

ALLOWED_ORIGINS: List[str] = [
    o.strip() for o in os.getenv("ALLOWED_ORIGINS", "*").split(",") if o.strip()
]

THINKING_KEYWORDS = [
    "why","how does","explain","analyze","reason","think","understand",
    "compare","difference","pros and cons","debate","argue","philosophy",
    "ethics","theory","hypothesis","research","complex","deep dive",
    "step by step","break down","implications","consequences","tradeoffs",
    "solve","algorithm","proof","derive","calculate","optimize",
    "should i","what if","is it true","help me understand","critically",
]
THINKING_MIN_LENGTH = 50

SUBSCRIPTION_LIMITS: Dict[str, Dict[str, Any]] = {
    "free":       {"messages_per_day":100,"documents":10,"max_file_mb":10,"agent_jobs":3,"websites":5,"code_runs_per_day":50,"api_keys":3,"memories":100,"connectors":3},
    "pro":        {"messages_per_day":500,"documents":50,"max_file_mb":25,"agent_jobs":10,"websites":25,"code_runs_per_day":100,"api_keys":10,"memories":500,"connectors":8},
    "premium":    {"messages_per_day":-1,"documents":100,"max_file_mb":50,"agent_jobs":20,"websites":100,"code_runs_per_day":-1,"api_keys":20,"memories":-1,"connectors":20},
    "enterprise": {"messages_per_day":-1,"documents":-1,"max_file_mb":200,"agent_jobs":-1,"websites":-1,"code_runs_per_day":-1,"api_keys":-1,"memories":-1,"connectors":-1},
}
PLAN_TIERS = ("free", "pro", "premium", "enterprise")
RATE_LIMIT_RESOURCES = [
    ("messages_per_day", "Messages per day", "Chat, multi-model, voice, and agent message budget."),
    ("documents", "Documents", "Processed documents a user can keep."),
    ("max_file_mb", "Max file size MB", "Largest upload size allowed."),
    ("agent_jobs", "Agent jobs", "Scheduled or saved agent jobs."),
    ("websites", "Websites", "Generated websites saved in the portal."),
    ("code_runs_per_day", "Code runs per day", "Code Runner and system code execution budget."),
    ("api_keys", "API keys", "Personal API keys a user can create."),
    ("memories", "Memories", "Saved long-term memories."),
    ("connectors", "Connectors", "Connected apps per user."),
]

ALLOWED_EXTENSIONS = {
    ".html",".css",".js",".py",".txt",".md",".json",".xml",".yaml",".yml",
    ".csv",".sql",".sh",".ts",".tsx",".jsx",".conf",".ini",".log",".svg",
    ".bat",".ps1",".rb",".go",".rs",".cpp",".c",".h",".java",".kt",".swift",".php",
}

_SERVER_START = datetime.now(timezone.utc)
_executor = ThreadPoolExecutor(max_workers=16)

# ── Model Registry ────────────────────────────────────────────────────────────
GROQ_MODELS = {
    "llama-3.3-70b-versatile": {"label":"Llama 3.3 70B","ctx":32768,"fast":False,"provider":"groq","censored":True},
    "llama-3.1-8b-instant":    {"label":"Llama 3.1 8B (Fast)","ctx":131072,"fast":True,"provider":"groq","censored":True},
}

HUGGINGFACE_MODELS = {
    # Protected per user request: do not remove without explicit permission.
    "dolphin-3.0-llama-3.1-8b": {
        "label":"Dolphin 3.0 Llama 3.1 8B (HF)",
        "model_name":"cognitivecomputations/Dolphin3.0-Llama3.1-8B",
        "ctx":131072,
        "fast":False,
        "provider":"huggingface",
        "censored":False,
        "protected":True,
    },
}

BUILTIN_MODELS = {**GROQ_MODELS, **HUGGINGFACE_MODELS}

DEFAULT_DB_MODELS = {
    "dolphin-mistral-24b-venice-hf": {
        "name":"Dolphin Mistral 24B Venice (HF Featherless)",
        "provider":"huggingface",
        "base_url":"https://router.huggingface.co/v1",
        "model_name":"dphn/Dolphin-Mistral-24B-Venice-Edition:featherless-ai",
        "context_length":32768,
        "max_output_tokens":4096,
        "temperature_default":0.15,
        "is_active":True,
        "is_default":False,
        "is_fast":False,
        "is_vision":False,
        "is_code":False,
        "description":"Uncensored, steerable Dolphin Mistral 24B Venice Edition via Hugging Face Inference Providers using Featherless AI.",
        "tags":["protected","uncensored","dolphin","huggingface","featherless-ai","venice","24b"],
    },
    "nvidia-moonshotai-kimi-k2-6": {
        "name":"Kimi K2.6 (NVIDIA)",
        "provider":"nvidia",
        "base_url":"https://integrate.api.nvidia.com/v1",
        "model_name":"moonshotai/kimi-k2.6",
        "context_length":262144,
        "max_output_tokens":16384,
        "temperature_default":1.0,
        "is_active":True,
        "is_default":False,
        "is_fast":False,
        "is_vision":False,
        "is_code":True,
        "description":"Moonshot AI Kimi K2.6 via NVIDIA NIM/OpenAI-compatible chat completions with thinking enabled.",
        "tags":["nvidia","moonshot","kimi","thinking","code","most-intelligent"],
    },
    "hf-router-zai-org-glm-5-1-together": {
        "name":"GLM 5.1 (HF Together)",
        "provider":"huggingface",
        "base_url":"https://router.huggingface.co/v1",
        "model_name":"zai-org/GLM-5.1:together",
        "context_length":202752,
        "max_output_tokens":25344,
        "temperature_default":0.7,
        "is_active":True,
        "is_default":False,
        "is_fast":False,
        "is_vision":False,
        "is_code":True,
        "description":"Z.ai GLM 5.1 through Hugging Face Router using Together. Thinking is disabled for normal chat content.",
        "tags":["huggingface","router","together","glm","zai-org","tools","structured","most-intelligent"],
    },
    "jazz-ai-testing": {
        "name":"Jazz AI V1.0",
        "provider":"local_generate",
        "base_url":"http://127.0.0.1:18080",
        "model_name":"JazzAI-V1.0-from-scratch-GPT2",
        "context_length":4096,
        "max_output_tokens":128,
        "temperature_default":0.0,
        "is_active":True,
        "is_default":False,
        "is_fast":True,
        "is_vision":False,
        "is_code":True,
        "description":"Private from-scratch GPT-style Jazz AI V1.0 runtime. Internal compatibility ID remains jazz-ai-testing.",
        "tags":["local","from-scratch","jazz-ai-v1","gpt2","coding","cpu","private","coach","hindi"],
    },
}

IMAGE_TO_TEXT_MODEL_ID = "hf-router-qwen-qwen2-5-vl-72b-instruct-ovhcloud"
IMAGE_TO_TEXT_MODEL_NAME = "Qwen/Qwen2.5-VL-72B-Instruct:ovhcloud"
DEFAULT_TEXT_MODEL_ID = "local-dolphin3-qwen25-05b"

_PROVIDER_DEFAULTS: Dict[str,str] = {
    "groq":        "https://api.groq.com/openai/v1",
    "openrouter":  "https://openrouter.ai/api/v1",
    "huggingface": "https://router.huggingface.co/v1",
    "nvidia":      "https://integrate.api.nvidia.com/v1",
    "local_generate": "http://127.0.0.1:18080",
    "ollama":      "http://localhost:11434/v1",
    "mistral":     "https://api.mistral.ai/v1",
    "cohere":      "https://api.cohere.ai/compatibility/v1",
    "together":    "https://api.together.xyz/v1",
    "custom":      "",
}

NVIDIA_KIMI_PROVIDER_FALLBACKS = {
    "moonshotai/kimi-k2.6": [
        "moonshotai/kimi-k2-instruct",
        "moonshotai/kimi-k2-instruct-0905",
        "moonshotai/kimi-k2-thinking",
        "moonshotai/kimi-k2.6",
    ],
}
NVIDIA_STREAM_TIMEOUTS = {
    "moonshotai/kimi-k2-instruct": int(os.getenv("NVIDIA_KIMI_INSTRUCT_TIMEOUT", "28")),
    "moonshotai/kimi-k2-instruct-0905": int(os.getenv("NVIDIA_KIMI_INSTRUCT_0905_TIMEOUT", "28")),
    "moonshotai/kimi-k2-thinking": int(os.getenv("NVIDIA_KIMI_THINKING_TIMEOUT", "12")),
    "moonshotai/kimi-k2.6": int(os.getenv("NVIDIA_KIMI_K26_TIMEOUT", "8")),
}

# ── Style hints for website builder ──────────────────────────────────────────
_STYLE_HINTS = {
    "modern":        "Clean, bold typography, whitespace, subtle shadows, CSS Grid",
    "retro":         "80s/90s aesthetic, pixel fonts, neon on dark, CRT scanlines",
    "minimal":       "Extreme whitespace, single accent color, fine typography",
    "glassmorphism": "Frosted glass, backdrop-filter blur, translucency",
    "brutalist":     "Raw feel, bold borders, high contrast, asymmetric layout",
    "dark":          "Dark background, glowing accents, modern dark UI patterns",
    "futuristic":    "Sci-fi aesthetic, neon glows, animated particles, holographic effects",
    "saas":          "Quiet SaaS product UI, dense sections, crisp feature grids, practical dashboards",
    "editorial":     "Premium magazine feel, immersive type scale, strong content rhythm",
    "commerce":      "Conversion-focused product storytelling, trust blocks, polished checkout cues",
    "app":           "Interactive app shell with realistic panels, controls, empty/loading states",
    "agency":        "Sharp creative studio style, case studies, bold service positioning",
}

# ══════════════════════════════════════════════════════════════════════════════
# §3  DATABASE
# ══════════════════════════════════════════════════════════════════════════════

_db: Optional[aiosqlite.Connection] = None

async def _init_db_connection() -> None:
    global _db
    _db = await aiosqlite.connect(DB_PATH, check_same_thread=False)
    _db.row_factory = aiosqlite.Row
    for p in ["PRAGMA journal_mode=WAL","PRAGMA synchronous=NORMAL",
              "PRAGMA foreign_keys=ON","PRAGMA cache_size=-32000",
              "PRAGMA temp_store=MEMORY","PRAGMA mmap_size=268435456"]:
        await _db.execute(p)
    await _db.commit()

def _conn() -> aiosqlite.Connection:
    if _db is None: raise RuntimeError("DB not initialised")
    return _db

def _utcnow() -> str: return datetime.now(timezone.utc).isoformat()
def _new_id() -> str: return str(uuid.uuid4())

async def db_fetchone(sql: str, params: tuple = ()) -> Optional[Dict]:
    async with _conn().execute(sql, params) as cur:
        row = await cur.fetchone()
        return dict(row) if row else None

async def db_fetchall(sql: str, params: tuple = ()) -> List[Dict]:
    async with _conn().execute(sql, params) as cur:
        rows = await cur.fetchall()
        return [dict(r) for r in rows]

async def db_count(sql: str, params: tuple = ()) -> int:
    row = await db_fetchone(sql, params)
    return list(row.values())[0] if row else 0

async def db_execute(sql: str, params: tuple = ()) -> int:
    async with _conn().execute(sql, params) as cur:
        await _conn().commit()
        return cur.lastrowid or 0

async def db_executemany(sql: str, param_list: List[tuple]) -> None:
    await _conn().executemany(sql, param_list)
    await _conn().commit()

async def _ensure_premium_subscription_tier() -> None:
    row = await db_fetchone("SELECT sql FROM sqlite_master WHERE type='table' AND name='users'")
    ddl = (row or {}).get("sql") or ""
    if "premium" in ddl:
        return
    logger.info("[MIGRATE] Rebuilding users table to allow premium subscriptions")
    conn = _conn()
    await conn.execute("PRAGMA foreign_keys=OFF")
    await conn.executescript("""
BEGIN;
DROP TABLE IF EXISTS users_new;
CREATE TABLE users_new (
    id TEXT PRIMARY KEY, email TEXT NOT NULL UNIQUE,
    password_hash TEXT NOT NULL, full_name TEXT DEFAULT '',
    role TEXT DEFAULT 'client' CHECK(role IN('admin','client')),
    subscription TEXT DEFAULT 'free' CHECK(subscription IN('free','pro','premium','enterprise')),
    subscription_expires_at TEXT, memory_enabled INTEGER DEFAULT 1,
    timezone TEXT DEFAULT 'UTC', is_active INTEGER DEFAULT 1, is_verified INTEGER DEFAULT 0,
    last_login_at TEXT, created_at TEXT NOT NULL, updated_at TEXT NOT NULL
);
INSERT INTO users_new(id,email,password_hash,full_name,role,subscription,subscription_expires_at,
    memory_enabled,timezone,is_active,is_verified,last_login_at,created_at,updated_at)
SELECT id,email,password_hash,full_name,role,subscription,subscription_expires_at,
    memory_enabled,timezone,is_active,is_verified,last_login_at,created_at,updated_at
FROM users;
DROP TABLE users;
ALTER TABLE users_new RENAME TO users;
CREATE INDEX IF NOT EXISTS idx_users_email ON users(email);
COMMIT;
""")
    await conn.commit()
    await conn.execute("PRAGMA foreign_keys=ON")
    await conn.commit()

async def _seed_default_ai_models(created_by: str) -> None:
    now = _utcnow()
    for mid, cfg in DEFAULT_DB_MODELS.items():
        exists = await db_fetchone(
            "SELECT id FROM ai_models WHERE id=? OR model_name=? OR name=? LIMIT 1",
            (mid, cfg["model_name"], cfg["name"]))
        if exists:
            continue
        await db_execute(
            "INSERT INTO ai_models(id,name,provider,base_url,model_name,encrypted_api_key,"
            "is_active,is_default,is_fast,is_vision,is_code,context_length,max_output_tokens,"
            "temperature_default,description,tags_json,created_by,created_at,updated_at)"
            " VALUES(?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)",
            (mid, cfg["name"], cfg["provider"], cfg["base_url"], cfg["model_name"],
             _encrypt({"key": ""}), int(cfg["is_active"]), int(cfg["is_default"]),
             int(cfg["is_fast"]), int(cfg["is_vision"]), int(cfg["is_code"]),
             int(cfg["context_length"]), int(cfg["max_output_tokens"]),
             float(cfg["temperature_default"]), cfg["description"],
             json.dumps(cfg["tags"]), created_by, now, now))
        logger.info("[INIT] Seeded default model: %s", cfg["name"])

async def _ensure_image_to_text_model(created_by: str) -> None:
    now = _utcnow()
    tags = ["internal", "image-to-text", "vision", "huggingface", "router", "protected", "qwen"]
    existing = await db_fetchone("SELECT is_active FROM ai_models WHERE id=? LIMIT 1", (IMAGE_TO_TEXT_MODEL_ID,))
    active = int(existing["is_active"]) if existing else 1
    await db_execute(
        "INSERT INTO ai_models(id,name,provider,base_url,model_name,encrypted_api_key,"
        "is_active,is_default,is_fast,is_vision,is_code,context_length,max_output_tokens,"
        "temperature_default,description,tags_json,created_by,created_at,updated_at)"
        " VALUES(?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)"
        " ON CONFLICT(id) DO UPDATE SET name=excluded.name,provider=excluded.provider,"
        "base_url=excluded.base_url,model_name=excluded.model_name,is_active=excluded.is_active,is_default=0,"
        "is_fast=0,is_vision=1,is_code=0,context_length=excluded.context_length,"
        "max_output_tokens=excluded.max_output_tokens,temperature_default=excluded.temperature_default,"
        "description=excluded.description,tags_json=excluded.tags_json,updated_at=excluded.updated_at",
        (
            IMAGE_TO_TEXT_MODEL_ID,
            "Qwen2.5 VL 72B Image-to-Text (HF internal)",
            "huggingface",
            _PROVIDER_DEFAULTS["huggingface"],
            IMAGE_TO_TEXT_MODEL_NAME,
            _encrypt({"key": ""}),
            active, 0, 0, 1, 0,
            32768,
            4096,
            0.15,
            "Internal image-to-text/OCR model. Hidden from user model picker; visible in Admin Models.",
            json.dumps(tags),
            created_by,
            now,
            now,
        ),
    )
    await db_execute(
        "INSERT INTO model_access(id,model_id,display_name,source,is_enabled,created_at,updated_at)"
        " VALUES(?,?,?,?,?,?,?)"
        " ON CONFLICT(model_id) DO UPDATE SET display_name=excluded.display_name,"
        "source='db',is_enabled=excluded.is_enabled,updated_at=excluded.updated_at",
        (_new_id(), IMAGE_TO_TEXT_MODEL_ID, "Qwen2.5 VL 72B Image-to-Text", "db", active, now, now),
    )

def _hf_router_model_id(model_id: str, provider: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", f"{model_id}-{provider}".lower()).strip("-")
    return ("hf-router-" + slug)[:120]

def _hf_router_label(model_id: str, provider: str) -> str:
    short = model_id.split("/")[-1].replace("-", " ")
    return f"{short} (HF {provider})"[:100]

async def _seed_hf_router_models(created_by: str, limit: int = None) -> int:
    """Seed current Hugging Face Router chat models that have live inference providers."""
    if not HF_TOKEN.strip():
        return 0
    limit = max(1, min(50, int(limit or os.getenv("HF_ROUTER_MODEL_LIMIT") or 30)))

    def _fetch_models() -> List[Dict[str, Any]]:
        req = urllib.request.Request(
            "https://router.huggingface.co/v1/models",
            headers={"Authorization": "Bearer " + HF_TOKEN.strip()})
        with urllib.request.urlopen(req, timeout=25) as r:
            data = json.loads(r.read())
        return data if isinstance(data, list) else data.get("data", [])

    try:
        rows = await asyncio.get_running_loop().run_in_executor(_executor, _fetch_models)
    except Exception as exc:
        logger.warning("[INIT] HF router model seed skipped: %s", exc)
        return 0

    created = 0
    now = _utcnow()
    for row in rows:
        if created >= limit:
            break
        model_id = str(row.get("id") or "").strip()
        arch = row.get("architecture") or {}
        if not model_id or "text" not in (arch.get("input_modalities") or []) or "text" not in (arch.get("output_modalities") or []):
            continue
        providers = [p for p in (row.get("providers") or []) if p.get("status") == "live" and p.get("provider")]
        if not providers:
            continue
        providers.sort(key=lambda p: (
            0 if p.get("supports_tools") else 1,
            float((p.get("pricing") or {}).get("output") or 999999),
            float((p.get("pricing") or {}).get("input") or 999999),
        ))
        provider = str(providers[0]["provider"])
        ctx = int(providers[0].get("context_length") or 32768)
        provider_model = f"{model_id}:{provider}"
        mid = _hf_router_model_id(model_id, provider)
        name = _hf_router_label(model_id, provider)
        exists = await db_fetchone(
            "SELECT id FROM ai_models WHERE id=? OR model_name=? OR name=? LIMIT 1",
            (mid, provider_model, name))
        if exists:
            continue
        tags = ["huggingface", "router", provider]
        if providers[0].get("supports_tools"):
            tags.append("tools")
        if providers[0].get("supports_structured_output"):
            tags.append("structured")
        await db_execute(
            "INSERT INTO ai_models(id,name,provider,base_url,model_name,encrypted_api_key,"
            "is_active,is_default,is_fast,is_vision,is_code,context_length,max_output_tokens,"
            "temperature_default,description,tags_json,created_by,created_at,updated_at)"
            " VALUES(?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)",
            (mid, name, "huggingface", _PROVIDER_DEFAULTS["huggingface"], provider_model,
             _encrypt({"key": ""}), 0, 0, int(ctx <= 32768),
             int("image" in (arch.get("input_modalities") or [])),
             int(any(x in model_id.lower() for x in ("coder", "code", "qwen3-coder"))),
             min(ctx, 2_000_000), min(32768, max(4096, ctx // 8)),
             0.7, f"Seeded from Hugging Face Router live inference providers using {provider}.",
             json.dumps(tags), created_by, now, now))
        created += 1
    if created:
        logger.info("[INIT] Seeded %d Hugging Face Router models", created)
    return created

@asynccontextmanager
async def db_transaction():
    await _conn().execute("BEGIN")
    try:
        yield; await _conn().commit()
    except Exception:
        await _conn().rollback(); raise

# ── Schema ────────────────────────────────────────────────────────────────────
_SCHEMA = """
CREATE TABLE IF NOT EXISTS users (
    id TEXT PRIMARY KEY, email TEXT NOT NULL UNIQUE,
    password_hash TEXT NOT NULL, full_name TEXT DEFAULT '',
    role TEXT DEFAULT 'client' CHECK(role IN('admin','client')),
    subscription TEXT DEFAULT 'free' CHECK(subscription IN('free','pro','premium','enterprise')),
    subscription_expires_at TEXT, memory_enabled INTEGER DEFAULT 1,
    timezone TEXT DEFAULT 'UTC', is_active INTEGER DEFAULT 1, is_verified INTEGER DEFAULT 0,
    last_login_at TEXT, created_at TEXT NOT NULL, updated_at TEXT NOT NULL
);
CREATE TABLE IF NOT EXISTS refresh_tokens (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL, token_hash TEXT NOT NULL UNIQUE,
    expires_at TEXT NOT NULL, revoked INTEGER DEFAULT 0, created_at TEXT NOT NULL,
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS email_verification_tokens (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL, token_hash TEXT NOT NULL UNIQUE,
    email TEXT NOT NULL, expires_at TEXT NOT NULL, used_at TEXT,
    created_at TEXT NOT NULL,
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS password_reset_tokens (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL, token_hash TEXT NOT NULL UNIQUE,
    email TEXT NOT NULL, expires_at TEXT NOT NULL, used_at TEXT,
    created_at TEXT NOT NULL,
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS chat_sessions (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL, title TEXT DEFAULT 'New Chat',
    model_id TEXT DEFAULT 'llama-3.3-70b-versatile', is_pinned INTEGER DEFAULT 0,
    is_archived INTEGER DEFAULT 0, turn_count INTEGER DEFAULT 0, last_message_at TEXT,
    parent_session_id TEXT, branch_from_msg_id TEXT,
    created_at TEXT NOT NULL, updated_at TEXT NOT NULL,
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS chat_history (
    id TEXT PRIMARY KEY, session_id TEXT NOT NULL, user_id TEXT NOT NULL,
    role TEXT NOT NULL CHECK(role IN('user','assistant','system','tool')),
    content TEXT NOT NULL, model_used TEXT, agent_used TEXT,
    thinking_content TEXT DEFAULT '',
    tool_calls_json TEXT DEFAULT '[]', rag_chunks_json TEXT DEFAULT '[]',
    tokens_input INTEGER DEFAULT 0, tokens_output INTEGER DEFAULT 0,
    latency_ms INTEGER DEFAULT 0, is_hidden INTEGER DEFAULT 0,
    parent_msg_id TEXT, edit_count INTEGER DEFAULT 0,
    mode TEXT DEFAULT 'normal', created_at TEXT NOT NULL,
    FOREIGN KEY(session_id) REFERENCES chat_sessions(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS message_edits (
    id TEXT PRIMARY KEY, message_id TEXT NOT NULL,
    old_content TEXT NOT NULL, new_content TEXT NOT NULL, edited_at TEXT NOT NULL,
    FOREIGN KEY(message_id) REFERENCES chat_history(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS summaries (
    id TEXT PRIMARY KEY, session_id TEXT NOT NULL, user_id TEXT NOT NULL,
    summary_text TEXT NOT NULL, covers_from_at TEXT NOT NULL, covers_to_at TEXT NOT NULL,
    turn_count INTEGER NOT NULL, tokens_saved INTEGER DEFAULT 0,
    model_used TEXT NOT NULL, created_at TEXT NOT NULL,
    FOREIGN KEY(session_id) REFERENCES chat_sessions(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS memories (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL,
    key TEXT NOT NULL, value TEXT NOT NULL,
    source TEXT DEFAULT 'auto' CHECK(source IN('auto','manual','agent')),
    confidence REAL DEFAULT 1.0, last_reinforced TEXT NOT NULL,
    reinforcement_count INTEGER DEFAULT 1, is_active INTEGER DEFAULT 1,
    created_at TEXT NOT NULL, updated_at TEXT NOT NULL,
    UNIQUE(user_id, key),
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS documents (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL,
    filename TEXT NOT NULL, original_name TEXT NOT NULL, file_type TEXT DEFAULT '',
    file_size_bytes INTEGER DEFAULT 0, page_count INTEGER, chunk_count INTEGER DEFAULT 0,
    char_count INTEGER DEFAULT 0, is_indexed INTEGER DEFAULT 0, index_error TEXT,
    collection_name TEXT DEFAULT 'documents', tags_json TEXT DEFAULT '[]',
    description TEXT DEFAULT '', uploaded_at TEXT NOT NULL, indexed_at TEXT,
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS embeddings_meta (
    id TEXT PRIMARY KEY, document_id TEXT NOT NULL, user_id TEXT NOT NULL,
    chunk_index INTEGER NOT NULL, chunk_text TEXT NOT NULL,
    chroma_id TEXT NOT NULL UNIQUE, embed_model TEXT NOT NULL,
    token_count INTEGER DEFAULT 0, created_at TEXT NOT NULL
);
CREATE TABLE IF NOT EXISTS api_keys (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL, name TEXT NOT NULL,
    key_hash TEXT NOT NULL UNIQUE, key_prefix TEXT NOT NULL,
    scopes_json TEXT DEFAULT '["chat","rag"]',
    rate_limit_rpm INTEGER DEFAULT 60, is_active INTEGER DEFAULT 1,
    last_used_at TEXT, usage_count INTEGER DEFAULT 0, expires_at TEXT,
    created_at TEXT NOT NULL,
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS agent_jobs (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL, name TEXT NOT NULL,
    description TEXT DEFAULT '', job_type TEXT NOT NULL,
    trigger_json TEXT DEFAULT '{}', prompt_template TEXT NOT NULL,
    tools_json TEXT DEFAULT '[]', model_id TEXT DEFAULT 'llama-3.3-70b-versatile',
    status TEXT DEFAULT 'idle', enabled INTEGER DEFAULT 1,
    max_retries INTEGER DEFAULT 2, retry_count INTEGER DEFAULT 0,
    timeout_seconds INTEGER DEFAULT 120, total_runs INTEGER DEFAULT 0,
    success_runs INTEGER DEFAULT 0, failed_runs INTEGER DEFAULT 0,
    last_run_at TEXT, last_run_status TEXT, next_run_at TEXT,
    metadata_json TEXT DEFAULT '{}', created_at TEXT NOT NULL, updated_at TEXT NOT NULL,
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS agent_job_logs (
    id TEXT PRIMARY KEY, job_id TEXT NOT NULL, user_id TEXT NOT NULL,
    run_number INTEGER NOT NULL, status TEXT NOT NULL, trigger_type TEXT DEFAULT 'cron',
    prompt_rendered TEXT NOT NULL, result_text TEXT,
    tool_calls_json TEXT DEFAULT '[]', tokens_used INTEGER DEFAULT 0,
    latency_ms INTEGER DEFAULT 0, error_message TEXT,
    started_at TEXT NOT NULL, finished_at TEXT,
    FOREIGN KEY(job_id) REFERENCES agent_jobs(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS connectors (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL, name TEXT NOT NULL,
    connector_type TEXT NOT NULL, encrypted_creds TEXT NOT NULL,
    is_active INTEGER DEFAULT 1, last_tested_at TEXT,
    last_test_ok INTEGER, test_error TEXT,
    metadata_json TEXT DEFAULT '{}', created_at TEXT NOT NULL, updated_at TEXT NOT NULL,
    UNIQUE(user_id, connector_type),
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS notifications (
    id TEXT PRIMARY KEY, user_id TEXT,
    title TEXT NOT NULL, message TEXT NOT NULL,
    type TEXT DEFAULT 'info', link TEXT,
    is_read INTEGER DEFAULT 0, is_broadcast INTEGER DEFAULT 0,
    actor_id TEXT, expires_at TEXT, created_at TEXT NOT NULL
);
CREATE TABLE IF NOT EXISTS coupons (
    id TEXT PRIMARY KEY, code TEXT NOT NULL UNIQUE COLLATE NOCASE,
    grants_tier TEXT NOT NULL, duration_days INTEGER DEFAULT 30,
    max_uses INTEGER DEFAULT 1, current_uses INTEGER DEFAULT 0,
    is_active INTEGER DEFAULT 1, created_by TEXT NOT NULL,
    expires_at TEXT, created_at TEXT NOT NULL
);
CREATE TABLE IF NOT EXISTS coupon_redemptions (
    id TEXT PRIMARY KEY, coupon_id TEXT NOT NULL, user_id TEXT NOT NULL,
    redeemed_at TEXT NOT NULL, UNIQUE(coupon_id, user_id)
);
CREATE TABLE IF NOT EXISTS billing_orders (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL, provider TEXT DEFAULT 'razorpay',
    tier TEXT NOT NULL, amount INTEGER NOT NULL, currency TEXT NOT NULL,
    provider_order_id TEXT UNIQUE, provider_payment_id TEXT,
    receipt TEXT, status TEXT DEFAULT 'created', raw_json TEXT DEFAULT '{}',
    created_at TEXT NOT NULL, updated_at TEXT NOT NULL,
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS websites (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL,
    title TEXT NOT NULL, description TEXT DEFAULT '',
    filename TEXT NOT NULL UNIQUE, style TEXT DEFAULT 'modern',
    prompt_used TEXT, html_size_bytes INTEGER DEFAULT 0,
    is_public INTEGER DEFAULT 0, view_count INTEGER DEFAULT 0,
    created_at TEXT NOT NULL, updated_at TEXT NOT NULL,
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS audit_log (
    id TEXT PRIMARY KEY, actor_id TEXT NOT NULL, target_id TEXT,
    action TEXT NOT NULL, detail_json TEXT DEFAULT '{}',
    ip_address TEXT, created_at TEXT NOT NULL
);
CREATE TABLE IF NOT EXISTS rate_limits (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL,
    resource TEXT NOT NULL, window_key TEXT NOT NULL,
    count INTEGER DEFAULT 0, updated_at TEXT NOT NULL,
    UNIQUE(user_id, resource, window_key)
);
CREATE TABLE IF NOT EXISTS code_run_logs (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL,
    language TEXT NOT NULL, cwd TEXT DEFAULT '',
    code TEXT NOT NULL, stdout TEXT DEFAULT '', stderr TEXT DEFAULT '',
    exit_code INTEGER DEFAULT 0, duration_ms INTEGER DEFAULT 0,
    created_at TEXT NOT NULL,
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS message_feedback (
    id TEXT PRIMARY KEY, message_id TEXT NOT NULL, user_id TEXT NOT NULL,
    rating INTEGER DEFAULT 0, feedback TEXT DEFAULT '',
    created_at TEXT NOT NULL, UNIQUE(message_id, user_id)
);
CREATE TABLE IF NOT EXISTS voice_sessions (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL,
    room_name TEXT NOT NULL, livekit_token TEXT,
    language TEXT DEFAULT 'auto', status TEXT DEFAULT 'active',
    started_at TEXT NOT NULL, ended_at TEXT,
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS computer_tasks (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL,
    task_type TEXT NOT NULL, command TEXT NOT NULL,
    result TEXT, status TEXT DEFAULT 'pending',
    stdout TEXT, stderr TEXT, exit_code INTEGER,
    created_at TEXT NOT NULL, completed_at TEXT,
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS ai_models (
    id TEXT PRIMARY KEY, name TEXT NOT NULL UNIQUE,
    provider TEXT NOT NULL, base_url TEXT NOT NULL,
    model_name TEXT NOT NULL, encrypted_api_key TEXT NOT NULL,
    is_active INTEGER DEFAULT 1, is_default INTEGER DEFAULT 0,
    is_fast INTEGER DEFAULT 0, is_vision INTEGER DEFAULT 0,
    is_code INTEGER DEFAULT 0, context_length INTEGER DEFAULT 4096,
    max_output_tokens INTEGER DEFAULT 4096,
    temperature_default REAL DEFAULT 0.7,
    description TEXT DEFAULT '', tags_json TEXT DEFAULT '[]',
    created_by TEXT NOT NULL, created_at TEXT NOT NULL, updated_at TEXT NOT NULL
);
CREATE TABLE IF NOT EXISTS mcp_servers (
    id TEXT PRIMARY KEY, name TEXT NOT NULL UNIQUE,
    description TEXT DEFAULT '', server_url TEXT NOT NULL,
    auth_type TEXT DEFAULT 'none' CHECK(auth_type IN('none','bearer','api_key','basic')),
    encrypted_creds TEXT DEFAULT '', tools_json TEXT DEFAULT '[]',
    is_active INTEGER DEFAULT 1, is_public INTEGER DEFAULT 0,
    created_by TEXT NOT NULL, last_pinged_at TEXT,
    ping_ok INTEGER DEFAULT 0, created_at TEXT NOT NULL, updated_at TEXT NOT NULL
);
CREATE TABLE IF NOT EXISTS slash_commands (
    id TEXT PRIMARY KEY, command TEXT NOT NULL UNIQUE,
    connector_type TEXT NOT NULL, description TEXT DEFAULT '',
    default_action TEXT NOT NULL, default_params TEXT DEFAULT '{}',
    is_active INTEGER DEFAULT 1, created_at TEXT NOT NULL
);
CREATE TABLE IF NOT EXISTS connector_logs (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL,
    connector_type TEXT NOT NULL, action TEXT NOT NULL,
    status TEXT NOT NULL CHECK(status IN('success','failure')),
    latency_ms INTEGER DEFAULT 0, error_msg TEXT,
    created_at TEXT NOT NULL
);
CREATE TABLE IF NOT EXISTS skills (
    id TEXT PRIMARY KEY,
    name TEXT NOT NULL UNIQUE,
    description TEXT DEFAULT '',
    activation_keywords TEXT DEFAULT '[]',
    instructions TEXT NOT NULL,
    resources_json TEXT DEFAULT '[]',
    script_metadata_json TEXT DEFAULT '{}',
    is_enabled INTEGER DEFAULT 1,
    created_by TEXT NOT NULL,
    created_at TEXT NOT NULL,
    updated_at TEXT NOT NULL
);
CREATE TABLE IF NOT EXISTS skill_resources (
    id TEXT PRIMARY KEY,
    skill_id TEXT NOT NULL,
    name TEXT NOT NULL,
    content TEXT NOT NULL,
    mime_type TEXT DEFAULT 'text/plain',
    created_at TEXT NOT NULL,
    FOREIGN KEY(skill_id) REFERENCES skills(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS skill_usage_logs (
    id TEXT PRIMARY KEY,
    skill_id TEXT NOT NULL,
    user_id TEXT NOT NULL,
    session_id TEXT,
    message_id TEXT,
    matched_by TEXT DEFAULT '',
    created_at TEXT NOT NULL,
    FOREIGN KEY(skill_id) REFERENCES skills(id) ON DELETE CASCADE,
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS smart_connectors (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL,
    app_name TEXT NOT NULL, connector_type TEXT NOT NULL,
    custom_url TEXT DEFAULT '', encrypted_conn_data TEXT NOT NULL,
    actions_json TEXT DEFAULT '[]', status TEXT DEFAULT 'active',
    last_tested_at TEXT, last_test_ok INTEGER, test_error TEXT,
    created_at TEXT NOT NULL, updated_at TEXT NOT NULL,
    UNIQUE(user_id, connector_type),
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS oauth_states (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL,
    connector_type TEXT NOT NULL, state TEXT NOT NULL UNIQUE,
    code_verifier TEXT DEFAULT '', redirect_uri TEXT DEFAULT '',
    scopes TEXT DEFAULT '', expires_at TEXT NOT NULL, created_at TEXT NOT NULL,
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS supabase_oauth_states (
    id TEXT PRIMARY KEY, state TEXT NOT NULL UNIQUE,
    code_verifier TEXT NOT NULL, redirect_uri TEXT NOT NULL,
    expires_at TEXT NOT NULL, created_at TEXT NOT NULL
);
CREATE TABLE IF NOT EXISTS user_last_result (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL UNIQUE,
    connector_type TEXT NOT NULL, result_json TEXT NOT NULL, updated_at TEXT NOT NULL,
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS platform_connectors (
    id TEXT PRIMARY KEY,
    connector_type TEXT NOT NULL UNIQUE,
    is_enabled INTEGER DEFAULT 0,
    display_name TEXT NOT NULL,
    category TEXT NOT NULL,
    icon TEXT DEFAULT '🔌',
    setup_status TEXT DEFAULT 'not_configured'
        CHECK(setup_status IN('ready','partial','not_configured')),
    admin_notes TEXT DEFAULT '',
    requires_oauth INTEGER DEFAULT 1,
    env_keys TEXT DEFAULT '[]',
    updated_at TEXT NOT NULL,
    created_at TEXT NOT NULL
);
CREATE TABLE IF NOT EXISTS sidebar_features (
    id TEXT PRIMARY KEY,
    feature_key TEXT NOT NULL UNIQUE,
    display_name TEXT NOT NULL,
    icon TEXT DEFAULT '',
    page_key TEXT NOT NULL,
    sort_order INTEGER DEFAULT 0,
    is_enabled INTEGER DEFAULT 1,
    created_at TEXT NOT NULL,
    updated_at TEXT NOT NULL
);
CREATE TABLE IF NOT EXISTS user_sidebar_feature_access (
    id TEXT PRIMARY KEY,
    user_id TEXT NOT NULL,
    feature_key TEXT NOT NULL,
    is_enabled INTEGER NOT NULL,
    updated_by TEXT,
    updated_at TEXT NOT NULL,
    UNIQUE(user_id, feature_key),
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS model_access (
    id TEXT PRIMARY KEY,
    model_id TEXT NOT NULL UNIQUE,
    display_name TEXT NOT NULL,
    source TEXT DEFAULT 'db',
    is_enabled INTEGER DEFAULT 1,
    created_at TEXT NOT NULL,
    updated_at TEXT NOT NULL
);
CREATE TABLE IF NOT EXISTS user_model_access (
    id TEXT PRIMARY KEY,
    user_id TEXT NOT NULL,
    model_id TEXT NOT NULL,
    is_enabled INTEGER NOT NULL,
    updated_by TEXT,
    updated_at TEXT NOT NULL,
    UNIQUE(user_id, model_id),
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS platform_settings (
    key TEXT PRIMARY KEY,
    value TEXT NOT NULL,
    updated_by TEXT,
    updated_at TEXT NOT NULL
);
CREATE TABLE IF NOT EXISTS plan_sidebar_feature_access (
    id TEXT PRIMARY KEY,
    plan TEXT NOT NULL,
    feature_key TEXT NOT NULL,
    is_enabled INTEGER NOT NULL,
    updated_by TEXT,
    updated_at TEXT NOT NULL,
    UNIQUE(plan, feature_key)
);
CREATE TABLE IF NOT EXISTS plan_model_access (
    id TEXT PRIMARY KEY,
    plan TEXT NOT NULL,
    model_id TEXT NOT NULL,
    is_enabled INTEGER NOT NULL,
    updated_by TEXT,
    updated_at TEXT NOT NULL,
    UNIQUE(plan, model_id)
);
CREATE TABLE IF NOT EXISTS livy_sessions (
    id TEXT PRIMARY KEY, user_id TEXT NOT NULL,
    livy_id INTEGER, livy_url TEXT NOT NULL, kind TEXT DEFAULT 'pyspark',
    state TEXT DEFAULT 'starting', app_id TEXT, log_lines TEXT DEFAULT '[]',
    created_at TEXT NOT NULL, updated_at TEXT NOT NULL,
    FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
);
CREATE TABLE IF NOT EXISTS livy_statements (
    id TEXT PRIMARY KEY, session_id TEXT NOT NULL, user_id TEXT NOT NULL,
    statement_id INTEGER, code TEXT NOT NULL, state TEXT DEFAULT 'waiting',
    output_json TEXT DEFAULT '{}', progress REAL DEFAULT 0.0,
    started_at TEXT NOT NULL, completed_at TEXT,
    FOREIGN KEY(session_id) REFERENCES livy_sessions(id) ON DELETE CASCADE
);
CREATE INDEX IF NOT EXISTS idx_users_email        ON users(email);
CREATE INDEX IF NOT EXISTS idx_sessions_user      ON chat_sessions(user_id,last_message_at DESC);
CREATE INDEX IF NOT EXISTS idx_history_session    ON chat_history(session_id,created_at);
CREATE INDEX IF NOT EXISTS idx_memories_user      ON memories(user_id,is_active,last_reinforced DESC);
CREATE INDEX IF NOT EXISTS idx_documents_user     ON documents(user_id,uploaded_at DESC);
CREATE INDEX IF NOT EXISTS idx_api_keys_hash      ON api_keys(key_hash);
CREATE INDEX IF NOT EXISTS idx_agent_jobs_user    ON agent_jobs(user_id,enabled);
CREATE INDEX IF NOT EXISTS idx_job_logs_job       ON agent_job_logs(job_id,started_at DESC);
CREATE INDEX IF NOT EXISTS idx_connectors_user    ON connectors(user_id,is_active);
CREATE INDEX IF NOT EXISTS idx_notifications_user ON notifications(user_id,is_read,created_at DESC);
CREATE INDEX IF NOT EXISTS idx_rate_limits        ON rate_limits(user_id,resource,window_key);
CREATE INDEX IF NOT EXISTS idx_code_logs_user     ON code_run_logs(user_id,created_at DESC);
CREATE INDEX IF NOT EXISTS idx_refresh_tokens     ON refresh_tokens(token_hash,revoked);
CREATE INDEX IF NOT EXISTS idx_email_ver_hash     ON email_verification_tokens(token_hash,used_at);
CREATE INDEX IF NOT EXISTS idx_email_ver_user     ON email_verification_tokens(user_id,used_at);
CREATE INDEX IF NOT EXISTS idx_pw_reset_hash      ON password_reset_tokens(token_hash,used_at);
CREATE INDEX IF NOT EXISTS idx_pw_reset_user      ON password_reset_tokens(user_id,used_at);
CREATE INDEX IF NOT EXISTS idx_ai_models_active   ON ai_models(is_active,is_default);
CREATE INDEX IF NOT EXISTS idx_oauth_states       ON oauth_states(state, expires_at);
CREATE INDEX IF NOT EXISTS idx_supabase_oauth_states ON supabase_oauth_states(state, expires_at);
CREATE INDEX IF NOT EXISTS idx_smart_conn_user    ON smart_connectors(user_id, status);
CREATE INDEX IF NOT EXISTS idx_connector_logs     ON connector_logs(connector_type,status,created_at DESC);
CREATE INDEX IF NOT EXISTS idx_user_sidebar_access_user ON user_sidebar_feature_access(user_id);
CREATE INDEX IF NOT EXISTS idx_model_access_enabled ON model_access(is_enabled,source);
CREATE INDEX IF NOT EXISTS idx_user_model_access_user ON user_model_access(user_id);
CREATE INDEX IF NOT EXISTS idx_plan_sidebar_access_plan ON plan_sidebar_feature_access(plan);
CREATE INDEX IF NOT EXISTS idx_plan_model_access_plan ON plan_model_access(plan);
CREATE INDEX IF NOT EXISTS idx_skills_enabled     ON skills(is_enabled,name);
CREATE INDEX IF NOT EXISTS idx_skill_usage_user   ON skill_usage_logs(user_id,created_at DESC);
"""

# ══════════════════════════════════════════════════════════════════════════════
# §4  SECURITY / AUTH
# ══════════════════════════════════════════════════════════════════════════════

pwd_ctx = CryptContext(schemes=["bcrypt"], deprecated="auto", bcrypt__rounds=BCRYPT_ROUNDS)
bearer  = HTTPBearer(auto_error=False)

def _hash_pw(pw: str) -> str: return pwd_ctx.hash(pw[:72].encode())
def _verify_pw(pw: str, h: str) -> bool:
    try: return pwd_ctx.verify(pw[:72].encode(), h)
    except Exception: return False

async def _hash_pw_async(pw: str) -> str:
    return await asyncio.get_running_loop().run_in_executor(_executor, _hash_pw, pw)

async def _verify_pw_async(pw: str, h: str) -> bool:
    return await asyncio.get_running_loop().run_in_executor(_executor, _verify_pw, pw, h)

def _make_access_token(user_id: str, role: str, email: str = "", subscription: str = "free") -> str:
    exp = datetime.now(timezone.utc) + timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    return jwt.encode({"sub": user_id, "role": role, "email": email, "subscription": subscription, "exp": exp},
                      JWT_SECRET, algorithm=JWT_ALGORITHM)

def _decode_jwt(token: str) -> Optional[Dict]:
    try: return jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGORITHM])
    except JWTError: return None

def _make_oauth_state_token(user_id: str, connector_type: str, code_verifier: str = "",
                            redirect_uri: str = "", client_state: str = "") -> str:
    exp = datetime.now(timezone.utc) + timedelta(minutes=30)
    return jwt.encode({
        "typ": "oauth_state",
        "sub": user_id,
        "connector_type": connector_type,
        "code_verifier": code_verifier or "",
        "redirect_uri": redirect_uri or "",
        "client_state": client_state or "",
        "nonce": secrets.token_urlsafe(12),
        "exp": exp,
    }, JWT_SECRET, algorithm=JWT_ALGORITHM)

def _decode_oauth_state_token(state: str) -> Optional[Dict]:
    try:
        payload = jwt.decode(state, JWT_SECRET, algorithms=[JWT_ALGORITHM])
    except JWTError:
        return None
    if payload.get("typ") != "oauth_state":
        return None
    if not payload.get("sub") or not payload.get("connector_type"):
        return None
    return payload

def _hash_token(raw: str) -> str: return hashlib.sha256(raw.encode()).hexdigest()

_EMAIL_RE = re.compile(r"^[A-Z0-9._%+\-]+@[A-Z0-9.\-]+\.[A-Z]{2,63}$", re.I)

def _env_bool(key: str, default: bool = False) -> bool:
    raw = os.getenv(key)
    if raw is None:
        return default
    return str(raw).strip().lower() not in ("0", "false", "no", "off", "")

def _env_int(key: str, default: int, min_value: int, max_value: int) -> int:
    try:
        value = int(os.getenv(key, str(default)) or default)
    except Exception:
        value = default
    return max(min_value, min(max_value, value))

def _email_verification_required() -> bool:
    return _env_bool("REQUIRE_EMAIL_VERIFICATION", True)

def _email_verification_return_link() -> bool:
    return _env_bool("EMAIL_VERIFICATION_RETURN_LINK", False)

def _email_verification_expiry_hours() -> int:
    return _env_int("EMAIL_VERIFICATION_EXPIRE_HOURS", 24, 1, 168)

def _password_reset_expiry_minutes() -> int:
    return _env_int("PASSWORD_RESET_EXPIRE_MINUTES", 30, 5, 1440)

def _password_reset_return_link() -> bool:
    return _env_bool("PASSWORD_RESET_RETURN_LINK", False)

def _supabase_public_config() -> Dict[str, Any]:
    url = os.getenv("NEXT_PUBLIC_SUPABASE_URL") or os.getenv("SUPABASE_URL") or ""
    key = os.getenv("NEXT_PUBLIC_SUPABASE_PUBLISHABLE_KEY") or os.getenv("SUPABASE_PUBLISHABLE_KEY") or ""
    return {
        "enabled": bool(url.strip() and key.strip()),
        "url": url.strip(),
        "publishable_key": key.strip(),
        "scope": "email_verification_only",
        "auth_mode": "jazz_jwt",
    }

def _normalize_email(email: str) -> str:
    return str(email or "").strip().lower()

def _validated_email(email: str) -> str:
    clean = _normalize_email(email)
    if not clean or len(clean) > 254 or not _EMAIL_RE.fullmatch(clean):
        raise HTTPException(400, "Enter a valid email address")
    return clean

def _password_issues(password: str) -> List[str]:
    pw = password or ""
    issues: List[str] = []
    if len(pw) < 8:
        issues.append("at least 8 characters")
    if len(pw) > 128:
        issues.append("128 characters or fewer")
    if re.search(r"\s", pw):
        issues.append("no spaces")
    if not re.search(r"[a-z]", pw):
        issues.append("one lowercase letter")
    if not re.search(r"[A-Z]", pw):
        issues.append("one uppercase letter")
    if not re.search(r"\d", pw):
        issues.append("one number")
    if not re.search(r"[^A-Za-z0-9\s]", pw):
        issues.append("one symbol")
    return issues

def _validated_password(password: str) -> str:
    issues = _password_issues(password)
    if issues:
        raise HTTPException(400, "Password must include " + ", ".join(issues))
    return password

async def _create_email_verification_link(user_id: str, email: str, request: Optional[Request] = None) -> str:
    raw = f"jev_{secrets.token_urlsafe(32)}"
    now = _utcnow()
    expires_at = (datetime.now(timezone.utc) + timedelta(hours=_email_verification_expiry_hours())).isoformat()
    await db_execute(
        "UPDATE email_verification_tokens SET used_at=? WHERE user_id=? AND used_at IS NULL",
        (now, user_id),
    )
    await db_execute(
        "INSERT INTO email_verification_tokens(id,user_id,token_hash,email,expires_at,created_at)"
        " VALUES(?,?,?,?,?,?)",
        (_new_id(), user_id, _hash_token(raw), email, expires_at, now),
    )
    return f"{_auth_link_base_url(request)}/auth/verify-email?token={urllib.parse.quote(raw)}"

def _mail_sender_from() -> str:
    raw = (
        os.getenv("SMTP_FROM")
        or os.getenv("EMAIL_FROM")
        or os.getenv("SMTP_USERNAME")
        or os.getenv("ADMIN_EMAIL")
        or "no-reply@jazzai.online"
    ).strip()
    name = (os.getenv("EMAIL_FROM_NAME") or "JAZZ AI").strip()
    if name and "<" not in raw:
        return f"{name} <{raw}>"
    return raw

def _has_primary_email_sender() -> bool:
    return bool((os.getenv("RESEND_API_KEY") or "").strip() or (os.getenv("SMTP_HOST") or "").strip())

def _supabase_mail_config() -> Tuple[str, str]:
    base = (
        os.getenv("SUPABASE_URL")
        or os.getenv("NEXT_PUBLIC_SUPABASE_URL")
        or ""
    ).strip().rstrip("/")
    key = (
        os.getenv("SUPABASE_PUBLISHABLE_KEY")
        or os.getenv("NEXT_PUBLIC_SUPABASE_PUBLISHABLE_KEY")
        or ""
    ).strip()
    if not base or not key:
        raise RuntimeError("Email delivery is not configured. Set RESEND_API_KEY, SMTP settings, or Supabase publishable auth keys.")
    return base, key

def _send_supabase_link_email_sync(to_email: str, link: str, flow: str) -> Dict[str, Any]:
    to_email = _validated_email(to_email)
    base, key = _supabase_mail_config()
    redirect = urllib.parse.quote(link, safe="")
    endpoint = f"{base}/auth/v1/otp?redirect_to={redirect}"
    payload = {
        "email": to_email,
        "create_user": True,
        "data": {"jazz_auth_flow": flow},
        "options": {"email_redirect_to": link},
    }
    req = urllib.request.Request(
        endpoint,
        data=json.dumps(payload).encode("utf-8"),
        headers={
            "apikey": key,
            "Authorization": f"Bearer {key}",
            "Content-Type": "application/json",
            "X-Client-Info": "jazz-ai-server/14",
        },
        method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=20) as resp:
            body = resp.read().decode("utf-8", "ignore")
            return {"provider": "supabase", "status": resp.status, "body": body[:300]}
    except urllib.error.HTTPError as exc:
        body = exc.read().decode("utf-8", "ignore")
        raise RuntimeError(f"Supabase email failed ({exc.code}): {body[:300]}")

def _supabase_user_from_access_token_sync(access_token: str) -> Dict[str, Any]:
    token = (access_token or "").strip()
    if not token:
        raise RuntimeError("Missing Supabase access token")
    base, key = _supabase_mail_config()
    req = urllib.request.Request(
        f"{base.rstrip('/')}/auth/v1/user",
        headers={
            "apikey": key,
            "Authorization": f"Bearer {token}",
            "X-Client-Info": "jazz-ai-server/14",
        },
        method="GET",
    )
    try:
        with urllib.request.urlopen(req, timeout=20) as resp:
            return json.loads(resp.read().decode("utf-8", "ignore") or "{}")
    except urllib.error.HTTPError as exc:
        body = exc.read().decode("utf-8", "ignore")
        raise RuntimeError(f"Supabase token verification failed ({exc.code}): {body[:300]}")

async def _supabase_user_from_access_token(access_token: str) -> Dict[str, Any]:
    return await asyncio.get_running_loop().run_in_executor(
        _executor, _supabase_user_from_access_token_sync, access_token
    )

def _send_email_sync(to_email: str, subject: str, text: str, html: str = "") -> Dict[str, Any]:
    to_email = _validated_email(to_email)
    sender = _mail_sender_from()
    resend_key = (os.getenv("RESEND_API_KEY") or "").strip()
    if resend_key:
        payload = {
            "from": sender,
            "to": [to_email],
            "subject": subject,
            "text": text,
        }
        if html:
            payload["html"] = html
        req = urllib.request.Request(
            "https://api.resend.com/emails",
            data=json.dumps(payload).encode("utf-8"),
            headers={
                "Authorization": f"Bearer {resend_key}",
                "Content-Type": "application/json",
            },
            method="POST",
        )
        with urllib.request.urlopen(req, timeout=20) as resp:
            body = resp.read().decode("utf-8", "ignore")
            return {"provider": "resend", "status": resp.status, "body": body[:300]}

    smtp_host = (os.getenv("SMTP_HOST") or "").strip()
    if smtp_host:
        smtp_port = _env_int("SMTP_PORT", 587, 1, 65535)
        username = (os.getenv("SMTP_USERNAME") or "").strip()
        password = os.getenv("SMTP_PASSWORD") or ""
        msg = EmailMessage()
        msg["From"] = sender
        msg["To"] = to_email
        msg["Subject"] = subject
        msg.set_content(text)
        if html:
            msg.add_alternative(html, subtype="html")
        use_ssl = _env_bool("SMTP_USE_SSL", False)
        use_tls = _env_bool("SMTP_USE_TLS", not use_ssl)
        if use_ssl:
            with smtplib.SMTP_SSL(smtp_host, smtp_port, timeout=20, context=ssl.create_default_context()) as server:
                if username:
                    server.login(username, password)
                server.send_message(msg)
        else:
            with smtplib.SMTP(smtp_host, smtp_port, timeout=20) as server:
                if use_tls:
                    server.starttls(context=ssl.create_default_context())
                if username:
                    server.login(username, password)
                server.send_message(msg)
        return {"provider": "smtp", "status": "sent"}

    raise RuntimeError("Email delivery is not configured. Set RESEND_API_KEY or SMTP_HOST/SMTP_USERNAME/SMTP_PASSWORD.")

async def _send_email_async(to_email: str, subject: str, text: str, html: str = "") -> Dict[str, Any]:
    return await asyncio.get_running_loop().run_in_executor(
        _executor, _send_email_sync, to_email, subject, text, html
    )

async def _send_supabase_link_email_async(to_email: str, link: str, flow: str) -> Dict[str, Any]:
    return await asyncio.get_running_loop().run_in_executor(
        _executor, _send_supabase_link_email_sync, to_email, link, flow
    )

async def _send_verification_email(email: str, link: str) -> Dict[str, Any]:
    if not _has_primary_email_sender():
        return await _send_supabase_link_email_async(email, link, "verify_email")
    subject = "Verify your JAZZ AI account"
    text = (
        "Welcome to JAZZ AI.\n\n"
        "Verify your email address with this secure link:\n"
        f"{link}\n\n"
        f"This link expires in {_email_verification_expiry_hours()} hours."
    )
    html = (
        "<div style='font-family:Arial,sans-serif;line-height:1.55;color:#111'>"
        "<h2>Verify your JAZZ AI account</h2>"
        "<p>Click the button below to verify your email address.</p>"
        f"<p><a href='{html_lib.escape(link)}' style='display:inline-block;background:#7c6ff7;color:#fff;padding:12px 18px;border-radius:8px;text-decoration:none;font-weight:700'>Verify email</a></p>"
        f"<p style='color:#555;font-size:13px'>This link expires in {_email_verification_expiry_hours()} hours.</p>"
        "</div>"
    )
    return await _send_email_async(email, subject, text, html)

async def _send_password_reset_email(email: str, link: str) -> Dict[str, Any]:
    if not _has_primary_email_sender():
        return await _send_supabase_link_email_async(email, link, "password_reset")
    subject = "Reset your JAZZ AI password"
    text = (
        "Reset your JAZZ AI password with this secure link:\n"
        f"{link}\n\n"
        f"This link expires in {_password_reset_expiry_minutes()} minutes."
    )
    html = (
        "<div style='font-family:Arial,sans-serif;line-height:1.55;color:#111'>"
        "<h2>Reset your JAZZ AI password</h2>"
        "<p>Click the button below to choose a new password.</p>"
        f"<p><a href='{html_lib.escape(link)}' style='display:inline-block;background:#7c6ff7;color:#fff;padding:12px 18px;border-radius:8px;text-decoration:none;font-weight:700'>Reset password</a></p>"
        f"<p style='color:#555;font-size:13px'>This link expires in {_password_reset_expiry_minutes()} minutes.</p>"
        "</div>"
    )
    return await _send_email_async(email, subject, text, html)

def _verification_response(email: str, link: str, email_sent: bool, delivery_error: str = "") -> Dict[str, Any]:
    resp: Dict[str, Any] = {
        "ok": True,
        "verification_required": True,
        "email": email,
        "email_sent": email_sent,
        "message": "Verification email sent. Check your inbox, then sign in." if email_sent else "Verification email could not be sent because email delivery is not configured.",
    }
    if _email_verification_return_link():
        resp["verification_link"] = link
        resp["dev_link"] = True
    if delivery_error:
        resp["delivery_error"] = delivery_error
    return resp

async def _create_password_reset_link(user_id: str, email: str, request: Optional[Request] = None) -> str:
    raw = f"jpr_{secrets.token_urlsafe(32)}"
    now = _utcnow()
    expires_at = (datetime.now(timezone.utc) + timedelta(minutes=_password_reset_expiry_minutes())).isoformat()
    await db_execute(
        "UPDATE password_reset_tokens SET used_at=? WHERE user_id=? AND used_at IS NULL",
        (now, user_id),
    )
    await db_execute(
        "INSERT INTO password_reset_tokens(id,user_id,token_hash,email,expires_at,created_at)"
        " VALUES(?,?,?,?,?,?)",
        (_new_id(), user_id, _hash_token(raw), email, expires_at, now),
    )
    return f"{_auth_link_base_url(request)}/?reset_token={urllib.parse.quote(raw)}"

def _password_reset_response(email: str, link: str = "", email_sent: bool = False, delivery_error: str = "") -> Dict[str, Any]:
    resp: Dict[str, Any] = {
        "ok": True,
        "email_sent": email_sent,
        "message": "If that account exists, a reset email has been sent.",
    }
    if link and _password_reset_return_link():
        resp["reset_link"] = link
        resp["dev_link"] = True
        resp["email"] = email
    if delivery_error:
        resp["delivery_error"] = delivery_error
    return resp

async def _reset_password_by_token(raw_token: str, new_password: str) -> str:
    _validated_password(new_password)
    token = str(raw_token or "").strip()
    if not token:
        raise HTTPException(400, "Missing reset token")
    row = await db_fetchone(
        "SELECT * FROM password_reset_tokens WHERE token_hash=? AND used_at IS NULL",
        (_hash_token(token),),
    )
    if not row:
        raise HTTPException(400, "Reset link is invalid or already used")
    try:
        expires_at = datetime.fromisoformat(row["expires_at"])
    except Exception:
        expires_at = datetime.now(timezone.utc) - timedelta(seconds=1)
    if expires_at < datetime.now(timezone.utc):
        raise HTTPException(400, "Reset link has expired. Please request a new one.")
    now = _utcnow()
    new_hash = await _hash_pw_async(new_password)
    await db_execute(
        "UPDATE users SET password_hash=?,is_verified=1,updated_at=? WHERE id=?",
        (new_hash, now, row["user_id"]),
    )
    await db_execute("UPDATE refresh_tokens SET revoked=1 WHERE user_id=? AND revoked=0", (row["user_id"],))
    await db_execute("UPDATE password_reset_tokens SET used_at=? WHERE id=?", (now, row["id"]))
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), row["user_id"], row["user_id"], "password_reset_link_used",
         json.dumps({"email": row["email"]}), now),
    )
    return row["email"]

async def _verify_email_token(raw_token: str) -> Tuple[bool, str]:
    token = str(raw_token or "").strip()
    if not token:
        return False, "Missing verification token"
    row = await db_fetchone(
        "SELECT * FROM email_verification_tokens WHERE token_hash=? AND used_at IS NULL",
        (_hash_token(token),),
    )
    if not row:
        return False, "Verification link is invalid or already used"
    try:
        expires_at = datetime.fromisoformat(row["expires_at"])
    except Exception:
        expires_at = datetime.now(timezone.utc) - timedelta(seconds=1)
    if expires_at < datetime.now(timezone.utc):
        return False, "Verification link has expired. Please request a new one."
    now = _utcnow()
    await db_execute("UPDATE users SET is_verified=1,updated_at=? WHERE id=?", (now, row["user_id"]))
    await db_execute("UPDATE email_verification_tokens SET used_at=? WHERE id=?", (now, row["id"]))
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), row["user_id"], row["user_id"], "email_verified", json.dumps({"email": row["email"]}), now),
    )
    return True, row["email"]

def _auth_html_page(title: str, message: str, ok: bool = True) -> HTMLResponse:
    color = "#34d399" if ok else "#f87171"
    safe_title = html_lib.escape(title)
    safe_msg = html_lib.escape(message)
    html = f"""<!doctype html><html><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1"><title>{safe_title}</title><style>body{{margin:0;min-height:100vh;display:grid;place-items:center;background:#07070f;color:#e4e4f4;font-family:Inter,system-ui,sans-serif}}.box{{width:min(460px,calc(100vw - 32px));border:1px solid #252540;border-radius:18px;background:#101020;padding:28px;box-shadow:0 20px 80px rgba(0,0,0,.45)}}h1{{margin:0 0 10px;font-size:24px}}p{{color:#a0a0c0;line-height:1.5}}a{{color:#8b7cf6}}</style></head><body><main class="box"><h1 style="color:{color}">{safe_title}</h1><p>{safe_msg}</p><p><a href="/">Back to JAZZ AI</a></p></main></body></html>"""
    return HTMLResponse(html, status_code=200 if ok else 400)

def _supabase_auth_base() -> str:
    base = (
        os.getenv("SUPABASE_URL")
        or os.getenv("NEXT_PUBLIC_SUPABASE_URL")
        or ""
    ).strip().rstrip("/")
    return base

def _supabase_oauth_client_id() -> str:
    return (os.getenv("SUPABASE_OAUTH_CLIENT_ID") or "").strip()

def _supabase_oauth_client_secret() -> str:
    return (os.getenv("SUPABASE_OAUTH_CLIENT_SECRET") or "").strip()

def _supabase_oauth_scopes() -> str:
    return (os.getenv("SUPABASE_OAUTH_SCOPES") or "openid email profile").strip() or "openid email profile"

def _supabase_oauth_auth_method() -> str:
    raw = (os.getenv("SUPABASE_OAUTH_AUTH_METHOD") or "auto").strip().lower()
    if raw not in {"auto", "none", "client_secret_basic", "client_secret_post"}:
        return "auto"
    return raw

def _supabase_oauth_redirect_uri(request: Optional[Request] = None, explicit: str = "") -> str:
    value = (explicit or "").strip()
    if value:
        return value
    return f"{_auth_link_base_url(request)}/auth/supabase/callback"

def _base64url(raw: bytes) -> str:
    return base64.urlsafe_b64encode(raw).decode("ascii").rstrip("=")

def _pkce_pair() -> Tuple[str, str]:
    verifier = _base64url(secrets.token_bytes(48))
    challenge = _base64url(hashlib.sha256(verifier.encode("ascii")).digest())
    return verifier, challenge

async def _supabase_oauth_authorize_url(request: Request, redirect_uri: str = "") -> Dict[str, str]:
    base = _supabase_auth_base()
    client_id = _supabase_oauth_client_id()
    if not base or not client_id:
        raise HTTPException(400, "Supabase OAuth is not configured. Add SUPABASE_OAUTH_CLIENT_ID in .env.")
    redirect = _supabase_oauth_redirect_uri(request, redirect_uri)
    verifier, challenge = _pkce_pair()
    state = _base64url(secrets.token_bytes(32))
    now = _utcnow()
    expires_at = (datetime.now(timezone.utc) + timedelta(minutes=10)).isoformat()
    await db_execute(
        "DELETE FROM supabase_oauth_states WHERE expires_at<?",
        (now,),
    )
    await db_execute(
        "INSERT INTO supabase_oauth_states(id,state,code_verifier,redirect_uri,expires_at,created_at)"
        " VALUES(?,?,?,?,?,?)",
        (_new_id(), state, verifier, redirect, expires_at, now),
    )
    params = {
        "response_type": "code",
        "client_id": client_id,
        "redirect_uri": redirect,
        "scope": _supabase_oauth_scopes(),
        "state": state,
        "code_challenge": challenge,
        "code_challenge_method": "S256",
    }
    return {
        "ok": True,
        "url": f"{base}/auth/v1/oauth/authorize?{urllib.parse.urlencode(params)}",
        "state": state,
        "redirect_uri": redirect,
    }

def _json_request_sync(url: str, *, method: str = "GET", token: str = "", apikey: str = "", timeout: int = 20) -> Dict[str, Any]:
    headers = {"Accept": "application/json", "X-Client-Info": "jazz-ai-server/14"}
    if token:
        headers["Authorization"] = f"Bearer {token}"
    if apikey:
        headers["apikey"] = apikey
    req = urllib.request.Request(url, headers=headers, method=method)
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        raw = resp.read().decode("utf-8", "ignore")
        return json.loads(raw or "{}")

def _form_request_sync(url: str, form: Dict[str, str], headers: Optional[Dict[str, str]] = None, timeout: int = 25) -> Dict[str, Any]:
    req = urllib.request.Request(
        url,
        data=urllib.parse.urlencode(form).encode("utf-8"),
        headers={
            "Accept": "application/json",
            "Content-Type": "application/x-www-form-urlencoded",
            "X-Client-Info": "jazz-ai-server/14",
            **(headers or {}),
        },
        method="POST",
    )
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        raw = resp.read().decode("utf-8", "ignore")
        return json.loads(raw or "{}")

def _decode_jwt_unverified(token: str) -> Dict[str, Any]:
    try:
        part = str(token or "").split(".")[1]
        part += "=" * (-len(part) % 4)
        raw = base64.urlsafe_b64decode(part.encode("ascii"))
        data = json.loads(raw.decode("utf-8", "ignore"))
        return data if isinstance(data, dict) else {}
    except Exception:
        return {}

def _exchange_supabase_oauth_code_sync(code: str, verifier: str, redirect_uri: str) -> Dict[str, Any]:
    base = _supabase_auth_base()
    client_id = _supabase_oauth_client_id()
    if not base or not client_id:
        raise RuntimeError("Supabase OAuth is not configured")
    client_secret = _supabase_oauth_client_secret()
    method = _supabase_oauth_auth_method()
    if method == "auto":
        method = "client_secret_basic" if client_secret else "none"
    form = {
        "grant_type": "authorization_code",
        "code": code,
        "redirect_uri": redirect_uri,
        "code_verifier": verifier,
        "client_id": client_id,
    }
    headers: Dict[str, str] = {}
    if method == "client_secret_basic" and client_secret:
        basic = base64.b64encode(f"{client_id}:{client_secret}".encode("utf-8")).decode("ascii")
        headers["Authorization"] = f"Basic {basic}"
    elif method == "client_secret_post" and client_secret:
        form["client_secret"] = client_secret
    token_url = f"{base}/auth/v1/oauth/token"
    try:
        return _form_request_sync(token_url, form, headers)
    except urllib.error.HTTPError as exc:
        body = exc.read().decode("utf-8", "ignore")
        if exc.code not in (404, 405):
            raise RuntimeError(f"Supabase OAuth token exchange failed ({exc.code}): {body[:300]}")
        legacy = f"{base}/auth/v1/token?grant_type=authorization_code"
        try:
            return _form_request_sync(legacy, form, headers)
        except urllib.error.HTTPError as exc2:
            body2 = exc2.read().decode("utf-8", "ignore")
            raise RuntimeError(f"Supabase OAuth token exchange failed ({exc2.code}): {body2[:300]}")

def _supabase_oauth_userinfo_sync(access_token: str, id_token: str = "") -> Dict[str, Any]:
    base = _supabase_auth_base()
    apikey = (
        os.getenv("SUPABASE_PUBLISHABLE_KEY")
        or os.getenv("NEXT_PUBLIC_SUPABASE_PUBLISHABLE_KEY")
        or ""
    ).strip()
    merged: Dict[str, Any] = {}
    if id_token:
        merged.update(_decode_jwt_unverified(id_token))
    if access_token:
        merged.update(_decode_jwt_unverified(access_token))
    if base and access_token:
        for path in ("/auth/v1/userinfo", "/auth/v1/user"):
            try:
                data = _json_request_sync(f"{base}{path}", token=access_token, apikey=apikey)
                if isinstance(data, dict):
                    merged.update(data)
                    break
            except Exception:
                continue
    return merged

async def _exchange_supabase_oauth_code(body: SupabaseOAuthExchange) -> Dict[str, Any]:
    state = str(body.state or "").strip()
    code = str(body.code or "").strip()
    redirect_uri = str(body.redirect_uri or "").strip()
    row = await db_fetchone("SELECT * FROM supabase_oauth_states WHERE state=?", (state,))
    if not row:
        raise HTTPException(400, "Supabase OAuth session expired. Start again.")
    try:
        expires_at = datetime.fromisoformat(row["expires_at"])
    except Exception:
        expires_at = datetime.now(timezone.utc) - timedelta(seconds=1)
    await db_execute("DELETE FROM supabase_oauth_states WHERE state=?", (state,))
    if expires_at < datetime.now(timezone.utc):
        raise HTTPException(400, "Supabase OAuth session expired. Start again.")
    if redirect_uri != row["redirect_uri"]:
        raise HTTPException(400, "Supabase OAuth redirect mismatch")
    try:
        tokens = await asyncio.get_running_loop().run_in_executor(
            _executor, _exchange_supabase_oauth_code_sync, code, row["code_verifier"], redirect_uri
        )
        profile = await asyncio.get_running_loop().run_in_executor(
            _executor,
            _supabase_oauth_userinfo_sync,
            tokens.get("access_token", ""),
            tokens.get("id_token", ""),
        )
        profile["_tokens"] = {k: bool(tokens.get(k)) for k in ("access_token", "id_token", "refresh_token")}
        return profile
    except HTTPException:
        raise
    except Exception as e:
        logger.warning("[AUTH] Supabase OAuth exchange failed: %s", e)
        raise HTTPException(400, str(e))

def _supabase_profile_email(profile: Dict[str, Any]) -> str:
    email = profile.get("email")
    if not email and isinstance(profile.get("user"), dict):
        email = profile["user"].get("email")
    if not email and isinstance(profile.get("user_metadata"), dict):
        email = profile["user_metadata"].get("email")
    return _validated_email(str(email or ""))

def _supabase_profile_name(profile: Dict[str, Any], email: str) -> str:
    meta = profile.get("user_metadata") if isinstance(profile.get("user_metadata"), dict) else {}
    name = (
        profile.get("name")
        or profile.get("full_name")
        or meta.get("full_name")
        or meta.get("name")
        or email.split("@")[0]
    )
    return str(name or email.split("@")[0]).strip()[:120]

async def _issue_jazz_auth_for_supabase_profile(profile: Dict[str, Any]) -> Dict[str, Any]:
    email = _supabase_profile_email(profile)
    full_name = _supabase_profile_name(profile, email)
    now = _utcnow()
    user = await db_fetchone("SELECT * FROM users WHERE email=?", (email,))
    if user:
        uid = user["id"]
        updates = ["is_verified=1", "is_active=1", "last_login_at=?", "updated_at=?"]
        vals: List[Any] = [now, now]
        if not (user.get("full_name") or "").strip() and full_name:
            updates.append("full_name=?")
            vals.append(full_name)
        vals.append(uid)
        await db_execute(f"UPDATE users SET {','.join(updates)} WHERE id=?", tuple(vals))
    else:
        uid = _new_id()
        password_hash = await _hash_pw_async(secrets.token_urlsafe(32))
        await db_execute(
            "INSERT INTO users(id,email,password_hash,full_name,role,subscription,is_active,is_verified,last_login_at,created_at,updated_at)"
            " VALUES(?,?,?,?,?,?,?,?,?,?,?)",
            (uid, email, password_hash, full_name, "client", "free", 1, 1, now, now, now),
        )
    user = await db_fetchone(
        "SELECT id,email,full_name,role,subscription,memory_enabled,is_verified FROM users WHERE id=?",
        (uid,),
    )
    raw = f"jzr_{secrets.token_urlsafe(48)}"
    h = _hash_token(raw)
    exp = (datetime.now(timezone.utc) + timedelta(minutes=REFRESH_TOKEN_EXPIRE_MINUTES)).isoformat()
    await db_execute(
        "INSERT INTO refresh_tokens(id,user_id,token_hash,expires_at,created_at) VALUES(?,?,?,?,?)",
        (_new_id(), uid, h, exp, now),
    )
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), uid, uid, "supabase_oauth_login", json.dumps({"email": email}), now),
    )
    return {
        "access_token": _make_access_token(user["id"], user["role"], user["email"], user["subscription"]),
        "refresh_token": raw,
        "token_type": "bearer",
        "user": {k: user[k] for k in ("id", "email", "full_name", "role", "subscription", "memory_enabled", "is_verified")},
    }

async def _setting_get(key: str, default: Any = None) -> Any:
    row = await db_fetchone("SELECT value FROM platform_settings WHERE key=?", (key,))
    if not row:
        return default
    try:
        return json.loads(row["value"])
    except Exception:
        return row["value"]

async def _setting_set(key: str, value: Any, actor_id: str = "") -> None:
    now = _utcnow()
    await db_execute(
        "INSERT INTO platform_settings(key,value,updated_by,updated_at) VALUES(?,?,?,?)"
        " ON CONFLICT(key) DO UPDATE SET value=excluded.value,"
        "updated_by=excluded.updated_by,updated_at=excluded.updated_at",
        (key, json.dumps(value), actor_id, now))

async def _subscription_limits_all() -> Dict[str, Dict[str, int]]:
    stored = await _setting_get("subscription_limits", {})
    if not isinstance(stored, dict):
        stored = {}
    merged: Dict[str, Dict[str, int]] = {}
    for plan in PLAN_TIERS:
        base = dict(SUBSCRIPTION_LIMITS.get(plan, {}))
        overrides = stored.get(plan, {})
        if isinstance(overrides, dict):
            for key, value in overrides.items():
                if key not in base:
                    continue
                try:
                    base[key] = max(-1, int(value))
                except Exception:
                    pass
        merged[plan] = base
    return merged

async def _subscription_limits_for_plan(plan: str) -> Dict[str, int]:
    limits = await _subscription_limits_all()
    return limits.get(plan, limits.get("free", {}))

def _prompt_policy_defaults() -> Dict[str, Any]:
    return {
        "assistant_instructions": "",
        "tool_routing_rules": "",
        "website_builder_instructions": "",
        "code_runner_policy": "",
        "admin_notes": "",
        "chat_runtime_applied": False,
    }

async def _prompt_policies() -> Dict[str, Any]:
    stored = await _setting_get("prompt_policies", {})
    if not isinstance(stored, dict):
        stored = {}
    defaults = _prompt_policy_defaults()
    defaults.update({k: stored.get(k, defaults[k]) for k in defaults})
    defaults["chat_runtime_applied"] = False
    return defaults

async def _platform_status() -> Dict[str, Any]:
    return {
        "maintenance": {
            "enabled": bool(await _setting_get("maintenance_enabled", False)),
            "message": await _setting_get("maintenance_message", "JAZZ AI is temporarily in maintenance mode."),
        },
        "banner": {
            "enabled": bool(await _setting_get("broadcast_banner_enabled", False)),
            "title": await _setting_get("broadcast_banner_title", ""),
            "message": await _setting_get("broadcast_banner_message", ""),
            "type": await _setting_get("broadcast_banner_type", "info"),
        },
    }

async def _maintenance_enabled_message() -> Tuple[bool, str]:
    status = await _platform_status()
    m = status.get("maintenance", {})
    return bool(m.get("enabled")), str(m.get("message") or "JAZZ AI is temporarily in maintenance mode.")

async def _get_current_user(request: Request,
                             creds: Optional[HTTPAuthorizationCredentials] = Depends(bearer)) -> Dict:
    token = None
    if creds: token = creds.credentials
    if not token and request:
        ak = request.headers.get("X-API-Key", "") or request.query_params.get("token","")
        if ak and ak.startswith("jz_"):
            h = _hash_token(ak)
            row = await db_fetchone(
                "SELECT ak.user_id, u.role, u.subscription, u.is_active, ak.scopes_json "
                "FROM api_keys ak JOIN users u ON u.id=ak.user_id "
                "WHERE ak.key_hash=? AND ak.is_active=1", (h,))
            if row and row["is_active"]:
                maintenance_on, maintenance_msg = await _maintenance_enabled_message()
                if maintenance_on and row["role"] != "admin":
                    raise HTTPException(status_code=503, detail=maintenance_msg)
                await db_execute("UPDATE api_keys SET last_used_at=?,usage_count=usage_count+1 WHERE key_hash=?",
                                 (_utcnow(), h))
                return {"id": row["user_id"], "sub": row["user_id"], "role": row["role"],
                        "subscription": row["subscription"]}
        if ak and not ak.startswith("jz_"):
            token = ak
    if not token: raise HTTPException(status_code=401, detail="Not authenticated")
    payload = _decode_jwt(token)
    if not payload: raise HTTPException(status_code=401, detail="Invalid token")
    user_id = payload.get("sub")
    user = await db_fetchone("SELECT * FROM users WHERE id=? AND is_active=1", (user_id,))
    if not user: raise HTTPException(status_code=401, detail="User not found")
    maintenance_on, maintenance_msg = await _maintenance_enabled_message()
    if maintenance_on and user.get("role") != "admin":
        raise HTTPException(status_code=503, detail=maintenance_msg)
    user["sub"] = user["id"]
    return user

async def _require_admin(user: Dict = Depends(_get_current_user)) -> Dict:
    if user.get("role") != "admin": raise HTTPException(status_code=403, detail="Admin only")
    return user

# ══════════════════════════════════════════════════════════════════════════════
# §5  RATE LIMITING
# ══════════════════════════════════════════════════════════════════════════════

_ENV_KEY_RE = re.compile(r"^[A-Z_][A-Z0-9_]*$")
_ENV_SECRET_RE = re.compile(r"(KEY|SECRET|TOKEN|PASSWORD|PASS|PRIVATE|CREDENTIAL|FERNET|JWT)", re.I)

def _validate_env_key(key: str) -> str:
    key = (key or "").strip().upper()
    if not _ENV_KEY_RE.match(key):
        raise HTTPException(400, "Environment key must use A-Z, 0-9, and underscores, and start with a letter or underscore")
    return key

def _format_env_line(key: str, value: str) -> str:
    value = "" if value is None else str(value)
    if "\n" in value or "\r" in value:
        raise HTTPException(400, "Environment values cannot contain new lines")
    if value == "":
        return f"{key}="
    if re.search(r"\s|#|\"|'", value):
        escaped = value.replace("\\", "\\\\").replace('"', '\\"')
        return f'{key}="{escaped}"'
    return f"{key}={value}"

def _env_read_pairs() -> Dict[str, str]:
    path = _env_file_path()
    if not path.exists():
        return {}
    values = dotenv_values(path)
    return {str(k): "" if v is None else str(v) for k, v in values.items() if k}

def _env_write_value(key: str, value: str) -> Path:
    key = _validate_env_key(key)
    path = _env_file_path()
    lines = path.read_text(encoding="utf-8").splitlines() if path.exists() else []
    new_line = _format_env_line(key, value)
    found = False
    out: List[str] = []
    pattern = re.compile(rf"^\s*(?:export\s+)?{re.escape(key)}\s*=")
    for line in lines:
        if pattern.match(line):
            if not found:
                out.append(new_line)
                found = True
            continue
        out.append(line)
    if not found:
        if out and out[-1].strip():
            out.append("")
        out.append(new_line)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text("\n".join(out).rstrip() + "\n", encoding="utf-8")
    os.environ[key] = "" if value is None else str(value)
    _reload_runtime_env()
    return path

def _env_delete_value(key: str) -> Path:
    key = _validate_env_key(key)
    path = _env_file_path()
    if not path.exists():
        return path
    pattern = re.compile(rf"^\s*(?:export\s+)?{re.escape(key)}\s*=")
    lines = [line for line in path.read_text(encoding="utf-8").splitlines() if not pattern.match(line)]
    path.write_text("\n".join(lines).rstrip() + ("\n" if lines else ""), encoding="utf-8")
    os.environ.pop(key, None)
    _reload_runtime_env()
    return path

async def _check_rate_limit(user: Dict, resource: str) -> None:
    sub   = user.get("subscription", "free")
    limit = (await _subscription_limits_for_plan(sub)).get(resource, 0)
    if limit == -1: return
    wkey  = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    uid   = user.get("id") or user.get("sub","")
    row   = await db_fetchone(
        "SELECT count FROM rate_limits WHERE user_id=? AND resource=? AND window_key=?",
        (uid, resource, wkey))
    count = row["count"] if row else 0
    if count >= limit: raise HTTPException(status_code=429, detail=f"Daily {resource} limit ({limit}) reached")
    await db_execute(
        "INSERT INTO rate_limits(id,user_id,resource,window_key,count,updated_at) VALUES(?,?,?,?,1,?) "
        "ON CONFLICT(user_id,resource,window_key) DO UPDATE SET count=count+1,updated_at=excluded.updated_at",
        (_new_id(), uid, resource, wkey, _utcnow()))

# ══════════════════════════════════════════════════════════════════════════════
# §6  LLM CLIENT
# ══════════════════════════════════════════════════════════════════════════════

def _groq_client() -> OpenAI:
    return OpenAI(api_key=GROQ_API_KEY, base_url="https://api.groq.com/openai/v1")

def _split_api_tokens(raw: str) -> List[str]:
    if not raw:
        return []
    return [t.strip() for t in re.split(r"[\s,;]+", str(raw)) if t.strip()]

def _hf_api_keys(primary: str = "") -> List[str]:
    keys: List[str] = []
    for raw in (
        primary,
        HF_TOKEN,
        HUGGINGFACE_API_KEY,
        HUGGINGFACEHUB_API_TOKEN,
        HF_TOKEN_BACKUP,
        HF_TOKEN_BACKUPS,
        HF_TOKENS,
    ):
        keys.extend(_split_api_tokens(raw))
    seen: Set[str] = set()
    return [k for k in keys if k and not (k in seen or seen.add(k))]

async def _get_model_client(model_id: str) -> Tuple[OpenAI, str]:
    model_id = _canonical_model_id(model_id)
    row = await db_fetchone("SELECT * FROM ai_models WHERE (id=? OR name=?) AND is_active=1 LIMIT 1",
                            (model_id, model_id))
    if row:
        try:
            creds = _decrypt(row["encrypted_api_key"])
            api_key = creds.get("key", creds.get("api_key", ""))
        except Exception:
            api_key = ""
        base_url = row["base_url"] or _PROVIDER_DEFAULTS.get(row["provider"], "")
        if row["provider"] == "huggingface" and not api_key:
            api_key = (_hf_api_keys("") or [""])[0]
        if row["provider"] == "groq" and not api_key:
            api_key = GROQ_API_KEY
        if row["provider"] == "nvidia" and not api_key:
            api_key = NVIDIA_API_KEY
        if row["provider"] == "nvidia":
            return OpenAI(api_key=api_key or "none", base_url=base_url, timeout=45.0), row["model_name"]
        return OpenAI(api_key=api_key or "none", base_url=base_url), row["model_name"]
    if model_id in HUGGINGFACE_MODELS:
        meta = HUGGINGFACE_MODELS[model_id]
        return OpenAI(api_key=(_hf_api_keys("") or ["none"])[0], base_url=_PROVIDER_DEFAULTS["huggingface"]), meta["model_name"]
    model_name = model_id if model_id in GROQ_MODELS else "llama-3.3-70b-versatile"
    return _groq_client(), model_name

async def _chat_completion_extra_body(model_id: str) -> Dict[str, Any]:
    mid = _canonical_model_id(model_id)
    row = await db_fetchone(
        "SELECT provider,model_name FROM ai_models WHERE (id=? OR name=? OR model_name=?) AND is_active=1 LIMIT 1",
        (mid, mid, mid))
    if not row:
        return {}
    provider = str(row.get("provider") or "").lower()
    model_name = str(row.get("model_name") or "").lower()
    if provider == "nvidia" and model_name == "moonshotai/kimi-k2.6":
        return {"chat_template_kwargs": {"thinking": True}}
    if provider == "huggingface" and model_name == "zai-org/glm-5.1:together":
        return {"chat_template_kwargs": {"thinking": False}}
    return {}

async def _llm_text(messages: List[Dict], model_id: str = "llama-3.3-70b-versatile",
                    max_tokens: int = 1024, temperature: float = 0.7) -> str:
    client, model_name = await _get_model_client(model_id)
    extra_body = await _chat_completion_extra_body(model_id)
    try:
        resp = await asyncio.get_running_loop().run_in_executor(
            _executor,
            lambda c=client, m=model_name, e=extra_body: c.chat.completions.create(
                model=m, messages=messages, max_tokens=max_tokens, temperature=temperature,
                **({"extra_body": e} if e else {})))
        return resp.choices[0].message.content or ""
    except Exception as exc:
        # Fallback to default groq model
        logger.warning("[LLM] %s failed: %s — falling back", model_id, exc)
        try:
            resp = await asyncio.get_running_loop().run_in_executor(
                _executor,
                lambda: _groq_client().chat.completions.create(
                    model="llama-3.3-70b-versatile",
                    messages=messages,
                    max_tokens=max_tokens,
                    temperature=temperature))
            return resp.choices[0].message.content or ""
        except Exception as fallback_exc:
            logger.error("[LLM] fallback failed: %s", fallback_exc)
            return ("I could not get a response from the selected model. "
                    "Please check the model API key in Admin > Env or switch models.")

def _needs_thinking(text: str) -> bool:
    t = text.lower()
    if len(text) < THINKING_MIN_LENGTH: return False
    return any(kw in t for kw in THINKING_KEYWORDS)

def _count_tokens(text: str) -> int: return math.ceil(len(text) / 4)

LOCAL_JAZZ_MODEL_ID = "local-dolphin3-qwen25-05b"
LOCAL_JAZZ_FAST_INPUT_TOKENS = int(os.getenv("LOCAL_JAZZ_FAST_INPUT_TOKENS") or "360")
LOCAL_JAZZ_FAST_MAX_TOKENS = int(os.getenv("LOCAL_JAZZ_FAST_MAX_TOKENS") or "160")

_MODEL_ALIASES = {
    "censored": "llama-3.3-70b-versatile",
    "fast": "llama-3.1-8b-instant",
    "uncensored": "local-dolphin3-qwen25-05b",
    "dolphin": "local-dolphin3-qwen25-05b",
    "jazz-ai": "local-dolphin3-qwen25-05b",
    "jazz": "local-dolphin3-qwen25-05b",
    "local-dolphin": "local-dolphin3-qwen25-05b",
    "jazz-ai-v1": "jazz-ai-testing",
    "jazz-ai-v1.0": "jazz-ai-testing",
    "jazz ai v1.0": "jazz-ai-testing",
    "qwen": "jazz-ai-testing",
    "gwen": "jazz-ai-testing",
    "hf-dolphin": "dolphin-mistral-24b-venice-hf",
    "venice": "dolphin-mistral-24b-venice-hf",
    "dolphin-venice": "dolphin-mistral-24b-venice-hf",
    "dphn/Dolphin-Mistral-24B-Venice-Edition:featherless-ai": "dolphin-mistral-24b-venice-hf",
    "dphn/Dolphin-Mistral-24B-Venice-Edition": "dolphin-mistral-24b-venice-hf",
    "dphn/dolphin-mistral-24b-venice-edition:featherless-ai": "dolphin-mistral-24b-venice-hf",
    "dphn/dolphin-mistral-24b-venice-edition": "dolphin-mistral-24b-venice-hf",
    "dolphin-3.0-llama-3.1-8b": "dolphin-mistral-24b-venice-hf",
    "dolphin-8b": "dolphin-mistral-24b-venice-hf",
    "cognitivecomputations/Dolphin3.0-Llama3.1-8B": "dolphin-mistral-24b-venice-hf",
    "dphn/Dolphin3.0-Llama3.1-8B": "dolphin-mistral-24b-venice-hf",
    "cognitivecomputations/dolphin3.0-llama3.1-8b": "dolphin-mistral-24b-venice-hf",
    "dphn/dolphin3.0-llama3.1-8b": "dolphin-mistral-24b-venice-hf",
    "kimi": "nvidia-moonshotai-kimi-k2-6",
    "kimi-k2.6": "nvidia-moonshotai-kimi-k2-6",
    "kimi-k2-6": "nvidia-moonshotai-kimi-k2-6",
    "moonshotai/kimi-k2.6": "nvidia-moonshotai-kimi-k2-6",
    "moonshotai/kimi-k2-6": "nvidia-moonshotai-kimi-k2-6",
    "nvidia-kimi": "nvidia-moonshotai-kimi-k2-6",
}

def _canonical_model_id(model_id: str) -> str:
    raw = (model_id or "llama-3.3-70b-versatile").strip()
    return _MODEL_ALIASES.get(raw) or _MODEL_ALIASES.get(raw.lower()) or raw

def _is_local_jazz_model(model_id: str) -> bool:
    return _canonical_model_id(model_id) == LOCAL_JAZZ_MODEL_ID

def _model_tags(row: Dict[str, Any]) -> List[str]:
    return [str(t).lower() for t in _safe_json_loads(row.get("tags_json") or row.get("tags"), [])]

def _model_is_internal(row: Dict[str, Any]) -> bool:
    tags = _model_tags(row)
    hay = " ".join(str(row.get(k, "")) for k in ("id", "name", "model_name", "description")).lower()
    return "internal" in tags or "hidden" in tags or "image-to-text" in tags or row.get("id") == IMAGE_TO_TEXT_MODEL_ID or "image-to-text" in hay

def _fit_text_to_token_budget(text: str, token_budget: int) -> Tuple[str, bool]:
    text = text or ""
    token_budget = max(128, int(token_budget or 128))
    if _count_tokens(text) <= token_budget:
        return text, False
    char_budget = max(600, token_budget * 4)
    marker = (
        "\n\n[Middle of this very long input was compacted automatically to fit the "
        "model context. Preserve and reason over the beginning and ending sections.]\n\n"
    )
    keep = max(200, char_budget - len(marker))
    head = max(200, int(keep * 0.65))
    tail = max(200, keep - head)
    if head + tail >= len(text):
        return text[:char_budget], True
    return text[:head].rstrip() + marker + text[-tail:].lstrip(), True

async def _model_display_name(model_id: str) -> str:
    mid = _canonical_model_id(model_id)
    if mid in BUILTIN_MODELS:
        return BUILTIN_MODELS[mid]["label"]
    row = await db_fetchone(
        "SELECT name,model_name FROM ai_models WHERE (id=? OR name=? OR model_name=?) AND is_active=1 LIMIT 1",
        (mid, mid, mid))
    if row:
        return row["name"] or row["model_name"] or mid
    return mid

async def _all_active_model_ids(user_id: Optional[str] = None) -> List[str]:
    allowed = set(await _allowed_model_ids_for_user(user_id))
    ids = [mid for mid in BUILTIN_MODELS.keys() if mid in allowed]
    rows = await db_fetchall(
        "SELECT id,name,model_name,description,tags_json FROM ai_models WHERE is_active=1 ORDER BY is_default DESC,is_fast DESC,name")
    ids.extend(r["id"] for r in rows if r["id"] in allowed and not _model_is_internal(r))
    if not ids:
        ids = ["llama-3.3-70b-versatile"]
    seen, out = set(), []
    for mid in ids:
        mid = _canonical_model_id(mid)
        if mid and mid not in seen:
            seen.add(mid); out.append(mid)
    return out

async def _model_supports_vision(model_id: str) -> bool:
    mid = _canonical_model_id(model_id)
    row = await db_fetchone(
        "SELECT is_vision FROM ai_models WHERE (id=? OR name=? OR model_name=?) AND is_active=1 LIMIT 1",
        (mid, mid, mid))
    return bool(row and row["is_vision"])

async def _vision_model_for_user(user_id: Optional[str], preferred_model_id: str = "") -> Optional[str]:
    preferred = _canonical_model_id(preferred_model_id)
    preferred_row = await db_fetchone(
        "SELECT id,name,model_name,description,tags_json FROM ai_models WHERE id=? LIMIT 1",
        (preferred,))
    if preferred and await _model_supports_vision(preferred) and not (preferred_row and _model_is_internal(preferred_row)):
        if not user_id or await _model_enabled_for_user(user_id, preferred):
            return preferred
    allowed = set(await _allowed_model_ids_for_user(user_id))
    rows = await db_fetchall(
        "SELECT id,name,model_name,description,tags_json FROM ai_models WHERE is_active=1 AND is_vision=1 ORDER BY is_default DESC,is_fast DESC,name")
    for row in rows:
        mid = _canonical_model_id(row["id"])
        if _model_is_internal(row):
            continue
        if not allowed or mid in allowed:
            return mid
    return None

async def _model_fallback_candidates(preferred_model_id: str,
                                     estimated_input_tokens: int = 0,
                                     user_id: Optional[str] = None) -> List[str]:
    preferred = _canonical_model_id(preferred_model_id)
    all_ids = await _all_active_model_ids(user_id)
    if preferred not in all_ids:
        if not all_ids:
            all_ids = [preferred]
        else:
            preferred = all_ids[0]
    preferred_budget = await _context_input_budget(preferred)
    # If the selected model is clearly too small, start with the largest fitting model.
    if estimated_input_tokens and estimated_input_tokens > int(preferred_budget * 0.88):
        scored = []
        for mid in all_ids:
            try:
                scored.append((await _context_input_budget(mid), mid))
            except Exception:
                scored.append((0, mid))
        fitting = [mid for budget, mid in sorted(scored, reverse=True) if budget > estimated_input_tokens + 1024]
        rest = [mid for _, mid in sorted(scored, reverse=True) if mid not in fitting]
        ordered = fitting + rest
    else:
        ordered = [preferred] + [mid for mid in all_ids if mid != preferred]
    seen, out = set(), []
    for mid in ordered:
        mid = _canonical_model_id(mid)
        if mid and mid not in seen:
            seen.add(mid); out.append(mid)
    return out[:6]

async def _auto_route_model(preferred_model_id: str, message: str,
                            user_id: Optional[str] = None) -> Tuple[str, str]:
    preferred = await _ensure_user_model(user_id, preferred_model_id) if user_id else _canonical_model_id(preferred_model_id)
    estimated = _count_tokens(message or "")
    preferred_budget = await _context_input_budget(preferred)
    if estimated <= int(preferred_budget * 0.82):
        return preferred, ""
    candidates = await _model_fallback_candidates(preferred, estimated, user_id)
    for mid in candidates:
        if await _context_input_budget(mid) > estimated + 1024:
            return mid, "long_input"
    return candidates[0] if candidates else preferred, "long_input"

def _message_content_to_text(content: Any) -> str:
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        parts: List[str] = []
        for item in content:
            if isinstance(item, dict):
                if item.get("type") == "text":
                    parts.append(str(item.get("text") or ""))
                elif item.get("text"):
                    parts.append(str(item.get("text")))
            elif item is not None:
                parts.append(str(item))
        return "\n".join(p for p in parts if p)
    return "" if content is None else str(content)

_LOCAL_HINDI_HINTS = {
    "hindi", "hinglish", "mai", "mein", "mujhe", "bata", "batao", "bol",
    "bolo", "kya", "kaise", "karu", "karun", "bheju", "bhejna", "ladki",
    "usne", "usko", "baat", "reply", "message", "pyaar", "pyar",
}
_LOCAL_FRENCH_HINTS = {"bonjour", "francais", "français", "parle", "merci", "salut"}
_LOCAL_SPANISH_HINTS = {"hola", "espanol", "español", "gracias", "habla"}
_LOCAL_GERMAN_HINTS = {"hallo", "deutsch", "danke", "sprich"}
_LOCAL_DATING_HINTS = {
    "dating", "date", "girl", "girls", "ladki", "crush", "flirt", "flirty",
    "text", "texting", "message", "dm", "whatsapp", "reply", "replied",
    "seen", "haha", "nice", "opener", "first message", "interested",
}
_LOCAL_REJECTION_HINTS = {
    "not interested", "no interest", "she rejected", "rejected", "reject",
    "mana kar diya", "interested nahi", "nahi interested", "not looking",
    "leave me", "stop texting",
}

def _local_norm(text: str) -> str:
    return re.sub(r"\s+", " ", (text or "").strip().lower())

def _local_has_devanagari(text: str) -> bool:
    return bool(re.search(r"[\u0900-\u097f]", text or ""))

def _local_word_hits(text: str, words: set) -> int:
    low = _local_norm(text)
    found = 0
    for word in words:
        if " " in word:
            found += 1 if word in low else 0
        else:
            found += 1 if re.search(rf"\b{re.escape(word)}\b", low) else 0
    return found

def _local_detect_language(text: str) -> str:
    low = _local_norm(text)
    if _local_has_devanagari(text) or _local_word_hits(low, _LOCAL_HINDI_HINTS) >= 1:
        return "hi"
    if _local_word_hits(low, _LOCAL_FRENCH_HINTS) >= 1:
        return "fr"
    if _local_word_hits(low, _LOCAL_SPANISH_HINTS) >= 1:
        return "es"
    if _local_word_hits(low, _LOCAL_GERMAN_HINTS) >= 1:
        return "de"
    return "en" if re.search(r"[a-z]{3,}", low) else "hi"

def _local_is_language_switch(text: str) -> bool:
    low = _local_norm(text)
    return bool(
        re.search(r"\b(hindi|hinglish)\b.*\b(bol|bolo|baat|talk|speak|reply)\b", low)
        or re.search(r"\b(bol|bolo|baat|talk|speak|reply)\b.*\b(hindi|hinglish)\b", low)
        or "hindi mai bol" in low
        or "hindi mein baat" in low
    )

def _local_is_dating_intent(text: str) -> bool:
    return _local_word_hits(text, _LOCAL_DATING_HINTS) >= 1

def _local_is_rejection(text: str) -> bool:
    return _local_word_hits(text, _LOCAL_REJECTION_HINTS) >= 1

def _local_is_first_message(text: str) -> bool:
    low = _local_norm(text)
    return (
        "first message" in low
        or "opener" in low
        or "pehla message" in low
        or "kya bheju" in low
        or "message kya" in low
    )

def _local_is_haha_reply(text: str) -> bool:
    low = _local_norm(text)
    return any(x in low for x in ("haha nice", "she replied haha", "she said haha", "reply haha", "nice"))

def _local_lines(title: str, options: List[str], reason: str, lang: str) -> str:
    numbered = "\n".join(f"{i}. {line}" for i, line in enumerate(options, 1))
    if lang == "hi":
        return f"{title}\n\n{numbered}\n\nKyu: {reason}"
    return f"{title}\n\n{numbered}\n\nWhy: {reason}"

def _local_generate_coach_answer(messages: List[Dict]) -> Optional[str]:
    text = _local_generate_last_user_text(messages)
    lang = _local_detect_language(text)
    low = _local_norm(text)
    if re.fullmatch(r"(hi|hello|hey|yo|hii|hiii|namaste|namaskar)", low):
        if lang == "hi":
            return "Haan, main Jazz AI V1.0 hoon. Batao, kis cheez mein help chahiye?"
        return "Hi, I am Jazz AI V1.0. Tell me what you want to work on."
    if re.search(r"\b(who are you|what are you|your name|model are you|kaun ho|tum kaun)\b", low):
        if lang == "hi":
            return "Main Jazz AI V1.0 hoon, ek private from-scratch local model jisme Hindi/Hinglish, coding, tools, aur dating/texting coach routing added hai."
        return "I am Jazz AI V1.0, a private from-scratch local model with added routing for Hindi/Hinglish, coding, tools, and dating/texting coaching."
    if _local_is_language_switch(text):
        return "Haan, ab main Hindi/Hinglish mein baat karunga. Batao, kis cheez mein help chahiye?"
    if lang == "fr" and not _local_is_dating_intent(text):
        return "Oui, je peux parler francais. Dis-moi ce que tu veux faire, et je te repondrai clairement."
    if lang == "es" and not _local_is_dating_intent(text):
        return "Si, puedo hablar espanol. Dime que necesitas y te respondo claro."
    if lang == "de" and not _local_is_dating_intent(text):
        return "Ja, ich kann Deutsch sprechen. Sag mir, wobei ich helfen soll."
    if not _local_is_dating_intent(text):
        return None
    if _local_is_rejection(text):
        if lang == "hi":
            return "Agar usne clearly bola ki interested nahi hai, graceful close kar: 'All good, thanks for being honest. Take care.' Phir move on. Pressure ya chase mat kar."
        return "If she clearly says she is not interested, close respectfully: 'All good, thanks for being honest. Take care.' Then move on. Do not pressure or chase."
    if _local_is_first_message(text):
        return _local_lines(
            "Ye bhej:",
            [
                "Teri vibe dekh ke lag raha hai tu sweet dikhti hai, par thodi trouble bhi hai.",
                "Wait, tu itni calm dikhti hai ya bas profile ka illusion hai?",
                "Tu woh type lagti hai jo innocent face bana ke sabse zyada chaos karti hai.",
            ],
            "Generic hi/hello boring hai. Light assumption curiosity create karti hai.",
            "hi",
        )
    if _local_is_haha_reply(text):
        if lang == "hi":
            return _local_lines(
                "Reply options:",
                [
                    "Haha nice? Bas itna hi? Mujhe laga tumhare paas thoda better comeback hoga.",
                    "Nice matlab impressed ho ya politely judge kar rahi ho?",
                    "Theek hai, ab tumhari turn. Ek honest assumption mere baare mein.",
                ],
                "Uske low-effort reply ko playful challenge mein convert karo.",
                "hi",
            )
        return _local_lines(
            "Reply options:",
            [
                "Haha nice? That's all? I expected a slightly better comeback from you.",
                "Nice as in impressed, or nice as in politely judging me?",
                "Okay, your turn. Make one honest assumption about me.",
            ],
            "Turn the low-effort reply into a playful challenge without chasing.",
            "en",
        )
    return _local_lines(
        "Dating/texting game plan:",
        [
            "Pehle boring question mat pooch. Curiosity hook ya light assumption se start kar.",
            "Uski profile/vibe se ek playful assumption bana.",
            "Reply aaye to 'why?' pooch aur uske answer par light tease kar.",
            "Energy low ho ya rejection ho to respectfully disengage.",
        ],
        "Curiosity -> assumption -> why -> light tease -> imagination game.",
        "hi" if lang == "hi" else "en",
    )

def _local_generate_system_context(last_user: str) -> str:
    lang = _local_detect_language(last_user)
    text = (
        "You are Jazz AI V1.0, a private from-scratch local LLM. "
        "Answer the current user directly. Default to Hindi/Hinglish when unclear. "
        "If the user clearly uses another language, reply in that same language. "
        "For dating or texting advice, be bold, playful, concise, and respectful. "
        "Use curiosity openers, light assumptions, question reframing, and imagination games. "
        "Never pressure after rejection; tell the user to disengage respectfully. "
        "Do not repeat role labels."
    )
    if lang == "hi":
        text += " Reply in natural Hindi/Hinglish."
    return text

def _messages_to_local_generate_prompt(messages: List[Dict]) -> str:
    lines: List[str] = ["System: " + _local_generate_system_context(_local_generate_last_user_text(messages))]
    for msg in messages[-8:]:
        role = str(msg.get("role") or "user").lower()
        if role == "system":
            continue
        text = _message_content_to_text(msg.get("content")).strip()
        if not text:
            continue
        if len(text) > 2000:
            text = text[-2000:]
        if role == "assistant":
            lines.append("Assistant: " + text)
        else:
            lines.append("User: " + text)
    if not lines or not lines[-1].startswith("Assistant:"):
        lines.append("Assistant:")
    return "\n".join(lines)

def _strip_local_generate_echo(text: str, prompt: str) -> str:
    out = (text or "").strip()
    prompt = (prompt or "").strip()
    if prompt and out.startswith(prompt):
        out = out[len(prompt):].strip()
    out = re.sub(r"^(?:Assistant:|assistant:)\s*", "", out).strip()
    return out or (text or "").strip()

def _local_generate_last_user_text(messages: List[Dict]) -> str:
    for msg in reversed(messages):
        if str(msg.get("role") or "").lower() == "user":
            text = _message_content_to_text(msg.get("content")).strip()
            if text:
                return text
    return ""

def _local_safe_math_eval(node: ast.AST) -> float:
    if isinstance(node, ast.Expression):
        return _local_safe_math_eval(node.body)
    if isinstance(node, ast.Constant) and isinstance(node.value, (int, float)):
        return node.value
    if isinstance(node, ast.UnaryOp):
        value = _local_safe_math_eval(node.operand)
        if isinstance(node.op, ast.UAdd):
            return value
        if isinstance(node.op, ast.USub):
            return -value
    if isinstance(node, ast.BinOp):
        left = _local_safe_math_eval(node.left)
        right = _local_safe_math_eval(node.right)
        if isinstance(node.op, ast.Add):
            return left + right
        if isinstance(node.op, ast.Sub):
            return left - right
        if isinstance(node.op, ast.Mult):
            return left * right
        if isinstance(node.op, ast.Div):
            if right == 0:
                raise ZeroDivisionError("division by zero")
            return left / right
        if isinstance(node.op, ast.FloorDiv):
            if right == 0:
                raise ZeroDivisionError("division by zero")
            return left // right
        if isinstance(node.op, ast.Mod):
            if right == 0:
                raise ZeroDivisionError("division by zero")
            return left % right
        if isinstance(node.op, ast.Pow):
            if abs(right) > 12:
                raise ValueError("exponent too large")
            return left ** right
    raise ValueError("unsupported math expression")

def _local_format_math_number(value: float) -> str:
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    if isinstance(value, float):
        return f"{value:.10g}"
    return str(value)

def _local_generate_simple_math_answer(messages: List[Dict]) -> Optional[str]:
    text = _local_generate_last_user_text(messages).strip().lower()
    if not text:
        return None
    text = text.replace("×", "*").replace("÷", "/").replace("−", "-")
    text = re.sub(r"\b(what\s+is|what's|calculate|compute|solve|please|answer|equals|equal\s+to)\b", " ", text)
    text = text.replace("?", " ").replace("=", " ")
    text = re.sub(r"\bplus\b", "+", text)
    text = re.sub(r"\bminus\b", "-", text)
    text = re.sub(r"\btimes\b|\bmultiplied\s+by\b", "*", text)
    text = re.sub(r"\bdivided\s+by\b", "/", text)
    expr = re.sub(r"\s+", " ", text).strip()
    if not re.fullmatch(r"[0-9\s\.\+\-\*\/%\(\)]{3,}", expr):
        return None
    if not re.search(r"[0-9]\s*[\+\-\*\/%]\s*[0-9]", expr):
        return None
    try:
        result = _local_safe_math_eval(ast.parse(expr, mode="eval"))
    except ZeroDivisionError:
        return "Cannot divide by zero."
    except Exception:
        return None
    return f"{expr} = {_local_format_math_number(result)}"

async def _local_generate_text_once(messages: List[Dict], row: Dict[str, Any],
                                    max_tokens: int = 512) -> Tuple[str, str, str]:
    base = (row["base_url"] or _PROVIDER_DEFAULTS["local_generate"]).rstrip("/")
    prompt = _messages_to_local_generate_prompt(messages)
    row_limit = int(row.get("max_output_tokens") or 48)
    token_limit = max(8, min(int(max_tokens or row_limit), row_limit, 128))
    math_answer = _local_generate_simple_math_answer(messages)
    if math_answer:
        return math_answer, _canonical_model_id(row["id"]), row["model_name"]

    coach_answer = _local_generate_coach_answer(messages)
    if coach_answer:
        return coach_answer, _canonical_model_id(row["id"]), row["model_name"]

    def _call() -> str:
        payload = json.dumps({
            "prompt": prompt,
            "max_new_tokens": token_limit,
            "temperature": float(row.get("temperature_default") or 0.7),
            "top_p": 0.9,
        }).encode("utf-8")
        req = urllib.request.Request(
            base + "/generate",
            data=payload,
            method="POST",
            headers={"Content-Type": "application/json"},
        )
        with urllib.request.urlopen(req, timeout=45) as resp:
            raw = resp.read().decode("utf-8", "replace")
        try:
            data = json.loads(raw)
        except Exception:
            return raw
        if isinstance(data, dict):
            for key in ("text", "generated_text", "response", "content", "output"):
                if data.get(key):
                    return str(data[key])
        return raw

    raw_text = await asyncio.get_running_loop().run_in_executor(_executor, _call)
    return _strip_local_generate_echo(raw_text, prompt), _canonical_model_id(row["id"]), row["model_name"]

async def _llm_text_once(messages: List[Dict], model_id: str,
                         max_tokens: int = 1024,
                         temperature: float = 0.7) -> Tuple[str, str, str]:
    nvidia_meta = await _nvidia_stream_model_meta(model_id)
    if nvidia_meta:
        chunks: List[str] = []
        provider_model = nvidia_meta["model_name"]
        async for kind, payload, meta in _stream_nvidia_text_once(messages, model_id, max_tokens, temperature):
            if kind == "delta":
                chunks.append(payload)
                provider_model = meta.get("provider_model") or provider_model
        return _strip_model_thinking_markup("".join(chunks)), _canonical_model_id(model_id), provider_model
    mid = _canonical_model_id(model_id)
    hf_row = await db_fetchone(
        "SELECT * FROM ai_models WHERE (id=? OR name=? OR model_name=?) AND is_active=1 LIMIT 1",
        (mid, mid, mid))
    if hf_row and str(hf_row["provider"]).lower() == "local_generate":
        return await _local_generate_text_once(messages, hf_row, max_tokens)
    if hf_row and str(hf_row["provider"]).lower() == "huggingface":
        try:
            creds = _decrypt(hf_row["encrypted_api_key"])
            primary_key = creds.get("key", creds.get("api_key", ""))
        except Exception:
            primary_key = ""
        keys = _hf_api_keys(primary_key) or ["none"]
        base_url = hf_row["base_url"] or _PROVIDER_DEFAULTS["huggingface"]
        extra_body = await _chat_completion_extra_body(mid)
        last_error: Optional[Exception] = None
        for idx, api_key in enumerate(keys):
            try:
                if idx > 0:
                    logger.info("[LLM] HF completion retrying %s with backup token #%s", mid, idx)
                client = OpenAI(api_key=api_key, base_url=base_url)
                resp = await asyncio.get_running_loop().run_in_executor(
                    _executor,
                    lambda c=client, m=hf_row["model_name"], e=extra_body: c.chat.completions.create(
                        model=m, messages=messages, max_tokens=max_tokens, temperature=temperature,
                        **({"extra_body": e} if e else {})))
                return _strip_model_thinking_markup(resp.choices[0].message.content or ""), mid, hf_row["model_name"]
            except Exception as exc:
                last_error = exc
                logger.warning("[LLM] HF completion token #%s failed for %s: %s", idx + 1, mid, exc)
        raise last_error or RuntimeError("No Hugging Face tokens available")
    if mid in HUGGINGFACE_MODELS:
        meta = HUGGINGFACE_MODELS[mid]
        keys = _hf_api_keys("") or ["none"]
        last_error: Optional[Exception] = None
        for idx, api_key in enumerate(keys):
            try:
                if idx > 0:
                    logger.info("[LLM] HF builtin completion retrying %s with backup token #%s", mid, idx)
                client = OpenAI(api_key=api_key, base_url=_PROVIDER_DEFAULTS["huggingface"])
                resp = await asyncio.get_running_loop().run_in_executor(
                    _executor,
                    lambda c=client, m=meta["model_name"]: c.chat.completions.create(
                        model=m, messages=messages, max_tokens=max_tokens, temperature=temperature))
                return _strip_model_thinking_markup(resp.choices[0].message.content or ""), mid, meta["model_name"]
            except Exception as exc:
                last_error = exc
                logger.warning("[LLM] HF builtin completion token #%s failed for %s: %s", idx + 1, mid, exc)
        raise last_error or RuntimeError("No Hugging Face tokens available")
    client, model_name = await _get_model_client(model_id)
    extra_body = await _chat_completion_extra_body(model_id)
    provider_messages = _messages_for_provider_model(messages, model_name)
    resp = await asyncio.get_running_loop().run_in_executor(
        _executor,
        lambda c=client, m=model_name, e=extra_body: c.chat.completions.create(
            model=m, messages=provider_messages, max_tokens=max_tokens, temperature=temperature,
            **({"extra_body": e} if e else {})))
    return _strip_model_thinking_markup(resp.choices[0].message.content or ""), _canonical_model_id(model_id), model_name

def _messages_for_provider_model(messages: List[Dict], provider_model: str) -> List[Dict]:
    """Provider-specific prompt nudge for models that otherwise emit hidden thinking."""
    model = str(provider_model or "").lower()
    if "qwen/qwen3" not in model:
        return messages
    out = [dict(m) for m in messages]
    for i in range(len(out) - 1, -1, -1):
        if out[i].get("role") != "user":
            continue
        content = out[i].get("content")
        if isinstance(content, str) and "/no_think" not in content.lower():
            out[i]["content"] = "/no_think " + content
            break
    return out

def _strip_model_thinking_markup(text: str) -> str:
    clean = str(text or "")
    clean = re.sub(r"(?is)<think>\s*</think>\s*", "", clean)
    clean = re.sub(r"(?is)<think>.*?</think>\s*", "", clean)
    return clean.strip()

def _image_text_is_weak(text: str) -> bool:
    cleaned = re.sub(r"\s+", " ", (text or "").strip()).strip(" .,:;!-").lower()
    if not cleaned:
        return True
    color_only = {
        "white", "black", "blue", "green", "red", "yellow", "orange", "purple",
        "gray", "grey", "blank", "empty", "plain white", "a white image",
        "the image is white", "white background"
    }
    if cleaned in color_only:
        return True
    words = re.findall(r"[A-Za-z0-9_]+", cleaned)
    if len(words) < 10 and not any(ch.isdigit() for ch in cleaned):
        return True
    return False

async def _image_to_text_model_id() -> Optional[str]:
    row = await db_fetchone(
        "SELECT id FROM ai_models WHERE id=? AND is_active=1 AND is_vision=1 LIMIT 1",
        (IMAGE_TO_TEXT_MODEL_ID,))
    if row:
        return row["id"]
    row = await db_fetchone(
        "SELECT id FROM ai_models WHERE is_active=1 AND is_vision=1 "
        "AND (tags_json LIKE '%image-to-text%' OR tags_json LIKE '%internal%') "
        "ORDER BY updated_at DESC LIMIT 1")
    return row["id"] if row else None

async def _image_to_text_model_candidates(user_id: Optional[str] = None,
                                          preferred_model_id: str = "") -> List[str]:
    candidates: List[str] = []
    internal = await _image_to_text_model_id()
    if internal:
        candidates.append(_canonical_model_id(internal))
    preferred = _canonical_model_id(preferred_model_id)
    if preferred and await _model_supports_vision(preferred):
        row = await db_fetchone(
            "SELECT id,name,model_name,description,tags_json FROM ai_models WHERE id=? LIMIT 1",
            (preferred,))
        if not (row and _model_is_internal(row)):
            if not user_id or await _model_enabled_for_user(user_id, preferred):
                candidates.append(preferred)
    rows = await db_fetchall(
        "SELECT id,name,model_name,description,tags_json FROM ai_models WHERE is_active=1 AND is_vision=1 ORDER BY is_default DESC,is_fast DESC,name")
    allowed = set(await _allowed_model_ids_for_user(user_id)) if user_id else set()
    for row in rows:
        mid = _canonical_model_id(row["id"])
        if _model_is_internal(row) and mid != internal:
            continue
        if allowed and mid not in allowed and mid != internal:
            continue
        candidates.append(mid)
    seen, out = set(), []
    for mid in candidates:
        if mid and mid not in seen:
            seen.add(mid); out.append(mid)
    return out

async def _extract_image_text(image_url: str, question: str = "",
                              preferred_model_id: str = "",
                              user_id: Optional[str] = None) -> Dict[str, Any]:
    candidates = await _image_to_text_model_candidates(user_id, preferred_model_id)
    if not candidates:
        raise RuntimeError("No active image-to-text or vision model is available")
    prompt = (
        "Convert the image into precise text for a downstream chat model. Return structured Markdown with:\n"
        "1. Image type and purpose.\n"
        "2. All visible OCR text, labels, headings, numbers, buttons, and captions.\n"
        "3. Layout/objects, arrows, sequence steps, tables, charts, code, or diagram relationships.\n"
        "4. Anything uncertain or too small to read.\n"
        "Be faithful and detailed. Do not answer the user's question directly unless it helps describe the image. "
        "Never reduce the image to a single color or generic phrase when there is visible content.\n\n"
        f"User question/context: {question or 'Describe the image accurately.'}"
    )
    messages = [
        {"role":"system","content":"You are a precise image-to-text/OCR extraction model for JAZZ AI."},
        {"role":"user","content":[
            {"type":"text","text":prompt},
            {"type":"image_url","image_url":{"url":image_url}},
        ]},
    ]
    last_error: Optional[Exception] = None
    best: Tuple[str, str, str] = ("", candidates[0], "")
    for mid in candidates[:4]:
        try:
            text, resolved_id, provider_model = await _llm_text_once(messages, mid, max_tokens=4096, temperature=0.12)
            text = (text or "").strip()
            best = (text, resolved_id or mid, provider_model)
            if text and not _image_text_is_weak(text):
                break
            logger.warning("[VISION] weak image extraction from %s: %r", mid, text[:120])
        except Exception as exc:
            last_error = exc
            logger.warning("[VISION] image extraction failed with %s: %s", mid, exc)
    text, resolved_id, provider_model = best
    if not text:
        raise RuntimeError(f"Image-to-text model returned empty output: {last_error}" if last_error else "Image-to-text model returned empty output")
    final_id = _canonical_model_id(resolved_id or candidates[0])
    return {
        "text": text,
        "model_id": final_id,
        "model_label": await _model_display_name(final_id),
        "provider_model": provider_model,
        "weak": _image_text_is_weak(text),
    }

def _chat_image_inputs(req: Any) -> List[Dict[str, str]]:
    items: List[Dict[str, str]] = []
    if getattr(req, "image_url", None):
        items.append({"name":"attached image", "mime":"image/*", "url":req.image_url})
    for att in getattr(req, "attachments", []) or []:
        if not isinstance(att, dict):
            continue
        url = att.get("image_url") or att.get("data_url") or att.get("dataUrl") or att.get("url")
        mime = str(att.get("type") or att.get("mime") or "")
        if not url:
            continue
        if not (str(url).startswith("data:image/") or mime.startswith("image/")):
            continue
        items.append({
            "name": str(att.get("name") or att.get("filename") or "attached image")[:160],
            "mime": mime or "image/*",
            "url": str(url),
        })
    seen, out = set(), []
    for item in items:
        key = hashlib.sha256(item["url"].encode("utf-8", "ignore")).hexdigest()
        if key in seen:
            continue
        seen.add(key); out.append(item)
    return out[:6]

async def _build_image_context(req: Any, user_id: str, model_id: str,
                               tool_log: Optional[List[Dict[str, Any]]] = None,
                               emit=None) -> Tuple[str, List[Dict[str, Any]]]:
    images = _chat_image_inputs(req)
    if not images:
        return req.message, []
    extracted: List[Dict[str, Any]] = []
    for idx, img in enumerate(images, 1):
        start = {"type":"tool_start","tool":"image_to_text","label":f"Reading image {idx}/{len(images)}"}
        if tool_log is not None:
            tool_log.append(start)
        if emit:
            await emit(start)
        try:
            result = await _extract_image_text(img["url"], req.message, model_id, user_id)
            label = "Image converted to text"
            if result.get("weak"):
                label = "Image converted to text with low confidence"
            done = {
                "type":"tool_result",
                "tool":"image_to_text",
                "label":label,
                "count":1,
                "model_id":result.get("model_id"),
                "model_label":result.get("model_label") or "Image-to-text model",
            }
            if tool_log is not None:
                tool_log.append({**done, "image_name": img["name"], "excerpt": result["text"][:2000]})
            if emit:
                await emit(done)
            extracted.append({"name":img["name"], "mime":img["mime"], **result})
        except Exception as exc:
            logger.warning("[VISION] image-to-text extraction failed: %s", exc)
            err = {"type":"tool_error","tool":"image_to_text","label":f"Image {idx} text extraction failed."}
            if tool_log is not None:
                tool_log.append({**err, "image_name": img["name"], "error":str(exc)[:500]})
            if emit:
                await emit(err)
            extracted.append({"name":img["name"], "mime":img["mime"], "text":f"[Extraction failed: {exc}]", "weak":True})
    parts = [
        req.message,
        "",
        "[Attached image analysis]",
        "Use the following extracted visual/OCR content as the source of truth for attached images. "
        "If the user asks what the image contains, answer from this content and mention uncertainty where extraction says so.",
    ]
    for idx, item in enumerate(extracted, 1):
        parts.append(f"\nImage {idx}: {item.get('name','attached image')} ({item.get('mime','image/*')})")
        parts.append(f"Extractor: {item.get('model_label','unknown')}")
        if item.get("weak"):
            parts.append("Confidence: low - verify carefully and avoid overclaiming.")
        parts.append(str(item.get("text") or "").strip())
    return "\n".join(parts).strip(), extracted

async def _llm_text_with_fallback(messages: List[Dict], model_id: str,
                                  max_tokens: int = 1024,
                                  temperature: float = 0.7,
                                  user_id: Optional[str] = None) -> Dict[str, Any]:
    estimated = sum(_count_tokens(m.get("content", "") if isinstance(m.get("content"), str) else "")
                    for m in messages)
    requested = await _ensure_user_model(user_id, model_id) if user_id else _canonical_model_id(model_id)
    candidates = await _model_fallback_candidates(requested, estimated, user_id)
    events: List[Dict[str, Any]] = []
    errors: List[Dict[str, str]] = []
    previous = requested
    for idx, mid in enumerate(candidates):
        label = await _model_display_name(mid)
        effective_max_tokens = min(max_tokens, LOCAL_JAZZ_FAST_MAX_TOKENS) if _is_local_jazz_model(mid) else max_tokens
        if idx == 0:
            events.append({"type":"tool_progress","tool":"model","label":f"Asking {label}..."})
        else:
            prev_label = await _model_display_name(previous)
            events.append({
                "type":"model_switch", "from_model_id":previous, "from_model_label":prev_label,
                "model_id":mid, "model_label":label, "selected":False,
                "label":f"{prev_label} did not answer. Trying {label}..."
            })
        try:
            text, resolved_id, provider_model = await _llm_text_once(messages, mid, effective_max_tokens, temperature)
            if text and text.strip():
                final_id = _canonical_model_id(resolved_id or mid)
                final_label = await _model_display_name(final_id)
                if final_id != requested or idx > 0:
                    events.append({
                        "type":"model_switch", "from_model_id":requested,
                        "model_id":final_id, "model_label":final_label,
                        "provider_model":provider_model, "selected":True,
                        "label":f"Using {final_label} for this response."
                    })
                return {
                    "content": text.strip(), "model_id": final_id, "model_label": final_label,
                    "provider_model": provider_model, "events": events, "errors": errors,
                }
            errors.append({"model_id":mid, "error":"empty response"})
        except Exception as exc:
            logger.warning("[LLM] %s failed during smart fallback: %s", mid, exc)
            errors.append({"model_id":mid, "error":str(exc)[:500]})
            events.append({
                "type":"tool_error",
                "tool":"model",
                "model_id":mid,
                "model_label":label,
                "label":_provider_failure_label(label, exc),
            })
        previous = mid
    msg = ("I could not get a response from any configured model. "
           "Please check Admin > Env/API keys or add another active model.")
    events.append({"type":"tool_error","tool":"model","label":"All configured models failed."})
    return {
        "content": msg, "model_id": requested,
        "model_label": await _model_display_name(requested),
        "provider_model": requested, "events": events, "errors": errors,
    }

async def _hf_stream_model_meta(model_id: str) -> Optional[Dict[str, Any]]:
    if not HAS_HF_INFERENCE_CLIENT:
        return None
    mid = _canonical_model_id(model_id)
    row = await db_fetchone(
        "SELECT * FROM ai_models WHERE (id=? OR name=? OR model_name=?) AND is_active=1 LIMIT 1",
        (mid, mid, mid))
    if not row or str(row["provider"]).lower() != "huggingface":
        return None
    try:
        creds = _decrypt(row["encrypted_api_key"])
        api_key = creds.get("key", creds.get("api_key", ""))
    except Exception:
        api_key = ""
    api_keys = _hf_api_keys(api_key)
    if not api_keys:
        return None
    return {
        "model_id": mid,
        "model_name": row["model_name"],
        "label": row["name"] or row["model_name"] or mid,
        "api_key": api_keys[0],
        "api_keys": api_keys,
        "max_output_tokens": int(row["max_output_tokens"] or 1024),
        "temperature": float(row["temperature_default"] or 0.2),
    }

async def _stream_hf_text_once(messages: List[Dict], model_id: str,
                               max_tokens: int = 1024,
                               temperature: float = 0.7):
    meta = await _hf_stream_model_meta(model_id)
    if not meta:
        raise RuntimeError("Hugging Face streaming is not available for this model")
    loop = asyncio.get_running_loop()
    queue: asyncio.Queue = asyncio.Queue()
    effective_max_tokens = max(16, min(int(max_tokens or 1024), int(meta["max_output_tokens"] or 1024)))
    effective_temperature = float(temperature if temperature is not None else meta["temperature"])

    def _run_stream():
        last_error: Optional[Exception] = None
        keys = [k for k in (meta.get("api_keys") or [meta.get("api_key")]) if k]
        for idx, api_key in enumerate(keys):
            try:
                if idx > 0:
                    logger.info("[LLM] HF stream retrying %s with backup token #%s", meta["model_id"], idx)
                client = InferenceClient(api_key=api_key)
                stream = client.chat.completions.create(
                    model=meta["model_name"],
                    messages=messages,
                    stream=True,
                    max_tokens=effective_max_tokens,
                    temperature=effective_temperature,
                )
                for chunk in stream:
                    choices = getattr(chunk, "choices", None) or []
                    if not choices:
                        continue
                    delta = getattr(choices[0].delta, "content", "") or ""
                    if delta:
                        loop.call_soon_threadsafe(queue.put_nowait, ("delta", delta))
                loop.call_soon_threadsafe(queue.put_nowait, ("done", None))
                return
            except Exception as exc:
                last_error = exc
                logger.warning("[LLM] HF token #%s failed for %s: %s", idx + 1, meta["model_id"], exc)
                if idx + 1 < len(keys):
                    continue
        loop.call_soon_threadsafe(queue.put_nowait, ("error", last_error or RuntimeError("No Hugging Face tokens available")))

    loop.run_in_executor(_executor, _run_stream)
    while True:
        kind, payload = await queue.get()
        if kind == "delta":
            yield payload, meta
        elif kind == "error":
            raise payload
        else:
            break

def _provider_failure_label(label: str, exc: Exception) -> str:
    err = str(exc or "")
    lower = err.lower()
    if "402" in err or "payment required" in lower or "depleted your monthly included credits" in lower:
        return (
            f"{label} is unavailable because Hugging Face Inference Provider credits are depleted. "
            "Trying the next working model..."
        )
    if "401" in err or "unauthorized" in lower or "invalid token" in lower:
        return f"{label} API key is not authorized. Trying the next working model..."
    if "404" in err or "not found" in lower:
        return f"{label} provider/model route was not found. Trying the next working model..."
    if "429" in err or "rate limit" in lower or "too many requests" in lower:
        return f"{label} is rate limited. Trying the next working model..."
    if "timeout" in lower or "timed out" in lower:
        return f"{label} timed out. Trying the next working model..."
    return f"{label} stream failed. Trying the next working model..."

async def _nvidia_stream_model_meta(model_id: str) -> Optional[Dict[str, Any]]:
    mid = _canonical_model_id(model_id)
    row = await db_fetchone(
        "SELECT * FROM ai_models WHERE (id=? OR name=? OR model_name=?) AND is_active=1 LIMIT 1",
        (mid, mid, mid))
    if not row or str(row["provider"]).lower() != "nvidia":
        return None
    try:
        creds = _decrypt(row["encrypted_api_key"])
        api_key = creds.get("key", creds.get("api_key", ""))
    except Exception:
        api_key = ""
    api_key = api_key or NVIDIA_API_KEY
    if not api_key:
        return None
    return {
        "model_id": mid,
        "model_name": row["model_name"],
        "label": row["name"] or row["model_name"] or mid,
        "api_key": api_key,
        "base_url": (row["base_url"] or _PROVIDER_DEFAULTS["nvidia"]).rstrip("/"),
        "max_output_tokens": int(row.get("max_output_tokens") or 4096),
        "temperature": float(row.get("temperature_default") or 0.7),
    }

def _nvidia_provider_model_attempts(model_name: str) -> List[str]:
    model = str(model_name or "").strip()
    attempts = NVIDIA_KIMI_PROVIDER_FALLBACKS.get(model, [model])
    seen: Set[str] = set()
    return [m for m in attempts if m and not (m in seen or seen.add(m))]

def _nvidia_stream_timeout(provider_model: str) -> int:
    return max(8, int(NVIDIA_STREAM_TIMEOUTS.get(provider_model, 45)))

def _nvidia_enable_thinking(provider_model: str) -> bool:
    return provider_model in {"moonshotai/kimi-k2.6", "moonshotai/kimi-k2-thinking"}

async def _stream_nvidia_text_once(messages: List[Dict], model_id: str,
                                   max_tokens: int = 1024,
                                   temperature: float = 0.7):
    meta = await _nvidia_stream_model_meta(model_id)
    if not meta:
        raise RuntimeError("NVIDIA streaming is not available for this model")
    loop = asyncio.get_running_loop()
    queue: asyncio.Queue = asyncio.Queue()
    effective_max_tokens = max(16, min(int(max_tokens or 1024), int(meta["max_output_tokens"] or 1024)))
    effective_temperature = float(temperature if temperature is not None else meta["temperature"])

    def _run_stream():
        last_error: Optional[Exception] = None
        attempts = _nvidia_provider_model_attempts(meta["model_name"])
        for idx, provider_model in enumerate(attempts):
            attempt_meta = {**meta, "provider_model": provider_model}
            if idx > 0:
                label = f"{meta['label']} is slow on NVIDIA. Trying {provider_model}..."
                loop.call_soon_threadsafe(queue.put_nowait, ("status", label, attempt_meta))
            try:
                logger.info("[LLM] NVIDIA stream attempt %s via %s", meta["model_id"], provider_model)
                timeout_seconds = _nvidia_stream_timeout(provider_model)
                payload = {
                    "model": provider_model,
                    "messages": messages,
                    "max_tokens": effective_max_tokens,
                    "temperature": effective_temperature,
                    "top_p": 1.0,
                    "stream": True,
                }
                if _nvidia_enable_thinking(provider_model):
                    payload["chat_template_kwargs"] = {"thinking": True}
                req = urllib.request.Request(
                    meta["base_url"] + "/chat/completions",
                    data=json.dumps(payload).encode("utf-8"),
                    method="POST",
                    headers={
                        "Authorization": "Bearer " + meta["api_key"],
                        "Accept": "text/event-stream",
                        "Content-Type": "application/json",
                    },
                )
                with urllib.request.urlopen(req, timeout=timeout_seconds) as resp:
                    for raw in resp:
                        line = raw.decode("utf-8", "replace").strip()
                        if not line:
                            continue
                        if line.startswith("data:"):
                            line = line[5:].strip()
                        if not line or line == "[DONE]":
                            continue
                        try:
                            data = json.loads(line)
                        except Exception:
                            continue
                        choices = data.get("choices") or []
                        if not choices:
                            continue
                        delta = choices[0].get("delta") or {}
                        reasoning = delta.get("reasoning_content") or delta.get("thinking") or ""
                        content = delta.get("content") or ""
                        if reasoning:
                            loop.call_soon_threadsafe(queue.put_nowait, ("thinking", reasoning, attempt_meta))
                        if content:
                            loop.call_soon_threadsafe(queue.put_nowait, ("delta", content, attempt_meta))
                loop.call_soon_threadsafe(queue.put_nowait, ("done", None, attempt_meta))
                return
            except Exception as exc:
                last_error = exc
                logger.warning("[LLM] NVIDIA provider model %s failed for %s: %s",
                               provider_model, meta["model_id"], exc)
                continue
        if last_error:
            loop.call_soon_threadsafe(queue.put_nowait, ("error", last_error, meta))
        else:
            loop.call_soon_threadsafe(queue.put_nowait, ("error", RuntimeError("No NVIDIA provider models configured"), meta))
        return

    loop.run_in_executor(_executor, _run_stream)
    while True:
        item = await queue.get()
        kind, payload = item[0], item[1]
        item_meta = item[2] if len(item) > 2 else meta
        if kind in ("delta", "thinking", "status"):
            yield kind, payload, item_meta
        elif kind == "error":
            raise payload
        else:
            break

async def _model_context_limit(model_id: str) -> int:
    resolved = _MODEL_ALIASES.get(model_id, model_id)
    if resolved in BUILTIN_MODELS:
        return int(BUILTIN_MODELS[resolved]["ctx"])
    row = await db_fetchone(
        "SELECT context_length FROM ai_models WHERE (id=? OR name=? OR model_name=?) AND is_active=1 LIMIT 1",
        (model_id, model_id, model_id))
    if row and row["context_length"]:
        return int(row["context_length"])
    return int(GROQ_MODELS["llama-3.3-70b-versatile"]["ctx"])

async def _context_input_budget(model_id: str) -> int:
    model_limit = await _model_context_limit(model_id)
    reserve = max(512, min(CONTEXT_OUTPUT_RESERVE_TOKENS, max(512, model_limit // 3)))
    return max(1024, min(MAX_CONTEXT_TOKENS, max(1024, model_limit - reserve)))

async def _model_request_input_budget(model_id: str) -> int:
    """Keep calls below provider/request TPM ceilings even when context windows are large."""
    mid = _canonical_model_id(model_id)
    if _is_local_jazz_model(mid):
        return max(192, min(768, LOCAL_JAZZ_FAST_INPUT_TOKENS))
    default_cap = int(os.getenv("MODEL_REQUEST_MAX_TOKENS") or "9000")
    if mid == "llama-3.1-8b-instant":
        default_cap = min(default_cap, 4200)
    elif mid == "llama-3.3-70b-versatile":
        default_cap = min(default_cap, 8500)
    return max(1024, min(await _context_input_budget(mid), default_cap))

# ══════════════════════════════════════════════════════════════════════════════
# §7  RAG  (ChromaDB)
# ══════════════════════════════════════════════════════════════════════════════

_chroma_client: Optional[Any] = None
_chroma_ef: Optional[Any] = None

def _init_chroma():
    global _chroma_client, _chroma_ef
    if not HAS_CHROMA: return
    try:
        _chroma_client = chromadb.PersistentClient(path="./chroma_db")
        _chroma_ef = embedding_functions.SentenceTransformerEmbeddingFunction(
            model_name="all-MiniLM-L6-v2")
        logger.info("[RAG] ChromaDB ready")
    except Exception as e:
        logger.warning("[RAG] ChromaDB init failed: %s", e); _chroma_client = None

def _get_collection(user_id: str, collection: str = "documents"):
    if not _chroma_client: return None
    name = f"u_{user_id[:8]}_{collection}"
    return _chroma_client.get_or_create_collection(name=name, embedding_function=_chroma_ef,
                                                    metadata={"hnsw:space":"cosine"})

def _chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    sentences = re.split(r'(?<=[.!?])\s+', text.replace("\n\n","  "))
    chunks, current = [], ""
    for sent in sentences:
        if len(current) + len(sent) + 1 > size and current:
            chunks.append(current.strip())
            words = current.split()
            carry = " ".join(words[max(0, len(words)-overlap//6):])
            current = carry + " " + sent
        else:
            current = (current + " " + sent).strip() if current else sent
    if current.strip(): chunks.append(current.strip())
    return [c for c in chunks if len(c) >= 40]

def _extract_text_sync(path: Path, mime: str = "") -> Tuple[str, int]:
    suf = path.suffix.lower()
    try:
        if suf == ".pdf" and HAS_FITZ:
            doc = fitz.open(str(path))
            return "\n".join(p.get_text() for p in doc), len(doc)
        if suf in (".docx",".doc") and HAS_DOCX:
            doc = DocxDocument(str(path))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip()), 0
        if suf in (".xlsx",".xls") and HAS_OPENPYXL:
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            lines = []
            for ws in wb.worksheets:
                lines.append(f"Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    r = " | ".join(str(c or "") for c in row)
                    if r.strip(): lines.append(r)
            return "\n".join(lines), 0
        if suf == ".csv":
            with open(str(path), encoding="utf-8", errors="replace") as f:
                return "\n".join(" | ".join(row) for row in csv.reader(f)), 0
        return path.read_text(encoding="utf-8", errors="replace"), 0
    except Exception as e:
        logger.warning("[RAG] extract_text: %s", e); return "", 0

async def _index_document(doc_id: str, user_id: str, path: Path, collection: str = "documents"):
    col = _get_collection(user_id, collection)
    text, pages = await asyncio.get_running_loop().run_in_executor(
        _executor, _extract_text_sync, path, "")
    if not text.strip():
        await db_execute("UPDATE documents SET index_error='No text extracted',is_indexed=0 WHERE id=?", (doc_id,))
        return
    chunks = _chunk_text(text)
    if not chunks:
        await db_execute("UPDATE documents SET index_error='No chunks',is_indexed=0 WHERE id=?", (doc_id,))
        return
    now = _utcnow()
    ids = [f"{doc_id}_c{i}" for i in range(len(chunks))]
    metas = [{"doc_id":doc_id,"chunk":i,"user_id":user_id,"original_name":path.name}
             for i in range(len(chunks))]
    if col and ids:
        try:
            for b in range(0, len(ids), 100):
                col.upsert(ids=ids[b:b+100], documents=chunks[b:b+100], metadatas=metas[b:b+100])
        except Exception as e: logger.warning("[RAG] chroma upsert: %s", e)
    await db_executemany(
        "INSERT OR REPLACE INTO embeddings_meta(id,document_id,user_id,chunk_index,chunk_text,chroma_id,embed_model,token_count,created_at)"
        " VALUES(?,?,?,?,?,?,?,?,?)",
        [(_new_id(), doc_id, user_id, i, chunks[i], ids[i], "all-MiniLM-L6-v2",
          _count_tokens(chunks[i]), now) for i in range(len(chunks))])
    await db_execute(
        "UPDATE documents SET is_indexed=1,indexed_at=?,chunk_count=?,char_count=?,page_count=?,index_error=NULL WHERE id=?",
        (now, len(chunks), len(text), pages, doc_id))
    logger.info("[RAG] Processed %d chunks for doc %s", len(chunks), doc_id)

def _bm25_score(query_terms: List[str], doc: str, avg_len: float = 400,
                k1: float = 1.5, b: float = 0.75) -> float:
    doc_lower = doc.lower(); words = doc_lower.split(); doc_len = len(words); score = 0.0
    for term in query_terms:
        tf = doc_lower.count(term)
        if tf == 0: continue
        idf = math.log(1 + 1.0 / (tf + 0.5))
        tf_norm = (tf * (k1+1)) / (tf + k1 * (1 - b + b * doc_len / max(avg_len,1)))
        score += idf * tf_norm
    return score

def _is_document_access_query(query: str) -> bool:
    q = (query or "").lower()
    return any(t in q for t in (
        "rag", "document", "documents", "doc ", "docs", "file", "files",
        "uploaded", "upload", "processed", "indexed", "knowledge", "source"
    ))

async def _processed_document_rows(user_id: str, limit: int = 8) -> List[Dict]:
    return await db_fetchall(
        "SELECT id,original_name,file_size_bytes,chunk_count,uploaded_at,indexed_at "
        "FROM documents WHERE user_id=? AND is_indexed=1 AND COALESCE(chunk_count,0)>0 "
        "ORDER BY uploaded_at DESC LIMIT ?",
        (user_id, limit))

def _format_rag_document_status(docs: List[Dict]) -> str:
    if not docs:
        return (
            "[RAG document status]\n"
            "RAG is enabled, but this user has no processed documents available yet."
        )
    lines = [
        "[RAG document status]",
        f"RAG is enabled. {len(docs)} processed document(s) are available for this user.",
        "If the user asks whether you can access uploaded documents, answer yes for these processed JAZZ documents. Do not claim you cannot access uploaded documents in JAZZ.",
    ]
    for d in docs:
        lines.append(
            f"- {d.get('original_name') or 'document'} "
            f"({int(d.get('chunk_count') or 0)} chunks, uploaded {d.get('uploaded_at') or 'unknown'})"
        )
    return "\n".join(lines)

async def _rag_sqlite_search(user_id: str, query: str, k: int = TOP_K_RETRIEVAL,
                             recent: bool = False) -> List[str]:
    rows = await db_fetchall(
        "SELECT em.chunk_text,em.chunk_index,d.original_name,d.uploaded_at "
        "FROM embeddings_meta em JOIN documents d ON d.id=em.document_id "
        "WHERE em.user_id=? AND d.user_id=? AND d.is_indexed=1 AND COALESCE(d.chunk_count,0)>0 "
        "ORDER BY d.uploaded_at DESC, em.chunk_index ASC LIMIT 600",
        (user_id, user_id))
    if not rows:
        return []
    terms = [
        t for t in re.findall(r"[a-zA-Z0-9_]{3,}", (query or "").lower())
        if t not in {"the", "and", "for", "you", "are", "can", "this", "that", "with", "from", "have", "using"}
    ]
    if recent or not terms or (_is_document_access_query(query) and len(terms) <= 2):
        chosen = rows[:k]
    else:
        scored = []
        for r in rows:
            text = str(r.get("chunk_text") or "")
            hay = (text + " " + str(r.get("original_name") or "")).lower()
            score = sum(hay.count(t) for t in terms)
            if score > 0:
                scored.append((score, r))
        chosen = [r for _, r in sorted(scored, key=lambda x: x[0], reverse=True)[:k]] or rows[:k]
    return [
        f"[Source: {r.get('original_name') or 'document'} | chunk {int(r.get('chunk_index') or 0) + 1}]\n{r.get('chunk_text') or ''}"
        for r in chosen
        if (r.get("chunk_text") or "").strip()
    ]

async def _rag_search(user_id: str, query: str, k: int = TOP_K_RETRIEVAL,
                      collection: str = "documents") -> List[str]:
    col = _get_collection(user_id, collection)
    if not col:
        return await _rag_sqlite_search(user_id, query, k)
    try:
        cnt = col.count()
        if cnt == 0:
            return await _rag_sqlite_search(user_id, query, k)
        res = col.query(query_texts=[query], n_results=min(k*2, cnt),
                        where={"user_id": user_id},
                        include=["documents","metadatas","distances"])
        if not res or not res["documents"] or not res["documents"][0]:
            return await _rag_sqlite_search(user_id, query, k)
        docs   = res["documents"][0]
        metas  = res.get("metadatas",[None])[0] or [{}]*len(docs)
        dists  = res.get("distances",[None])[0] or [0.5]*len(docs)
        candidates = [{"text":d,"meta":m,"sem_rank":i+1,"sem_score":1-dv}
                      for i,(d,m,dv) in enumerate(zip(docs,metas,dists))]
        query_terms = [t.lower() for t in query.split() if len(t) > 3]
        avg_len = sum(len(c["text"].split()) for c in candidates) / max(len(candidates),1)
        for c in candidates:
            c["bm25"] = _bm25_score(query_terms, c["text"], avg_len)
        bm25_sorted = sorted(candidates, key=lambda x: x["bm25"], reverse=True)
        for rank, c in enumerate(bm25_sorted): c["bm25_rank"] = rank+1
        kk = 60
        for c in candidates:
            c["rrf"] = (1/(kk+c["sem_rank"])) + (1/(kk+c["bm25_rank"]))
        fused = sorted(candidates, key=lambda x: x["rrf"], reverse=True)[:k]
        parts = []
        for c in fused[:RERANK_TOP_K]:
            fname = c["meta"].get("original_name", c["meta"].get("filename","document"))
            parts.append(f"[Source: {fname}]\n{c['text']}")
        return parts or await _rag_sqlite_search(user_id, query, k)
    except Exception as e:
        logger.warning("[RAG] search error: %s", e)
        return await _rag_sqlite_search(user_id, query, k)

# ══════════════════════════════════════════════════════════════════════════════
# §8  MEMORY
# ══════════════════════════════════════════════════════════════════════════════

_NOISY_MEMORY_KEYS = {
    "user_input", "user_query", "user_question", "user_score", "user_attachment",
    "ai_answer", "ai_rating", "assistant_answer", "assistant_response",
    "processed_document", "document_uploaded", "uploaded_file", "source", "sources",
}
_NOISY_MEMORY_PREFIXES = ("ai_", "assistant_", "bot_", "response_", "answer_")
_LOW_VALUE_MEMORY_RE = re.compile(r"^(yes|no|ok|okay|hi|hello|thanks|thank you|done|[0-9]+/?[0-9]*|[0-9.]+\s*(kb|mb|gb)?)$", re.I)
_MEMORY_TERM_STOPWORDS = {
    "the", "and", "for", "with", "this", "that", "what", "when", "where", "why",
    "how", "you", "are", "can", "should", "would", "could", "please", "about",
}

def _normalize_memory_key(key: str) -> str:
    return re.sub(r"[^a-z0-9_]+", "_", (key or "").strip().lower()).strip("_")[:48]

def _memory_is_high_quality(key: str, value: str, source: str = "auto") -> bool:
    key = _normalize_memory_key(key)
    value = re.sub(r"\s+", " ", (value or "").strip())
    if not key or not value or len(value) < 3:
        return False
    if source != "manual":
        if key in _NOISY_MEMORY_KEYS or key.startswith(_NOISY_MEMORY_PREFIXES):
            return False
        if key.startswith("user_") and key not in {
            "user_name", "user_language", "user_location", "user_timezone",
            "user_role", "user_goal", "user_preference", "user_project",
        }:
            return False
        if _LOW_VALUE_MEMORY_RE.match(value):
            return False
        if len(value) > 260:
            return False
        lower = value.lower()
        if any(x in lower for x in (
            "whether i am", "as an ai", "i am a very strong llm",
            "found web result", "http://", "https://", "uploaded ",
            "conversion can be done", "online latex editor",
        )):
            return False
    return True

def _memory_relevance_score(memory: Dict[str, Any], query_terms: set) -> float:
    key = str(memory.get("key") or "")
    value = str(memory.get("value") or "")
    confidence = float(memory.get("confidence") or 0.0)
    if not query_terms:
        return confidence
    hay_terms = set(re.findall(r"[a-z0-9]{3,}", f"{key} {value}".lower()))
    overlap = len(query_terms & hay_terms)
    boost = 0.35 if any(x in key.lower() for x in ("name", "preference", "goal", "project", "language", "location", "role")) else 0
    return confidence + overlap * 1.25 + boost

async def _get_memories(user_id: str, limit: int = 30, query: str = "") -> List[Dict]:
    rows = await db_fetchall(
        "SELECT key,value,confidence,source,last_reinforced FROM memories WHERE user_id=? AND is_active=1 "
        "ORDER BY confidence DESC, last_reinforced DESC LIMIT ?",
        (user_id, max(limit * 4, 40)))
    rows = [r for r in rows if _memory_is_high_quality(r.get("key", ""), r.get("value", ""), r.get("source", "auto"))]
    query_terms = {
        t for t in re.findall(r"[a-z0-9]{3,}", (query or "").lower())
        if t not in _MEMORY_TERM_STOPWORDS
    }
    rows.sort(key=lambda r: _memory_relevance_score(r, query_terms), reverse=True)
    return rows[:limit]

def _format_memories(memories: List[Dict]) -> str:
    if not memories: return ""
    lines = [
        "[User memory context]",
        "These are durable user facts/preferences from JAZZ memory. Use them only when relevant. "
        "Do not invent beyond them; if current user text contradicts memory, prefer the current message.",
    ]
    for m in memories:
        lines.append(f"  • {m['key']}: {m['value']}")
    return "\n".join(lines)

async def _upsert_memory(user_id: str, key: str, value: str, source: str = "auto",
                          confidence: float = 0.8) -> None:
    key = _normalize_memory_key(key)
    value = re.sub(r"\s+", " ", (value or "").strip())[:300]
    if not _memory_is_high_quality(key, value, source):
        return
    now = _utcnow()
    row = await db_fetchone("SELECT id,reinforcement_count FROM memories WHERE user_id=? AND key=?",
                            (user_id, key))
    if row:
        new_conf = min(1.0, confidence + 0.05 * row["reinforcement_count"])
        await db_execute(
            "UPDATE memories SET value=?,confidence=?,reinforcement_count=reinforcement_count+1,"
            "last_reinforced=?,updated_at=? WHERE id=?",
            (value, new_conf, now, now, row["id"]))
    else:
        current = await db_count("SELECT COUNT(*) as c FROM memories WHERE user_id=? AND is_active=1",(user_id,))
        if current >= MAX_MEMORIES:
            oldest = await db_fetchone("SELECT id FROM memories WHERE user_id=? AND source='auto' AND is_active=1"
                                       " ORDER BY confidence ASC, last_reinforced ASC LIMIT 1",(user_id,))
            if oldest: await db_execute("UPDATE memories SET is_active=0 WHERE id=?",(oldest["id"],))
        await db_execute(
            "INSERT INTO memories(id,user_id,key,value,source,confidence,last_reinforced,"
            "reinforcement_count,is_active,created_at,updated_at) VALUES(?,?,?,?,?,?,?,1,1,?,?)",
            (_new_id(), user_id, key, value, source, confidence, now, now, now))

async def _auto_extract_memories(user_id: str, user_msg: str, ai_msg: str):
    try:
        mem_row = await db_fetchone("SELECT memory_enabled FROM users WHERE id=?",(user_id,))
        if not mem_row or not mem_row.get("memory_enabled"): return
        prompt = (
            "Extract 0-3 durable memories about the USER only. Save long-term facts, preferences, goals, "
            "identity/profile details, projects, tools they use, constraints, language preference, or recurring needs. "
            "Do NOT save the user's one-off question, the assistant answer, ratings of the AI, uploaded file names, "
            "web search results, document contents, temporary tasks, yes/no replies, or generic facts not about the user. "
            "Never use keys like user_input, user_query, user_question, ai_answer, ai_rating, processed_document. "
            "Return ONLY valid JSON array: [{\"key\":\"snake_case_key\",\"value\":\"short durable user fact\",\"confidence\":0.8}]\n"
            f"User said: {user_msg[:400]}\nAI replied: {ai_msg[:300]}\n"
            "If nothing memorable, return []."
        )
        result = await _llm_text([{"role":"user","content":prompt}],
                                  model_id="llama-3.1-8b-instant", max_tokens=300, temperature=0.1)
        raw = re.sub(r"```json|```","",result.strip()).strip()
        items = json.loads(raw)
        if not isinstance(items, list): return
        for item in items[:4]:
            if not isinstance(item, dict): continue
            key = str(item.get("key","")).strip()[:40]
            val = str(item.get("value","")).strip()[:300]
            conf = float(item.get("confidence", 0.8))
            if key and val and _memory_is_high_quality(key, val, "auto"):
                await _upsert_memory(user_id, key, val, "auto", conf)
    except Exception as e:
        logger.debug("[MEMORY] auto-extract: %s", e)

# ══════════════════════════════════════════════════════════════════════════════
# §9  CONTEXT BUILDER
# ══════════════════════════════════════════════════════════════════════════════

_SYSTEM_PROMPT = """You are JAZZ — a sharp, production-grade AI assistant with real tool access.

CAPABILITIES:
• Connected apps (Gmail, Calendar, Slack, GitHub, Notion, etc.) via /command syntax
• Execute code, run shell commands, search documents via RAG
• Memory of user preferences and past context
• Deep reasoning mode for complex questions

RESPONSE GUIDELINES:
• Be direct, precise, and genuinely helpful
• Use markdown for structure (tables, code blocks, lists)
• For connector results, present data in clean tables or structured lists
• For code, always include the language tag in code blocks"""

_ANTI_HALLUCINATION_PROMPT = """[Evidence and anti-hallucination rules]
- Do not invent facts, names, files, document contents, web sources, tool results, prices, dates, or citations.
- If the needed evidence is not in the current message, JAZZ memory, retrieved documents, attached files/images, web results, or tool output, say what is missing and ask for it.
- For uploaded documents, answer only from retrieved/processed document context. If no relevant excerpt is provided, say you cannot verify it from the processed docs yet.
- For web answers, cite only URLs provided in the live web search context. Never fabricate links.
- For images/files, mention uncertainty when extraction is partial or low-confidence.
- Prefer current user text over older memory if they conflict."""

_LOCAL_JAZZ_FAST_SYSTEM_PROMPT = (
    "You are Jazz AI, a direct assistant. Answer in the user's language. "
    "For greetings and simple questions, reply in one short natural answer. "
    "Use provided tool, web, image, or document context when it is present. "
    "Do not invent the user's name or facts. If evidence is missing, say so clearly. "
    "Use JAZZ memory only as user-specific context, and prefer the current message if it conflicts."
)

_GENERIC_PROFILE_NAMES = {"jazz", "jazzai", "jazz ai", "admin", "user", "test", "client", "account"}

def _clean_person_name(value: str) -> str:
    raw = (value or "").strip()
    if not raw:
        return ""
    raw = re.sub(r"\S+@\S+", " ", raw)
    raw = re.sub(r"[^A-Za-zÀ-ÖØ-öø-ÿ' -]+", " ", raw)
    parts = [p.strip(" '-") for p in re.split(r"\s+", raw) if p.strip(" '-")]
    parts = [p for p in parts if len(p) > 1 and p.lower() not in _GENERIC_PROFILE_NAMES]
    if not parts:
        return ""
    return " ".join(p[:1].upper() + p[1:] for p in parts[:3])[:120]

def _name_from_email(email: str) -> str:
    local = (email or "").split("@", 1)[0].split("+", 1)[0]
    local = re.sub(r"\d+", " ", local)
    local = re.sub(r"[._-]+", " ", local).strip()
    return _clean_person_name(local)

def _display_name_from_user_row(row: Optional[Dict[str, Any]]) -> str:
    if not row:
        return ""
    full_name = (row.get("full_name") or "").strip()
    cleaned_full = _clean_person_name(full_name)
    if cleaned_full:
        return cleaned_full
    email = (row.get("email") or "").strip()
    if not email or "@" not in email:
        return ""
    return _name_from_email(email)

def _is_name_recall_query(text: str) -> bool:
    t = re.sub(r"\s+", " ", (text or "").lower()).strip()
    if not t:
        return False
    phrases = (
        "what is my name", "what's my name", "who am i", "who i am",
        "do you know my name", "remember my name", "recall my name",
        "tell me my name", "try to remember", "do you remember me",
    )
    return any(p in t for p in phrases) or (("my name" in t or "my identity" in t) and any(v in t for v in ("know", "remember", "recall", "tell", "guess")))

def _name_from_memory_row(row: Dict[str, Any]) -> str:
    key = str(row.get("key") or "").lower()
    value = str(row.get("value") or "")
    if "name" in key or "identity" in key:
        return _clean_person_name(value)
    m = re.search(r"\b(?:my name is|i am|i'm|call me)\s+([A-Za-z][A-Za-zÀ-ÖØ-öø-ÿ' -]{1,60})", value, re.I)
    return _clean_person_name(m.group(1)) if m else ""

def _name_from_history_text(text: str) -> str:
    patterns = (
        r"\bmy name is\s+([A-Za-z][A-Za-zÀ-ÖØ-öø-ÿ' -]{1,60})",
        r"\bcall me\s+([A-Za-z][A-Za-zÀ-ÖØ-öø-ÿ' -]{1,60})",
        r"\bi am\s+([A-Za-z][A-Za-zÀ-ÖØ-öø-ÿ' -]{1,60})",
        r"\bi'm\s+([A-Za-z][A-Za-zÀ-ÖØ-öø-ÿ' -]{1,60})",
    )
    for pat in patterns:
        m = re.search(pat, text or "", re.I)
        if m:
            name = _clean_person_name(m.group(1))
            if name:
                return name
    return ""

async def _resolve_user_name(user_id: str, user_claims: Optional[Dict[str, Any]] = None) -> Tuple[str, str]:
    row = await db_fetchone("SELECT email,full_name FROM users WHERE id=?", (user_id,))
    name = _display_name_from_user_row(row)
    if name and (row or {}).get("full_name"):
        return name, "your profile"
    mems = await _get_memories(user_id, 40)
    for mem in mems:
        name = _name_from_memory_row(mem)
        if name:
            return name, "memory"
    rows = await db_fetchall(
        "SELECT content FROM chat_history WHERE user_id=? AND role='user' AND is_hidden=0 "
        "ORDER BY created_at DESC LIMIT 120",
        (user_id,))
    for item in rows:
        name = _name_from_history_text(item.get("content") or "")
        if name:
            return name, "previous chat"
    if row and row.get("email"):
        name = _name_from_email(row["email"])
        if name:
            return name, "your email"
    email = (user_claims or {}).get("email", "")
    name = _name_from_email(email)
    return (name, "your email") if name else ("", "")

async def _identity_recall_reply(user_id: str, message: str, user_claims: Optional[Dict[str, Any]] = None) -> Optional[str]:
    if not _is_name_recall_query(message):
        return None
    name, source = await _resolve_user_name(user_id, user_claims)
    if name:
        return f"Your name is {name}. I got it from {source}."
    return "I do not have a reliable name saved for you yet. Add it in your profile or tell me “my name is ...”, and I will remember it."

_PRIVATE_FACT_PATTERNS: Tuple[Tuple[re.Pattern, str], ...] = (
    (re.compile(r"\bpassport(?:\s+(?:number|no|id))?\b", re.I), "passport number"),
    (re.compile(r"\b(?:aadhaar|aadhar)\b", re.I), "Aadhaar number"),
    (re.compile(r"\bpan(?:\s+(?:card|number|no))?\b", re.I), "PAN number"),
    (re.compile(r"\b(?:ssn|social security)\b", re.I), "SSN"),
    (re.compile(r"\b(?:bank|account|iban|routing)\s+(?:number|no|details?)\b", re.I), "bank details"),
    (re.compile(r"\b(?:credit|debit)\s+card\b", re.I), "card details"),
    (re.compile(r"\b(?:otp|one[-\s]?time password|verification code)\b", re.I), "verification code"),
    (re.compile(r"\b(?:password|passcode|pin)\b", re.I), "password or PIN"),
    (re.compile(r"\b(?:private|personal)\s+(?:file|document|folder|data)\b", re.I), "private file or document"),
)

def _is_personal_fact_lookup(text: str) -> bool:
    t = re.sub(r"\s+", " ", (text or "").strip().lower())
    if not t or "my" not in t:
        return False
    lookup_verbs = (
        "what is", "what's", "tell me", "show me", "give me", "find my",
        "do you know", "do u know", "remember", "recall", "try to remember",
    )
    return any(v in t for v in lookup_verbs)

async def _verified_context_missing_reply(user_id: str, message: str, user_claims: Optional[Dict[str, Any]] = None) -> Optional[str]:
    """Deterministic guard for private/user-specific facts the model must not guess."""
    if not _is_personal_fact_lookup(message):
        return None
    t = re.sub(r"\s+", " ", (message or "").strip().lower())
    if _is_name_recall_query(t):
        return None
    user_row = await db_fetchone("SELECT email,full_name FROM users WHERE id=?", (user_id,))
    if re.search(r"\bmy\s+(?:email|mail|email id|mail id)\b", t):
        email = (user_row or {}).get("email") or (user_claims or {}).get("email") or ""
        return f"Your email is {email}." if email else "I do not have a verified email for this account."

    matched_labels = [label for pattern, label in _PRIVATE_FACT_PATTERNS if pattern.search(t)]
    if not matched_labels:
        return None

    mems = await _get_memories(user_id, 8, query=message)
    for mem in mems:
        haystack = f"{mem.get('key','')} {mem.get('value','')}".lower()
        if any(label.split()[0].lower() in haystack for label in matched_labels):
            return None
    label = matched_labels[0]
    return (
        f"I do not have your {label} in verified account context, JAZZ memory, "
        "attached files, processed documents, or tool output, so I cannot answer it reliably. "
        "Send that information or upload the relevant document if you want me to use it."
    )

def _authenticated_user_context(row: Optional[Dict[str, Any]]) -> str:
    name = _display_name_from_user_row(row)
    email = (row.get("email") or "").strip() if row else ""
    lines = ["[Current authenticated user identity]"]
    lines.append(f"Name: {name}" if name else "Name: not provided")
    if email:
        lines.append(f"Email: {email}")
    lines.append("Use only this identity when addressing the user. If the name is not provided, do not invent a name. This current identity overrides older memories, examples, and prior chat summaries.")
    return "\n".join(lines)

_THINKING_SYSTEM = """You are JAZZ — an advanced AI with deep reasoning capabilities.
When given complex questions, think through the problem carefully before answering.

Format your response as:
<thinking>
[Your step-by-step reasoning process here. Be thorough.]
</thinking>

<answer>
[Your final, well-reasoned response here. Use markdown for clarity.]
</answer>"""

async def _build_context(session_id: str, user_id: str, new_msg: str,
                          use_rag: bool, model_id: str,
                          web_results: Optional[List[Dict]] = None,
                          plan_mode: bool = False,
                          skill_context: str = "",
                          mcp_context: str = "") -> List[Dict]:
    local_fast = _is_local_jazz_model(model_id)
    budget = await _model_request_input_budget(model_id)
    sys_parts = (
        [_LOCAL_JAZZ_FAST_SYSTEM_PROMPT]
        if local_fast
        else [_SYSTEM_PROMPT, f"Today: {datetime.now(timezone.utc).strftime('%A, %B %d, %Y %H:%M UTC')}"]
    )
    sys_parts.append(_ANTI_HALLUCINATION_PROMPT)
    user_row = await db_fetchone("SELECT email,full_name FROM users WHERE id=?", (user_id,))
    if local_fast:
        user_name = _display_name_from_user_row(user_row)
        sys_parts.append(
            f"Current user name: {user_name}."
            if user_name else
            "Current user name is not provided; do not invent one."
        )
    else:
        sys_parts.append(_authenticated_user_context(user_row))
    mems = await _get_memories(user_id, query=new_msg)
    if mems:
        if local_fast:
            compact = [
                f"- {str(m.get('key') or '')[:48]}: {str(m.get('value') or '')[:120]}"
                for m in mems[:3]
            ]
            sys_parts.append("[User memory]\n" + "\n".join(compact))
        else:
            sys_parts.append(_format_memories(mems))

    # Active connectors context
    sc_rows = await db_fetchall(
        "SELECT connector_type,name FROM connectors WHERE user_id=? AND is_active=1", (user_id,))
    if sc_rows:
        conn_list = ", ".join(f"{r['connector_type']}({r['name']})" for r in sc_rows)
        sys_parts.append(f"[Active connectors: {conn_list}]")
    if plan_mode:
        sys_parts.append(
            "Plan mode is enabled. Start with a concise plan that names the intended "
            "steps, then carry out the answer. Keep the plan practical."
        )
    if web_results:
        lines = []
        for i, item in enumerate(web_results[:6], 1):
            title = item.get("title") or "Untitled"
            url = item.get("url") or ""
            snippet = item.get("snippet") or ""
            lines.append(f"{i}. {title}\nURL: {url}\nSnippet: {snippet}")
        sys_parts.append("[Live web search results]\n" + "\n\n".join(lines) + "\n\nCite the URLs you use.")
    if skill_context:
        sys_parts.append(skill_context)
    if mcp_context:
        sys_parts.append(mcp_context)

    rag_context = ""
    use_rag_now = bool(use_rag) and (
        not local_fast
        or bool(skill_context)
        or _is_document_access_query(new_msg or "")
    )
    if use_rag_now:
        processed_docs = await _processed_document_rows(user_id)
        sys_parts.append(_format_rag_document_status(processed_docs))
        chunks = await _rag_search(user_id, (new_msg or "")[:4000])
        if not chunks and processed_docs and _is_document_access_query(new_msg or ""):
            chunks = await _rag_sqlite_search(user_id, new_msg or "", TOP_K_RETRIEVAL, recent=True)
        if chunks:
            rag_str = "\n\n---\n\n".join(chunks[:RERANK_TOP_K])
            rag_context = (
                "\n[Retrieved document context]\n"
                "Use these processed JAZZ document excerpts when answering. "
                "When relevant, mention the source filename.\n"
                f"{rag_str}"
            )

    # Rolling summary
    summary = await db_fetchone(
        "SELECT summary_text FROM summaries WHERE session_id=? ORDER BY created_at DESC LIMIT 1",
        (session_id,))
    if summary and not local_fast:
        sys_parts.append(f"[Conversation Summary]\n{summary['summary_text']}")

    system_content = "\n\n".join(sys_parts) + rag_context
    msgs = [{"role":"system","content":system_content}]
    budget -= _count_tokens(system_content)

    safe_new_msg = new_msg
    new_msg_compacted = False
    if new_msg:
        user_budget = max(256, int(max(512, budget) * 0.78))
        safe_new_msg, new_msg_compacted = _fit_text_to_token_budget(new_msg, user_budget)
        if new_msg_compacted:
            notice = (
                "\n\n[Long input handling]\nThe current user message was automatically "
                "compacted with beginning and ending sections preserved because it was "
                "larger than the selected model context."
            )
            msgs[0]["content"] += notice
            budget -= _count_tokens(notice)
        budget -= _count_tokens(safe_new_msg)

    history_limit = 8 if local_fast else 60
    history = await db_fetchall(
        "SELECT role,content FROM chat_history WHERE session_id=? AND is_hidden=0 "
        "ORDER BY created_at DESC LIMIT ?", (session_id, history_limit))
    window = []
    for h in history:
        tok = _count_tokens(h["content"])
        if budget - tok < 200: break
        window.append({"role":h["role"],"content":h["content"]}); budget -= tok
    window.reverse()
    msgs.extend(window)
    if safe_new_msg:
        msgs.append({"role":"user","content":safe_new_msg})
    return msgs

async def _context_window_status(session_id: str, user_id: str, model_id: Optional[str] = None) -> Dict[str, Any]:
    session = await db_fetchone(
        "SELECT id,model_id,turn_count FROM chat_sessions WHERE id=? AND user_id=?",
        (session_id, user_id))
    if not session:
        raise HTTPException(404, "Session not found")
    mid = model_id or session["model_id"] or "llama-3.3-70b-versatile"
    context_limit = await _model_context_limit(mid)
    request_budget = await _context_input_budget(mid)
    messages = await _build_context(session_id, user_id, "", False, mid)
    used_tokens = sum(_count_tokens(m.get("content","") if isinstance(m.get("content"), str) else "")
                      for m in messages)
    summary_rows = await db_fetchall(
        "SELECT summary_text,tokens_saved,created_at FROM summaries WHERE session_id=? ORDER BY created_at",
        (session_id,))
    visible_rows = await db_fetchall(
        "SELECT content FROM chat_history WHERE session_id=? AND is_hidden=0",
        (session_id,))
    hidden_messages = await db_count(
        "SELECT COUNT(*) as c FROM chat_history WHERE session_id=? AND is_hidden=1",
        (session_id,))
    visible_tokens = sum(_count_tokens(r["content"]) for r in visible_rows)
    summary_tokens = sum(_count_tokens(r["summary_text"]) for r in summary_rows)
    pct = round(min(100.0, (used_tokens / max(1, context_limit)) * 100), 1)
    return {
        "session_id": session_id,
        "model_id": mid,
        "context_limit_tokens": context_limit,
        "request_budget_tokens": request_budget,
        "used_tokens": used_tokens,
        "remaining_tokens": max(0, context_limit - used_tokens),
        "percent_used": pct,
        "percent_left": round(max(0.0, 100.0 - pct), 1),
        "visible_tokens": visible_tokens,
        "summary_tokens": summary_tokens,
        "visible_messages": len(visible_rows),
        "hidden_messages": hidden_messages,
        "summary_count": len(summary_rows),
        "tokens_saved": sum(int(r["tokens_saved"] or 0) for r in summary_rows),
        "auto_compact_enabled": True,
        "auto_compact_threshold_percent": round(CONTEXT_AUTO_COMPACT_PCT * 100, 1),
        "compacted": bool(hidden_messages or summary_rows),
    }

async def _compact_session_context(session_id: str, user_id: str, model_id: str,
                                   reason: str = "auto") -> Dict[str, Any]:
    history = await db_fetchall(
        "SELECT id,role,content,created_at FROM chat_history WHERE session_id=? AND is_hidden=0 "
        "ORDER BY created_at ASC",
        (session_id,))
    keep_recent = max(4, CONTEXT_COMPACT_KEEP_RECENT)
    if len(history) <= keep_recent + 4:
        status = await _context_window_status(session_id, user_id, model_id)
        status["compact_result"] = "not_needed"
        return status

    compact_rows = history[:-keep_recent]
    prior = await db_fetchone(
        "SELECT summary_text FROM summaries WHERE session_id=? ORDER BY created_at DESC LIMIT 1",
        (session_id,))
    prior_text = prior["summary_text"] if prior else ""
    conv = "\n".join(
        f"{r['role'].upper()}: {str(r['content'])[:1200]}" for r in compact_rows
    )
    prompt = (
        "Create an updated rolling summary for this chat. Preserve facts, decisions, "
        "open tasks, user preferences, connector results, filenames, and code choices. "
        "Do not include filler.\n\n"
        f"Existing summary:\n{prior_text or '(none)'}\n\n"
        f"Messages to compact:\n{conv}"
    )
    summary = await _llm_text(
        [{"role":"system","content":"You maintain compact, faithful chat memory for long-context conversations."},
         {"role":"user","content":prompt}],
        model_id="llama-3.1-8b-instant", max_tokens=700, temperature=0.2)
    now = _utcnow()
    tokens_before = sum(_count_tokens(r["content"]) for r in compact_rows)
    tokens_saved = max(0, tokens_before - _count_tokens(summary))
    await db_execute(
        "INSERT INTO summaries(id,session_id,user_id,summary_text,covers_from_at,covers_to_at,"
        "turn_count,tokens_saved,model_used,created_at) VALUES(?,?,?,?,?,?,?,?,?,?)",
        (_new_id(), session_id, user_id, summary, compact_rows[0]["created_at"],
         compact_rows[-1]["created_at"], len(compact_rows), tokens_saved, model_id, now))
    ids = tuple(r["id"] for r in compact_rows)
    if ids:
        placeholders = ",".join("?" * len(ids))
        await db_execute(f"UPDATE chat_history SET is_hidden=1 WHERE id IN ({placeholders})", ids)
    status = await _context_window_status(session_id, user_id, model_id)
    status["compact_result"] = "compacted"
    status["compact_reason"] = reason
    return status

def _safe_json_loads(raw: Any, default: Any) -> Any:
    if raw is None:
        return default
    if isinstance(raw, (dict, list)):
        return raw
    try:
        return json.loads(raw)
    except Exception:
        return default

def _coerce_str_list(raw: Any) -> List[str]:
    """Accept JSON arrays, comma strings, sets/tuples, and single values as string lists."""
    if raw is None:
        return []
    if isinstance(raw, str):
        s = raw.strip()
        if not s:
            return []
        parsed = _safe_json_loads(s, None)
        if isinstance(parsed, list):
            raw = parsed
        elif isinstance(parsed, dict):
            raw = list(parsed.values())
        else:
            raw = s.split(",")
    elif isinstance(raw, dict):
        raw = list(raw.values())
    elif isinstance(raw, (set, tuple)):
        raw = list(raw)
    if isinstance(raw, list):
        return [str(x).strip() for x in raw if str(x).strip()]
    return [str(raw).strip()] if str(raw).strip() else []

def _coerce_tool_options(raw: Any) -> Dict[str, Any]:
    if raw is None or raw == "":
        return {}
    if isinstance(raw, str):
        raw = _safe_json_loads(raw, {})
    if isinstance(raw, list):
        raw = {"mcp_tools": raw}
    if not isinstance(raw, dict):
        return {}
    out = dict(raw)
    for key in ("connector_ids", "skill_ids"):
        out[key] = _coerce_str_list(out.get(key))
    mcp_tools = out.get("mcp_tools", [])
    if isinstance(mcp_tools, str):
        mcp_tools = _safe_json_loads(mcp_tools, [])
    if isinstance(mcp_tools, dict):
        mcp_tools = [mcp_tools]
    out["mcp_tools"] = mcp_tools if isinstance(mcp_tools, list) else []
    return out

def _normalise_skill_row(row: Dict) -> Dict:
    out = dict(row)
    out["activation_keywords"] = _safe_json_loads(out.get("activation_keywords"), [])
    out["resources"] = _safe_json_loads(out.get("resources_json"), [])
    out["script_metadata"] = _safe_json_loads(out.get("script_metadata_json"), {})
    out["is_enabled"] = bool(out.get("is_enabled"))
    out.pop("resources_json", None)
    out.pop("script_metadata_json", None)
    return out

def _markdown_section(markdown: str, title: str) -> str:
    pat = re.compile(rf"^##+\s+{re.escape(title)}\s*$([\s\S]*?)(?=^##+\s+|\Z)", re.I | re.M)
    m = pat.search(markdown or "")
    return m.group(1).strip() if m else ""

def _parse_skill_markdown(markdown: str, name_override: str = "") -> Dict[str, Any]:
    text = (markdown or "").strip()
    if not text:
        raise ValueError("markdown is required")
    front: Dict[str, str] = {}
    if text.startswith("---"):
        end = text.find("\n---", 3)
        if end > 0:
            raw_front = text[3:end].strip()
            text = text[end + 4:].strip()
            for line in raw_front.splitlines():
                if ":" in line:
                    k, v = line.split(":", 1)
                    front[k.strip().lower()] = v.strip().strip('"\'')
    h1 = re.search(r"^#\s+(.+)$", text, flags=re.M)
    name = (name_override or front.get("name") or (h1.group(1).strip() if h1 else "") or "Imported Skill").strip()
    description = (front.get("description") or _markdown_section(text, "Description") or "").strip()
    kw_raw = front.get("activation_keywords") or front.get("keywords") or _markdown_section(text, "Activation Keywords") or _markdown_section(text, "Keywords")
    keywords = []
    if kw_raw:
        cleaned = re.sub(r"^[\s*-]+", "", kw_raw, flags=re.M)
        parts = re.split(r"[,\n]", cleaned)
        keywords = [p.strip(" -\t") for p in parts if p.strip(" -\t")]
    instructions = _markdown_section(text, "Instructions") or _markdown_section(text, "Skill Instructions")
    if not instructions:
        instructions = re.sub(r"^#\s+.+$", "", text, count=1, flags=re.M).strip()
    if len(instructions.strip()) < 5:
        raise ValueError("SKILL.md needs instructions")
    return {
        "name": name[:100],
        "description": description[:1000],
        "activation_keywords": keywords[:50],
        "instructions": instructions,
        "resources": [],
        "script_metadata": {"source":"SKILL.md"},
    }

_SKILL_MD_TEMPLATE = """---
name: My Skill
description: What this skill helps JAZZ do.
activation_keywords: keyword one, keyword two
---

# My Skill

## Description
Short explanation of when this skill should be used.

## Activation Keywords
- keyword one
- keyword two

## Instructions
Give JAZZ clear, reusable operating instructions. Include constraints, workflow steps,
examples, and what output format should be used.
"""

async def _select_skills_for_message(message: str, explicit_ids: List[str] = None) -> Tuple[List[Dict], List[str]]:
    explicit_ids = [str(x) for x in (explicit_ids or []) if str(x).strip()]
    skills: List[Dict] = []
    matched_by: List[str] = []
    if explicit_ids:
        placeholders = ",".join("?" for _ in explicit_ids)
        rows = await db_fetchall(
            f"SELECT * FROM skills WHERE is_enabled=1 AND id IN ({placeholders}) ORDER BY name",
            tuple(explicit_ids))
        skills.extend([_normalise_skill_row(r) for r in rows])
        matched_by.extend(["selected"] * len(rows))

    text = (message or "").lower()
    rows = await db_fetchall("SELECT * FROM skills WHERE is_enabled=1 ORDER BY name")
    seen = {s["id"] for s in skills}
    for r in rows:
        if r["id"] in seen:
            continue
        skill = _normalise_skill_row(r)
        candidates = [skill.get("name",""), skill.get("description","")] + skill.get("activation_keywords", [])
        if any(c and str(c).lower() in text for c in candidates):
            skills.append(skill)
            matched_by.append("auto")
            seen.add(skill["id"])
        if len(skills) >= 5:
            break
    return skills[:5], matched_by[:5]

def _format_skill_context(skills: List[Dict]) -> str:
    if not skills:
        return ""
    blocks = []
    for s in skills:
        resources = s.get("resources") or []
        res_text = ""
        if resources:
            parts = []
            for r in resources[:5]:
                parts.append(f"- {r.get('name','resource')}: {str(r.get('content',''))[:1200]}")
            res_text = "\nResources:\n" + "\n".join(parts)
        blocks.append(
            f"[Skill: {s.get('name','Unnamed')}]\n"
            f"Description: {s.get('description','')}\n"
            f"Instructions:\n{s.get('instructions','')[:4000]}"
            f"{res_text}"
        )
    return "[Active Skills]\n" + "\n\n".join(blocks)

async def _log_skill_usage(skills: List[Dict], matched_by: List[str], user_id: str,
                           session_id: str = None, message_id: str = None) -> None:
    now = _utcnow()
    for i, s in enumerate(skills):
        await db_execute(
            "INSERT INTO skill_usage_logs(id,skill_id,user_id,session_id,message_id,matched_by,created_at)"
            " VALUES(?,?,?,?,?,?,?)",
            (_new_id(), s["id"], user_id, session_id, message_id, matched_by[i] if i < len(matched_by) else "", now))

def _chat_tool_options(req: ChatReq) -> Dict[str, Any]:
    opts = _coerce_tool_options(getattr(req, "tools", {}))
    if not opts.get("skill_ids"):
        opts["skill_ids"] = _coerce_str_list(getattr(req, "skill_ids", []))
    return opts

def _format_mcp_context(mcp_tools: Any) -> str:
    if not isinstance(mcp_tools, list) or not mcp_tools:
        return ""
    lines = []
    for t in mcp_tools[:12]:
        if isinstance(t, dict):
            server = t.get("server_name") or t.get("server_id") or "MCP"
            name = t.get("name") or t.get("tool_name") or ""
            if name:
                lines.append(f"- {server}: {name}")
    return "[Selected MCP tools available to the user]\n" + "\n".join(lines) if lines else ""

async def _format_agent_tool_context(user_id: str, tool_opts: Dict[str, Any]) -> str:
    selected = set(_coerce_str_list((tool_opts or {}).get("connector_ids")))
    rows = await db_fetchall(
        "SELECT connector_type,name FROM connectors WHERE user_id=? AND is_active=1 "
        "UNION SELECT connector_type,app_name as name FROM smart_connectors WHERE user_id=? AND status='active'",
        (user_id, user_id))
    connectors = []
    seen = set()
    for r in rows:
        ctype = r["connector_type"]
        if selected and ctype not in selected:
            continue
        if ctype in seen:
            continue
        seen.add(ctype)
        connectors.append(f"- {ctype}: connected as {r.get('name') or ctype}")
    if "n8n_local" in selected or not selected:
        connectors.append("- n8n_local: local Docker n8n. Use connector_action with action status/open to get portal and webhook URLs.")
    mcp = _format_mcp_context((tool_opts or {}).get("mcp_tools", []))
    return (
        "[Tool Agent Context]\n"
        "You may call connector_action for connected apps when the user asks you to use tools. "
        "Use exact connector ids and compact JSON params. Prefer read/list/search actions before write/send actions unless the user clearly asks to write.\n"
        + ("Connected connectors:\n" + "\n".join(connectors) if connectors else "No connected connectors were selected.")
        + ("\n\n" + mcp if mcp else "")
    )

async def _maybe_summarize(session_id: str, user_id: str, model_id: str):
    count = await db_count("SELECT COUNT(*) as c FROM chat_history WHERE session_id=? AND role='user'",(session_id,))
    try:
        status = await _context_window_status(session_id, user_id, model_id)
        should_compact = (
            (count >= SUMMARY_TRIGGER and count % SUMMARY_TRIGGER == 0) or
            (status["percent_used"] >= CONTEXT_AUTO_COMPACT_PCT * 100)
        )
        if should_compact:
            await _compact_session_context(session_id, user_id, model_id, reason="auto")
    except Exception as e:
        logger.warning("[CTX] summarize: %s", e)

async def _auto_title(session_id: str, user_msg: str, ai_msg: str = "") -> None:
    try:
        ctx = f"User: {user_msg[:300]}"
        if ai_msg: ctx += f"\nAI: {ai_msg[:200]}"
        title = await _llm_text(
            [{"role":"user","content":f"Generate a short 4-6 word chat title:\n{ctx}\nReturn ONLY the title, no quotes."}],
            model_id="llama-3.1-8b-instant", max_tokens=20, temperature=0.3)
        title = title.strip().strip('"\'').strip()[:80]
        if title: await db_execute("UPDATE chat_sessions SET title=?,updated_at=? WHERE id=?",(title,_utcnow(),session_id))
    except Exception as e:
        logger.debug("[TITLE] %s", e)

# ══════════════════════════════════════════════════════════════════════════════
# §10  OAUTH 2.0 — FULL PKCE FLOW (FIXED — All 25+ connectors)
# ══════════════════════════════════════════════════════════════════════════════

_pending_oauth: Dict[str, Dict] = {}
_completed_oauth: Dict[str, Dict] = {}
_GOOGLE_CONNECTOR_TYPES = {"gmail", "google_calendar", "google_drive", "google_sheets", "bigquery", "google_meet"}

def _app_base_url(request: Optional[Request] = None) -> str:
    """Return the public base URL used for OAuth redirects."""
    if APP_BASE_URL:
        return APP_BASE_URL.strip().rstrip("/")
    if request is not None:
        return str(request.base_url).rstrip("/")
    return "https://imperceptibly-hymnlike-leesa.ngrok-free.dev"

def _auth_link_base_url(request: Optional[Request] = None) -> str:
    """Return the public website URL used in verification/password emails."""
    configured = (os.getenv("AUTH_LINK_BASE_URL") or os.getenv("PUBLIC_APP_URL") or AUTH_LINK_BASE_URL or "").strip()
    if configured:
        return configured.rstrip("/")
    return _app_base_url(request)

def _oauth_callback_connector_id(app_id: str) -> str:
    # Google OAuth clients usually share one redirect URI for all Google APIs.
    if app_id in _GOOGLE_CONNECTOR_TYPES:
        return os.getenv("GOOGLE_OAUTH_CALLBACK_CONNECTOR", "gmail").strip() or "gmail"
    return app_id

def _oauth_redirect_uri(app_id: str, request: Optional[Request] = None) -> str:
    return f"{_app_base_url(request)}/connectors/oauth/{_oauth_callback_connector_id(app_id)}/callback"

def _oauth_redirect_uris() -> Dict[str, str]:
    providers = _oauth_providers()
    uris = {app_id: _oauth_redirect_uri(app_id) for app_id in providers.keys()}
    uris["google_shared"] = _oauth_redirect_uri("gmail")
    return uris

def _oauth_providers() -> Dict[str, Dict]:
    """Build provider config resolved at request time so env vars are ready."""
    return {
        "gmail": {
            "auth_url":    "https://accounts.google.com/o/oauth2/v2/auth",
            "token_url":   "https://oauth2.googleapis.com/token",
            "client_id":   GOOGLE_CLIENT_ID,
            "client_secret": GOOGLE_CLIENT_SECRET,
            "scope":       "https://www.googleapis.com/auth/gmail.modify",
            "extra_auth_params": {"access_type":"offline","prompt":"consent"},
        },
        "google_calendar": {
            "auth_url":    "https://accounts.google.com/o/oauth2/v2/auth",
            "token_url":   "https://oauth2.googleapis.com/token",
            "client_id":   GOOGLE_CLIENT_ID,
            "client_secret": GOOGLE_CLIENT_SECRET,
            "scope":       "https://www.googleapis.com/auth/calendar",
            "extra_auth_params": {"access_type":"offline","prompt":"consent"},
        },
        "google_drive": {
            "auth_url":    "https://accounts.google.com/o/oauth2/v2/auth",
            "token_url":   "https://oauth2.googleapis.com/token",
            "client_id":   GOOGLE_CLIENT_ID,
            "client_secret": GOOGLE_CLIENT_SECRET,
            "scope":       "https://www.googleapis.com/auth/drive.readonly",
            "extra_auth_params": {"access_type":"offline","prompt":"consent"},
        },
        "google_sheets": {
            "auth_url":    "https://accounts.google.com/o/oauth2/v2/auth",
            "token_url":   "https://oauth2.googleapis.com/token",
            "client_id":   GOOGLE_CLIENT_ID,
            "client_secret": GOOGLE_CLIENT_SECRET,
            "scope":       "https://www.googleapis.com/auth/spreadsheets",
            "extra_auth_params": {"access_type":"offline","prompt":"consent"},
        },
        "bigquery": {
            "auth_url":    "https://accounts.google.com/o/oauth2/v2/auth",
            "token_url":   "https://oauth2.googleapis.com/token",
            "client_id":   GOOGLE_CLIENT_ID,
            "client_secret": GOOGLE_CLIENT_SECRET,
            "scope":       "https://www.googleapis.com/auth/bigquery",
            "extra_auth_params": {"access_type":"offline","prompt":"consent"},
        },
        "google_meet": {
            "auth_url":    "https://accounts.google.com/o/oauth2/v2/auth",
            "token_url":   "https://oauth2.googleapis.com/token",
            "client_id":   GOOGLE_CLIENT_ID,
            "client_secret": GOOGLE_CLIENT_SECRET,
            "scope":       "https://www.googleapis.com/auth/calendar",
            "extra_auth_params": {"access_type":"offline","prompt":"consent"},
        },
        "github": {
            "auth_url":    "https://github.com/login/oauth/authorize",
            "token_url":   "https://github.com/login/oauth/access_token",
            "client_id":   GITHUB_CLIENT_ID,
            "client_secret": GITHUB_CLIENT_SECRET,
            "scope":       "repo,issues,read:user",
            "token_response_type": "form",
        },
        "gitlab": {
            "auth_url":    "https://gitlab.com/oauth/authorize",
            "token_url":   "https://gitlab.com/oauth/token",
            "client_id":   GITLAB_CLIENT_ID,
            "client_secret": GITLAB_CLIENT_SECRET,
            "scope":       "api",
        },
        "notion": {
            "auth_url":    "https://api.notion.com/v1/oauth/authorize",
            "token_url":   "https://api.notion.com/v1/oauth/token",
            "client_id":   NOTION_CLIENT_ID,
            "client_secret": NOTION_CLIENT_SECRET,
            "scope":       "",
            "extra_auth_params": {"owner":"user"},
            "token_auth_type": "basic",
        },
        "slack": {
            "auth_url":    "https://slack.com/oauth/v2/authorize",
            "token_url":   "https://slack.com/api/oauth.v2.access",
            "client_id":   SLACK_CLIENT_ID,
            "client_secret": SLACK_CLIENT_SECRET,
            "scope":       "channels:read,chat:write,im:history,users:read",
        },
        "jira": {
            "auth_url":    "https://auth.atlassian.com/authorize",
            "token_url":   "https://auth.atlassian.com/oauth/token",
            "client_id":   ATLASSIAN_CLIENT_ID,
            "client_secret": ATLASSIAN_CLIENT_SECRET,
            "scope":       "read:jira-work write:jira-work offline_access",
            "extra_auth_params": {"audience":"api.atlassian.com","prompt":"consent"},
        },
        "linear": {
            "auth_url":    "https://linear.app/oauth/authorize",
            "token_url":   "https://api.linear.app/oauth/token",
            "client_id":   LINEAR_CLIENT_ID,
            "client_secret": LINEAR_CLIENT_SECRET,
            "scope":       "read,write",
        },
        "asana": {
            "auth_url":    "https://app.asana.com/-/oauth_authorize",
            "token_url":   "https://app.asana.com/-/oauth_token",
            "client_id":   ASANA_CLIENT_ID,
            "client_secret": ASANA_CLIENT_SECRET,
            "scope":       "default",
        },
        "hubspot": {
            "auth_url":    "https://app.hubspot.com/oauth/authorize",
            "token_url":   "https://api.hubapi.com/oauth/v1/token",
            "client_id":   HUBSPOT_CLIENT_ID,
            "client_secret": HUBSPOT_CLIENT_SECRET,
            "scope":       "contacts crm.objects.deals.read",
        },
        "airtable": {
            "auth_url":    "https://airtable.com/oauth2/v1/authorize",
            "token_url":   "https://airtable.com/oauth2/v1/token",
            "client_id":   AIRTABLE_CLIENT_ID,
            "client_secret": AIRTABLE_CLIENT_SECRET,
            "scope":       "data.records:read data.records:write schema.bases:read",
            "pkce":        True,
        },
        "dropbox": {
            "auth_url":    "https://www.dropbox.com/oauth2/authorize",
            "token_url":   "https://api.dropboxapi.com/oauth2/token",
            "client_id":   DROPBOX_CLIENT_ID,
            "client_secret": DROPBOX_CLIENT_SECRET,
            "scope":       "files.content.read files.content.write",
            "extra_auth_params": {"token_access_type":"offline"},
        },
        "zoom": {
            "auth_url":    "https://zoom.us/oauth/authorize",
            "token_url":   "https://zoom.us/oauth/token",
            "client_id":   ZOOM_CLIENT_ID,
            "client_secret": ZOOM_CLIENT_SECRET,
            "scope":       "meeting:read:admin meeting:write:admin",
        },
        "outlook": {
            "auth_url":    "https://login.microsoftonline.com/common/oauth2/v2.0/authorize",
            "token_url":   "https://login.microsoftonline.com/common/oauth2/v2.0/token",
            "client_id":   MICROSOFT_CLIENT_ID,
            "client_secret": MICROSOFT_CLIENT_SECRET,
            "scope":       "https://graph.microsoft.com/Mail.ReadWrite offline_access",
            "extra_auth_params": {"response_mode":"query"},
        },
        "onedrive": {
            "auth_url":    "https://login.microsoftonline.com/common/oauth2/v2.0/authorize",
            "token_url":   "https://login.microsoftonline.com/common/oauth2/v2.0/token",
            "client_id":   MICROSOFT_CLIENT_ID,
            "client_secret": MICROSOFT_CLIENT_SECRET,
            "scope":       "https://graph.microsoft.com/Files.ReadWrite offline_access",
            "extra_auth_params": {"response_mode":"query"},
        },
        "microsoft_teams": {
            "auth_url":    "https://login.microsoftonline.com/common/oauth2/v2.0/authorize",
            "token_url":   "https://login.microsoftonline.com/common/oauth2/v2.0/token",
            "client_id":   MICROSOFT_CLIENT_ID,
            "client_secret": MICROSOFT_CLIENT_SECRET,
            "scope":       "https://graph.microsoft.com/Chat.ReadWrite offline_access",
            "extra_auth_params": {"response_mode":"query"},
        },
        "power_bi": {
            "auth_url":    "https://login.microsoftonline.com/common/oauth2/v2.0/authorize",
            "token_url":   "https://login.microsoftonline.com/common/oauth2/v2.0/token",
            "client_id":   MICROSOFT_CLIENT_ID,
            "client_secret": MICROSOFT_CLIENT_SECRET,
            "scope":       "https://analysis.windows.net/powerbi/api/.default offline_access",
            "extra_auth_params": {"response_mode":"query"},
        },
        "salesforce": {
            "auth_url":    "https://login.salesforce.com/services/oauth2/authorize",
            "token_url":   "https://login.salesforce.com/services/oauth2/token",
            "client_id":   SALESFORCE_CLIENT_ID,
            "client_secret": SALESFORCE_CLIENT_SECRET,
            "scope":       "api refresh_token",
        },
        "linkedin": {
            "auth_url":    "https://www.linkedin.com/oauth/v2/authorization",
            "token_url":   "https://www.linkedin.com/oauth/v2/accessToken",
            "client_id":   LINKEDIN_CLIENT_ID,
            "client_secret": LINKEDIN_CLIENT_SECRET,
            "scope":       "r_liteprofile r_emailaddress",
        },
        "box": {
            "auth_url":    "https://account.box.com/api/oauth2/authorize",
            "token_url":   "https://api.box.com/oauth2/token",
            "client_id":   BOX_CLIENT_ID,
            "client_secret": BOX_CLIENT_SECRET,
            "scope":       "",
        },
        "clickup": {
            "auth_url":    "https://app.clickup.com/api",
            "token_url":   "https://api.clickup.com/api/v2/oauth/token",
            "client_id":   CLICKUP_CLIENT_ID,
            "client_secret": CLICKUP_CLIENT_SECRET,
            "scope":       "",
        },
        "discord": {
            "auth_url":    "https://discord.com/api/oauth2/authorize",
            "token_url":   "https://discord.com/api/oauth2/token",
            "client_id":   DISCORD_CLIENT_ID,
            "client_secret": DISCORD_CLIENT_SECRET,
            "scope":       "identify guilds",
        },
    }

def _pkce_pair() -> Tuple[str, str]:
    verifier  = base64.urlsafe_b64encode(secrets.token_bytes(32)).rstrip(b"=").decode()
    challenge = base64.urlsafe_b64encode(
        hashlib.sha256(verifier.encode()).digest()).rstrip(b"=").decode()
    return verifier, challenge

async def _oauth_post_form(url: str, data: Dict,
                            basic_auth: Optional[Tuple[str,str]] = None) -> Dict:
    """POST form-encoded data, returns JSON dict (handles GitHub form response too)."""
    headers = {"Accept":"application/json","Content-Type":"application/x-www-form-urlencoded"}
    encoded = urllib.parse.urlencode(data).encode()
    req = urllib.request.Request(url, data=encoded, method="POST")
    req.add_header("Accept","application/json")
    req.add_header("Content-Type","application/x-www-form-urlencoded")
    if basic_auth:
        cred = base64.b64encode(f"{basic_auth[0]}:{basic_auth[1]}".encode()).decode()
        req.add_header("Authorization", f"Basic {cred}")
    if HAS_AIOHTTP:
        auth_obj = aiohttp.BasicAuth(*basic_auth) if basic_auth else None
        async with aiohttp.ClientSession() as sess:
            async with sess.post(url, data=data, headers={"Accept":"application/json"},
                                  auth=auth_obj) as resp:
                ct = resp.headers.get("Content-Type","")
                raw = await resp.text()
                if "json" in ct: return json.loads(raw)
                return dict(urllib.parse.parse_qsl(raw))
    else:
        try:
            with urllib.request.urlopen(req, timeout=20) as r:
                raw = r.read().decode()
                ct  = r.headers.get("Content-Type","")
                if "json" in ct: return json.loads(raw)
                return dict(urllib.parse.parse_qsl(raw))
        except urllib.error.HTTPError as e:
            raw = e.read().decode("utf-8","replace")
            try: return json.loads(raw)
            except Exception: return {"error":"http_error","error_description":raw[:200]}

_OAUTH_SUCCESS_HTML = """<!DOCTYPE html><html><head>
<style>*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#07070f;color:#e4e4f4;font-family:system-ui,sans-serif;
  display:flex;align-items:center;justify-content:center;min-height:100vh}}
.card{{background:#0c0c1a;border:1px solid #252540;border-radius:16px;
  padding:32px 40px;text-align:center;max-width:400px;width:90%}}
h2{{color:#34d399;font-size:22px;margin-bottom:10px}}
p{{color:#9090b0;font-size:14px;line-height:1.5}}</style></head><body>
<div class="card"><h2>✅ {name} Connected!</h2>
<p>Authorization successful. You can close this window and return to JAZZ.</p>
<script>setTimeout(()=>window.close(),2000)</script></div></body></html>"""

_OAUTH_ERROR_HTML = """<!DOCTYPE html><html><head>
<style>*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#07070f;color:#e4e4f4;font-family:system-ui,sans-serif;
  display:flex;align-items:center;justify-content:center;min-height:100vh}}
.card{{background:#0c0c1a;border:1px solid #452020;border-radius:16px;
  padding:32px 40px;text-align:center;max-width:450px;width:90%}}
h2{{color:#f87171;font-size:20px;margin-bottom:10px}}
p{{color:#9090b0;font-size:14px;line-height:1.5;margin-bottom:8px}}
code{{background:#111;padding:4px 8px;border-radius:4px;font-size:12px;display:block;
  margin-top:10px;word-break:break-all;color:#fbbf24}}</style></head><body>
<div class="card"><h2>❌ Connection Failed</h2>
<p>{message}</p><code>{detail}</code>
<script>setTimeout(()=>window.close(),6000)</script></div></body></html>"""

_OAUTH_SETUP_HTML = """<!DOCTYPE html><html><head>
<style>*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#07070f;color:#e4e4f4;font-family:system-ui,sans-serif;
  display:flex;align-items:center;justify-content:center;min-height:100vh;padding:16px}}
.card{{background:#0c0c1a;border:1px solid #252540;border-radius:16px;
  padding:32px 40px;max-width:520px;width:100%}}
h2{{color:#fbbf24;font-size:20px;margin-bottom:12px}}
p{{color:#9090b0;font-size:14px;line-height:1.7;margin-bottom:10px}}
code{{background:#111;padding:3px 7px;border-radius:4px;font-size:12px;color:#7c6ff7;word-break:break-all}}
.steps{{margin:16px 0;padding:16px;background:#0a0a16;border-radius:8px;border:1px solid #1a1a30}}
.step{{font-size:13px;color:#9090b0;margin-bottom:8px}}</style></head><body>
<div class="card"><h2>⚙️ {name} OAuth Not Configured</h2>
<p>To enable <b>{name}</b> OAuth, set these environment variables and restart JAZZ:</p>
<div class="steps">
<div class="step">• <code>{id_var}</code> — your {name} App Client ID</div>
<div class="step">• <code>{secret_var}</code> — your {name} App Client Secret</div>
<div class="step">• <code>APP_BASE_URL</code> — your server base URL (e.g. https://yourdomain.com)</div>
</div>
<p>Also add this <b>Redirect URI</b> in your {name} developer console:</p>
<code>{redirect_uri}</code>
</div></body></html>"""

def _user_id_from_token_str(token: str) -> Optional[str]:
    payload = _decode_jwt(token)
    return payload.get("sub") if payload else None

@asynccontextmanager
async def lifespan(app: FastAPI):
    await _init_db_connection()
    await _conn().executescript(_SCHEMA)
    await _conn().commit()
    await _ensure_premium_subscription_tier()
    _init_chroma()
    await _seed_slash_commands()
    await _seed_platform_connectors()
    await _seed_sidebar_features()

    # Create admin user
    admin = await db_fetchone("SELECT id FROM users WHERE email=?", (ADMIN_EMAIL,))
    if not admin:
        now = _utcnow()
        await db_execute(
            "INSERT INTO users(id,email,password_hash,full_name,role,subscription,"
            "is_active,is_verified,created_at,updated_at) VALUES(?,?,?,'Admin','admin','enterprise',1,1,?,?)",
            (_new_id(), ADMIN_EMAIL, _hash_pw(ADMIN_PASSWORD), now, now))
        logger.info("[INIT] Admin user created: %s", ADMIN_EMAIL)
        admin = await db_fetchone("SELECT id FROM users WHERE email=?", (ADMIN_EMAIL,))
    if admin:
        await _seed_default_ai_models(admin["id"])
        if os.getenv("HF_ROUTER_MODEL_SEED", "1").lower() not in ("0", "false", "no"):
            await _seed_hf_router_models(admin["id"])
        await _ensure_image_to_text_model(admin["id"])
        await _seed_model_access()

    # Schedule existing jobs
    jobs = await db_fetchall("SELECT * FROM agent_jobs WHERE enabled=1")
    for j in jobs: _schedule_job(j)
    _scheduler.start()
    logger.info("[INIT] JAZZ AI v14 started — %d agent jobs scheduled", len(jobs))
    yield
    _scheduler.shutdown(wait=False)
    if _db: await _db.close()

# ══════════════════════════════════════════════════════════════════════════════
# §11  CONNECTOR ACTIONS
# ══════════════════════════════════════════════════════════════════════════════

async def _get_connector_creds_record(user_id: str, connector_type: str) -> Tuple[Optional[Dict], str, str]:
    smart = await db_fetchone(
        "SELECT id,encrypted_conn_data FROM smart_connectors "
        "WHERE user_id=? AND connector_type=? AND status='active' ORDER BY updated_at DESC LIMIT 1",
        (user_id, connector_type))
    if smart:
        try: return _decrypt(smart["encrypted_conn_data"]), "smart", smart["id"]
        except Exception: pass

    row = await db_fetchone(
        "SELECT id,encrypted_creds FROM connectors WHERE user_id=? AND connector_type=? AND is_active=1",
        (user_id, connector_type))
    if row:
        try: return _decrypt(row["encrypted_creds"]), "legacy", row["id"]
        except Exception: pass
    return None, "", ""

async def _get_connector_creds(user_id: str, connector_type: str) -> Optional[Dict]:
    creds, _, _ = await _get_connector_creds_record(user_id, connector_type)
    return creds

def _n8n_local_info(action: str = "status") -> Dict[str, Any]:
    base_url = os.getenv("N8N_LOCAL_URL", "http://45.79.124.28:5678/").rstrip("/") + "/"
    info: Dict[str, Any] = {
        "status": "available",
        "service": "n8n Local",
        "portal_url": base_url,
        "webhook_base_url": base_url.rstrip("/") + "/webhook/",
        "action": action or "status",
    }
    try:
        r = subprocess.run(
            ["docker", "inspect", "-f", "{{.State.Running}}", "jazz-n8n"],
            capture_output=True, text=True, timeout=5)
        info["container"] = "running" if r.stdout.strip() == "true" else "stopped"
    except Exception:
        info["container"] = "unknown"
    return info

def _connector_error_text(error_obj: Any) -> str:
    if isinstance(error_obj, dict):
        msg = error_obj.get("message") or error_obj.get("error_description") or error_obj.get("status")
        code = error_obj.get("code") or error_obj.get("status_code")
        if code and msg:
            return f"{code}: {msg}"
        if msg:
            return str(msg)
        try:
            return json.dumps(error_obj)[:500]
        except Exception:
            return str(error_obj)
    return str(error_obj)

def _is_oauth_token_error(result: Any) -> bool:
    if not isinstance(result, dict) or "error" not in result:
        return False
    err = result.get("error")
    if isinstance(err, dict):
        code = int(err.get("code") or err.get("status_code") or 0)
        msg = str(err.get("message") or err.get("status") or "").lower()
        return code == 401 or "invalid credentials" in msg or "expired" in msg
    msg = str(err).lower()
    return "401" in msg or "invalid token" in msg or "expired" in msg or "unauthorized" in msg

async def _refresh_connector_creds(user_id: str, connector_type: str, creds: Dict,
                                   source: str, connector_id: str) -> Optional[Dict]:
    refresh_token = creds.get("refresh_token") or creds.get("raw", {}).get("refresh_token")
    if not refresh_token:
        return None
    prov = _oauth_providers().get(connector_type) or {}
    token_url = prov.get("token_url")
    if not token_url:
        return None
    data = {
        "grant_type": "refresh_token",
        "refresh_token": refresh_token,
        "client_id": prov.get("client_id", ""),
        "client_secret": prov.get("client_secret", ""),
    }
    basic_auth = None
    if prov.get("token_auth_type") == "basic":
        basic_auth = (prov.get("client_id", ""), prov.get("client_secret", ""))
        data.pop("client_secret", None)
    tok = await _oauth_post_form(token_url, data, basic_auth=basic_auth)
    if not tok or tok.get("error") or not tok.get("access_token"):
        return None
    updated = dict(creds)
    updated["access_token"] = tok.get("access_token", "")
    updated["refresh_token"] = tok.get("refresh_token") or refresh_token
    updated["expires_in"] = tok.get("expires_in", updated.get("expires_in", 0))
    updated["token_type"] = tok.get("token_type", updated.get("token_type", "Bearer"))
    updated["scope"] = tok.get("scope", updated.get("scope", ""))
    updated["raw"] = {k:v for k,v in tok.items() if k not in ("access_token","refresh_token")}
    encrypted = _encrypt(updated)
    now = _utcnow()
    if source == "smart":
        await db_execute("UPDATE smart_connectors SET encrypted_conn_data=?,updated_at=? WHERE id=? AND user_id=?",
                         (encrypted, now, connector_id, user_id))
    elif source == "legacy":
        await db_execute("UPDATE connectors SET encrypted_creds=?,updated_at=? WHERE id=? AND user_id=?",
                         (encrypted, now, connector_id, user_id))
    return updated

async def _connector_api_call(user_id: str, connector_type: str, action: str, params: Dict = None) -> Dict:
    if connector_type == "n8n_local":
        return _n8n_local_info(action)
    creds, source, connector_id = await _get_connector_creds_record(user_id, connector_type)
    if not creds:
        err = f"{connector_type} not connected. Connect via /connectors/oauth/{connector_type}/init"
        await db_execute(
            "INSERT INTO connector_logs(id,user_id,connector_type,action,status,error_msg,created_at)"
            " VALUES(?,?,?,?,'failure',?,?)",
            (_new_id(), user_id, connector_type or "unknown", action or "unknown", err[:300], _utcnow()))
        return {"error": err}
    params = params or {}
    try:
        result = await _dispatch_connector_action(connector_type, action, creds, params)
        if _is_oauth_token_error(result):
            refreshed = await _refresh_connector_creds(user_id, connector_type, creds, source, connector_id)
            if refreshed:
                result = await _dispatch_connector_action(connector_type, action, refreshed, params)
        if isinstance(result, dict) and result.get("error"):
            err = _connector_error_text(result.get("error"))
            await db_execute(
                "INSERT INTO connector_logs(id,user_id,connector_type,action,status,error_msg,created_at)"
                " VALUES(?,?,?,?,'failure',?,?)",
                (_new_id(), user_id, connector_type, action, err[:300], _utcnow()))
            return {"error": err, "detail": result.get("error")}
        await db_execute(
            "INSERT INTO connector_logs(id,user_id,connector_type,action,status,created_at)"
            " VALUES(?,?,?,?,'success',?)",
            (_new_id(), user_id, connector_type, action, _utcnow()))
        return result
    except Exception as e:
        logger.warning("[CONNECTOR] %s.%s error: %s", connector_type, action, e)
        await db_execute(
            "INSERT INTO connector_logs(id,user_id,connector_type,action,status,error_msg,created_at)"
            " VALUES(?,?,?,?,'failure',?,?)",
            (_new_id(), user_id, connector_type, action, str(e)[:300], _utcnow()))
        return {"error": str(e)}

async def _dispatch_connector_action(connector_type: str, action: str, creds: Dict, params: Dict) -> Dict:
    token = creds.get("access_token","")
    h = {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}

    async def _get(url): return await _http_get(url, h)
    async def _post(url, body): return await _http_post_json(url, body, h)

    if connector_type in ("zapier", "make", "n8n", "webhook"):
        webhook_url = (params.get("webhook_url") or creds.get("webhook_url") or
                       creds.get("url") or creds.get("api_key") or creds.get("auth_value") or "").strip()
        if action in ("status", "info"):
            return {"status":"configured" if webhook_url else "missing_webhook_url", "connector":connector_type}
        if not webhook_url:
            return {"error": f"{connector_type} webhook URL is required"}
        payload = params.get("payload")
        if payload is None:
            payload = {k:v for k,v in params.items() if k not in ("webhook_url",)}
        if not isinstance(payload, (dict, list)):
            payload = {"text": str(payload)}
        return await _http_request_json("POST", webhook_url, payload, {"Content-Type":"application/json"})

    if connector_type == "http":
        base_url = (params.get("base_url") or creds.get("base_url") or "").rstrip("/")
        path = str(params.get("path") or params.get("endpoint") or "").strip()
        url = params.get("url") or (base_url + ("" if path.startswith("/") else "/") + path if base_url else "")
        if not url:
            return {"error":"HTTP connector requires base_url or url"}
        method = (action if action in ("GET","POST","PUT","PATCH","DELETE") else params.get("method","GET")).upper()
        headers = dict(params.get("headers") or {})
        auth_type = (creds.get("auth_type") or params.get("auth_type") or "").lower()
        auth_value = creds.get("auth_value") or creds.get("api_key") or params.get("auth_value") or ""
        if auth_type == "bearer" and auth_value:
            headers["Authorization"] = "Bearer " + auth_value.replace("Bearer ", "")
        elif auth_type == "basic" and auth_value:
            headers["Authorization"] = "Basic " + auth_value
        elif auth_type in ("api_key", "apikey") and auth_value:
            headers[creds.get("auth_header") or params.get("auth_header") or "X-API-Key"] = auth_value
        return await _http_request_json(method, url, params.get("body") or params.get("payload"), headers)

    if connector_type == "graphql":
        endpoint = params.get("endpoint") or creds.get("endpoint") or creds.get("base_url") or ""
        if not endpoint:
            return {"error":"GraphQL endpoint is required"}
        headers = {"Content-Type":"application/json"}
        auth_value = creds.get("auth_value") or creds.get("api_key") or ""
        if auth_value:
            headers["Authorization"] = auth_value if auth_value.lower().startswith(("bearer ", "basic ")) else "Bearer " + auth_value
        body = {"query": params.get("query") or params.get("mutation") or "", "variables": params.get("variables") or {}}
        if not body["query"]:
            return {"error":"GraphQL query is required"}
        return await _http_request_json("POST", endpoint, body, headers)

    # ── Gmail ──────────────────────────────────────────────────────────────────
    if connector_type == "gmail":
        base = "https://gmail.googleapis.com/gmail/v1/users/me"
        def _gmail_body_from_payload(payload: Dict[str, Any]) -> str:
            stack = [payload] if payload else []
            while stack:
                part = stack.pop(0) or {}
                stack.extend(part.get("parts") or [])
                if part.get("mimeType") not in ("text/plain", "text/html"):
                    continue
                data = part.get("body", {}).get("data", "")
                if data:
                    import base64 as b64
                    text = b64.urlsafe_b64decode(data + "==").decode("utf-8", "replace")
                    if part.get("mimeType") == "text/html":
                        text = _strip_html(text)
                    return text
            return ""

        async def _gmail_metadata(messages: List[Dict[str, Any]], count: int, include_body: bool = False) -> List[Dict[str, Any]]:
            details = []
            for m in messages[:count]:
                mid = m.get("id") if isinstance(m, dict) else ""
                if not mid:
                    continue
                try:
                    if include_body:
                        d = await _get(f"{base}/messages/{urllib.parse.quote(mid)}?format=full")
                    else:
                        d = await _get(
                            f"{base}/messages/{urllib.parse.quote(mid)}?format=metadata"
                            "&metadataHeaders=Subject&metadataHeaders=From&metadataHeaders=Date&metadataHeaders=To"
                        )
                    if isinstance(d, dict) and d.get("error"):
                        details.append({"id":mid,"error":_connector_error_text(d.get("error"))})
                        continue
                    hdrs = {
                        str(h2.get("name", "")).lower(): h2.get("value", "")
                        for h2 in d.get("payload", {}).get("headers", [])
                    }
                    body_text = _gmail_body_from_payload(d.get("payload", {})) if include_body else ""
                    details.append({"id":mid,"subject":hdrs.get("subject",""),"from":hdrs.get("from",""),
                                    "to":hdrs.get("to",""),"date":hdrs.get("date",""),
                                    "snippet":d.get("snippet",""),"body":body_text[:3000]})
                except Exception:
                    details.append({"id":mid})
            return details

        if action == "list_emails":
            count = max(1, min(50, int(params.get("count", 10) or 10)))
            query = (params.get("query") or "").strip()
            include_body = bool(params.get("include_body") or params.get("read"))
            searches = []
            if query:
                searches.append({"q": query, "maxResults": count})
            searches.extend([
                {"q": "in:inbox", "maxResults": count},
                {"labelIds": "INBOX", "maxResults": count},
                {"q": "category:primary", "maxResults": count},
                {"q": "newer_than:365d", "maxResults": count},
                {"q": "in:anywhere", "maxResults": count},
                {"maxResults": count},
            ])
            last_error = None
            seen_urls = set()
            for qs in searches:
                url = f"{base}/messages?{urllib.parse.urlencode(qs)}"
                if url in seen_urls:
                    continue
                seen_urls.add(url)
                result = await _get(url)
                if isinstance(result, dict) and result.get("error"):
                    last_error = result
                    continue
                msgs = result.get("messages", []) if isinstance(result, dict) else []
                if msgs:
                    return await _gmail_metadata(msgs, count, include_body)
            thread_result = await _get(f"{base}/threads?maxResults={count}")
            if isinstance(thread_result, dict) and thread_result.get("threads"):
                thread_messages = []
                for t in thread_result.get("threads", [])[:count]:
                    tid = t.get("id")
                    if tid:
                        thread = await _get(f"{base}/threads/{urllib.parse.quote(tid)}?format=metadata")
                        if isinstance(thread, dict):
                            thread_messages.extend({"id": m.get("id")} for m in thread.get("messages", []) if m.get("id"))
                if thread_messages:
                    return await _gmail_metadata(thread_messages, count, include_body)
            if last_error:
                return last_error
            return {"messages": [], "note": "Gmail returned zero messages for this connected account. Reconnect Gmail and confirm the same Google account has inbox mail."}
        if action == "send_email":
            import base64 as b64
            from email.mime.multipart import MIMEMultipart
            from email.mime.text import MIMEText
            from email.mime.base import MIMEBase
            from email import encoders as _encoders
            to = params.get("to", "").strip()
            subject = params.get("subject", "").strip()
            body_text = params.get("body", "").strip()
            if not to:
                return {"error": "Recipient address required. Example: /gmail send to:user@example.com subject:Hello body:Your message"}
            try:
                profile = await _get(f"{base}/profile")
                from_addr = profile.get("emailAddress", "") if isinstance(profile, dict) else ""
            except Exception:
                from_addr = ""
            attachment_path = params.get("attachment_path", "").strip()
            attachment_data = params.get("attachment_data", "")
            attachment_name = params.get("attachment_name", "attachment")
            if attachment_path or attachment_data:
                msg = MIMEMultipart()
                if from_addr: msg["From"] = from_addr
                msg["To"] = to
                msg["Subject"] = subject
                msg.attach(MIMEText(body_text, "plain", "utf-8"))
                if attachment_path:
                    try:
                        file_bytes = Path(attachment_path).read_bytes()
                        fname = Path(attachment_path).name
                        part = MIMEBase("application", "octet-stream")
                        part.set_payload(file_bytes)
                        _encoders.encode_base64(part)
                        part.add_header("Content-Disposition", f'attachment; filename="{fname}"')
                        msg.attach(part)
                    except Exception as e:
                        return {"error": f"Could not read attachment: {e}"}
                elif attachment_data:
                    try:
                        file_bytes = b64.b64decode(attachment_data)
                        part = MIMEBase("application", "octet-stream")
                        part.set_payload(file_bytes)
                        _encoders.encode_base64(part)
                        part.add_header("Content-Disposition", f'attachment; filename="{attachment_name}"')
                        msg.attach(part)
                    except Exception as e:
                        return {"error": f"Could not attach data: {e}"}
                raw_bytes = msg.as_bytes()
            else:
                from_line = f"From: {from_addr}\r\n" if from_addr else ""
                raw_str = f"{from_line}To: {to}\r\nSubject: {subject}\r\nContent-Type: text/plain; charset=utf-8\r\n\r\n{body_text}"
                raw_bytes = raw_str.encode()
            encoded = b64.urlsafe_b64encode(raw_bytes).decode().rstrip("=")
            return await _post(f"{base}/messages/send", {"raw": encoded})
        if action == "create_draft":
            import base64 as b64
            to = params.get("to","").strip(); subject = params.get("subject","").strip()
            body_text = params.get("body","").strip()
            try:
                profile = await _get(f"{base}/profile")
                from_addr = profile.get("emailAddress","") if isinstance(profile, dict) else ""
            except Exception: from_addr = ""
            from_line = f"From: {from_addr}\r\n" if from_addr else ""
            raw_str = f"{from_line}To: {to}\r\nSubject: {subject}\r\nContent-Type: text/plain; charset=utf-8\r\n\r\n{body_text}"
            encoded = b64.urlsafe_b64encode(raw_str.encode()).decode().rstrip("=")
            return await _post(f"{base}/drafts", {"message":{"raw":encoded}})
        if action == "search_emails":
            q = urllib.parse.quote(params.get("query",""))
            count = int(params.get("count", 10) or 10)
            result = await _get(f"{base}/messages?q={q}&maxResults={count}")
            if isinstance(result, dict) and result.get("error"):
                return result
            return await _gmail_metadata(result.get("messages", []), count, bool(params.get("include_body") or params.get("read")))
        if action == "read_email":
            mid = params.get("message_id","")
            result = await _get(f"{base}/messages/{mid}?format=full")
            if isinstance(result, dict) and result.get("error"):
                return result
            payload = result.get("payload",{})
            hdrs = {hh["name"]:hh["value"] for hh in payload.get("headers",[])}
            body_text = ""
            for part in payload.get("parts",[payload]):
                if part.get("mimeType") == "text/plain":
                    data = part.get("body",{}).get("data","")
                    if data:
                        import base64 as b64
                        body_text = b64.urlsafe_b64decode(data+"==").decode("utf-8","replace"); break
            return {"id":mid,"subject":hdrs.get("Subject",""),"from":hdrs.get("From",""),
                    "date":hdrs.get("Date",""),"body":body_text[:3000]}
        if action == "reply_email":
            import base64 as b64
            mid = params.get("message_id","")
            body_text = params.get("body","")
            # Fetch original to get threadId, subject, and sender
            orig = await _get(f"{base}/messages/{mid}?format=metadata&metadataHeaders=Subject&metadataHeaders=From&metadataHeaders=Message-ID")
            thread_id = orig.get("threadId","") if isinstance(orig, dict) else ""
            orig_hdrs = {}
            if isinstance(orig, dict):
                orig_hdrs = {h["name"]: h["value"] for h in orig.get("payload",{}).get("headers",[])}
            to_addr = orig_hdrs.get("From","")
            subj = orig_hdrs.get("Subject","")
            if not subj.lower().startswith("re:"):
                subj = "Re: " + subj
            msg_id_ref = orig_hdrs.get("Message-ID","")
            headers = f"To: {to_addr}\r\nSubject: {subj}\r\nContent-Type: text/plain; charset=utf-8\r\n"
            if msg_id_ref:
                headers += f"In-Reply-To: {msg_id_ref}\r\nReferences: {msg_id_ref}\r\n"
            raw_msg = headers + f"\r\n{body_text}"
            encoded = b64.urlsafe_b64encode(raw_msg.encode()).decode().rstrip("=")
            payload_body: Dict = {"raw": encoded}
            if thread_id:
                payload_body["threadId"] = thread_id
            result = await _post(f"{base}/messages/send", payload_body)
            return result if isinstance(result, dict) and result.get("error") else {"status": "replied", "thread_id": thread_id, "to": to_addr, "subject": subj}

    # ── Google Calendar ────────────────────────────────────────────────────────
    if connector_type == "google_calendar":
        base = "https://www.googleapis.com/calendar/v3/calendars/primary"
        cal_base = "https://www.googleapis.com/calendar/v3"
        if action == "list_events":
            now_iso = datetime.now(timezone.utc).isoformat()
            count = int(params.get("count", 10) or 10)
            qs = f"maxResults={count}&orderBy=startTime&singleEvents=true&timeMin={urllib.parse.quote(now_iso)}"
            if params.get("query"): qs += f"&q={urllib.parse.quote(params['query'])}"
            result = await _get(f"{base}/events?{qs}")
            if isinstance(result, dict) and result.get("error"): return result
            items = result.get("items", [])
            return [{"id":e["id"],"summary":e.get("summary",""),"description":e.get("description",""),
                     "start":e.get("start",{}).get("dateTime",e.get("start",{}).get("date","")),
                     "end":e.get("end",{}).get("dateTime",e.get("end",{}).get("date","")),
                     "location":e.get("location",""),
                     "attendees":[a.get("email") for a in e.get("attendees",[])]} for e in items]
        if action == "get_event":
            return await _get(f"{base}/events/{params.get('event_id','')}")
        if action == "create_event":
            body = {"summary":params.get("summary",params.get("title","Event")),
                    "description":params.get("description",""),
                    "start":{"dateTime":params.get("start"),"timeZone":params.get("timezone","UTC")},
                    "end":{"dateTime":params.get("end"),"timeZone":params.get("timezone","UTC")}}
            if params.get("location"): body["location"] = params["location"]
            if params.get("attendees"): body["attendees"] = [{"email":e} for e in params["attendees"]]
            return await _post(f"{base}/events", body)
        if action == "update_event":
            eid = params.get("event_id","")
            patch: Dict = {}
            if params.get("summary") or params.get("title"): patch["summary"] = params.get("summary", params.get("title"))
            if params.get("description"): patch["description"] = params["description"]
            if params.get("start"): patch["start"] = {"dateTime":params["start"],"timeZone":params.get("timezone","UTC")}
            if params.get("end"):   patch["end"]   = {"dateTime":params["end"],  "timeZone":params.get("timezone","UTC")}
            if params.get("location"): patch["location"] = params["location"]
            return await _http_patch_json(f"{base}/events/{eid}", patch, h)
        if action == "delete_event":
            eid = params.get("event_id","")
            await _http_delete(f"{base}/events/{eid}", h)
            return {"status":"deleted","event_id":eid}
        if action == "list_calendars":
            return await _get(f"{cal_base}/users/me/calendarList")

    # ── Google Drive ────────────────────────────────────────────────────────────
    if connector_type == "google_drive":
        base = "https://www.googleapis.com/drive/v3"
        upload_base = "https://www.googleapis.com/upload/drive/v3"
        if action == "list_files":
            count = int(params.get("count", 20) or 20)
            qs = f"pageSize={count}&fields=files(id,name,mimeType,size,modifiedTime)"
            if params.get("folder_id"): qs += f"&q={urllib.parse.quote(repr(params['folder_id'])+' in parents')}"
            return await _get(f"{base}/files?{qs}")
        if action == "search_files":
            safe_q = params.get("query","").replace("'","\'")
            q = f"name contains '{safe_q}'"
            if params.get("mime_type"):
                safe_mime = params["mime_type"].replace("'","\'")
                q += f" and mimeType='{safe_mime}'"
            return await _get(f"{base}/files?q={urllib.parse.quote(q)}&pageSize={int(params.get('count',10) or 10)}&fields=files(id,name,mimeType,size,modifiedTime,webViewLink)")
        if action == "get_file":
            fid = params.get("file_id","")
            return await _get(f"{base}/files/{fid}?fields=id,name,mimeType,size,webViewLink,description")
        if action == "delete_file":
            await _http_delete(f"{base}/files/{params.get('file_id','')}", h)
            return {"status":"deleted","file_id":params.get("file_id","")}
        if action == "upload_file":
            import base64 as _b64
            file_data = params.get("file_data","")
            file_name = params.get("file_name","upload")
            mime_type = params.get("mime_type","application/octet-stream")
            if not file_data: return {"error":"file_data (base64) required for upload_file"}
            try: raw_bytes = _b64.b64decode(file_data)
            except Exception as e: return {"error":f"Invalid base64 file_data: {e}"}
            meta: Dict = {"name":file_name}
            if params.get("folder_id"): meta["parents"] = [params["folder_id"]]
            boundary = "jazz_boundary_" + secrets.token_hex(8)
            body_parts = (
                f"--{boundary}\r\nContent-Type: application/json; charset=UTF-8\r\n\r\n".encode()
                + json.dumps(meta).encode()
                + f"\r\n--{boundary}\r\nContent-Type: {mime_type}\r\n\r\n".encode()
                + raw_bytes + f"\r\n--{boundary}--".encode())
            def _upload_call():
                req = urllib.request.Request(
                    f"{upload_base}/files?uploadType=multipart", data=body_parts, method="POST",
                    headers={**h, "Content-Type": f"multipart/related; boundary={boundary}"})
                try:
                    with urllib.request.urlopen(req, timeout=30) as r: return json.loads(r.read())
                except urllib.error.HTTPError as e:
                    raw = e.read()
                    try: return json.loads(raw)
                    except Exception: return {"error":f"HTTP {e.code}","detail":raw.decode("utf-8","replace")[:300]}
            return await asyncio.get_running_loop().run_in_executor(_executor, _upload_call)
        if action == "download":
            fid = params.get("file_id","")
            # Returns file metadata with a direct download link
            meta_result = await _get(f"{base}/files/{fid}?fields=id,name,mimeType,size,webViewLink,exportLinks")
            if isinstance(meta_result, dict) and not meta_result.get("error"):
                meta_result["download_url"] = f"{base}/files/{fid}?alt=media"
            return meta_result

    # ── Google Sheets ────────────────────────────────────────────────────────────
    if connector_type == "google_sheets":
        sheets_base = "https://sheets.googleapis.com/v4/spreadsheets"
        if action == "list_files":
            q = urllib.parse.quote("mimeType='application/vnd.google-apps.spreadsheet'")
            return await _http_get(f"https://www.googleapis.com/drive/v3/files?q={q}&pageSize=20&fields=files(id,name,modifiedTime)", h)
        if action in ("get_values", "read_range"):
            sid = params.get("spreadsheet_id",""); rng = urllib.parse.quote(params.get("range","Sheet1"))
            return await _get(f"{sheets_base}/{sid}/values/{rng}")
        if action in ("append_values", "write_range"):
            sid = params.get("spreadsheet_id",""); rng = urllib.parse.quote(params.get("range","Sheet1"))
            return await _post(f"{sheets_base}/{sid}/values/{rng}:append?valueInputOption=USER_ENTERED&insertDataOption=INSERT_ROWS",
                {"values": params.get("values",[])})
        if action == "update_values":
            sid = params.get("spreadsheet_id",""); rng = urllib.parse.quote(params.get("range","Sheet1!A1"))
            return await _post(f"{sheets_base}/{sid}/values/{rng}?valueInputOption=USER_ENTERED",
                {"values": params.get("values",[])})
        if action == "create_spreadsheet":
            return await _post(sheets_base, {"properties":{"title":params.get("title","New Spreadsheet")}})
        if action in ("get_spreadsheet", "get_metadata"):
            sid = params.get("spreadsheet_id","")
            return await _get(f"{sheets_base}/{sid}?fields=spreadsheetId,properties,sheets.properties")

    # ── Google Meet ──────────────────────────────────────────────────────────────
    if connector_type == "google_meet":
        cal_base = "https://www.googleapis.com/calendar/v3/calendars/primary"
        if action in ("create_meeting","create_event"):
            req_id = secrets.token_hex(8)
            body = {"summary":params.get("summary",params.get("title","Meeting")),
                    "description":params.get("description",""),
                    "start":{"dateTime":params.get("start"),"timeZone":params.get("timezone","UTC")},
                    "end":  {"dateTime":params.get("end"),  "timeZone":params.get("timezone","UTC")},
                    "conferenceData":{"createRequest":{"requestId":req_id,
                        "conferenceSolutionKey":{"type":"hangoutsMeet"}}}}
            if params.get("attendees"): body["attendees"] = [{"email":e} for e in params["attendees"]]
            result = await _http_post_json(f"{cal_base}/events?conferenceDataVersion=1", body, h)
            if isinstance(result, dict) and not result.get("error"):
                meet_link = next((ep.get("uri","") for ep in result.get("conferenceData",{}).get("entryPoints",[])
                                  if ep.get("entryPointType")=="video"), "")
                result["meet_link"] = meet_link
            return result
        if action == "list_meetings":
            now_iso = datetime.now(timezone.utc).isoformat()
            qs = f"maxResults=10&orderBy=startTime&singleEvents=true&timeMin={urllib.parse.quote(now_iso)}&q=meet"
            result = await _http_get(f"{cal_base}/events?{qs}", h)
            items = result.get("items",[]) if isinstance(result, dict) else []
            out_list = []
            for e in items:
                meet_link = next((ep.get("uri","") for ep in e.get("conferenceData",{}).get("entryPoints",[])
                                  if ep.get("entryPointType")=="video"), "")
                out_list.append({"id":e["id"],"summary":e.get("summary",""),
                                 "start":e.get("start",{}).get("dateTime",""),"meet_link":meet_link})
            return out_list

    # ── GitHub (full integration — repos, issues, PRs, commits, files, search) ───
    if connector_type == "github":
        base = "https://api.github.com"
        gh_h = {"Authorization":f"Bearer {token}","Accept":"application/vnd.github+json",
                "X-GitHub-Api-Version":"2022-11-28"}

        # ── profile & repos ─────────────────────────────────────────────────────
        if action == "get_user":
            return await _http_get(f"{base}/user", gh_h)
        if action == "list_repos":
            per = max(1, min(50, int(params.get("per_page", 20) or 20)))
            sort = params.get("sort", "updated")  # updated|created|pushed|full_name
            return await _http_get(f"{base}/user/repos?sort={sort}&per_page={per}&type=all", gh_h)
        if action == "get_repo":
            repo = params.get("repo","")
            return await _http_get(f"{base}/repos/{repo}", gh_h)
        if action == "list_branches":
            repo = params.get("repo","")
            return await _http_get(f"{base}/repos/{repo}/branches?per_page=30", gh_h)
        if action == "list_contributors":
            repo = params.get("repo","")
            return await _http_get(f"{base}/repos/{repo}/contributors?per_page=20", gh_h)

        # ── issues ──────────────────────────────────────────────────────────────
        if action == "list_issues":
            repo = params.get("repo","")
            state = params.get("state","open")
            label = params.get("label","")
            per = max(1, min(50, int(params.get("per_page", 20) or 20)))
            qs = f"state={state}&per_page={per}"
            if label: qs += f"&labels={urllib.parse.quote(label)}"
            return await _http_get(f"{base}/repos/{repo}/issues?{qs}", gh_h)
        if action == "get_issue":
            repo = params.get("repo",""); num = params.get("number",1)
            return await _http_get(f"{base}/repos/{repo}/issues/{num}", gh_h)
        if action == "create_issue":
            repo = params.get("repo","")
            body: Dict = {"title":params.get("title",""),"body":params.get("body","")}
            if params.get("labels"): body["labels"] = params["labels"]
            if params.get("assignees"): body["assignees"] = params["assignees"]
            return await _http_post_json(f"{base}/repos/{repo}/issues", body, gh_h)
        if action == "close_issue":
            repo = params.get("repo",""); num = params.get("number",1)
            return await _http_patch_json(f"{base}/repos/{repo}/issues/{num}", {"state":"closed"}, gh_h)
        if action == "comment_issue":
            repo = params.get("repo",""); num = params.get("number",1)
            return await _http_post_json(f"{base}/repos/{repo}/issues/{num}/comments",
                {"body":params.get("body","")}, gh_h)

        # ── pull requests ────────────────────────────────────────────────────────
        if action == "list_prs":
            repo = params.get("repo","")
            state = params.get("state","open")  # open|closed|all
            per = max(1, min(50, int(params.get("per_page", 20) or 20)))
            return await _http_get(f"{base}/repos/{repo}/pulls?state={state}&per_page={per}&sort=updated", gh_h)
        if action == "get_pr":
            repo = params.get("repo",""); num = params.get("number",1)
            return await _http_get(f"{base}/repos/{repo}/pulls/{num}", gh_h)
        if action == "create_pr":
            repo = params.get("repo","")
            body: Dict = {
                "title": params.get("title",""),
                "body":  params.get("body",""),
                "head":  params.get("head",""),  # branch to merge FROM
                "base":  params.get("base","main"),  # branch to merge INTO
                "draft": bool(params.get("draft", False)),
            }
            return await _http_post_json(f"{base}/repos/{repo}/pulls", body, gh_h)
        if action == "merge_pr":
            repo = params.get("repo",""); num = params.get("number",1)
            merge_body: Dict = {
                "merge_method": params.get("method","squash"),  # squash|merge|rebase
                "commit_title": params.get("title",""),
            }
            return await _http_post_json(f"{base}/repos/{repo}/pulls/{num}/merge", merge_body, gh_h)
        if action == "pr_files":
            repo = params.get("repo",""); num = params.get("number",1)
            return await _http_get(f"{base}/repos/{repo}/pulls/{num}/files?per_page=30", gh_h)
        if action == "pr_reviews":
            repo = params.get("repo",""); num = params.get("number",1)
            return await _http_get(f"{base}/repos/{repo}/pulls/{num}/reviews", gh_h)
        if action == "review_pr":
            repo = params.get("repo",""); num = params.get("number",1)
            event = params.get("event","COMMENT")  # APPROVE|REQUEST_CHANGES|COMMENT
            review_body: Dict = {"body": params.get("body",""), "event": event}
            if params.get("comments"):
                review_body["comments"] = params["comments"]
            return await _http_post_json(f"{base}/repos/{repo}/pulls/{num}/reviews", review_body, gh_h)

        # ── commits ─────────────────────────────────────────────────────────────
        if action == "list_commits":
            repo = params.get("repo","")
            branch = params.get("branch","") or params.get("sha","")
            per = max(1, min(50, int(params.get("per_page", 20) or 20)))
            qs = f"per_page={per}"
            if branch: qs += f"&sha={urllib.parse.quote(branch)}"
            if params.get("author"): qs += f"&author={urllib.parse.quote(params['author'])}"
            return await _http_get(f"{base}/repos/{repo}/commits?{qs}", gh_h)
        if action == "get_commit":
            repo = params.get("repo",""); sha = params.get("sha","")
            return await _http_get(f"{base}/repos/{repo}/commits/{sha}", gh_h)

        # ── files & content ──────────────────────────────────────────────────────
        if action == "get_file":
            repo = params.get("repo",""); path = params.get("path","")
            ref = params.get("ref","") or params.get("branch","")
            url = f"{base}/repos/{repo}/contents/{path}"
            if ref: url += f"?ref={urllib.parse.quote(ref)}"
            result = await _http_get(url, gh_h)
            if isinstance(result, dict) and result.get("content") and not result.get("error"):
                try:
                    import base64 as _b64
                    decoded = _b64.b64decode(result["content"].replace("\n","")).decode("utf-8","replace")
                    result["decoded_content"] = decoded[:8000]
                except Exception:
                    pass
            return result
        if action == "list_files":
            repo = params.get("repo",""); path = params.get("path","")
            ref = params.get("ref","") or params.get("branch","")
            url = f"{base}/repos/{repo}/contents/{path}"
            if ref: url += f"?ref={urllib.parse.quote(ref)}"
            return await _http_get(url, gh_h)
        if action == "create_or_update_file":
            repo = params.get("repo",""); path = params.get("path","")
            import base64 as _b64
            content_b64 = _b64.b64encode(params.get("content","").encode()).decode()
            file_body: Dict = {
                "message": params.get("message","Update file"),
                "content": content_b64,
            }
            if params.get("sha"): file_body["sha"] = params["sha"]  # required for updates
            if params.get("branch"): file_body["branch"] = params["branch"]
            return await _http_post_json(f"{base}/repos/{repo}/contents/{path}", file_body, gh_h)

        # ── search ───────────────────────────────────────────────────────────────
        if action == "search_code":
            q = params.get("query","")
            repo = params.get("repo","")
            if repo: q = f"{q} repo:{repo}"
            per = max(1, min(30, int(params.get("per_page", 10) or 10)))
            return await _http_get(f"{base}/search/code?q={urllib.parse.quote(q)}&per_page={per}", gh_h)
        if action == "search_issues":
            q = params.get("query","")
            repo = params.get("repo","")
            if repo: q = f"{q} repo:{repo}"
            per = max(1, min(30, int(params.get("per_page", 10) or 10)))
            return await _http_get(f"{base}/search/issues?q={urllib.parse.quote(q)}&per_page={per}", gh_h)
        if action == "search_repos":
            q = params.get("query","")
            per = max(1, min(30, int(params.get("per_page", 10) or 10)))
            return await _http_get(f"{base}/search/repositories?q={urllib.parse.quote(q)}&per_page={per}&sort=stars", gh_h)
        if action == "search_commits":
            q = params.get("query","")
            repo = params.get("repo","")
            if repo: q = f"{q} repo:{repo}"
            per = max(1, min(30, int(params.get("per_page", 10) or 10)))
            return await _http_get(f"{base}/search/commits?q={urllib.parse.quote(q)}&per_page={per}", gh_h)

        # ── workflows & actions ──────────────────────────────────────────────────
        if action == "list_workflows":
            repo = params.get("repo","")
            return await _http_get(f"{base}/repos/{repo}/actions/workflows", gh_h)
        if action == "list_workflow_runs":
            repo = params.get("repo",""); wid = params.get("workflow_id","")
            url = f"{base}/repos/{repo}/actions/runs?per_page=10"
            if wid: url = f"{base}/repos/{repo}/actions/workflows/{wid}/runs?per_page=10"
            return await _http_get(url, gh_h)
        if action == "trigger_workflow":
            repo = params.get("repo",""); wid = params.get("workflow_id","")
            ref = params.get("ref","main")
            trigger_body: Dict = {"ref": ref, "inputs": params.get("inputs",{})}
            return await _http_post_json(f"{base}/repos/{repo}/actions/workflows/{wid}/dispatches", trigger_body, gh_h)

        # ── releases ─────────────────────────────────────────────────────────────
        if action == "list_releases":
            repo = params.get("repo","")
            return await _http_get(f"{base}/repos/{repo}/releases?per_page=10", gh_h)
        if action == "create_release":
            repo = params.get("repo","")
            rel_body: Dict = {
                "tag_name":   params.get("tag",""),
                "name":       params.get("name",""),
                "body":       params.get("body",""),
                "draft":      bool(params.get("draft",False)),
                "prerelease": bool(params.get("prerelease",False)),
            }
            return await _http_post_json(f"{base}/repos/{repo}/releases", rel_body, gh_h)

        # ── gists ────────────────────────────────────────────────────────────────
        if action == "list_gists":
            return await _http_get(f"{base}/gists?per_page=20", gh_h)
        if action == "create_gist":
            files = {params.get("filename","snippet.txt"): {"content": params.get("content","")}}
            gist_body: Dict = {
                "description": params.get("description",""),
                "public": bool(params.get("public",False)),
                "files": files,
            }
            return await _http_post_json(f"{base}/gists", gist_body, gh_h)

        # ── notifications & activity ─────────────────────────────────────────────
        if action == "list_notifications":
            return await _http_get(f"{base}/notifications?per_page=20&all=false", gh_h)
        if action == "list_starred":
            return await _http_get(f"{base}/user/starred?per_page=20&sort=updated", gh_h)
        if action == "star_repo":
            repo = params.get("repo","")
            req_obj = urllib.request.Request(f"{base}/user/starred/{repo}", method="PUT", headers=gh_h)
            req_obj.add_header("Content-Length","0")
            try:
                urllib.request.urlopen(req_obj, timeout=15)
                return {"starred":True,"repo":repo}
            except Exception as exc:
                return {"error":str(exc)}

        return {"error":f"Unknown GitHub action: {action}"}

    # ── Notion ─────────────────────────────────────────────────────────────────
    if connector_type == "notion":
        base = "https://api.notion.com/v1"
        nh = {**h, "Notion-Version":"2022-06-28"}
        if action == "search":
            body: Dict = {"query": params.get("query","")}
            if params.get("filter"): body["filter"] = {"value": params["filter"], "property": "object"}
            return await _http_post_json(f"{base}/search", body, nh)
        if action == "get_page":
            pid = params.get("id","") or params.get("page_id","")
            return await _http_get(f"{base}/pages/{pid}", nh)
        if action == "get_page_content":
            pid = params.get("id","") or params.get("page_id","")
            return await _http_get(f"{base}/blocks/{pid}/children?page_size=50", nh)
        if action == "create_page":
            parent: Dict = {}
            if params.get("database_id"): parent = {"database_id": params["database_id"]}
            elif params.get("page_id"): parent = {"page_id": params["page_id"]}
            body = {"parent": parent, "properties": {
                "Name": {"title": [{"text": {"content": params.get("title","")}}]}}}
            if params.get("content"):
                body["children"] = [{"object":"block","type":"paragraph",
                    "paragraph":{"rich_text":[{"type":"text","text":{"content":params["content"]}}]}}]
            return await _http_post_json(f"{base}/pages", body, nh)
        if action == "update_page":
            pid = params.get("id","") or params.get("page_id","")
            props: Dict = {}
            if params.get("title"): props["Name"] = {"title": [{"text": {"content": params["title"]}}]}
            return await _http_patch_json(f"{base}/pages/{pid}", {"properties": props}, nh)
        if action == "list_databases":
            result = await _http_post_json(f"{base}/search", {"filter":{"value":"database","property":"object"}}, nh)
            dbs = result.get("results",[]) if isinstance(result,dict) else []
            return [{"id":d.get("id",""),"title":((d.get("title") or [{}])[0].get("plain_text",""))} for d in dbs]

    # ── Slack ──────────────────────────────────────────────────────────────────
    if connector_type == "slack":
        bot_token = creds.get("bot_token", SLACK_BOT_TOKEN) or token
        sh = {"Authorization":f"Bearer {bot_token}","Content-Type":"application/json"}
        base = "https://slack.com/api"
        async def _slack_resolve_channel(name_or_id: str) -> str:
            """Resolve #channel-name to channel ID if needed."""
            name = name_or_id.lstrip("#").lower()
            if not name: return ""
            # If it looks like a Slack channel ID (C + alphanumeric) return as-is
            if re.match(r"^[CG][A-Z0-9]{6,}$", name_or_id): return name_or_id
            ch_list = await _http_get(f"{base}/conversations.list?limit=200&exclude_archived=true", sh)
            for ch in (ch_list.get("channels",[]) if isinstance(ch_list,dict) else []):
                if ch.get("name","").lower() == name:
                    return ch.get("id","")
            return name_or_id
        if action == "send_message":
            channel = await _slack_resolve_channel(params.get("channel","general"))
            text = params.get("text","")
            if not text: return {"error": "text is required for send_message"}
            result = await _http_post_json(f"{base}/chat.postMessage", {"channel":channel,"text":text}, sh)
            if isinstance(result,dict) and not result.get("ok"):
                return {"error": result.get("error","slack_error"), "detail": result}
            return {"ok": True, "channel": channel, "ts": result.get("ts","")} if isinstance(result,dict) else result
        if action == "list_channels":
            result = await _http_get(f"{base}/conversations.list?limit=100&exclude_archived=true", sh)
            channels = result.get("channels",[]) if isinstance(result,dict) else []
            return [{"id":c.get("id",""),"name":c.get("name",""),"is_private":c.get("is_private",False),
                     "num_members":c.get("num_members",0)} for c in channels]
        if action == "get_messages":
            channel = await _slack_resolve_channel(params.get("channel",""))
            count = int(params.get("count",20) or 20)
            result = await _http_get(f"{base}/conversations.history?channel={channel}&limit={count}", sh)
            if isinstance(result,dict) and not result.get("ok"):
                return {"error": result.get("error","channel_not_found"), "detail": "Ensure the bot is added to the channel"}
            msgs = result.get("messages",[]) if isinstance(result,dict) else []
            return [{"ts":m.get("ts",""),"user":m.get("user",""),"text":m.get("text",""),"type":m.get("type","")} for m in msgs]
        if action == "search":
            q = params.get("query","")
            result = await _http_get(f"{base}/search.messages?query={urllib.parse.quote(q)}&count=10", sh)
            if isinstance(result,dict) and not result.get("ok"):
                return {"error": result.get("error","search_failed")}
            return result.get("messages",{}).get("matches",[]) if isinstance(result,dict) else result

    # ── Jira ───────────────────────────────────────────────────────────────────
    if connector_type == "jira":
        cloud_id = creds.get("cloud_id","")
        if not cloud_id:
            # Fetch and cache cloud_id from Atlassian accessible resources
            try:
                res = await _http_get("https://api.atlassian.com/oauth/token/accessible-resources", h)
                if isinstance(res, list) and res:
                    cloud_id = res[0].get("id","")
            except Exception:
                pass
        base = f"https://api.atlassian.com/ex/jira/{cloud_id}/rest/api/3"
        if action == "list_issues":
            jql = params.get("jql","assignee=currentUser() ORDER BY updated DESC")
            return await _get(f"{base}/search?jql={urllib.parse.quote(jql)}&maxResults={int(params.get('count',20))}")
        if action == "get_issue":
            return await _get(f"{base}/issue/{params.get('issue_key','')}")
        if action == "create_issue":
            body = {"fields":{"project":{"key":params.get("project","PROJ")},
                              "summary":params.get("title",""),"issuetype":{"name":params.get("issue_type","Task")},
                              "description":{"version":1,"type":"doc","content":[{"type":"paragraph","content":[{"type":"text","text":params.get("description","")}]}]}}}
            return await _post(f"{base}/issue", body)
        if action == "update_issue":
            issue_key = params.get("issue_key","")
            fields: Dict = {}
            if params.get("summary"): fields["summary"] = params["summary"]
            if params.get("status"):
                trans = await _get(f"{base}/issue/{issue_key}/transitions")
                if isinstance(trans, dict):
                    for t in trans.get("transitions",[]):
                        if params["status"].lower() in t.get("name","").lower():
                            await _post(f"{base}/issue/{issue_key}/transitions", {"transition":{"id":t["id"]}})
                            break
            if fields:
                return await _http_patch_json(f"{base}/issue/{issue_key}", {"fields":fields}, h)
            return {"status":"updated","issue_key":issue_key}
        if action == "search":
            q = params.get("query","")
            jql = f'text ~ "{q}" ORDER BY updated DESC' if q else "ORDER BY updated DESC"
            return await _get(f"{base}/search?jql={urllib.parse.quote(jql)}&maxResults=10")

    # ── Linear ─────────────────────────────────────────────────────────────────
    if connector_type == "linear":
        lin_base = "https://api.linear.app/graphql"
        def _gql(query: str, variables: Dict = None) -> Dict:
            return {"query": query, "variables": variables or {}}
        if action == "list_issues":
            count = int(params.get("count", 20) or 20)
            state_clause = f', filter: {{state: {{name: {{eq: "{params["state"]}"}}}}}}'.rstrip() if params.get("state") else ""
            q = _gql(f'{{ issues(first:{count}{state_clause}){{nodes{{id title state{{name}} assignee{{name}} priority url createdAt}}}}}}')
            return await _http_post_json(lin_base, q, h)
        if action == "create_issue":
            teams = await _http_post_json(lin_base, _gql("{ teams { nodes { id name } } }"), h)
            team_id = ""
            if isinstance(teams, dict):
                nodes = teams.get("data", {}).get("teams", {}).get("nodes", [])
                team_id = nodes[0].get("id", "") if nodes else ""
            q = _gql("mutation CreateIssue($title: String!, $desc: String, $teamId: String!) { issueCreate(input: {title: $title, description: $desc, teamId: $teamId}) { success issue { id title url } } }",
                     {"title": params.get("title", ""), "desc": params.get("description", ""), "teamId": team_id})
            return await _http_post_json(lin_base, q, h)
        if action == "update_issue":
            q = _gql("mutation UpdateIssue($id: String!, $title: String, $desc: String) { issueUpdate(id: $id, input: {title: $title, description: $desc}) { success issue { id title } } }",
                     {"id": params.get("issue_id", ""), "title": params.get("title"), "desc": params.get("description")})
            return await _http_post_json(lin_base, q, h)
        if action == "list_projects":
            q = _gql("{ projects { nodes { id name description url } } }")
            return await _http_post_json(lin_base, q, h)
        # default: list issues
        q = _gql("{ issues(first:20) { nodes { id title state{name} assignee{name} priority url } } }")
        return await _http_post_json(lin_base, q, h)

    # ── Asana ──────────────────────────────────────────────────────────────────
    if connector_type == "asana":
        base = "https://app.asana.com/api/1.0"
        asana_h = {**h, "Accept": "application/json"}
        if action == "list_tasks":
            workspace = params.get("workspace","")
            if not workspace:
                # Auto-discover workspace
                try:
                    me = await _http_get(f"{base}/users/me?opt_fields=workspaces", asana_h)
                    workspaces = (me.get("data",{}) if isinstance(me,dict) else {}).get("workspaces",[])
                    workspace = workspaces[0].get("gid","") if workspaces else ""
                except Exception: pass
            qs = f"assignee=me&opt_fields=name,completed,due_on,notes,projects.name"
            if workspace: qs += f"&workspace={workspace}"
            result = await _http_get(f"{base}/tasks?{qs}", asana_h)
            items = result.get("data",[]) if isinstance(result,dict) else []
            return [{"id":t.get("gid",""),"name":t.get("name",""),"completed":t.get("completed",False),
                     "due_on":t.get("due_on",""),"notes":t.get("notes","")} for t in items]
        if action == "create_task":
            workspace = params.get("workspace","")
            if not workspace:
                try:
                    me = await _http_get(f"{base}/users/me?opt_fields=workspaces", asana_h)
                    ws_list = (me.get("data",{}) if isinstance(me,dict) else {}).get("workspaces",[])
                    workspace = ws_list[0].get("gid","") if ws_list else ""
                except Exception: pass
            task_data: Dict = {"name": params.get("name",""), "notes": params.get("notes",""), "assignee": "me"}
            if workspace: task_data["workspace"] = workspace
            if params.get("project"): task_data["projects"] = [params["project"]]
            if params.get("due_on"): task_data["due_on"] = params["due_on"]
            result = await _http_post_json(f"{base}/tasks", {"data": task_data}, asana_h)
            data = result.get("data",{}) if isinstance(result,dict) else {}
            return {"id": data.get("gid",""), "name": data.get("name",""), "status": "created"}
        if action == "list_projects":
            workspace = params.get("workspace","")
            if not workspace:
                try:
                    me = await _http_get(f"{base}/users/me?opt_fields=workspaces", asana_h)
                    ws_list = (me.get("data",{}) if isinstance(me,dict) else {}).get("workspaces",[])
                    workspace = ws_list[0].get("gid","") if ws_list else ""
                except Exception: pass
            result = await _http_get(f"{base}/projects?workspace={workspace}&opt_fields=name,color", asana_h)
            items = result.get("data",[]) if isinstance(result,dict) else []
            return [{"id":p.get("gid",""),"name":p.get("name","")} for p in items]

    # ── HubSpot ────────────────────────────────────────────────────────────────
    if connector_type == "hubspot":
        base = "https://api.hubapi.com"
        if action == "list_contacts":
            count = int(params.get("count",20) or 20)
            return await _get(f"{base}/crm/v3/objects/contacts?limit={count}&properties=email,firstname,lastname,phone,company")
        if action == "create_contact":
            return await _post(f"{base}/crm/v3/objects/contacts",
                {"properties":{"email":params.get("email",""),"firstname":params.get("first",params.get("firstname","")),
                               "lastname":params.get("last",params.get("lastname","")),"phone":params.get("phone",""),
                               "company":params.get("company","")}})
        if action == "search":
            q = params.get("query","")
            obj = params.get("object","contacts")
            body = {"query":q,"limit":10,"properties":["email","firstname","lastname","name","dealname"]}
            return await _post(f"{base}/crm/v3/objects/{obj}/search", body)
        if action == "list_deals":
            count = int(params.get("count",20) or 20)
            return await _get(f"{base}/crm/v3/objects/deals?limit={count}&properties=dealname,amount,closedate,dealstage")
        if action == "create_deal":
            return await _post(f"{base}/crm/v3/objects/deals",
                {"properties":{"dealname":params.get("name",""),"amount":str(params.get("amount","0")),
                               "closedate":params.get("close_date",""),"pipeline":params.get("pipeline","default"),
                               "dealstage":params.get("stage","appointmentscheduled")}})

    # ── Airtable ───────────────────────────────────────────────────────────────
    if connector_type == "airtable":
        base = "https://api.airtable.com/v0"
        if action == "list_bases":
            return await _http_get("https://api.airtable.com/v0/meta/bases", h)
        if action == "list_records":
            return await _get(f"{base}/{params.get('base_id','')}/{params.get('table','')}")

    # ── Dropbox ────────────────────────────────────────────────────────────────
    if connector_type == "dropbox":
        dbx_h = {"Authorization":f"Bearer {token}","Content-Type":"application/json"}
        dbx_api = "https://api.dropboxapi.com/2"
        if action == "list_files":
            path = params.get("path","")
            return await _http_post_json(f"{dbx_api}/files/list_folder",
                {"path": path, "limit": int(params.get("count",50) or 50)}, dbx_h)
        if action == "get_file":
            path = params.get("path","")
            meta = await _http_post_json(f"{dbx_api}/files/get_metadata", {"path":path}, dbx_h)
            meta["download_url"] = f"https://content.dropboxapi.com/2/files/download"
            return meta
        if action == "search":
            q = params.get("query","")
            return await _http_post_json(f"{dbx_api}/files/search_v2",
                {"query":q,"options":{"max_results":20}}, dbx_h)
        if action == "delete":
            path = params.get("path","")
            result = await _http_post_json(f"{dbx_api}/files/delete_v2", {"path":path}, dbx_h)
            return {"status":"deleted","path":path} if not (isinstance(result,dict) and result.get("error")) else result
        if action == "create_folder":
            path = params.get("path","")
            return await _http_post_json(f"{dbx_api}/files/create_folder_v2", {"path":path,"autorename":True}, dbx_h)

    # ── Zoom ───────────────────────────────────────────────────────────────────
    if connector_type == "zoom":
        base = "https://api.zoom.us/v2"
        if action == "list_meetings":
            meeting_type = params.get("type","scheduled")
            return await _get(f"{base}/users/me/meetings?type={meeting_type}")
        if action == "create_meeting":
            body: Dict = {
                "topic": params.get("topic","Meeting"),
                "type": 2,  # Scheduled
                "duration": int(params.get("duration",60) or 60),
                "agenda": params.get("agenda",""),
                "settings": {"host_video": True, "participant_video": True, "join_before_host": False}
            }
            if params.get("start_time"): body["start_time"] = params["start_time"]
            if params.get("timezone"): body["timezone"] = params["timezone"]
            result = await _post(f"{base}/users/me/meetings", body)
            if isinstance(result,dict) and not result.get("error"):
                result["join_url"] = result.get("join_url","")
            return result
        if action == "get_meeting":
            mid = params.get("meeting_id","")
            return await _get(f"{base}/meetings/{mid}")
        if action == "delete_meeting":
            mid = params.get("meeting_id","")
            await _http_delete(f"{base}/meetings/{mid}", h)
            return {"status":"deleted","meeting_id":mid}
        if action == "list_recordings":
            return await _get(f"{base}/users/me/recordings?page_size=10")

    # ── Microsoft Graph (Outlook / OneDrive / Teams / Power BI) ───────────────
    if connector_type in ("outlook","microsoft_teams","onedrive","power_bi"):
        base = "https://graph.microsoft.com/v1.0"
        if connector_type == "outlook":
            if action == "list_emails":
                count = int(params.get("count",20) or 20)
                select = "$select=id,subject,from,receivedDateTime,bodyPreview,isRead"
                return await _get(f"{base}/me/messages?$top={count}&$orderby=receivedDateTime%20desc&{select}")
            if action == "search_emails":
                q = params.get("query","")
                return await _get(f'{base}/me/messages?$search="{urllib.parse.quote(q)}"&$top=10&$select=id,subject,from,receivedDateTime,bodyPreview')
            if action == "send_email":
                to = params.get("to","")
                if not to: return {"error":"to address required"}
                body = {"message":{"subject":params.get("subject","(no subject)"),
                                   "body":{"contentType":"Text","content":params.get("body","")},
                                   "toRecipients":[{"emailAddress":{"address":to}}]}}
                await _post(f"{base}/me/sendMail", body)
                return {"status":"sent","to":to}
            if action == "read_email":
                mid = params.get("message_id","")
                return await _get(f"{base}/me/messages/{mid}")
            if action == "delete_email":
                mid = params.get("message_id","")
                await _http_delete(f"{base}/me/messages/{mid}", h)
                return {"status":"deleted"}
        if connector_type == "onedrive":
            if action == "list_files":
                return await _get(f"{base}/me/drive/root/children")
        if connector_type == "microsoft_teams":
            if action == "list_channels":
                return await _get(f"{base}/me/joinedTeams")
        if connector_type == "power_bi":
            if action == "list_datasets":
                wid = creds.get("workspace_id","")
                return await _get(f"https://api.powerbi.com/v1.0/myorg/groups/{wid}/datasets")

    # ── Salesforce ─────────────────────────────────────────────────────────────
    if connector_type == "salesforce":
        instance_url = (creds.get("instance_url") or "").rstrip("/")
        if not instance_url:
            return {"error":"Salesforce instance_url not found in credentials. Re-connect Salesforce."}
        base = f"{instance_url}/services/data/v57.0"
        sf_h = {**h, "Authorization":f"Bearer {token}"}
        if action == "list_contacts":
            count = int(params.get("count",20) or 20)
            q = f"SELECT Id,Name,Email,Phone,Account.Name FROM Contact LIMIT {count}"
            return await _http_get(f"{base}/query?q={urllib.parse.quote(q)}", sf_h)
        if action == "list_accounts":
            count = int(params.get("count",20) or 20)
            q = f"SELECT Id,Name,Industry,Phone,Website FROM Account LIMIT {count}"
            return await _http_get(f"{base}/query?q={urllib.parse.quote(q)}", sf_h)
        if action == "query":
            soql = params.get("query","SELECT Id,Name FROM Contact LIMIT 10")
            return await _http_get(f"{base}/query?q={urllib.parse.quote(soql)}", sf_h)
        if action == "get_record":
            obj = params.get("object","Contact"); rid = params.get("record_id","")
            return await _http_get(f"{base}/sobjects/{obj}/{rid}", sf_h)
        if action == "create_record":
            obj = params.get("object","Contact")
            return await _http_post_json(f"{base}/sobjects/{obj}", params.get("fields",{}), sf_h)

    # ── LinkedIn ───────────────────────────────────────────────────────────────
    if connector_type == "linkedin":
        if action == "get_profile":
            return await _http_get("https://api.linkedin.com/v2/me",
                {"Authorization":f"Bearer {token}","X-Restli-Protocol-Version":"2.0.0"})

    # ── Discord ────────────────────────────────────────────────────────────────
    if connector_type == "discord":
        base = "https://discord.com/api/v10"
        if action == "get_guilds":
            return await _get(f"{base}/users/@me/guilds")
        if action == "send_message":
            return await _post(f"{base}/channels/{params.get('channel_id','')}/messages",
                {"content":params.get("text","")})
        if action == "list_channels":
            guild_id = params.get("guild_id","")
            if not guild_id: return {"error":"guild_id required for list_channels"}
            result = await _get(f"{base}/guilds/{guild_id}/channels")
            if isinstance(result, list):
                return [{"id":c.get("id",""),"name":c.get("name",""),"type":c.get("type",0)} for c in result]
            return result
        if action == "get_messages":
            channel_id = params.get("channel_id","")
            count = int(params.get("count",20) or 20)
            result = await _get(f"{base}/channels/{channel_id}/messages?limit={count}")
            if isinstance(result, list):
                return [{"id":m.get("id",""),"content":m.get("content",""),
                         "author":m.get("author",{}).get("username",""),
                         "timestamp":m.get("timestamp","")} for m in result]
            return result

    # ── Stripe ─────────────────────────────────────────────────────────────────
    if connector_type == "stripe":
        stripe_key = creds.get("api_key","") or os.getenv("STRIPE_SECRET_KEY","")
        sh = {"Authorization":f"Bearer {stripe_key}","Content-Type":"application/x-www-form-urlencoded"}
        sbase = "https://api.stripe.com/v1"
        async def _sget(url): return await _http_get(url, sh)
        async def _spost(url, data):
            def _call():
                body = urllib.parse.urlencode(data).encode()
                req = urllib.request.Request(url, data=body, headers=sh, method="POST")
                try:
                    with urllib.request.urlopen(req, timeout=20) as r: return json.loads(r.read())
                except urllib.error.HTTPError as e:
                    raw = e.read()
                    try: return json.loads(raw)
                    except Exception: return {"error":f"HTTP {e.code}","detail":raw.decode("utf-8","replace")[:300]}
            return await asyncio.get_running_loop().run_in_executor(_executor, _call)
        if action == "get_balance":
            return await _sget(f"{sbase}/balance")
        if action == "list_charges":
            count = int(params.get("count",10) or 10)
            return await _sget(f"{sbase}/charges?limit={count}")
        if action == "list_customers":
            count = int(params.get("count",10) or 10)
            q = params.get("query","")
            if q: return await _sget(f"{sbase}/customers/search?query={urllib.parse.quote(q)}&limit={count}")
            return await _sget(f"{sbase}/customers?limit={count}")
        if action == "create_payment":
            data = {"amount":str(int(float(params.get("amount",0))*100)),
                    "currency":params.get("currency","usd"),
                    "customer":params.get("customer",""),
                    "description":params.get("description","")}
            return await _spost(f"{sbase}/payment_intents", data)
        if action == "list_subscriptions":
            return await _sget(f"{sbase}/subscriptions?limit={params.get('count',10)}")
        if action == "list_invoices":
            return await _sget(f"{sbase}/invoices?limit={params.get('count',10)}")

    # ── BigQuery ────────────────────────────────────────────────────────────────
    if connector_type == "bigquery":
        bq_base = "https://bigquery.googleapis.com/bigquery/v2"
        project = params.get("project_id","") or creds.get("project_id","")
        if action == "list_datasets":
            if not project: return {"error":"project_id required"}
            return await _get(f"{bq_base}/projects/{project}/datasets")
        if action == "list_tables":
            dataset = params.get("dataset_id","")
            if not project or not dataset: return {"error":"project_id and dataset_id required"}
            return await _get(f"{bq_base}/projects/{project}/datasets/{dataset}/tables")
        if action == "run_query":
            if not project: return {"error":"project_id required"}
            body = {"query":params.get("query",""),"useLegacySql":False,
                    "maxResults":int(params.get("max_results",100) or 100)}
            return await _post(f"{bq_base}/projects/{project}/queries", body)

    # ── Box ─────────────────────────────────────────────────────────────────────
    if connector_type == "box":
        bx_base = "https://api.box.com/2.0"
        if action == "list_files":
            folder_id = params.get("folder_id","0")
            return await _get(f"{bx_base}/folders/{folder_id}/items?limit=50&fields=id,name,type,size,modified_at")
        if action == "get_file":
            fid = params.get("file_id","")
            return await _get(f"{bx_base}/files/{fid}?fields=id,name,size,shared_link,modified_at")
        if action == "upload":
            return {"error":"Box upload requires multipart — use the Box web interface or SDK directly"}
        if action == "search":
            q = urllib.parse.quote(params.get("query",""))
            return await _get(f"{bx_base}/search?query={q}&limit=20&fields=id,name,type,size")

    return {"error": f"Unknown action '{action}' for connector '{connector_type}'"}

# ── HTTP helpers ──────────────────────────────────────────────────────────────
async def _http_get(url: str, headers: Dict) -> Any:
    def _call():
        req = urllib.request.Request(url, headers=headers)
        try:
            with urllib.request.urlopen(req, timeout=20) as r:
                return json.loads(r.read())
        except urllib.error.HTTPError as e:
            raw = e.read()
            try: return json.loads(raw)
            except Exception: return {"error":f"HTTP {e.code}","detail":raw.decode("utf-8","replace")[:200]}
    return await asyncio.get_running_loop().run_in_executor(_executor, _call)

def _strip_html(text: str) -> str:
    text = re.sub(r"<[^>]+>", " ", text or "")
    return re.sub(r"\s+", " ", urllib.parse.unquote(text)).strip()

async def _web_search(query: str, count: int = 5) -> List[Dict]:
    query = (query or "").strip()
    if not query:
        return []
    count = max(1, min(10, int(count or 5)))

    brave_key = os.getenv("BRAVE_SEARCH_API_KEY", "").strip()
    tavily_key = os.getenv("TAVILY_API_KEY", "").strip()
    serpapi_key = os.getenv("SERPAPI_API_KEY", "").strip()

    def _call() -> List[Dict]:
        try:
            if brave_key:
                url = "https://api.search.brave.com/res/v1/web/search?q=" + urllib.parse.quote(query) + f"&count={count}"
                req = urllib.request.Request(url, headers={
                    "Accept": "application/json",
                    "X-Subscription-Token": brave_key,
                    "User-Agent": "JAZZ-AI/14",
                })
                with urllib.request.urlopen(req, timeout=20) as r:
                    data = json.loads(r.read())
                return [
                    {"title": x.get("title", ""), "url": x.get("url", ""), "snippet": _strip_html(x.get("description", ""))}
                    for x in (data.get("web", {}).get("results", [])[:count])
                ]
            if tavily_key:
                body = json.dumps({"api_key": tavily_key, "query": query, "max_results": count}).encode()
                req = urllib.request.Request("https://api.tavily.com/search", data=body,
                                             headers={"Content-Type": "application/json"}, method="POST")
                with urllib.request.urlopen(req, timeout=20) as r:
                    data = json.loads(r.read())
                return [
                    {"title": x.get("title", ""), "url": x.get("url", ""), "snippet": _strip_html(x.get("content", ""))}
                    for x in data.get("results", [])[:count]
                ]
            if serpapi_key:
                url = "https://serpapi.com/search.json?engine=google&q=" + urllib.parse.quote(query) + "&api_key=" + urllib.parse.quote(serpapi_key)
                with urllib.request.urlopen(url, timeout=20) as r:
                    data = json.loads(r.read())
                return [
                    {"title": x.get("title", ""), "url": x.get("link", ""), "snippet": _strip_html(x.get("snippet", ""))}
                    for x in data.get("organic_results", [])[:count]
                ]

            url = "https://duckduckgo.com/html/?q=" + urllib.parse.quote(query)
            req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
            with urllib.request.urlopen(req, timeout=20) as r:
                html = r.read().decode("utf-8", "replace")
            results = []
            pattern = re.compile(
                r'class="result__a" href="(?P<href>[^"]+)".*?>(?P<title>.*?)</a>.*?class="result__snippet".*?>(?P<snippet>.*?)</a>',
                re.S,
            )
            for m in pattern.finditer(html):
                href = m.group("href")
                parsed = urllib.parse.urlparse(href)
                qs = urllib.parse.parse_qs(parsed.query)
                if "uddg" in qs:
                    href = qs["uddg"][0]
                results.append({
                    "title": _strip_html(m.group("title")),
                    "url": href,
                    "snippet": _strip_html(m.group("snippet")),
                })
                if len(results) >= count:
                    break
            return results
        except Exception as e:
            logger.warning("[WEB_SEARCH] %s", e)
            return [{"title": "Web search failed", "url": "", "snippet": str(e)}]

    return await asyncio.get_running_loop().run_in_executor(_executor, _call)

async def _http_post_json(url: str, body: Dict, headers: Dict) -> Any:
    def _call():
        data = json.dumps(body).encode()
        req = urllib.request.Request(url, data=data, headers=headers, method="POST")
        req.add_header("Content-Type","application/json")
        try:
            with urllib.request.urlopen(req, timeout=20) as r:
                raw = r.read()
                if raw: return json.loads(raw)
                return {"status":"ok"}
        except urllib.error.HTTPError as e:
            raw = e.read()
            try: return json.loads(raw)
            except Exception: return {"error":f"HTTP {e.code}","detail":raw.decode("utf-8","replace")[:300]}
    return await asyncio.get_running_loop().run_in_executor(_executor, _call)

async def _http_request_json(method: str, url: str, body: Any = None, headers: Dict = None) -> Any:
    def _call():
        data = None if body is None else json.dumps(body).encode()
        req_headers = dict(headers or {})
        if data is not None:
            req_headers.setdefault("Content-Type", "application/json")
        req = urllib.request.Request(url, data=data, headers=req_headers, method=(method or "GET").upper())
        try:
            with urllib.request.urlopen(req, timeout=30) as r:
                raw = r.read().decode("utf-8", "replace")
                if not raw:
                    return {"status":"ok", "status_code":r.status}
                try:
                    return json.loads(raw)
                except Exception:
                    return {"status":"ok", "status_code":r.status, "text":raw[:4000]}
        except urllib.error.HTTPError as e:
            raw = e.read().decode("utf-8", "replace")
            try:
                detail = json.loads(raw)
            except Exception:
                detail = raw[:1000]
            return {"error":f"HTTP {e.code}", "detail":detail}
    return await asyncio.get_running_loop().run_in_executor(_executor, _call)

async def _http_delete(url: str, headers: Dict) -> None:
    def _call():
        req = urllib.request.Request(url, headers=headers, method="DELETE")
        try: urllib.request.urlopen(req, timeout=20)
        except urllib.error.HTTPError: pass
    await asyncio.get_running_loop().run_in_executor(_executor, _call)

async def _http_patch_json(url: str, body: Dict, headers: Dict) -> Any:
    def _call():
        data = json.dumps(body).encode()
        req = urllib.request.Request(url, data=data,
            headers={**headers, "Content-Type":"application/json"}, method="PATCH")
        try:
            with urllib.request.urlopen(req, timeout=20) as r:
                raw = r.read()
                return json.loads(raw) if raw else {"status":"ok"}
        except urllib.error.HTTPError as e:
            raw = e.read()
            try: return json.loads(raw)
            except Exception: return {"error":f"HTTP {e.code}","detail":raw.decode("utf-8","replace")[:300]}
    return await asyncio.get_running_loop().run_in_executor(_executor, _call)

# ══════════════════════════════════════════════════════════════════════════════
# §12  AGENT / TOOL SYSTEM
# ══════════════════════════════════════════════════════════════════════════════

TOOL_REGISTRY = {
    "rag_search":       {"desc":"Search user's documents","params":["query"]},
    "code_exec":        {"desc":"Execute Python/JS/bash code","params":["code","language"]},
    "shell":            {"desc":"Run shell command","params":["command"]},
    "connector_action": {"desc":"Call a connected app","params":["connector","action","params"]},
    "file_read":        {"desc":"Read a file","params":["path"]},
    "file_write":       {"desc":"Write a file","params":["path","content"]},
    "file_list":        {"desc":"List directory","params":["path"]},
    "system_info":      {"desc":"Get system information","params":[]},
    "web_fetch":        {"desc":"Fetch a URL","params":["url"]},
}

def _run_shell_sync(command: str, timeout: int = CODE_EXEC_TIMEOUT, cwd: str = None) -> Dict:
    t0 = time.time()
    try:
        r = subprocess.run(command, shell=True, capture_output=True, text=True,
                           timeout=timeout, cwd=cwd or os.getcwd())
        return {"stdout":r.stdout[:CODE_EXEC_MAX_OUT],"stderr":r.stderr[:CODE_EXEC_MAX_OUT],
                "exit_code":r.returncode,"duration_ms":int((time.time()-t0)*1000)}
    except subprocess.TimeoutExpired:
        return {"stdout":"","stderr":f"Timeout ({timeout}s)","exit_code":-1,"duration_ms":timeout*1000}
    except Exception as e:
        return {"stdout":"","stderr":str(e),"exit_code":-1,"duration_ms":0}

def _run_code_sync(code: str, language: str = "python", cwd: str = None) -> Dict:
    cmd_map = {
        "python":     [sys.executable,"-c",code],
        "javascript": ["node","-e",code],
        "bash":       ["bash","-c",code],
        "shell":      ["bash","-c",code],
    }
    cmd = cmd_map.get(language)
    if not cmd: return _run_shell_sync(code, cwd=cwd)
    t0 = time.time()
    try:
        r = subprocess.run(cmd, capture_output=True, text=True, timeout=CODE_EXEC_TIMEOUT,
                           cwd=cwd or os.getcwd())
        return {"stdout":r.stdout[:CODE_EXEC_MAX_OUT],"stderr":r.stderr[:CODE_EXEC_MAX_OUT],
                "exit_code":r.returncode,"duration_ms":int((time.time()-t0)*1000)}
    except subprocess.TimeoutExpired:
        return {"stdout":"","stderr":"Timeout","exit_code":-1,"duration_ms":CODE_EXEC_TIMEOUT*1000}
    except FileNotFoundError as e:
        return {"stdout":"","stderr":f"Runtime not found: {e}","exit_code":-1,"duration_ms":0}

async def _run_tool(tool: str, params: Dict, user: Dict) -> str:
    uid = user.get("id") or user.get("sub","")
    try:
        if tool == "rag_search":
            chunks = await _rag_search(uid, params.get("query",""))
            return "\n\n".join(chunks) if chunks else "No relevant documents found."
        if tool == "code_exec":
            result = await asyncio.get_running_loop().run_in_executor(
                _executor, lambda: _run_code_sync(params.get("code",""), params.get("language","python")))
            return f"stdout: {result['stdout']}\nstderr: {result['stderr']}\nexit: {result['exit_code']}"
        if tool == "shell":
            result = await asyncio.get_running_loop().run_in_executor(
                _executor, lambda: _run_shell_sync(params.get("command","")))
            return result["stdout"][-3000:] + (result["stderr"][-500:] if result["stderr"] else "")
        if tool == "connector_action":
            result = await _connector_api_call(uid, params.get("connector",""), params.get("action",""), params.get("params",{}))
            return json.dumps(result, indent=2)[:3000]
        if tool == "file_read":
            p = Path(params.get("path",".")).expanduser().resolve()
            if p.exists() and p.is_file(): return p.read_text(encoding="utf-8", errors="replace")[:5000]
            return "File not found"
        if tool == "file_write":
            p = Path(params.get("path",".")).expanduser().resolve()
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(params.get("content",""), encoding="utf-8")
            return f"Written {len(params.get('content',''))} chars to {p}"
        if tool == "file_list":
            p = Path(params.get("path",".")).expanduser().resolve()
            if p.is_dir():
                items = [{"name":f.name,"type":"dir" if f.is_dir() else "file",
                           "size":f.stat().st_size if f.is_file() else None} for f in p.iterdir()]
                return json.dumps(sorted(items, key=lambda x:(x["type"]!="dir",x["name"])), indent=2)
            return "Not a directory"
        if tool == "system_info":
            return json.dumps({"os":platform.system(),"version":platform.version(),
                               "python":sys.version,"cwd":os.getcwd(),"hostname":platform.node()})
        if tool == "web_fetch":
            url = params.get("url","")
            result = await _http_get(url, {"User-Agent":"JAZZ-AI/14.0"})
            return json.dumps(result, indent=2)[:3000]
        return f"Tool '{tool}' not available"
    except Exception as e:
        return f"Tool error: {e}"

_AGENT_SYSTEM = """You are a reasoning agent with FULL system access. Work step by step.
Available tools: code_exec(code,language), shell(command), file_read(path), file_write(path,content), file_list(path), rag_search(query), connector_action(connector,action,params), system_info(), web_fetch(url)

Each step respond ONLY with JSON:
{"thought":"what I'm thinking","action":"tool_name|FINISH","args":{},"final_answer":"only when FINISH"}
Max 12 steps."""

async def _agent_loop(prompt: str, user: Dict, tools: List[str],
                      model_id: str, max_steps: int = 12,
                      tool_context: str = "") -> Tuple[str, List[Dict]]:
    client, model_name = await _get_model_client(model_id)
    system = _AGENT_SYSTEM + ("\n\n" + tool_context if tool_context else "")
    messages = [{"role":"system","content":system},{"role":"user","content":prompt}]
    tool_calls = []
    for step in range(max_steps):
        try:
            resp = await asyncio.get_running_loop().run_in_executor(
                _executor, lambda c=client, m=model_name: c.chat.completions.create(
                    model=m, messages=messages, max_tokens=800, temperature=0.2))
            raw = resp.choices[0].message.content.strip()
            clean = re.sub(r"```json|```","",raw).strip()
            step_json = json.loads(clean)
        except Exception as e:
            logger.warning("[AGENT] step %d parse error: %s", step, e); break

        thought = step_json.get("thought","")
        action  = step_json.get("action","FINISH")
        args    = step_json.get("args",{})
        messages.append({"role":"assistant","content":raw})

        if action == "FINISH":
            return step_json.get("final_answer","Task complete."), tool_calls

        t0 = time.time()
        obs = await _run_tool(action, args, user)
        tool_calls.append({"step":step,"tool":action,"args":args,
                           "success":not obs.startswith("Tool error"),
                           "latency_ms":int((time.time()-t0)*1000)})
        messages.append({"role":"user","content":f"[Observation]: {obs}"})

    return "Agent reached max steps.", tool_calls

# ══════════════════════════════════════════════════════════════════════════════
# §13  SCHEDULER
# ══════════════════════════════════════════════════════════════════════════════

_scheduler = BackgroundScheduler(timezone="UTC")

async def _run_agent_job(job_id: str):
    job = await db_fetchone("SELECT * FROM agent_jobs WHERE id=?", (job_id,))
    if not job or job["status"] == "running": return
    user = await db_fetchone("SELECT * FROM users WHERE id=?", (job["user_id"],))
    if not user: return
    user["sub"] = user["id"]
    log_id = _new_id(); run_num = (job["total_runs"] or 0)+1; now_iso = _utcnow()
    await db_execute("UPDATE agent_jobs SET status='running' WHERE id=?", (job_id,))
    await db_execute(
        "INSERT INTO agent_job_logs(id,job_id,user_id,run_number,status,trigger_type,prompt_rendered,started_at)"
        " VALUES(?,?,?,?,'running','cron',?,?)",
        (log_id, job_id, job["user_id"], run_num, job["prompt_template"], now_iso))
    tools = json.loads(job["tools_json"] or "[]")
    result_text = None; tool_log = []; run_status = "failure"; error_msg = None
    t0 = time.time()
    for attempt in range(job.get("max_retries",2)+1):
        try:
            result_text, tool_log = await _agent_loop(
                job["prompt_template"], user, tools, job.get("model_id","llama-3.3-70b-versatile"),
                max_steps=12)
            run_status = "success"; break
        except Exception as e:
            error_msg = str(e)
            await asyncio.sleep(min(2**attempt, 30))
    elapsed = int((time.time()-t0)*1000); now = _utcnow()
    await db_execute(
        "UPDATE agent_job_logs SET status=?,result_text=?,tool_calls_json=?,latency_ms=?,error_message=?,finished_at=? WHERE id=?",
        (run_status, result_text, json.dumps(tool_log), elapsed, error_msg, now, log_id))
    col = "success_runs=success_runs+1" if run_status == "success" else "failed_runs=failed_runs+1"
    await db_execute(
        f"UPDATE agent_jobs SET status='idle',total_runs=total_runs+1,{col},"
        "last_run_at=?,last_run_status=?,retry_count=0,updated_at=? WHERE id=?",
        (now, run_status, now, job_id))

def _sync_run_job(job_id: str) -> None:
    try:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        try: loop.run_until_complete(_run_agent_job(job_id))
        finally: loop.close()
    except Exception as e:
        logger.error("[SCHEDULER] Job %s failed: %s", job_id, e)

def _schedule_job(job: Dict):
    jid = job["id"]; jtype = job.get("job_type","cron")
    trigger_cfg = json.loads(job.get("trigger_json") or "{}")
    try:
        if jtype == "cron":
            expr = trigger_cfg.get("cron","0 9 * * *")
            parts = expr.split()
            if len(parts) >= 5:
                trig = CronTrigger(minute=parts[0],hour=parts[1],day=parts[2],month=parts[3],day_of_week=parts[4])
            else: trig = IntervalTrigger(hours=24)
        elif jtype == "manual": return
        else:
            trig = IntervalTrigger(seconds=trigger_cfg.get("interval_seconds",3600))
        _scheduler.add_job(_sync_run_job, args=[jid], trigger=trig, id=jid, replace_existing=True)
    except Exception as e:
        logger.warning("[SCHEDULER] Failed to schedule job %s: %s", jid, e)

# ══════════════════════════════════════════════════════════════════════════════
# §14  SLASH COMMANDS
# ══════════════════════════════════════════════════════════════════════════════

_DEFAULT_SLASH_COMMANDS = [
    # ── Communication ──────────────────────────────────────────────────────────
    ("gmail",    "gmail",           "Access Gmail inbox",            "list_emails",   {}),
    ("mail",     "gmail",           "Alias for /gmail",             "list_emails",   {}),
    ("calendar", "google_calendar", "Manage Google Calendar",        "list_events",   {}),
    ("cal",      "google_calendar", "Alias for /calendar",          "list_events",   {}),
    ("slack",    "slack",           "Send/read Slack messages",      "list_channels", {}),
    # ── Storage & Docs ─────────────────────────────────────────────────────────
    ("drive",    "google_drive",    "Search Google Drive",           "list_files",    {}),
    ("sheets",   "google_sheets",   "Google Sheets",                 "list_files",    {}),
    # ── GitHub (full) ──────────────────────────────────────────────────────────
    ("github",   "github",          "GitHub repos, PRs, issues, code","list_repos",  {}),
    ("gh",       "github",          "Alias for /github",            "list_repos",    {}),
    ("repos",    "github",          "List your GitHub repos",        "list_repos",    {}),
    ("prs",      "github",          "List pull requests",            "list_prs",      {}),
    ("issues",   "github",          "List GitHub issues",            "list_issues",   {}),
    ("commits",  "github",          "List recent commits",           "list_commits",  {}),
    # ── Project Mgmt ───────────────────────────────────────────────────────────
    ("notion",   "notion",          "Search Notion pages",          "search",        {}),
    ("jira",     "jira",            "Manage Jira issues",           "list_issues",   {}),
    ("asana",    "asana",           "Manage Asana tasks",           "list_tasks",    {}),
    ("linear",   "linear",          "Manage Linear issues",         "list_issues",   {}),
    # ── CRM / Business ─────────────────────────────────────────────────────────
    ("hubspot",  "hubspot",         "CRM contacts and deals",        "list_contacts", {}),
    ("airtable", "airtable",        "Airtable bases",               "list_bases",    {}),
    ("zoom",     "zoom",            "Create Zoom meetings",          "list_meetings", {}),
    ("discord",  "discord",         "Discord channels and messages", "list_channels", {}),
    ("teams",    "microsoft_teams", "Microsoft Teams channels",      "list_channels", {}),
    ("outlook",  "outlook",         "Outlook mail",                  "list_emails",   {}),
    ("onedrive", "onedrive",        "OneDrive files",                "list_files",    {}),
    ("excel",    "excel",           "Microsoft Excel files",         "list_sheets",   {}),
    ("gitlab",   "gitlab",          "GitLab projects and issues",    "list_projects", {}),
    ("trello",   "trello",          "Trello boards and cards",       "list_boards",   {}),
    ("monday",   "monday",          "Monday.com boards",             "list_boards",   {}),
    ("clickup",  "clickup",         "ClickUp tasks",                 "list_tasks",    {}),
    ("salesforce","salesforce",     "Salesforce records",            "query",         {}),
    ("stripe",   "stripe",          "Stripe payments",               "list_charges",  {}),
    ("powerbi",  "power_bi",        "Power BI datasets",             "list_datasets", {}),
    ("bigquery", "bigquery",        "BigQuery datasets",             "list_datasets", {}),
    ("dropbox",  "dropbox",         "Dropbox files",                 "list_files",    {}),
    ("box",      "box",             "Box files",                     "list_files",    {}),
    ("s3",       "s3",              "Amazon S3 buckets",             "list_buckets",  {}),
    ("meet",     "google_meet",     "Google Meet meetings",          "list_meetings", {}),
    ("zapier",   "zapier",          "Trigger Zapier webhooks",       "trigger_zap",   {}),
    ("make",     "make",            "Trigger Make scenarios",        "trigger_scenario", {}),
    ("n8n",      "n8n",             "Trigger n8n Cloud workflows",   "trigger_workflow", {}),
    ("n8n_local","n8n_local",       "Open or check local n8n",       "status",        {}),
    ("indeed",   "indeed",          "Indeed jobs",                   "search_jobs",   {}),
    ("linkedin", "linkedin",        "LinkedIn jobs and profile",     "search_jobs",   {}),
    ("shopify",  "shopify",         "Shopify products and orders",   "list_products", {}),
    ("woocommerce","woocommerce",   "WooCommerce products",          "list_products", {}),
    ("http",     "http",            "Call a configured HTTP API",    "GET",           {}),
    ("graphql",  "graphql",         "Run a GraphQL query",           "query",         {}),
    ("webhook",  "webhook",         "Send a webhook payload",        "send",          {}),
    # ── Web ────────────────────────────────────────────────────────────────────
    ("web",      "web_search",      "Search the live web",          "search",        {}),
    ("search",   "web_search",      "Search the live web",          "search",        {}),
]

async def _seed_slash_commands():
    for cmd, ctype, desc, action, params in _DEFAULT_SLASH_COMMANDS:
        exists = await db_fetchone("SELECT id FROM slash_commands WHERE command=?", (cmd,))
        if not exists:
            await db_execute(
                "INSERT INTO slash_commands(id,command,connector_type,description,default_action,default_params,created_at)"
                " VALUES(?,?,?,?,?,?,?)",
                (_new_id(), cmd, ctype, desc, action, json.dumps(params), _utcnow()))

# ── Platform Connector Catalog (mirrors frontend CONNECTOR_CATALOG) ────────────
_PLATFORM_CONNECTOR_CATALOG = [
    # Communication
    ("gmail",           "Gmail",            "Communication", "📧", 1, ["GOOGLE_CLIENT_ID","GOOGLE_CLIENT_SECRET"]),
    ("slack",           "Slack",            "Communication", "💬", 1, ["SLACK_CLIENT_ID","SLACK_CLIENT_SECRET"]),
    ("discord",         "Discord",          "Communication", "🎮", 0, ["DISCORD_CLIENT_ID","DISCORD_CLIENT_SECRET"]),
    ("microsoft_teams", "Microsoft Teams",  "Communication", "🔷", 1, ["MICROSOFT_CLIENT_ID","MICROSOFT_CLIENT_SECRET"]),
    # Productivity
    ("google_calendar", "Google Calendar",  "Productivity",  "📅", 1, ["GOOGLE_CLIENT_ID","GOOGLE_CLIENT_SECRET"]),
    ("google_drive",    "Google Drive",     "Productivity",  "💾", 1, ["GOOGLE_CLIENT_ID","GOOGLE_CLIENT_SECRET"]),
    ("google_sheets",   "Google Sheets",    "Productivity",  "📊", 1, ["GOOGLE_CLIENT_ID","GOOGLE_CLIENT_SECRET"]),
    ("notion",          "Notion",           "Productivity",  "📝", 1, ["NOTION_CLIENT_ID","NOTION_CLIENT_SECRET"]),
    ("excel",           "Microsoft Excel",  "Productivity",  "📗", 0, []),
    ("outlook",         "Outlook",          "Productivity",  "📨", 1, ["MICROSOFT_CLIENT_ID","MICROSOFT_CLIENT_SECRET"]),
    ("onedrive",        "OneDrive",         "Productivity",  "☁️", 1, ["MICROSOFT_CLIENT_ID","MICROSOFT_CLIENT_SECRET"]),
    # Developer
    ("github",          "GitHub",           "Developer",     "🐙", 1, ["GITHUB_CLIENT_ID","GITHUB_CLIENT_SECRET"]),
    ("gitlab",          "GitLab",           "Developer",     "🦊", 1, ["GITLAB_CLIENT_ID","GITLAB_CLIENT_SECRET"]),
    ("jira",            "Jira",             "Developer",     "🔵", 1, ["ATLASSIAN_CLIENT_ID","ATLASSIAN_CLIENT_SECRET"]),
    ("linear",          "Linear",           "Developer",     "📐", 1, ["LINEAR_CLIENT_ID","LINEAR_CLIENT_SECRET"]),
    # Project Mgmt
    ("asana",           "Asana",            "Project Mgmt",  "✅", 1, ["ASANA_CLIENT_ID","ASANA_CLIENT_SECRET"]),
    ("trello",          "Trello",           "Project Mgmt",  "📋", 0, []),
    ("monday",          "Monday.com",       "Project Mgmt",  "📅", 0, []),
    ("clickup",         "ClickUp",          "Project Mgmt",  "🎯", 1, ["CLICKUP_CLIENT_ID","CLICKUP_CLIENT_SECRET"]),
    # CRM / Sales
    ("hubspot",         "HubSpot",          "CRM / Sales",   "🧡", 1, ["HUBSPOT_CLIENT_ID","HUBSPOT_CLIENT_SECRET"]),
    ("salesforce",      "Salesforce",       "CRM / Sales",   "☁️", 1, ["SALESFORCE_CLIENT_ID","SALESFORCE_CLIENT_SECRET"]),
    ("stripe",          "Stripe",           "CRM / Sales",   "💳", 0, ["STRIPE_SECRET_KEY"]),
    # Analytics
    ("power_bi",        "Power BI",         "Analytics",     "📈", 1, ["MICROSOFT_CLIENT_ID","MICROSOFT_CLIENT_SECRET"]),
    ("airtable",        "Airtable",         "Analytics",     "🗂️", 1, ["AIRTABLE_CLIENT_ID","AIRTABLE_CLIENT_SECRET"]),
    ("bigquery",        "BigQuery",         "Analytics",     "🔬", 1, ["GOOGLE_CLIENT_ID","GOOGLE_CLIENT_SECRET"]),
    # Storage
    ("dropbox",         "Dropbox",          "Storage",       "📦", 1, ["DROPBOX_CLIENT_ID","DROPBOX_CLIENT_SECRET"]),
    ("box",             "Box",              "Storage",       "📫", 1, ["BOX_CLIENT_ID","BOX_CLIENT_SECRET"]),
    ("s3",              "Amazon S3",        "Storage",       "🪣", 0, []),
    # Meetings
    ("zoom",            "Zoom",             "Meetings",      "📹", 1, ["ZOOM_CLIENT_ID","ZOOM_CLIENT_SECRET"]),
    ("google_meet",     "Google Meet",      "Meetings",      "📺", 1, ["GOOGLE_CLIENT_ID","GOOGLE_CLIENT_SECRET"]),
    # Automation
    ("zapier",          "Zapier",           "Automation",    "⚡", 0, []),
    ("make",            "Make",             "Automation",    "🔧", 0, []),
    ("n8n",             "n8n Cloud",        "Automation",    "https://i.ibb.co/jPBvQxWj/download-8.png", 0, []),
    ("n8n_local",       "n8n Local",        "Automation",    "https://i.ibb.co/jPBvQxWj/download-8.png", 0, []),
    # Recruiting
    ("indeed",          "Indeed",           "Recruiting",    "💼", 0, []),
    ("linkedin",        "LinkedIn",         "Recruiting",    "🔗", 1, ["LINKEDIN_CLIENT_ID","LINKEDIN_CLIENT_SECRET"]),
    # E-commerce
    ("shopify",         "Shopify",          "E-commerce",    "🛍️", 0, []),
    ("woocommerce",     "WooCommerce",      "E-commerce",    "🛒", 0, []),
]

def _calc_setup_status(env_keys: list) -> str:
    """Check if all required env keys are set."""
    if not env_keys:
        return "ready"
    filled = [k for k in env_keys if os.getenv(k, "").strip()]
    if len(filled) == len(env_keys):
        return "ready"
    if filled:
        return "partial"
    return "not_configured"

async def _seed_platform_connectors():
    now = _utcnow()
    for (ctype, name, cat, icon, requires_oauth, env_keys) in _PLATFORM_CONNECTOR_CATALOG:
        exists = await db_fetchone(
            "SELECT id, setup_status FROM platform_connectors WHERE connector_type=?", (ctype,))
        status = _calc_setup_status(env_keys)
        if not exists:
            await db_execute(
                "INSERT INTO platform_connectors"
                "(id,connector_type,is_enabled,display_name,category,icon,setup_status,"
                "requires_oauth,env_keys,admin_notes,updated_at,created_at)"
                " VALUES(?,?,0,?,?,?,?,?,?,?,?,?)",
                (_new_id(), ctype, name, cat, icon, status,
                 requires_oauth, json.dumps(env_keys), "", now, now))
        else:
            # Refresh catalog metadata and setup status in case env or labels changed.
            await db_execute(
                "UPDATE platform_connectors SET display_name=?,category=?,icon=?,requires_oauth=?,"
                "env_keys=?,setup_status=?,updated_at=? WHERE connector_type=?",
                (name, cat, icon, requires_oauth, json.dumps(env_keys), status, now, ctype))

_SIDEBAR_FEATURE_CATALOG = [
    ("agents",     "Agent Jobs",      "🤖",  "agents",     20),
    ("playground", "Playground",      "⚡",  "playground", 30),
    ("multimodel", "Multi-Model",     "MM",  "multimodel", 40),
    ("websites",   "Website Builder", "🌐",  "websites",   50),
    ("code",       "Code Runner",     "💻",  "code",       60),
]
_WORKSPACE_TOOL_FEATURE_KEYS = {"documents", "memory", "connectors", "mcp", "livy", "apikeys"}
_SIDEBAR_FEATURE_KEYS = {k for k, *_ in _SIDEBAR_FEATURE_CATALOG}

async def _seed_sidebar_features():
    now = _utcnow()
    workspace_keys = tuple(sorted(_WORKSPACE_TOOL_FEATURE_KEYS))
    if workspace_keys:
        placeholders = ",".join("?" for _ in workspace_keys)
        await db_execute(
            f"DELETE FROM sidebar_features WHERE feature_key IN ({placeholders})",
            workspace_keys)
        await db_execute(
            f"DELETE FROM user_sidebar_feature_access WHERE feature_key IN ({placeholders})",
            workspace_keys)
        await db_execute(
            f"DELETE FROM plan_sidebar_feature_access WHERE feature_key IN ({placeholders})",
            workspace_keys)
    for key, name, icon, page_key, sort_order in _SIDEBAR_FEATURE_CATALOG:
        exists = await db_fetchone(
            "SELECT id FROM sidebar_features WHERE feature_key=?", (key,))
        if not exists:
            await db_execute(
                "INSERT INTO sidebar_features"
                "(id,feature_key,display_name,icon,page_key,sort_order,is_enabled,created_at,updated_at)"
                " VALUES(?,?,?,?,?,?,1,?,?)",
                (_new_id(), key, name, icon, page_key, sort_order, now, now))
        else:
            await db_execute(
                "UPDATE sidebar_features SET display_name=?,icon=?,page_key=?,sort_order=?,updated_at=?"
                " WHERE feature_key=?",
                (name, icon, page_key, sort_order, now, key))

async def _sidebar_feature_key_set() -> Set[str]:
    await _seed_sidebar_features()
    rows = await db_fetchall("SELECT feature_key FROM sidebar_features")
    keys = {str(r.get("feature_key") or "").strip() for r in rows}
    keys = {k for k in keys if k not in _WORKSPACE_TOOL_FEATURE_KEYS}
    return {k for k in keys if k} or set(_SIDEBAR_FEATURE_KEYS)

async def _sidebar_feature_rows() -> List[Dict]:
    await _seed_sidebar_features()
    return await db_fetchall(
        "SELECT feature_key,display_name,icon,page_key,sort_order,is_enabled"
        f" FROM sidebar_features WHERE feature_key NOT IN ({','.join('?' for _ in _WORKSPACE_TOOL_FEATURE_KEYS)})"
        " ORDER BY sort_order,display_name",
        tuple(sorted(_WORKSPACE_TOOL_FEATURE_KEYS)))

async def _sidebar_feature_access_for_user(user_id: str) -> Dict[str, Any]:
    rows = await _sidebar_feature_rows()
    user_row = await db_fetchone("SELECT subscription FROM users WHERE id=?", (user_id,))
    plan = (user_row or {}).get("subscription") or "free"
    overrides = {
        r["feature_key"]: bool(r["is_enabled"])
        for r in await db_fetchall(
            "SELECT feature_key,is_enabled FROM user_sidebar_feature_access WHERE user_id=?",
            (user_id,))
    }
    plan_access = {
        r["feature_key"]: bool(r["is_enabled"])
        for r in await db_fetchall(
            "SELECT feature_key,is_enabled FROM plan_sidebar_feature_access WHERE plan=?",
            (plan,))
    }
    features = []
    enabled = []
    for row in rows:
        key = row["feature_key"]
        global_enabled = bool(row["is_enabled"])
        plan_enabled = plan_access.get(key, True)
        override = overrides.get(key)
        effective = global_enabled and plan_enabled and (override is not False)
        item = dict(row)
        item["global_enabled"] = global_enabled
        item["plan_enabled"] = plan_enabled
        item["plan"] = plan
        item["user_override"] = override
        item["is_enabled"] = effective
        features.append(item)
        if effective:
            enabled.append(key)
    return {"features": features, "enabled": enabled}

async def _seed_model_access() -> None:
    now = _utcnow()
    for mid, cfg in BUILTIN_MODELS.items():
        exists = await db_fetchone("SELECT id FROM model_access WHERE model_id=?", (mid,))
        if not exists:
            await db_execute(
                "INSERT INTO model_access(id,model_id,display_name,source,is_enabled,created_at,updated_at)"
                " VALUES(?,?,?,?,1,?,?)",
                (_new_id(), mid, cfg.get("label", mid), "builtin", now, now))
        else:
            await db_execute(
                "UPDATE model_access SET display_name=?,source='builtin',updated_at=? WHERE model_id=?",
                (cfg.get("label", mid), now, mid))

    rows = await db_fetchall("SELECT id,name,is_active FROM ai_models")
    for row in rows:
        exists = await db_fetchone("SELECT id FROM model_access WHERE model_id=?", (row["id"],))
        if not exists:
            await db_execute(
                "INSERT INTO model_access(id,model_id,display_name,source,is_enabled,created_at,updated_at)"
                " VALUES(?,?,?,?,?,?,?)",
                (_new_id(), row["id"], row.get("name") or row["id"], "db",
                 1 if row.get("is_active") else 0, now, now))
        else:
            await db_execute(
                "UPDATE model_access SET display_name=?,source='db',updated_at=? WHERE model_id=?",
                (row.get("name") or row["id"], now, row["id"]))

async def _model_access_rows(include_inactive_db: bool = True) -> List[Dict]:
    await _seed_model_access()
    access = {
        r["model_id"]: r
        for r in await db_fetchall("SELECT model_id,display_name,source,is_enabled FROM model_access")
    }
    rows: List[Dict[str, Any]] = []
    for mid, cfg in BUILTIN_MODELS.items():
        a = access.get(mid, {})
        rows.append({
            "id": mid,
            "name": cfg.get("label", mid),
            "provider": cfg.get("provider", "groq"),
            "base_url": _PROVIDER_DEFAULTS.get(cfg.get("provider", "groq"), ""),
            "model_name": mid,
            "is_active": bool(a.get("is_enabled", 1)),
            "access_enabled": bool(a.get("is_enabled", 1)),
            "is_default": mid == "llama-3.3-70b-versatile",
            "is_fast": bool(cfg.get("fast")),
            "is_vision": False,
            "is_code": False,
            "context_length": cfg.get("ctx", 32768),
            "max_output_tokens": 4096,
            "temperature_default": 0.7,
            "description": "Built-in JAZZ model.",
            "tags": [cfg.get("provider", "groq"), "builtin"],
            "protected": bool(cfg.get("protected", False)),
            "source": "builtin",
            "is_builtin": True,
            "created_at": "",
            "updated_at": "",
        })

    where = "" if include_inactive_db else " WHERE is_active=1"
    db_rows = await db_fetchall(
        "SELECT id,name,provider,base_url,model_name,is_active,is_default,is_fast,is_vision,is_code,"
        "context_length,max_output_tokens,temperature_default,description,tags_json,created_at,updated_at "
        f"FROM ai_models{where} ORDER BY created_at DESC")
    for row in db_rows:
        a = access.get(row["id"], {})
        row["tags"] = _safe_json_loads(row.get("tags_json"), [])
        row["protected"] = "protected" in [str(t).lower() for t in row["tags"]]
        row["access_enabled"] = bool(a.get("is_enabled", row.get("is_active", 1)))
        row["source"] = "db"
        row["is_builtin"] = False
        rows.append(row)
    return rows

async def _model_access_for_user(user_id: str) -> Dict[str, Any]:
    rows = await _model_access_rows(include_inactive_db=False)
    user_row = await db_fetchone("SELECT subscription FROM users WHERE id=?", (user_id,))
    plan = (user_row or {}).get("subscription") or "free"
    overrides = {
        r["model_id"]: bool(r["is_enabled"])
        for r in await db_fetchall(
            "SELECT model_id,is_enabled FROM user_model_access WHERE user_id=?",
            (user_id,))
    }
    plan_access = {
        r["model_id"]: bool(r["is_enabled"])
        for r in await db_fetchall(
            "SELECT model_id,is_enabled FROM plan_model_access WHERE plan=?",
            (plan,))
    }
    models = []
    enabled = []
    for row in rows:
        if _model_is_internal(row):
            continue
        mid = row["id"]
        global_enabled = bool(row.get("access_enabled", row.get("is_active", True)))
        plan_enabled = plan_access.get(mid, True)
        override = overrides.get(mid)
        effective = global_enabled and plan_enabled and (override is not False)
        item = dict(row)
        item["global_enabled"] = global_enabled
        item["plan_enabled"] = plan_enabled
        item["plan"] = plan
        item["user_override"] = override
        item["is_enabled"] = effective
        models.append(item)
        if effective:
            enabled.append(mid)
    return {"models": models, "enabled": enabled}

async def _model_enabled_for_user(user_id: Optional[str], model_id: str) -> bool:
    mid = _canonical_model_id(model_id)
    if not user_id:
        row = await db_fetchone("SELECT is_enabled FROM model_access WHERE model_id=?", (mid,))
        return bool(row["is_enabled"]) if row else True
    access = await _model_access_for_user(user_id)
    return mid in set(access.get("enabled", []))

async def _allowed_model_ids_for_user(user_id: Optional[str]) -> List[str]:
    if user_id:
        return [m for m in (await _model_access_for_user(user_id)).get("enabled", [])]
    return [m["id"] for m in await _model_access_rows(include_inactive_db=False) if m.get("access_enabled")]

async def _ensure_user_model(user_id: str, model_id: str) -> str:
    mid = _canonical_model_id(model_id)
    if await _model_enabled_for_user(user_id, mid):
        return mid
    allowed = await _allowed_model_ids_for_user(user_id)
    return allowed[0] if allowed else "llama-3.3-70b-versatile"

_PASTED_CONTENT_BLOCK_RE = re.compile(
    r"<JAZZ_PASTED_CONTENT\b[^>]*>.*?</JAZZ_PASTED_CONTENT>",
    re.IGNORECASE | re.DOTALL,
)
_SLASH_RE = re.compile(r"(?<![\w:/.-])/([a-zA-Z0-9_-]+)\b\s*(.*)?$", re.DOTALL)

def _strip_pasted_content_for_tools(message: str) -> str:
    return _PASTED_CONTENT_BLOCK_RE.sub(" ", message or "").strip()

def _parse_slash(message: str) -> Optional[Tuple[str, str]]:
    m = _SLASH_RE.match(_strip_pasted_content_for_tools(message))
    if m: return m.group(1).lower(), (m.group(2) or "").strip()
    return None

_TOOL_INTENT_RE = re.compile(
    r"\b(use|run|call|open|launch|read|check|list|show|search|find|send|create|trigger|post|summarize|summary|draft|reply)\b",
    re.IGNORECASE)

async def _infer_slash_command(message: str) -> Optional[Tuple[str, str]]:
    """Let natural language such as 'read gmail inbox' use the same connector path."""
    text = _strip_pasted_content_for_tools(message)
    if not text or not _TOOL_INTENT_RE.search(text):
        return None
    low = text.lower()
    rows = await db_fetchall(
        "SELECT command,connector_type FROM slash_commands WHERE is_active=1 ORDER BY LENGTH(command) DESC")
    platform = {
        r["connector_type"]: (r.get("display_name") or "")
        for r in await db_fetchall("SELECT connector_type,display_name FROM platform_connectors")
    }
    for r in rows:
        cmd = str(r["command"]).lower()
        ctype = str(r["connector_type"]).lower()
        display = platform.get(ctype, "").lower()
        terms = {"web", "web search", "internet"} if ctype == "web_search" else {cmd, ctype.replace("_", " ")}
        if display:
            terms.add(display)
        if ctype == "google_calendar":
            terms.update({"calendar", "google calendar"})
        if ctype == "microsoft_teams":
            terms.update({"teams", "microsoft teams"})
        if ctype == "power_bi":
            terms.update({"power bi", "powerbi"})
        if any(re.search(r"\b" + re.escape(term) + r"\b", low) for term in terms if term):
            return cmd, text
    return None

def _extract_first_int(text: str, default: int = 10) -> int:
    m = re.search(r"\b(\d{1,3})\b", text or "")
    if not m:
        return default
    n = int(m.group(1))
    return max(1, min(50, n))

def _apply_slash_rest(ctype: str, action: str, params: Dict, rest: str) -> Tuple[str, Dict]:
    rest_l = (rest or "").strip().lower()
    out = dict(params or {})
    if not rest_l:
        return action, out

    # ── Gmail ──────────────────────────────────────────────────────────────────
    if ctype == "gmail":
        out["count"] = _extract_first_int(rest_l, int(out.get("count", 10) or 10))
        if any(k in rest_l for k in ("reply", "respond")):
            # /gmail reply <message_id> <body text>
            action = "reply_email"
            reply_m = re.match(r"(?:reply|respond)\s+([a-f0-9]{6,})\s*(.*)", rest, flags=re.IGNORECASE | re.DOTALL)
            if reply_m:
                out["message_id"] = reply_m.group(1).strip()
                out["body"] = reply_m.group(2).strip()
        elif any(k in rest_l for k in ("search", "find")):
            action = "search_emails"
            q = re.sub(r"^\s*(search|find)\s+", "", rest, flags=re.IGNORECASE)
            out["query"] = q.strip() or "in:inbox"
        elif any(k in rest_l for k in ("read", "open", "show")):
            mid_match = re.search(r"\b([a-f0-9]{10,})\b", rest_l)
            if mid_match:
                action = "read_email"; out["message_id"] = mid_match.group(1)
            else:
                action = "list_emails"; out["include_body"] = True
        elif any(k in rest_l for k in ("send", "compose", "write")):
            action = "send_email"
            # handles "to: x@y", "to x@y", "to : x@y" (spaces around colon)
            to_m = re.search(r"\bto(?:\s*:\s*|\s+)([^\s,;]+@[^\s,;]+)", rest, flags=re.IGNORECASE)
            if to_m:
                out["to"] = to_m.group(1).strip().rstrip(",;")
            # subject — stops before body/message/saying/attach keyword
            subj_m = re.search(
                r"\bsubject(?:\s*:\s*|\s+)(.+?)(?=\s+(?:body|message|saying|attach)(?:\s*:\s*|\s+)|\s*$)",
                rest, flags=re.IGNORECASE | re.DOTALL)
            if subj_m:
                out["subject"] = subj_m.group(1).strip()
            # explicit body: wins over natural "saying X"
            body_explicit = re.search(
                r"\b(?:body|message)(?:\s*:\s*)(.+?)(?=\s+subject(?:\s*:\s*|\s+)|\s+attach(?:ment)?(?:\s*:\s*|\s+)|\s*$)",
                rest, flags=re.IGNORECASE | re.DOTALL)
            body_saying = re.search(
                r"\bsaying(?:\s+)(.+?)(?=\s+subject(?:\s*:\s*|\s+)|\s+(?:body|message)(?:\s*:\s*|\s+)|\s+attach(?:ment)?(?:\s*:\s*|\s+)|\s*$)",
                rest, flags=re.IGNORECASE | re.DOTALL)
            body_m = body_explicit or body_saying
            if body_m:
                out["body"] = body_m.group(1).strip()
            # optional attachment path — "attach: /path/file.pdf"
            attach_m = re.search(r"\battach(?:ment)?(?:\s*:\s*|\s+)([^\s]+)", rest, flags=re.IGNORECASE)
            if attach_m:
                out["attachment_path"] = attach_m.group(1).strip()
        else:
            action = "list_emails"
            out["count"] = _extract_first_int(rest_l, int(out.get("count", 10) or 10))
            if "inbox" in rest_l: out["query"] = "in:inbox"
            elif "unread" in rest_l: out["query"] = "is:unread"
            elif "sent" in rest_l: out["query"] = "in:sent"

    # ── Google Calendar ─────────────────────────────────────────────────────────
    elif ctype == "google_calendar":
        if any(k in rest_l for k in ("create", "add", "new", "schedule")):
            action = "create_event"
            out["summary"] = re.sub(r"^\s*(create|add|new|schedule)\s+", "", rest, flags=re.IGNORECASE).strip()
        elif any(k in rest_l for k in ("update", "edit", "change", "move", "reschedule")):
            action = "update_event"
            eid_m = re.search(r"\b([a-zA-Z0-9_]{20,})\b", rest)
            if eid_m: out["event_id"] = eid_m.group(1)
        elif any(k in rest_l for k in ("delete", "remove", "cancel")):
            action = "delete_event"
            eid_m = re.search(r"\b([a-zA-Z0-9_]{20,})\b", rest)
            if eid_m: out["event_id"] = eid_m.group(1)
        elif any(k in rest_l for k in ("get", "show", "view", "open")):
            action = "get_event"
            eid_m = re.search(r"\b([a-zA-Z0-9_]{20,})\b", rest)
            if eid_m: out["event_id"] = eid_m.group(1)
        else:
            action = "list_events"
            out["count"] = _extract_first_int(rest_l, 10)
            if rest_l.strip(): out["query"] = rest.strip()

    # ── Slack ───────────────────────────────────────────────────────────────────
    elif ctype == "slack":
        if any(k in rest_l for k in ("send", "message", "post", "say")):
            action = "send_message"
            ch_m = re.search(r"#?([\w-]+)", rest)
            if ch_m:
                out["channel"] = "#" + ch_m.group(1)
                remaining = rest[ch_m.end():].strip()
                if remaining:
                    out["text"] = remaining
            else:
                body_m = re.search(r'([^" ]+)', rest)
                if body_m: out["text"] = body_m.group(1)
        elif any(k in rest_l for k in ("history", "read", "messages")):
            action = "get_messages"
            ch_m = re.search(r"#([\w-]+)", rest_l)
            if ch_m: out["channel"] = ch_m.group(1)
        else:
            action = "list_channels"

    # ── GitHub (natural language routing) ────────────────────────────────────────
    elif ctype == "github":
        # Extract repo pattern: owner/repo or just repo name
        repo_m = re.search(r"\b([\w.-]+/[\w.-]+)\b", rest)
        if repo_m: out["repo"] = repo_m.group(1)

        # PR actions
        if any(k in rest_l for k in ("pr", "pull request", "pull-request", "merge request")):
            if any(k in rest_l for k in ("create", "new", "open")):
                action = "create_pr"
            elif any(k in rest_l for k in ("merge",)):
                action = "merge_pr"
                num_m = re.search(r"#?(\d+)", rest_l)
                if num_m: out["number"] = int(num_m.group(1))
            elif any(k in rest_l for k in ("review", "approve")):
                action = "pr_reviews"
                num_m = re.search(r"#?(\d+)", rest_l)
                if num_m: out["number"] = int(num_m.group(1))
            elif any(k in rest_l for k in ("file", "diff", "change")):
                action = "pr_files"
                num_m = re.search(r"#?(\d+)", rest_l)
                if num_m: out["number"] = int(num_m.group(1))
            else:
                action = "list_prs"
                if "closed" in rest_l: out["state"] = "closed"
                elif "all" in rest_l: out["state"] = "all"
        # Issue actions
        elif any(k in rest_l for k in ("issue", "bug", "task")):
            if any(k in rest_l for k in ("create", "new", "open", "add")):
                action = "create_issue"
                title_m = re.search(r'([^" ]+)', rest)
                if title_m: out["title"] = title_m.group(1)
            elif any(k in rest_l for k in ("close",)):
                action = "close_issue"
                num_m = re.search(r"#?(\d+)", rest_l)
                if num_m: out["number"] = int(num_m.group(1))
            else:
                action = "list_issues"
                if "closed" in rest_l: out["state"] = "closed"
        # Commit actions
        elif any(k in rest_l for k in ("commit", "history", "log")):
            action = "list_commits"
            out["per_page"] = _extract_first_int(rest_l, 15)
        # File/code actions
        elif any(k in rest_l for k in ("file", "content", "read", "show", "get")):
            action = "get_file"
            path_m = re.search(r'([^" ]+)', rest)
            if path_m: out["path"] = path_m.group(1)
        # Search actions
        elif any(k in rest_l for k in ("search", "find", "look")):
            q = re.sub(r"^\s*(search|find|look for?)\s+", "", rest, flags=re.IGNORECASE).strip()
            if any(k in rest_l for k in ("code", "function", "class")):
                action = "search_code"; out["query"] = q
            elif any(k in rest_l for k in ("repo", "repos", "repository")):
                action = "search_repos"; out["query"] = q
            else:
                action = "search_issues"; out["query"] = q
        # Workflow/Actions
        elif any(k in rest_l for k in ("workflow", "action", "ci", "pipeline")):
            action = "list_workflow_runs"
        # Release
        elif any(k in rest_l for k in ("release", "tag", "version")):
            action = "list_releases"
        # Notification
        elif any(k in rest_l for k in ("notification", "inbox", "mention")):
            action = "list_notifications"
        # Default: list repos
        else:
            action = "list_repos"
            out["per_page"] = _extract_first_int(rest_l, 20)

    # ── Notion ──────────────────────────────────────────────────────────────────
    elif ctype == "notion":
        if any(k in rest_l for k in ("create", "new", "add")):
            action = "create_page"
            out["title"] = re.sub(r"^\s*(create|new|add)\s+", "", rest, flags=re.IGNORECASE).strip()
        else:
            action = "search"
            out["query"] = rest.strip()

    # ── Jira ────────────────────────────────────────────────────────────────────
    elif ctype == "jira":
        if any(k in rest_l for k in ("create", "new", "add")):
            action = "create_issue"
            out["title"] = re.sub(r"^\s*(create|new|add)\s+", "", rest, flags=re.IGNORECASE).strip()
        else:
            action = "list_issues"

    # ── Asana ───────────────────────────────────────────────────────────────────
    elif ctype == "asana":
        if any(k in rest_l for k in ("create", "new", "add")):
            action = "create_task"
            out["name"] = re.sub(r"^\s*(create|new|add)\s+", "", rest, flags=re.IGNORECASE).strip()
        else:
            action = "list_tasks"

    elif ctype in ("zapier", "make", "n8n", "webhook"):
        if rest.strip():
            out["payload"] = {"text": rest.strip()}

    elif ctype == "n8n_local":
        action = "open" if any(k in rest_l for k in ("open", "launch", "portal")) else "status"

    elif ctype == "http":
        method_m = re.search(r"\b(GET|POST|PUT|PATCH|DELETE)\b", rest, flags=re.IGNORECASE)
        if method_m:
            action = method_m.group(1).upper()
        url_m = re.search(r"https?://\S+", rest)
        if url_m:
            out["url"] = url_m.group(0).rstrip(".,)")
        if rest.strip() and action in ("POST", "PUT", "PATCH"):
            out["payload"] = {"text": rest.strip()}

    elif ctype == "graphql":
        if rest.strip():
            out["query"] = rest.strip()

    return action, out

async def _execute_slash(user_id: str, command: str, rest: str) -> Optional[Dict]:
    sc_cmd = await db_fetchone("SELECT * FROM slash_commands WHERE command=? AND is_active=1", (command,))
    if not sc_cmd: return None
    ctype = sc_cmd["connector_type"]
    action = sc_cmd["default_action"]
    params: Dict = json.loads(sc_cmd.get("default_params") or "{}")
    if ctype == "web_search":
        result = await _web_search(rest or command, _extract_first_int(rest, 5))
        return {"type":"tool","app":ctype,"action":action,"result":result,"command":f"/{command}"}
    action, params = _apply_slash_rest(ctype, action, params, rest)
    result = await _connector_api_call(user_id, ctype, action, params)
    out = {"type":"tool","app":ctype,"action":action,"result":result,"command":f"/{command}"}
    if isinstance(result, dict) and result.get("error"):
        out["error"] = result.get("error")
    return out

# ══════════════════════════════════════════════════════════════════════════════
# §15  VOICE (LiveKit + Whisper STT + PlayAI TTS)
# ══════════════════════════════════════════════════════════════════════════════

def _lk_available() -> bool:
    return bool(LIVEKIT_API_KEY and LIVEKIT_API_SECRET and LIVEKIT_URL)

def _lk_token(identity: str, name: str, room: str, ttl_s: int = 3600) -> str:
    if HAS_LIVEKIT:
        try:
            token = (livekit_api.AccessToken(LIVEKIT_API_KEY, LIVEKIT_API_SECRET)
                     .with_identity(identity).with_name(name)
                     .with_grants(livekit_api.VideoGrants(
                         room_join=True, room=room,
                         can_publish=True, can_subscribe=True, can_publish_data=True))
                     .to_jwt())
            return token
        except Exception: pass
    # Fallback: manual JWT
    import hmac as _hmac
    now = int(time.time())
    payload = {"iss":LIVEKIT_API_KEY,"sub":identity,"iat":now,"nbf":now,"exp":now+ttl_s,
               "video":{"roomJoin":True,"room":room,"canPublish":True,"canSubscribe":True}}
    header = base64.urlsafe_b64encode(json.dumps({"alg":"HS256","typ":"JWT"}).encode()).rstrip(b"=").decode()
    body_b = base64.urlsafe_b64encode(json.dumps(payload).encode()).rstrip(b"=").decode()
    sig_in = f"{header}.{body_b}".encode()
    sig = _hmac.new(LIVEKIT_API_SECRET.encode(), sig_in, hashlib.sha256).digest()
    sig_b64 = base64.urlsafe_b64encode(sig).rstrip(b"=").decode()
    return f"{header}.{body_b}.{sig_b64}"

async def _stt_transcribe(audio_bytes: bytes, mime_type: str = "audio/webm") -> str:
    ext_map = {"audio/webm":"webm","audio/mp4":"mp4","audio/wav":"wav",
               "audio/ogg":"ogg","audio/mpeg":"mp3","audio/m4a":"m4a","video/webm":"webm"}
    ext = ext_map.get(mime_type, "webm")
    client = _groq_client()
    try:
        audio_io = io.BytesIO(audio_bytes); audio_io.name = f"audio.{ext}"
        result = await asyncio.get_running_loop().run_in_executor(
            None, lambda: client.audio.transcriptions.create(
                model="whisper-large-v3-turbo", file=audio_io, response_format="text"))
        return str(result).strip()
    except Exception as e:
        raise HTTPException(500, f"Transcription failed: {e}")

# ══════════════════════════════════════════════════════════════════════════════
# §16  WEBSOCKET MANAGER
# ══════════════════════════════════════════════════════════════════════════════

class ConnectionManager:
    def __init__(self): self._connections: Dict[str, Set[WebSocket]] = defaultdict(set)
    async def connect(self, user_id: str, ws: WebSocket):
        await ws.accept(); self._connections[user_id].add(ws)
    def disconnect(self, user_id: str, ws: WebSocket):
        self._connections[user_id].discard(ws)
        if not self._connections[user_id]: del self._connections[user_id]
    async def send_to_user(self, user_id: str, data: Dict):
        dead = []
        for ws in list(self._connections.get(user_id,[])):
            try: await ws.send_json(data)
            except Exception: dead.append(ws)
        for ws in dead: self.disconnect(user_id, ws)
    async def broadcast(self, data: Dict):
        for uid in list(self._connections.keys()): await self.send_to_user(uid, data)
    def online_users(self) -> List[str]: return list(self._connections.keys())

ws_manager = ConnectionManager()

async def _push_notification(user_id: Optional[str], title: str, message: str,
                              ntype: str = "info", link: str = None, is_broadcast: bool = False) -> None:
    nid = _new_id()
    await db_execute(
        "INSERT INTO notifications(id,user_id,title,message,type,link,is_broadcast,created_at)"
        " VALUES(?,?,?,?,?,?,?,?)",
        (nid, user_id, title, message, ntype, link, int(is_broadcast), _utcnow()))
    payload = {"type":"notification","id":nid,"title":title,"message":message,"ntype":ntype}
    if is_broadcast: await ws_manager.broadcast(payload)
    elif user_id: await ws_manager.send_to_user(user_id, payload)

# ══════════════════════════════════════════════════════════════════════════════
# §17  WEBSITE BUILDER
# ══════════════════════════════════════════════════════════════════════════════

_WEB_SYSTEM = """You are an elite product designer and senior frontend engineer. Generate a COMPLETE, SELF-CONTAINED single-file HTML website or web app prototype.
All CSS must live in one <style> tag. All JavaScript must live in one <script> tag. Use semantic HTML, accessible controls, responsive CSS Grid/Flexbox, polished mobile layouts, and realistic production copy.
Fully functional, visually stunning, fully responsive. No Lorem Ipsum. Return ONLY raw HTML - no markdown fences."""

def _clean_generated_html(result: str, title: str = "Website") -> str:
    html = re.sub(r"^```(?:html)?\s*", "", (result or "").strip(), flags=re.IGNORECASE)
    html = re.sub(r"```\s*$", "", html.strip()).strip()
    if "<html" not in html.lower():
        safe_title = html_lib.escape(title or "Website")
        html = f"<!doctype html><html lang=\"en\"><head><meta charset=\"utf-8\"><meta name=\"viewport\" content=\"width=device-width,initial-scale=1\"><title>{safe_title}</title></head><body>{html}</body></html>"
    if "<meta name=\"viewport\"" not in html.lower():
        html = re.sub(r"(<head[^>]*>)", r"\1<meta name=\"viewport\" content=\"width=device-width,initial-scale=1\">", html, count=1, flags=re.IGNORECASE)
    return html

def _website_quality_spec(req: Dict[str, Any]) -> str:
    features = ", ".join([str(x).strip() for x in req.get("features", []) if str(x).strip()][:12])
    sections = ", ".join([str(x).strip() for x in req.get("pages", []) if str(x).strip()][:12])
    return (
        f"Site type: {req.get('site_type') or 'choose the best fit'}\n"
        f"Audience: {req.get('audience') or 'target users for the brief'}\n"
        f"Sections/pages: {sections or 'choose the best sections for the brief'}\n"
        f"Requested features: {features or 'add useful interactions appropriate to the site'}\n"
        f"Theme mode: {req.get('theme_mode') or 'auto'}\n"
        f"Interactivity level: {req.get('interactivity') or 'rich'}\n"
        f"SEO keywords: {req.get('seo_keywords') or 'derive from brief'}\n"
        f"Palette: {req.get('color_palette') or 'domain-appropriate, not one-note'}\n"
    )

async def _build_website(description: str, title: str, style: str, model_id: str,
                         pages: Optional[List[str]] = None,
                         color_palette: str = "",
                         extra_instructions: str = "",
                         site_type: str = "",
                         audience: str = "",
                         features: Optional[List[str]] = None,
                         theme_mode: str = "auto",
                         interactivity: str = "rich",
                         seo_keywords: str = "") -> str:
    hint = _STYLE_HINTS.get(style, _STYLE_HINTS["modern"])
    policies = await _prompt_policies()
    admin_guidance = str(policies.get("website_builder_instructions") or "").strip()
    quality = _website_quality_spec({
        "pages": pages or [], "color_palette": color_palette, "site_type": site_type,
        "audience": audience, "features": features or [], "theme_mode": theme_mode,
        "interactivity": interactivity, "seo_keywords": seo_keywords,
    })
    prompt = (
        f"Build a complete website.\nTitle: {title}\nDescription: {description}\n"
        f"Style: {style} - {hint}\n"
        f"{quality}"
        f"Extra instructions: {extra_instructions or 'none'}\n"
        f"Admin builder guidance: {admin_guidance or 'none'}\n"
        "Quality bar: make it feel closer to a Lovable/Emergent-style generated app: complete, coherent, interactive, mobile-perfect, visually specific, with real content and no placeholder sections.\n"
        "Return ONLY raw HTML."
    )
    result = await _llm_text(
        [{"role":"system","content":_WEB_SYSTEM},{"role":"user","content":prompt}],
        model_id=model_id, max_tokens=8192, temperature=0.72)
    return _clean_generated_html(result, title)

async def _iterate_website(existing_html: str, change_request: str, title: str,
                           model_id: str, preserve_brand: bool = True) -> str:
    prompt = (
        "Revise this existing single-file website. Return the full updated HTML file, not a diff.\n"
        f"Title: {title}\n"
        f"Change request: {change_request}\n"
        f"Preserve brand/content unless explicitly changed: {'yes' if preserve_brand else 'no'}\n"
        "Keep it self-contained, responsive, accessible, and interactive. Do not remove working sections unless asked.\n\n"
        "CURRENT HTML:\n"
        f"{existing_html[:180000]}"
    )
    result = await _llm_text(
        [{"role":"system","content":_WEB_SYSTEM},{"role":"user","content":prompt}],
        model_id=model_id, max_tokens=8192, temperature=0.55)
    return _clean_generated_html(result, title)

# ══════════════════════════════════════════════════════════════════════════════
# §18  MCP SERVER ENGINE
# ══════════════════════════════════════════════════════════════════════════════

async def _mcp_list_tools(server_id: str) -> List[Dict]:
    row = await db_fetchone("SELECT * FROM mcp_servers WHERE id=? AND is_active=1", (server_id,))
    if not row: return []
    headers = {"Content-Type":"application/json"}
    if row.get("encrypted_creds"):
        try:
            creds = _decrypt(row["encrypted_creds"])
            if row["auth_type"] == "bearer": headers["Authorization"] = f"Bearer {creds.get('token','')}"
            elif row["auth_type"] == "api_key": headers["X-API-Key"] = creds.get("api_key","")
        except Exception: pass
    payload = {"jsonrpc":"2.0","id":1,"method":"tools/list","params":{}}
    try:
        result = await _http_post_json(row["server_url"].rstrip("/")+"/mcp", payload, headers)
        return result.get("result",{}).get("tools",[])
    except Exception: return []

async def _mcp_call_tool(server_id: str, tool_name: str, tool_args: Dict) -> Dict:
    row = await db_fetchone("SELECT * FROM mcp_servers WHERE id=? AND is_active=1", (server_id,))
    if not row: raise ValueError(f"MCP server {server_id} not found")
    headers = {"Content-Type":"application/json"}
    if row.get("encrypted_creds"):
        try:
            creds = _decrypt(row["encrypted_creds"])
            if row["auth_type"] == "bearer": headers["Authorization"] = f"Bearer {creds.get('token','')}"
        except Exception: pass
    payload = {"jsonrpc":"2.0","id":1,"method":"tools/call",
               "params":{"name":tool_name,"arguments":tool_args}}
    result = await _http_post_json(row["server_url"].rstrip("/")+"/mcp", payload, headers)
    if "error" in result: raise ValueError(result["error"].get("message","MCP error"))
    return result.get("result",{})

# ══════════════════════════════════════════════════════════════════════════════
# §19  APACHE LIVY ENGINE
# ══════════════════════════════════════════════════════════════════════════════

def _livy_headers() -> Dict:
    headers = {"Content-Type":"application/json","Accept":"application/json"}
    if LIVY_USER and LIVY_PASSWORD:
        cred = base64.b64encode(f"{LIVY_USER}:{LIVY_PASSWORD}".encode()).decode()
        headers["Authorization"] = f"Basic {cred}"
    return headers

def _livy_configured() -> bool:
    return bool((LIVY_URL or "").strip())

def _livy_reachable_quick(timeout: float = 0.6) -> bool:
    if not _livy_configured():
        return False
    try:
        parsed = urllib.parse.urlparse(LIVY_URL)
        host = parsed.hostname
        port = parsed.port or (443 if parsed.scheme == "https" else 80)
        if not host:
            return False
        with socket.create_connection((host, port), timeout=timeout):
            return True
    except Exception:
        return False

async def _livy_request(method: str, path: str, body: Optional[Dict] = None) -> Any:
    url = LIVY_URL.rstrip("/") + path; headers = _livy_headers()
    def _call():
        data = json.dumps(body).encode() if body else None
        req = urllib.request.Request(url, data, headers, method=method)
        try:
            with urllib.request.urlopen(req, timeout=30) as r: return json.loads(r.read())
        except urllib.error.HTTPError as e:
            raw = e.read()
            try: return json.loads(raw)
            except Exception: raise ValueError(f"Livy HTTP {e.code}: {raw.decode('utf-8','replace')[:300]}")
    return await asyncio.get_running_loop().run_in_executor(_executor, _call)

# ══════════════════════════════════════════════════════════════════════════════
# §20  PYDANTIC MODELS
# ══════════════════════════════════════════════════════════════════════════════

class AuthIn(BaseModel):
    email: str
    password: str
    full_name: str = ""

class EmailOnlyIn(BaseModel):
    email: str

class PasswordResetIn(BaseModel):
    token: str
    new_password: str = Field(..., min_length=8)

class RefreshIn(BaseModel):
    refresh_token: str

class SupabaseOAuthExchange(BaseModel):
    code: str = Field(..., min_length=8)
    state: str = Field(..., min_length=8)
    redirect_uri: str = Field(..., min_length=8)

class SupabaseEmailConfirmIn(BaseModel):
    access_token: str = Field(..., min_length=20, max_length=12000)

class ProfileUpdate(BaseModel):
    full_name: Optional[str] = None
    memory_enabled: Optional[bool] = None
    timezone: Optional[str] = None

class ChangePasswordReq(BaseModel):
    current_password: str
    new_password: str = Field(..., min_length=8)

class ChatReq(BaseModel):
    message: str = Field(..., min_length=1, max_length=1_000_000)
    session_id: Optional[str] = None
    model_id: str = "llama-3.3-70b-versatile"
    use_rag: bool = True
    web_search: bool = False
    plan_mode: bool = False
    skill_ids: List[str] = Field(default_factory=list)
    force_thinking: Optional[bool] = None
    image_url: Optional[str] = None
    attachments: List[Dict[str, Any]] = Field(default_factory=list)
    agent_mode: bool = False
    tools: Any = Field(default_factory=dict)

    @field_validator("skill_ids", mode="before")
    @classmethod
    def _skill_ids_to_list(cls, value):
        return _coerce_str_list(value)

    @field_validator("attachments", mode="before")
    @classmethod
    def _attachments_to_list(cls, value):
        if value is None:
            return []
        if isinstance(value, list):
            return [v for v in value if isinstance(v, dict)][:12]
        return []

    @field_validator("tools", mode="before")
    @classmethod
    def _tools_to_options(cls, value):
        return _coerce_tool_options(value)

class FileGenerateReq(BaseModel):
    prompt: str = Field("", max_length=20_000)
    content: str = Field(..., min_length=1, max_length=1_000_000)
    file_type: str = Field(
        "txt",
        pattern="^(txt|text|md|markdown|csv|xlsx|excel|pdf|tex|latex|docx|word|pptx|ppt|powerpoint|slides|presentation|html|json|zip|archive)$",
    )
    filename: Optional[str] = Field(None, max_length=120)

class LatexCompileReq(BaseModel):
    source: str = Field(..., min_length=1, max_length=1_000_000)
    filename: Optional[str] = Field(None, max_length=120)
    output: str = Field("pdf", pattern="^(pdf|tex|analysis)$")

class AgentRunReq(BaseModel):
    prompt: str = Field(..., min_length=5)
    tools: List[str] = []
    model_id: str = "llama-3.3-70b-versatile"
    session_id: Optional[str] = None
    timeout_seconds: int = Field(90, ge=10, le=300)

class ApiKeyCreate(BaseModel):
    name: str = Field(..., min_length=1, max_length=60)
    scopes_json: str = '["chat","rag"]'
    rate_limit_rpm: int = Field(60, ge=1, le=1000)
    expires_days: Optional[int] = None

class AgentJobCreate(BaseModel):
    name: str = Field(..., min_length=1, max_length=100)
    description: str = ""
    job_type: str = "cron"
    trigger_json: str = "{}"
    prompt_template: str = Field(..., min_length=5)
    tools_json: str = '["code_exec","shell","rag_search"]'
    model_id: str = "llama-3.3-70b-versatile"
    max_retries: int = Field(2, ge=0, le=5)
    timeout_seconds: int = Field(120, ge=10, le=600)

    @field_validator("trigger_json", "tools_json", mode="before")
    @classmethod
    def _jsonish_to_string(cls, value):
        if value is None:
            return "{}"
        if isinstance(value, str):
            return value
        return json.dumps(value)

class AgentJobUpdate(BaseModel):
    name: Optional[str] = None
    prompt_template: Optional[str] = None
    trigger_json: Optional[str] = None
    tools_json: Optional[str] = None
    enabled: Optional[bool] = None
    model_id: Optional[str] = None

    @field_validator("trigger_json", "tools_json", mode="before")
    @classmethod
    def _optional_jsonish_to_string(cls, value):
        if value is None or isinstance(value, str):
            return value
        return json.dumps(value)

class AIModelCreate(BaseModel):
    name: str = Field(..., min_length=1, max_length=100)
    provider: str = "groq"
    base_url: Optional[str] = None
    model_name: str = Field(..., min_length=1, max_length=200)
    api_key: str = ""
    is_active: bool = True
    is_default: bool = False
    is_fast: bool = False
    is_vision: bool = False
    is_code: bool = False
    context_length: int = Field(4096, ge=512, le=2_000_000)
    max_output_tokens: int = Field(4096, ge=128, le=128_000)
    temperature_default: float = Field(0.7, ge=0.0, le=2.0)
    description: str = ""

class AIModelUpdate(BaseModel):
    name: Optional[str] = None
    provider: Optional[str] = None
    base_url: Optional[str] = None
    model_name: Optional[str] = None
    api_key: Optional[str] = None
    is_active: Optional[bool] = None
    is_default: Optional[bool] = None
    is_fast: Optional[bool] = None
    is_vision: Optional[bool] = None
    is_code: Optional[bool] = None
    context_length: Optional[int] = Field(None, ge=512, le=2_000_000)
    max_output_tokens: Optional[int] = Field(None, ge=128, le=128_000)
    temperature_default: Optional[float] = Field(None, ge=0.0, le=2.0)
    description: Optional[str] = None

class AIModelDumpImport(BaseModel):
    raw: str = Field(..., min_length=2, max_length=2_000_000)
    dry_run: bool = False

class MCPServerCreate(BaseModel):
    name: str = Field(..., min_length=1, max_length=80)
    description: str = ""
    server_url: str = Field(..., min_length=5)
    auth_type: str = "none"
    creds: Dict = {}
    is_public: bool = False

class MCPToolCallReq(BaseModel):
    server_id: str
    tool_name: str
    tool_args: Dict = {}

class SkillCreate(BaseModel):
    name: str = Field(..., min_length=1, max_length=100)
    description: str = ""
    activation_keywords: List[str] = []
    instructions: str = Field(..., min_length=5)
    resources: List[Dict[str, Any]] = []
    script_metadata: Dict[str, Any] = {}
    is_enabled: bool = True

    @field_validator("activation_keywords", mode="before")
    @classmethod
    def _keywords_to_list(cls, value):
        return _coerce_str_list(value)

class SkillUpdate(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    activation_keywords: Optional[List[str]] = None
    instructions: Optional[str] = None
    resources: Optional[List[Dict[str, Any]]] = None
    script_metadata: Optional[Dict[str, Any]] = None
    is_enabled: Optional[bool] = None

    @field_validator("activation_keywords", mode="before")
    @classmethod
    def _optional_keywords_to_list(cls, value):
        return None if value is None else _coerce_str_list(value)

class WebsiteCreate(BaseModel):
    title: str = Field(..., min_length=1, max_length=100)
    description: str = Field(..., min_length=5)
    style: str = "modern"
    model_id: str = "llama-3.3-70b-versatile"
    pages: Optional[List[str]] = None
    color_palette: Optional[str] = None
    extra_instructions: Optional[str] = None
    site_type: str = "landing"
    audience: str = ""
    features: List[str] = Field(default_factory=list)
    theme_mode: str = "auto"
    interactivity: str = "rich"
    seo_keywords: str = ""

    @field_validator("pages", "features", mode="before")
    @classmethod
    def _website_lists(cls, value):
        return _coerce_str_list(value)

class WebsiteIterateReq(BaseModel):
    change_request: str = Field(..., min_length=3, max_length=20000)
    model_id: str = "llama-3.3-70b-versatile"
    save_as_new: bool = False
    preserve_brand: bool = True

class CodeRunReq(BaseModel):
    code: str = Field(..., min_length=1)
    language: str = "python"
    cwd: Optional[str] = None
    save_log: bool = True

class ShellReq(BaseModel):
    command: Optional[str] = None
    cmd: Optional[str] = None
    timeout: int = Field(30, ge=1, le=300)
    cwd: Optional[str] = None

class MemoryCreate(BaseModel):
    key: str = Field(..., min_length=1, max_length=40)
    value: str = Field(..., min_length=1, max_length=300)

class VoiceTokenReq(BaseModel):
    room_name: str = "jazz-voice"
    identity: Optional[str] = None
    ttl_seconds: int = Field(3600, ge=60, le=86400)

class LiveVoiceSessionReq(BaseModel):
    room_name: Optional[str] = None
    language: str = "auto"
    ttl_seconds: int = Field(3600, ge=60, le=86400)

class VoiceChatReq(BaseModel):
    transcript: str
    session_id: Optional[str] = None
    model_id: str = "llama-3.3-70b-versatile"
    return_audio: bool = False
    voice: str = TTS_VOICE

class NotifCreate(BaseModel):
    user_id: Optional[str] = None
    title: str
    message: str
    type: str = "info"
    link: Optional[str] = None
    is_broadcast: bool = False

class SubscriptionUpdate(BaseModel):
    user_id: str
    tier: str = Field(..., pattern="^(free|pro|premium|enterprise)$")

class BillingOrderReq(BaseModel):
    tier: str = Field(..., pattern="^(pro|premium|enterprise)$")

class RazorpayVerifyReq(BaseModel):
    razorpay_order_id: str = Field(..., min_length=5, max_length=80)
    razorpay_payment_id: str = Field(..., min_length=5, max_length=80)
    razorpay_signature: str = Field(..., min_length=20, max_length=256)

class CouponCreate(BaseModel):
    code: Optional[str] = None
    grants_tier: str = Field(..., pattern="^(free|pro|premium|enterprise)$")
    duration_days: int = 30
    max_uses: int = 1
    expires_days: Optional[int] = None

class MultiModelChatReq(BaseModel):
    message: str = Field(..., min_length=1, max_length=1_000_000)
    model_ids: List[str] = Field(..., min_length=2)
    session_id: Optional[str] = None
    use_rag: bool = True

    @field_validator("model_ids", mode="before")
    @classmethod
    def _model_ids_to_list(cls, value):
        return _coerce_str_list(value)

class EnvVarCreate(BaseModel):
    key: str = Field(..., min_length=1, max_length=128)
    value: str = ""

class EnvVarUpdate(BaseModel):
    value: str = ""

class ConnectorActionReq(BaseModel):
    action: str
    params: Dict = {}

class LivySessionCreate(BaseModel):
    kind: str = "pyspark"
    name: str = "jazz-session"

class LivyStatementCreate(BaseModel):
    code: str = Field(..., min_length=1)
    kind: str = "pyspark"

class LivyBatchCreate(BaseModel):
    file: str
    class_name: str = ""
    args: List[str] = []
    conf: Dict = {}

# ══════════════════════════════════════════════════════════════════════════════
# §21  APP FACTORY
# ══════════════════════════════════════════════════════════════════════════════

app = FastAPI(title="JAZZ AI", version="14.0", lifespan=lifespan)
app.add_middleware(CORSMiddleware, allow_origins=ALLOWED_ORIGINS,
                   allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

if SITES_DIR.exists():
    app.mount("/preview", StaticFiles(directory=str(SITES_DIR)), name="sites")

@app.middleware("http")
async def _request_logger(request: Request, call_next):
    t0 = time.time(); resp = await call_next(request)
    ms = int((time.time()-t0)*1000); path = request.url.path
    if path not in ("/health","/favicon.ico") and not path.startswith("/static") and not path.startswith("/preview"):
        logger.info("%s %s → %d (%dms)", request.method, path, resp.status_code, ms)
    resp.headers["ngrok-skip-browser-warning"] = "true"
    return resp

FRONTEND_PATH = Path("index.html")

def _serve_frontend_response():
    if FRONTEND_PATH.exists():
        return FileResponse(
            FRONTEND_PATH,
            media_type="text/html",
            headers={"Cache-Control": "no-store, max-age=0", "Pragma": "no-cache"},
        )
    return HTMLResponse(
        "<h1>JAZZ AI Server is Running</h1>",
        headers={"Cache-Control": "no-store, max-age=0", "Pragma": "no-cache"},
    )

@app.get("/", response_class=HTMLResponse)
async def serve_frontend():
    return _serve_frontend_response()

@app.get("/c/{session_id:path}", response_class=HTMLResponse)
async def serve_chat_session(session_id: str):
    return _serve_frontend_response()

@app.get("/auth/supabase/callback", response_class=HTMLResponse)
async def serve_supabase_auth_callback():
    return _serve_frontend_response()

# ══════════════════════════════════════════════════════════════════════════════
# §22  AUTH ROUTES
# ══════════════════════════════════════════════════════════════════════════════

@app.get("/platform/status")
async def platform_status():
    return await _platform_status()

@app.get("/auth/supabase-config")
async def auth_supabase_config():
    return _supabase_public_config()

@app.get("/auth/supabase/status")
async def auth_supabase_status(request: Request):
    base = _supabase_auth_base()
    client_id = _supabase_oauth_client_id()
    redirect_uri = _supabase_oauth_redirect_uri(request)
    return {
        "enabled": bool(base and client_id),
        "project_url": base,
        "redirect_uri": redirect_uri,
        "scopes": _supabase_oauth_scopes(),
        "auth_method": _supabase_oauth_auth_method(),
    }

@app.get("/auth/supabase/start-url")
async def auth_supabase_start_url(request: Request, redirect_uri: str = ""):
    return await _supabase_oauth_authorize_url(request, redirect_uri)

@app.post("/auth/supabase/exchange")
async def auth_supabase_exchange(body: SupabaseOAuthExchange):
    profile = await _exchange_supabase_oauth_code(body)
    return await _issue_jazz_auth_for_supabase_profile(profile)

@app.post("/auth/supabase/confirm-email")
async def auth_supabase_confirm_email(body: SupabaseEmailConfirmIn):
    try:
        profile = await _supabase_user_from_access_token(body.access_token)
    except Exception as e:
        logger.warning("[AUTH] Supabase email token verification failed: %s", e)
        raise HTTPException(400, "Verification link is invalid or expired")
    email = _validated_email(
        profile.get("email")
        or (profile.get("user_metadata") or {}).get("email")
        or ""
    )
    is_confirmed = bool(
        profile.get("email_confirmed_at")
        or profile.get("confirmed_at")
        or (profile.get("user_metadata") or {}).get("email_verified")
    )
    if not is_confirmed:
        raise HTTPException(400, "Email is not verified by Supabase yet")
    row = await db_fetchone("SELECT id,is_verified FROM users WHERE email=?", (email,))
    if not row:
        return {
            "ok": True,
            "email": email,
            "jazz_user_found": False,
            "message": "Email verified. Create your JAZZ account, then sign in.",
        }
    now = _utcnow()
    if not int(row.get("is_verified") or 0):
        await db_execute("UPDATE users SET is_verified=1,updated_at=? WHERE id=?", (now, row["id"]))
        await db_execute(
            "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
            (_new_id(), row["id"], row["id"], "supabase_email_verified",
             json.dumps({"email": email}), now),
        )
    return {
        "ok": True,
        "email": email,
        "jazz_user_found": True,
        "message": "Email verified. Sign in now.",
    }

@app.post("/auth/register")
async def auth_register(request: Request, body: AuthIn):
    maintenance_on, maintenance_msg = await _maintenance_enabled_message()
    if maintenance_on:
        raise HTTPException(503, maintenance_msg)
    email = _validated_email(body.email)
    _validated_password(body.password)
    if await db_fetchone("SELECT id FROM users WHERE email=?", (email,)):
        raise HTTPException(400, "Email already registered. Sign in or resend verification.")
    uid = _new_id(); now = _utcnow()
    full_name = (body.full_name or email.split("@")[0]).strip()[:120]
    is_verified = 0 if _email_verification_required() else 1
    password_hash = await _hash_pw_async(body.password)
    await db_execute(
        "INSERT INTO users(id,email,password_hash,full_name,role,subscription,is_active,is_verified,created_at,updated_at)"
       " VALUES(?,?,?,?,?,?,?,?,?,?)",
        (
            uid,
            email,
            password_hash,
            full_name,
            "client",
            "free",
            1,
            is_verified,
            now,
            now,
        ),
    )
    if _email_verification_required():
        link = await _create_email_verification_link(uid, email, request)
        try:
            delivery = await _send_verification_email(email, link)
            logger.info("[AUTH] Verification email sent for %s via %s", email, delivery.get("provider"))
            return _verification_response(email, link, True)
        except Exception as e:
            logger.error("[AUTH] Verification email failed for %s: %s", email, e)
            logger.info("[AUTH] Verification link for %s: %s", email, link)
            return _verification_response(email, link, False, str(e))
    raw = f"jzr_{secrets.token_urlsafe(48)}"; h = _hash_token(raw)
    exp = (datetime.now(timezone.utc)+timedelta(minutes=REFRESH_TOKEN_EXPIRE_MINUTES)).isoformat()
    await db_execute("INSERT INTO refresh_tokens(id,user_id,token_hash,expires_at,created_at) VALUES(?,?,?,?,?)",
                     (_new_id(), uid, h, exp, now))
    return {"access_token":_make_access_token(uid,"client",email,"free"),
            "refresh_token":raw,"token_type":"bearer",
            "user":{"id":uid,"email":email,"full_name":full_name,"role":"client","subscription":"free","is_verified":is_verified}}

@app.post("/auth/signup")
async def auth_signup(request: Request, body: AuthIn):
    return await auth_register(request, body)

@app.post("/auth/resend-verification")
async def auth_resend_verification(request: Request, body: EmailOnlyIn):
    email = _validated_email(body.email)
    user = await db_fetchone("SELECT id,email,is_verified FROM users WHERE email=?", (email,))
    generic = {"ok": True, "message": "If this account needs verification, a fresh link has been created."}
    if not user:
        return generic
    if int(user.get("is_verified") or 0):
        return {"ok": True, "already_verified": True, "message": "Email is already verified. You can sign in."}
    link = await _create_email_verification_link(user["id"], email, request)
    try:
        delivery = await _send_verification_email(email, link)
        logger.info("[AUTH] Verification email sent for %s via %s", email, delivery.get("provider"))
        return _verification_response(email, link, True)
    except Exception as e:
        logger.error("[AUTH] Verification email failed for %s: %s", email, e)
        logger.info("[AUTH] Verification link for %s: %s", email, link)
        return _verification_response(email, link, False, str(e))

@app.post("/auth/forgot-password")
async def auth_forgot_password(request: Request, body: EmailOnlyIn):
    email = _validated_email(body.email)
    user = await db_fetchone("SELECT id,email,is_active FROM users WHERE email=?", (email,))
    if not user or not int(user.get("is_active") or 0):
        return _password_reset_response(email)
    link = await _create_password_reset_link(user["id"], email, request)
    try:
        delivery = await _send_password_reset_email(email, link)
        logger.info("[AUTH] Password reset email sent for %s via %s", email, delivery.get("provider"))
        return _password_reset_response(email, link, True)
    except Exception as e:
        logger.error("[AUTH] Password reset email failed for %s: %s", email, e)
        logger.info("[AUTH] Password reset link for %s: %s", email, link)
        return _password_reset_response(email, link, False, str(e))

@app.get("/auth/reset-password")
async def auth_reset_password_page(token: str = ""):
    if not token:
        return _auth_html_page("Reset password", "Open the reset link from your email, or request a new one from the sign-in screen.", True)
    link = f"/?reset_token={urllib.parse.quote(token)}"
    html = f"""<!doctype html><html><head><meta charset="utf-8"><meta http-equiv="refresh" content="0; url={html_lib.escape(link)}"><title>Reset password</title></head><body><a href="{html_lib.escape(link)}">Continue to reset password</a></body></html>"""
    return HTMLResponse(html)

@app.post("/auth/reset-password")
async def auth_reset_password(body: PasswordResetIn):
    email = await _reset_password_by_token(body.token, body.new_password)
    return {"ok": True, "email": email, "message": "Password reset. Sign in with your new password."}

@app.get("/auth/verify-email")
async def auth_verify_email(token: str = ""):
    ok, message = await _verify_email_token(token)
    if ok:
        return _auth_html_page("Email verified", f"{message} is verified. You can sign in now.", True)
    return _auth_html_page("Verification failed", message, False)

@app.post("/auth/login")
async def auth_login(body: AuthIn):
    email = _validated_email(body.email)
    user = await db_fetchone("SELECT * FROM users WHERE email=? AND is_active=1", (email,))
    if not user or not await _verify_pw_async(body.password, user["password_hash"]):
        raise HTTPException(401, "Invalid credentials")
    if _email_verification_required() and not int(user.get("is_verified") or 0):
        raise HTTPException(
            403,
            {
                "code": "email_not_verified",
                "message": "Please verify your email before signing in. Use Resend verification if the link expired.",
            },
        )
    maintenance_on, maintenance_msg = await _maintenance_enabled_message()
    if maintenance_on and user["role"] != "admin":
        raise HTTPException(503, maintenance_msg)
    now = _utcnow()
    await db_execute("UPDATE users SET last_login_at=? WHERE id=?", (now, user["id"]))
    raw = f"jzr_{secrets.token_urlsafe(48)}"; h = _hash_token(raw)
    exp = (datetime.now(timezone.utc)+timedelta(minutes=REFRESH_TOKEN_EXPIRE_MINUTES)).isoformat()
    await db_execute("INSERT INTO refresh_tokens(id,user_id,token_hash,expires_at,created_at) VALUES(?,?,?,?,?)",
                     (_new_id(), user["id"], h, exp, now))
    return {"access_token":_make_access_token(user["id"],user["role"],user["email"],user["subscription"]),
            "refresh_token":raw,"token_type":"bearer",
            "user":{k:user[k] for k in ("id","email","full_name","role","subscription","memory_enabled","is_verified")}}

@app.post("/auth/refresh")
async def auth_refresh(body: RefreshIn):
    h = _hash_token(body.refresh_token)
    row = await db_fetchone(
        "SELECT rt.*,u.role,u.email,u.subscription FROM refresh_tokens rt JOIN users u ON u.id=rt.user_id "
        "WHERE rt.token_hash=? AND rt.revoked=0", (h,))
    if not row: raise HTTPException(401, "Invalid refresh token")
    if datetime.fromisoformat(row["expires_at"]) < datetime.now(timezone.utc):
        raise HTTPException(401, "Refresh token expired")
    maintenance_on, maintenance_msg = await _maintenance_enabled_message()
    if maintenance_on and row["role"] != "admin":
        raise HTTPException(503, maintenance_msg)
    await db_execute("UPDATE refresh_tokens SET revoked=1 WHERE id=?", (row["id"],))
    raw = f"jzr_{secrets.token_urlsafe(48)}"; nh = _hash_token(raw)
    exp = (datetime.now(timezone.utc)+timedelta(minutes=REFRESH_TOKEN_EXPIRE_MINUTES)).isoformat()
    await db_execute("INSERT INTO refresh_tokens(id,user_id,token_hash,expires_at,created_at) VALUES(?,?,?,?,?)",
                     (_new_id(), row["user_id"], nh, exp, _utcnow()))
    return {"access_token":_make_access_token(row["user_id"],row["role"],row["email"],row["subscription"]),
            "refresh_token":raw}

@app.get("/auth/me")
async def auth_me(user: Dict = Depends(_get_current_user)):
    row = await db_fetchone(
        "SELECT id,email,full_name,role,subscription,memory_enabled,timezone,is_verified,created_at FROM users WHERE id=?",
        (user.get("id") or user.get("sub"),))
    return row or {}

@app.post("/auth/logout")
async def auth_logout(body: RefreshIn):
    h = _hash_token(body.refresh_token)
    await db_execute("UPDATE refresh_tokens SET revoked=1 WHERE token_hash=?", (h,))
    return {"ok":True}

@app.post("/auth/change-password")
async def auth_change_password(body: ChangePasswordReq, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    row = await db_fetchone("SELECT password_hash FROM users WHERE id=?", (uid,))
    if not row or not await _verify_pw_async(body.current_password, row["password_hash"]):
        raise HTTPException(400, "Wrong current password")
    _validated_password(body.new_password)
    new_hash = await _hash_pw_async(body.new_password)
    await db_execute("UPDATE users SET password_hash=?,updated_at=? WHERE id=?",
                     (new_hash, _utcnow(), uid))
    return {"ok":True}

@app.put("/auth/profile")
async def auth_update_profile(req: ProfileUpdate, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    updates, vals = [], []
    if req.full_name      is not None: updates.append("full_name=?"); vals.append(req.full_name)
    if req.memory_enabled is not None: updates.append("memory_enabled=?"); vals.append(int(req.memory_enabled))
    if req.timezone       is not None: updates.append("timezone=?"); vals.append(req.timezone)
    if updates:
        updates.append("updated_at=?"); vals.append(_utcnow()); vals.append(uid)
        await db_execute(f"UPDATE users SET {','.join(updates)} WHERE id=?", tuple(vals))
    return {"ok":True}

# ══════════════════════════════════════════════════════════════════════════════
# §23  CHAT SESSIONS
# ══════════════════════════════════════════════════════════════════════════════

@app.get("/sessions")
async def list_sessions(limit: int = 60, offset: int = 0, archived: bool = False, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    limit = max(1, min(200, int(limit or 60)))
    offset = max(0, int(offset or 0))
    rows = await db_fetchall(
        "SELECT id,title,model_id,turn_count,last_message_at,is_pinned,is_archived,created_at"
        " FROM chat_sessions WHERE user_id=? AND is_archived=?"
        " ORDER BY is_pinned DESC,last_message_at DESC LIMIT ? OFFSET ?",
        (uid, 1 if archived else 0, limit, offset))
    return {"sessions":rows}

@app.get("/sessions/search")
async def search_sessions(q: str = "", limit: int = 50, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    query = (q or "").strip()
    if not query:
        return {"sessions":[]}
    like = f"%{query.lower()}%"
    limit = max(1, min(100, int(limit or 50)))
    rows = await db_fetchall(
        "SELECT s.id,s.title,s.model_id,s.turn_count,s.last_message_at,s.is_pinned,s.is_archived,s.created_at,"
        " CASE WHEN LOWER(COALESCE(s.title,'')) LIKE ? THEN 'title' ELSE 'history' END AS match_source,"
        " COALESCE((SELECT h.content FROM chat_history h"
        "   WHERE h.session_id=s.id AND h.is_hidden=0 AND LOWER(COALESCE(h.content,'')) LIKE ?"
        "   ORDER BY h.created_at DESC LIMIT 1),'') AS match_text,"
        " (SELECT COUNT(*) FROM chat_history h"
        "   WHERE h.session_id=s.id AND h.is_hidden=0 AND LOWER(COALESCE(h.content,'')) LIKE ?) AS match_count"
        " FROM chat_sessions s"
        " WHERE s.user_id=? AND s.is_archived=0 AND (LOWER(COALESCE(s.title,'')) LIKE ?"
        "   OR EXISTS (SELECT 1 FROM chat_history h WHERE h.session_id=s.id AND h.is_hidden=0"
        "      AND LOWER(COALESCE(h.content,'')) LIKE ?))"
        " ORDER BY s.is_pinned DESC, s.last_message_at DESC, s.created_at DESC LIMIT ?",
        (like, like, like, uid, like, like, limit))
    out = []
    ql = query.lower()
    for r in rows:
        item = dict(r)
        text = item.pop("match_text", "") or ""
        snippet = ""
        if text:
            idx = text.lower().find(ql)
            start = max(0, idx - 70) if idx >= 0 else 0
            end = min(len(text), start + 190)
            snippet = text[start:end].strip()
            if start > 0:
                snippet = "..." + snippet
            if end < len(text):
                snippet += "..."
        item["match_snippet"] = snippet
        out.append(item)
    return {"sessions":out, "query":query}

@app.post("/sessions")
async def create_session(body: dict, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    sid = _new_id(); now = _utcnow()
    model = body.get("model_id","llama-3.3-70b-versatile")
    await db_execute(
        "INSERT INTO chat_sessions(id,user_id,title,model_id,created_at,updated_at) VALUES(?,?,?,?,?,?)",
        (sid, uid, body.get("title","New Chat"), model, now, now))
    return await db_fetchone("SELECT * FROM chat_sessions WHERE id=?", (sid,))

@app.post("/sessions/archive-all")
async def archive_all_sessions(user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    now = _utcnow()
    count = await db_count("SELECT COUNT(*) as c FROM chat_sessions WHERE user_id=? AND is_archived=0", (uid,))
    await db_execute(
        "UPDATE chat_sessions SET is_archived=1,is_pinned=0,updated_at=? WHERE user_id=? AND is_archived=0",
        (now, uid))
    return {"ok": True, "archived": count}

@app.post("/sessions/unarchive-all")
async def unarchive_all_sessions(user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    now = _utcnow()
    count = await db_count("SELECT COUNT(*) as c FROM chat_sessions WHERE user_id=? AND is_archived=1", (uid,))
    await db_execute(
        "UPDATE chat_sessions SET is_archived=0,updated_at=? WHERE user_id=? AND is_archived=1",
        (now, uid))
    return {"ok": True, "restored": count}

@app.get("/sessions/{sid}")
async def get_session(sid: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    s = await db_fetchone("SELECT * FROM chat_sessions WHERE id=? AND user_id=?", (sid, uid))
    if not s: raise HTTPException(404, "Session not found")
    return s

@app.patch("/sessions/{sid}")
async def update_session(sid: str, body: dict, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    s = await db_fetchone("SELECT id,model_id FROM chat_sessions WHERE id=? AND user_id=?", (sid, uid))
    if not s: raise HTTPException(404)
    fields = {k:v for k,v in body.items() if k in ("title","model_id","is_pinned","is_archived")}
    if not fields: return {"ok":True}
    fields["updated_at"] = _utcnow()
    set_clause = ", ".join(f"{k}=?" for k in fields)
    await db_execute(f"UPDATE chat_sessions SET {set_clause} WHERE id=?", tuple(fields.values())+(sid,))
    return {"ok":True}

@app.patch("/sessions/{sid}/pin")
async def pin_session_alias(sid: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    s = await db_fetchone("SELECT id,is_pinned FROM chat_sessions WHERE id=? AND user_id=?", (sid, uid))
    if not s: raise HTTPException(404, "Session not found")
    await db_execute("UPDATE chat_sessions SET is_pinned=?,updated_at=? WHERE id=?",
                     (0 if int(s.get("is_pinned") or 0) else 1, _utcnow(), sid))
    return {"ok":True}

@app.patch("/sessions/{sid}/rename")
async def rename_session_alias(sid: str, title: str = "", user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    s = await db_fetchone("SELECT id,model_id FROM chat_sessions WHERE id=? AND user_id=?", (sid, uid))
    if not s: raise HTTPException(404, "Session not found")
    clean_title = (title or "").strip()
    if not clean_title: raise HTTPException(400, "title required")
    await db_execute("UPDATE chat_sessions SET title=?,updated_at=? WHERE id=?", (clean_title[:140], _utcnow(), sid))
    return {"ok":True}

@app.post("/sessions/{sid}/auto-title")
async def auto_title_session_alias(sid: str, body: dict = None, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    s = await db_fetchone("SELECT id FROM chat_sessions WHERE id=? AND user_id=?", (sid, uid))
    if not s: raise HTTPException(404, "Session not found")
    rows = await db_fetchall(
        "SELECT role,content FROM chat_history WHERE session_id=? AND is_hidden=0 ORDER BY created_at DESC LIMIT 6",
        (sid,))
    user_msg = next((r.get("content","") for r in rows if r.get("role") == "user"), "")
    assistant_msg = next((r.get("content","") for r in rows if r.get("role") == "assistant"), "")
    asyncio.create_task(_auto_title(sid, user_msg, assistant_msg))
    return {"ok":True}

@app.delete("/sessions/{sid}")
async def delete_session(sid: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    await db_execute("DELETE FROM chat_sessions WHERE id=? AND user_id=?", (sid, uid))
    return {"ok":True}

@app.get("/sessions/{sid}/messages")
async def get_messages(sid: str, limit: int = 60, offset: int = 0, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    s = await db_fetchone("SELECT id,model_id FROM chat_sessions WHERE id=? AND user_id=?", (sid, uid))
    if not s: raise HTTPException(404)
    msgs = await db_fetchall(
        "SELECT id,role,content,model_used,agent_used,thinking_content,tool_calls_json,"
        "rag_chunks_json,tokens_input,tokens_output,latency_ms,mode,created_at"
        " FROM chat_history WHERE session_id=? AND is_hidden=0 ORDER BY created_at ASC LIMIT ? OFFSET ?",
        (sid, limit, offset))
    summaries = await db_fetchall(
        "SELECT summary_text,covers_from_at,covers_to_at FROM summaries WHERE session_id=? ORDER BY created_at",
        (sid,))
    context = await _context_window_status(sid, uid, s["model_id"])
    return {"messages":msgs,"summaries":summaries,"context":context}

@app.get("/sessions/{sid}/context")
async def get_session_context(sid: str, model_id: Optional[str] = None,
                              user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    return await _context_window_status(sid, uid, model_id)

@app.post("/sessions/{sid}/compact")
async def compact_session_context(sid: str, body: dict = None,
                                  user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    session = await db_fetchone("SELECT id,model_id FROM chat_sessions WHERE id=? AND user_id=?", (sid, uid))
    if not session: raise HTTPException(404)
    model_id = (body or {}).get("model_id") or session["model_id"] or "llama-3.3-70b-versatile"
    return await _compact_session_context(sid, uid, model_id, reason="manual")

@app.post("/sessions/{sid}/branch")
async def branch_session(sid: str, body: dict, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    orig = await db_fetchone("SELECT * FROM chat_sessions WHERE id=? AND user_id=?", (sid, uid))
    if not orig: raise HTTPException(404)
    new_sid = _new_id(); now = _utcnow()
    msg_id = body.get("from_message_id")
    await db_execute(
        "INSERT INTO chat_sessions(id,user_id,title,model_id,parent_session_id,branch_from_msg_id,created_at,updated_at)"
        " VALUES(?,?,?,?,?,?,?,?)",
        (new_sid, uid, f"Branch of {orig['title']}", orig["model_id"], sid, msg_id, now, now))
    return {"session_id":new_sid}

@app.get("/sessions/{sid}/export")
async def export_session(sid: str, format: str = "markdown", user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    msgs = await db_fetchall(
        "SELECT role,content,model_used,latency_ms,created_at FROM chat_history WHERE session_id=? AND user_id=? AND is_hidden=0 ORDER BY created_at",
        (sid, uid))
    session = await db_fetchone("SELECT title FROM chat_sessions WHERE id=?", (sid,))
    if not msgs: raise HTTPException(404, "No messages found")
    title = (session or {}).get("title","Chat") if session else "Chat"
    if format == "json":
        data = json.dumps({"session_id":sid,"title":title,"messages":msgs}, indent=2, default=str)
        return StreamingResponse(io.StringIO(data), media_type="application/json",
                                 headers={"Content-Disposition":f'attachment; filename="chat_{sid[:8]}.json"'})
    elif format == "txt":
        lines = [f"# {title}\n"]
        for m in msgs: lines.append(f"[{m['role'].upper()} — {m['created_at']}]\n{m['content']}\n")
        return StreamingResponse(io.StringIO("\n".join(lines)), media_type="text/plain",
                                 headers={"Content-Disposition":f'attachment; filename="chat_{sid[:8]}.txt"'})
    else:  # markdown (default)
        lines = [f"# {title}\n"]
        for m in msgs:
            lines.append(f"**{m['role'].capitalize()}** ({m['created_at']}):\n{m['content']}\n")
        return StreamingResponse(io.StringIO("\n".join(lines)), media_type="text/markdown",
                                 headers={"Content-Disposition":f'attachment; filename="chat_{sid[:8]}.md"'})

# ══════════════════════════════════════════════════════════════════════════════
# §24  CHAT — STREAMING SSE
# ══════════════════════════════════════════════════════════════════════════════

async def _stream_text(text: str, chunk_size: int = 8):
    for i in range(0, len(text), chunk_size):
        yield json.dumps({"type":"delta","content":text[i:i+chunk_size]})
        await asyncio.sleep(0.01)

_EMPTY_MODEL_REPLY = (
    "I did not receive text from the selected model on that attempt. "
    "I retried automatically, so please send the message again if this appears."
)

async def _recover_empty_reply(messages: List[Dict], model_id: str) -> str:
    try:
        retry_result = await _llm_text_with_fallback(messages, model_id, max_tokens=1024, temperature=0.3)
        retry = retry_result.get("content", "")
        if retry and retry.strip():
            return retry.strip()
    except Exception as exc:
        logger.warning("[CHAT] empty-response retry failed: %s", exc)
    return _EMPTY_MODEL_REPLY

async def _get_or_create_session(user_id: str, session_id: Optional[str], model_id: str) -> str:
    if session_id:
        row = await db_fetchone("SELECT id FROM chat_sessions WHERE id=? AND user_id=?", (session_id, user_id))
        if row: return session_id
    sid = _new_id(); now = _utcnow()
    await db_execute(
        "INSERT INTO chat_sessions(id,user_id,title,model_id,created_at,updated_at) VALUES(?,?,?,?,?,?)",
        (sid, user_id, "New Chat", model_id, now, now))
    return sid

@app.post("/chat/stream")
async def chat_stream_post(req: ChatReq, user: Dict = Depends(_get_current_user)) -> StreamingResponse:
    await _check_rate_limit(user, "messages_per_day")
    uid = user.get("id") or user.get("sub","")

    async def _gen():
        t0 = time.time()
        requested_model_id = _canonical_model_id(req.model_id)
        active_model_id, route_reason = await _auto_route_model(requested_model_id, req.message, uid)
        sid = await _get_or_create_session(uid, req.session_id, active_model_id)
        tool_opts = _chat_tool_options(req)
        use_rag = bool(tool_opts.get("rag", req.use_rag))
        use_web = bool(tool_opts.get("web_search", req.web_search))
        plan_mode = bool(tool_opts.get("plan_mode", req.plan_mode))
        explicit_skill_ids = tool_opts.get("skill_ids", req.skill_ids) or []
        mcp_context = _format_mcp_context(tool_opts.get("mcp_tools", []))
        tool_log: List[Dict[str, Any]] = []
        image_inputs = _chat_image_inputs(req)
        effective_message = req.message
        await db_execute("UPDATE chat_sessions SET model_id=?,updated_at=? WHERE id=?",
                         (active_model_id, _utcnow(), sid))
        yield f"data: {json.dumps({'type':'start','session_id':sid,'mode':'normal','web_search':use_web,'plan_mode':plan_mode,'model_id':active_model_id,'model_label':await _model_display_name(active_model_id)})}\n\n"
        try:
            if route_reason and active_model_id != requested_model_id:
                route_label = f"Long input detected. Switched to {await _model_display_name(active_model_id)}."
                ev = {
                    "type":"model_switch", "from_model_id":requested_model_id,
                    "model_id":active_model_id,
                    "model_label":await _model_display_name(active_model_id),
                    "selected":True,
                    "label":route_label
                }
                tool_log.append(ev)
                yield f"data: {json.dumps(ev)}\n\n"

            identity_reply = await _identity_recall_reply(uid, req.message, user)
            if identity_reply:
                ev = {"type":"tool_result","tool":"identity","label":"Recovered your name from account context","count":1}
                tool_log.append(ev)
                yield f"data: {json.dumps(ev)}\n\n"
                async for chunk in _stream_text(identity_reply):
                    yield f"data: {chunk}\n\n"
                elapsed = int((time.time()-t0)*1000)
                now = _utcnow(); mid = _new_id()
                await db_execute(
                    "INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,created_at) VALUES(?,?,?,'user',?,?,?)",
                    (_new_id(), sid, uid, req.message, active_model_id, now))
                await db_execute(
                    "INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,latency_ms,tool_calls_json,mode,created_at)"
                    " VALUES(?,?,?,'assistant',?,?,?,?,'normal',?)",
                    (mid, sid, uid, identity_reply, active_model_id, elapsed, json.dumps(tool_log), now))
                await db_execute("UPDATE chat_sessions SET model_id=?,last_message_at=?,turn_count=turn_count+1,updated_at=? WHERE id=?",
                                 (active_model_id, now, now, sid))
                context = await _context_window_status(sid, uid, active_model_id)
                yield f"data: {json.dumps({'type':'done','message_id':mid,'content':identity_reply,'latency_ms':elapsed,'tokens':{'input':_count_tokens(req.message),'output':_count_tokens(identity_reply)},'mode':'normal','context':context,'model_id':active_model_id,'model_label':await _model_display_name(active_model_id)})}\n\n"
                return

            missing_context_reply = await _verified_context_missing_reply(uid, req.message, user)
            if missing_context_reply:
                ev = {"type":"tool_result","tool":"evidence_guard","label":"No verified context found for that private detail","count":0}
                tool_log.append(ev)
                yield f"data: {json.dumps(ev)}\n\n"
                async for chunk in _stream_text(missing_context_reply):
                    yield f"data: {chunk}\n\n"
                elapsed = int((time.time()-t0)*1000)
                now = _utcnow(); mid = _new_id()
                await db_execute(
                    "INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,created_at) VALUES(?,?,?,'user',?,?,?)",
                    (_new_id(), sid, uid, req.message, active_model_id, now))
                await db_execute(
                    "INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,latency_ms,tool_calls_json,mode,created_at)"
                    " VALUES(?,?,?,'assistant',?,?,?,?,'guard',?)",
                    (mid, sid, uid, missing_context_reply, active_model_id, elapsed, json.dumps(tool_log), now))
                await db_execute("UPDATE chat_sessions SET model_id=?,last_message_at=?,turn_count=turn_count+1,updated_at=? WHERE id=?",
                                 (active_model_id, now, now, sid))
                context = await _context_window_status(sid, uid, active_model_id)
                yield f"data: {json.dumps({'type':'done','message_id':mid,'content':missing_context_reply,'latency_ms':elapsed,'tokens':{'input':_count_tokens(req.message),'output':_count_tokens(missing_context_reply)},'mode':'guard','context':context,'model_id':active_model_id,'model_label':await _model_display_name(active_model_id)})}\n\n"
                return

            latex_req = _latex_compile_request_from_message(req.message)
            if latex_req:
                start_ev = {"type":"tool_start","tool":"latex","label":"Compiling LaTeX to PDF" if latex_req.output == "pdf" else "Preparing LaTeX source"}
                tool_log.append(start_ev)
                yield f"data: {json.dumps(start_ev)}\n\n"
                result = await asyncio.get_running_loop().run_in_executor(_executor, _build_latex_artifact, latex_req, uid)
                done_ev = {
                    "type":"tool_result",
                    "tool":"latex",
                    "label":"LaTeX PDF ready" if result.get("compiled") is not False else "LaTeX PDF ready (interpreted fallback)",
                    "count":1,
                    "file":result,
                }
                tool_log.append(done_ev)
                yield f"data: {json.dumps(done_ev)}\n\n"
                elapsed = int((time.time()-t0)*1000)
                now = _utcnow(); mid = _new_id()
                reply = (
                    f"Done - I created {result.get('original') or 'the PDF'} and added it to the Files panel."
                    if result.get("compiled") is not False else
                    f"Done - I created {result.get('original') or 'the PDF'} and added it to the Files panel. The server has no TeX engine installed, so this is an interpreted PDF fallback."
                )
                await db_execute(
                    "INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,created_at) VALUES(?,?,?,'user',?,?,?)",
                    (_new_id(), sid, uid, req.message, active_model_id, now))
                await db_execute(
                    "INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,latency_ms,tool_calls_json,mode,created_at)"
                    " VALUES(?,?,?,'assistant',?,?,?,?,'tool',?)",
                    (mid, sid, uid, reply, active_model_id, elapsed, json.dumps(tool_log), now))
                await db_execute("UPDATE chat_sessions SET last_message_at=?,turn_count=turn_count+1,updated_at=? WHERE id=?",
                                 (now, now, sid))
                context = await _context_window_status(sid, uid, active_model_id)
                yield f"data: {json.dumps({'type':'done','message_id':mid,'content':reply,'latency_ms':elapsed,'tokens':{'input':_count_tokens(req.message),'output':_count_tokens(reply)},'mode':'tool','context':context,'model_id':active_model_id,'model_label':await _model_display_name(active_model_id)})}\n\n"
                return

            skills, skill_matches = await _select_skills_for_message(req.message, explicit_skill_ids)
            skill_context = _format_skill_context(skills)
            for i, skill in enumerate(skills):
                ev = {"type":"skill_loaded","skill_id":skill["id"],"name":skill["name"],
                      "matched_by":skill_matches[i] if i < len(skill_matches) else ""}
                tool_log.append(ev)
                yield f"data: {json.dumps(ev)}\n\n"

            # Slash command
            slash = _parse_slash(req.message) or await _infer_slash_command(req.message)
            if slash:
                cmd, rest = slash
                start_ev = {"type":"tool_start","tool":"slash","label":f"Running /{cmd}"}
                tool_log.append(start_ev)
                yield f"data: {json.dumps(start_ev)}\n\n"
                result = await _execute_slash(uid, cmd, rest)
                if result is None:
                    err_ev = {"type":"tool_error","tool":cmd,"label":f"/{cmd} is not a recognised slash command or is inactive."}
                    tool_log.append(err_ev)
                    yield f"data: {json.dumps(err_ev)}\n\n"
                    yield f"data: {json.dumps({'type':'error','message':f'Unknown slash command: /{cmd}'})}\n\n"
                    return
                if result and not result.get("error"):
                    payload = json.dumps({"type":"tool","app":result.get("app",cmd),"action":result.get("action",""),"result":result.get("result",{})})
                    r = result.get("result", {})
                    count = len(r) if isinstance(r, list) else len(r.get("results", [])) if isinstance(r, dict) and isinstance(r.get("results"), list) else None
                    done_ev = {"type":"tool_result","tool":result.get("app",cmd),"action":result.get("action",""),"label":f"Found {count} item(s)" if count is not None else "Tool completed","count":count}
                    tool_log.append(done_ev)
                    yield f"data: {json.dumps(done_ev)}\n\n"
                    async for ev in _stream_text(payload): yield f"data: {ev}\n\n"
                    elapsed = int((time.time()-t0)*1000)
                    now = _utcnow(); mid = _new_id()
                    await db_execute(
                        "INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,created_at) VALUES(?,?,?,'user',?,?,?)",
                        (_new_id(), sid, uid, req.message, active_model_id, now))
                    await db_execute(
                        "INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,latency_ms,tool_calls_json,mode,created_at)"
                        " VALUES(?,?,?,'assistant',?,?,?,?,'tool',?)",
                        (mid, sid, uid, payload, active_model_id, elapsed, json.dumps(tool_log), now))
                    await db_execute("UPDATE chat_sessions SET last_message_at=?,turn_count=turn_count+1,updated_at=? WHERE id=?",
                                     (now, now, sid))
                    if skills:
                        await _log_skill_usage(skills, skill_matches, uid, sid, mid)
                    context = await _context_window_status(sid, uid, active_model_id)
                    yield f"data: {json.dumps({'type':'done','message_id':mid,'latency_ms':elapsed,'tokens':{'input':_count_tokens(req.message),'output':_count_tokens(payload)},'mode':'tool','context':context,'model_id':active_model_id,'model_label':await _model_display_name(active_model_id)})}\n\n"
                    return
                if result and result.get("error"):
                    err_ev = {"type":"tool_error","tool":cmd,"label":str(result.get("error"))}
                    tool_log.append(err_ev)
                    yield f"data: {json.dumps(err_ev)}\n\n"
                    yield f"data: {json.dumps({'type':'error','message':result.get('error')})}\n\n"
                    return

            web_results = None
            if use_web:
                start_ev = {"type":"tool_start","tool":"web_search","label":"Searching web"}
                tool_log.append(start_ev)
                yield f"data: {json.dumps(start_ev)}\n\n"
                web_results = await _web_search(req.message[:500], 6)
                result_ev = {"type":"tool_result","tool":"web_search","label":f"Found {len(web_results)} web result(s)","count":len(web_results),"results":web_results[:6]}
                tool_log.append(result_ev)
                yield f"data: {json.dumps(result_ev)}\n\n"

            # Vision
            if image_inputs:
                extracted_images: List[Dict[str, Any]] = []
                for img_idx, img in enumerate(image_inputs, 1):
                    img_start = {"type":"tool_start","tool":"image_to_text","label":f"Reading image {img_idx}/{len(image_inputs)}"}
                    tool_log.append(img_start)
                    yield f"data: {json.dumps(img_start)}\n\n"
                    try:
                        image_text_result = await _extract_image_text(img["url"], req.message, active_model_id, uid)
                        image_text = image_text_result["text"]
                        img_done = {
                            "type":"tool_result",
                            "tool":"image_to_text",
                            "label":"Image converted to text with low confidence" if image_text_result.get("weak") else "Image converted to text",
                            "count":1,
                            "model_id":image_text_result.get("model_id"),
                            "model_label":image_text_result.get("model_label") or "Image-to-text model",
                        }
                        tool_log.append({**img_done, "image_name":img.get("name"), "excerpt": image_text[:2000]})
                        yield f"data: {json.dumps(img_done)}\n\n"
                        extracted_images.append({"name":img.get("name"), "mime":img.get("mime"), **image_text_result})
                    except Exception as exc:
                        logger.warning("[VISION] image-to-text extraction failed: %s", exc)
                        err_ev = {"type":"tool_error","tool":"image_to_text","label":f"Image {img_idx} text extraction failed."}
                        tool_log.append({**err_ev, "image_name":img.get("name"), "error":str(exc)[:500]})
                        yield f"data: {json.dumps(err_ev)}\n\n"
                        extracted_images.append({"name":img.get("name"), "mime":img.get("mime"), "text":f"[Extraction failed: {exc}]", "weak":True})
                context_parts = [
                    req.message,
                    "",
                    "[Attached image analysis]",
                    "Use the extracted visual/OCR content below as the source of truth for attached images. "
                    "Answer from this content and mention uncertainty where extraction is low-confidence.",
                ]
                for img_idx, item in enumerate(extracted_images, 1):
                    context_parts.append(f"\nImage {img_idx}: {item.get('name') or 'attached image'} ({item.get('mime') or 'image/*'})")
                    context_parts.append(f"Extractor: {item.get('model_label') or 'unknown'}")
                    if item.get("weak"):
                        context_parts.append("Confidence: low - verify carefully and avoid overclaiming.")
                    context_parts.append(str(item.get("text") or "").strip())
                effective_message = "\n".join(context_parts).strip()
                messages = await _build_context(sid, uid, effective_message, use_rag, active_model_id, web_results, plan_mode, skill_context, mcp_context)

            # Agent mode
            agent_tools = tool_opts.get("mcp_tools") or []
            if not isinstance(agent_tools, list):
                agent_tools = []
            if req.agent_mode or bool(tool_opts.get("agent_mode")):
                start_ev = {"type":"tool_start","tool":"agent","label":"Tool Agent is planning actions"}
                tool_log.append(start_ev)
                yield f"data: {json.dumps(start_ev)}\n\n"
                agent_context = await _format_agent_tool_context(uid, tool_opts)
                result_text, agent_tool_log = await _agent_loop(
                    effective_message, user, agent_tools, active_model_id, tool_context=agent_context)
                for step_ev in agent_tool_log:
                    tool_log.append({"type":"agent_step", **step_ev})
                async for ev in _stream_text(result_text): yield f"data: {ev}\n\n"
                tok_in = _count_tokens(req.message); tok_out = _count_tokens(result_text)
                elapsed = int((time.time()-t0)*1000)
                now = _utcnow(); mid = _new_id()
                await db_execute(
                    "INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,created_at) VALUES(?,?,?,'user',?,?,?)",
                    (_new_id(), sid, uid, req.message, active_model_id, now))
                await db_execute(
                    "INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,latency_ms,tool_calls_json,mode,created_at)"
                    " VALUES(?,?,?,'assistant',?,?,?,?,'agent',?)",
                    (mid, sid, uid, result_text, active_model_id, elapsed, json.dumps(tool_log), now))
                await db_execute("UPDATE chat_sessions SET model_id=?,last_message_at=?,turn_count=turn_count+1,updated_at=? WHERE id=?",
                                 (active_model_id, now, now, sid))
                context = await _context_window_status(sid, uid, active_model_id)
                yield f"data: {json.dumps({'type':'done','message_id':mid,'content':result_text,'latency_ms':elapsed,'tokens':{'input':tok_in,'output':tok_out},'mode':'agent','context':context,'model_id':active_model_id,'model_label':await _model_display_name(active_model_id),'tool_calls':agent_tool_log})}\n\n"
                asyncio.create_task(_maybe_summarize(sid, uid, active_model_id))
                return

            # Thinking mode
            use_thinking = req.force_thinking if req.force_thinking is not None else _needs_thinking(req.message)
            if use_thinking:
                yield f"data: {json.dumps({'type':'thinking_start'})}\n\n"
                messages = await _build_context(sid, uid, effective_message, use_rag, active_model_id, web_results, plan_mode, skill_context, mcp_context)
                think_sys = [{"role":"system","content":_THINKING_SYSTEM}] + \
                            [m for m in messages if m["role"] != "system"]
                llm_result = await _llm_text_with_fallback(think_sys, active_model_id, max_tokens=4096, temperature=0.6, user_id=uid)
                for ev in llm_result.get("events", []):
                    tool_log.append(ev)
                    yield f"data: {json.dumps(ev)}\n\n"
                final_model_id = llm_result.get("model_id") or active_model_id
                content = llm_result.get("content") or ""
                thinking_match = re.search(r'<thinking>(.*?)</thinking>', content, re.DOTALL)
                answer_match = re.search(r'<answer>(.*?)</answer>', content, re.DOTALL)
                thinking_content = thinking_match.group(1).strip() if thinking_match else ""
                full_response = answer_match.group(1).strip() if answer_match else content
                if not full_response.strip():
                    logger.warning("[CHAT] thinking response was empty; recovering with fallback model")
                    full_response = await _recover_empty_reply(messages, final_model_id)
                yield f"data: {json.dumps({'type':'thinking_content','content':thinking_content})}\n\n"
                yield f"data: {json.dumps({'type':'thinking_done'})}\n\n"
                async for ev in _stream_text(full_response): yield f"data: {ev}\n\n"
                tok_in = _count_tokens(req.message); tok_out = _count_tokens(full_response)
                elapsed = int((time.time()-t0)*1000)
                now = _utcnow(); mid = _new_id()
                await db_execute(
                    "INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,latency_ms,thinking_content,mode,created_at)"
                    " VALUES(?,?,?,'user',?,?,0,'','thinking',?)",
                    (_new_id(), sid, uid, req.message, final_model_id, now))
                await db_execute(
                    "INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,latency_ms,thinking_content,tool_calls_json,mode,created_at)"
                    " VALUES(?,?,?,'assistant',?,?,?,?,?,'thinking',?)",
                    (mid, sid, uid, full_response, final_model_id, elapsed, thinking_content, json.dumps(tool_log), now))
                await db_execute("UPDATE chat_sessions SET model_id=?,last_message_at=?,turn_count=turn_count+1,updated_at=? WHERE id=?",
                                 (final_model_id, now, now, sid))
                if skills:
                    await _log_skill_usage(skills, skill_matches, uid, sid, mid)
                context = await _context_window_status(sid, uid, final_model_id)
                yield f"data: {json.dumps({'type':'done','message_id':mid,'content':full_response,'latency_ms':elapsed,'tokens':{'input':tok_in,'output':tok_out},'mode':'thinking','context':context,'model_id':final_model_id,'model_label':await _model_display_name(final_model_id)})}\n\n"
                if len(req.message.strip()) >= 30:
                    asyncio.create_task(_auto_extract_memories(uid, req.message, full_response))
                asyncio.create_task(_maybe_summarize(sid, uid, final_model_id))
                return

            # Normal streaming
            if not image_inputs:
                messages = await _build_context(sid, uid, req.message, use_rag, active_model_id, web_results, plan_mode, skill_context, mcp_context)
            tok_in = sum(_count_tokens(m.get("content","") if isinstance(m.get("content"), str) else "")
                        for m in messages)
            full_response = ""
            final_model_id = active_model_id
            streamed_direct = False
            fallback_model_id = active_model_id
            nvidia_meta = await _nvidia_stream_model_meta(active_model_id)
            if nvidia_meta and not image_inputs:
                start_ev = {"type":"tool_progress","tool":"model","label":f"Asking {nvidia_meta['label']}..."}
                tool_log.append(start_ev)
                yield f"data: {json.dumps(start_ev)}\n\n"
                thinking_content = ""
                try:
                    async for kind, payload, meta in _stream_nvidia_text_once(messages, active_model_id, max_tokens=2048, temperature=0.7):
                        if kind == "status":
                            status_ev = {"type":"tool_progress","tool":"model","label":payload}
                            tool_log.append(status_ev)
                            yield f"data: {json.dumps(status_ev)}\n\n"
                        elif kind == "thinking":
                            thinking_content += payload
                            yield f"data: {json.dumps({'type':'thinking_content','content':thinking_content[-12000:]})}\n\n"
                        else:
                            full_response += payload
                            yield f"data: {json.dumps({'type':'delta','content':payload})}\n\n"
                    streamed_direct = bool(full_response.strip())
                    final_model_id = nvidia_meta["model_id"]
                    if not streamed_direct:
                        logger.warning("[CHAT] NVIDIA stream returned empty response for %s", active_model_id)
                except Exception as exc:
                    logger.warning("[CHAT] NVIDIA direct stream failed for %s: %s", active_model_id, exc)
                    err_ev = {
                        "type":"tool_error","tool":"model",
                        "model_id":active_model_id,
                        "model_label":nvidia_meta["label"],
                        "label":_provider_failure_label(nvidia_meta["label"], exc),
                    }
                    tool_log.append(err_ev)
                    yield f"data: {json.dumps(err_ev)}\n\n"
                    candidates = await _model_fallback_candidates(active_model_id, tok_in, uid)
                    fallback_model_id = next((c for c in candidates if c != active_model_id), active_model_id)
            hf_meta = await _hf_stream_model_meta(active_model_id)
            if not streamed_direct and hf_meta and not image_inputs:
                start_ev = {"type":"tool_progress","tool":"model","label":f"Asking {hf_meta['label']}..."}
                tool_log.append(start_ev)
                yield f"data: {json.dumps(start_ev)}\n\n"
                try:
                    async for delta, meta in _stream_hf_text_once(messages, active_model_id, max_tokens=768, temperature=0.2):
                        full_response += delta
                        yield f"data: {json.dumps({'type':'delta','content':delta})}\n\n"
                    streamed_direct = bool(full_response.strip())
                    final_model_id = hf_meta["model_id"]
                    if not streamed_direct:
                        logger.warning("[CHAT] HF stream returned empty response for %s", active_model_id)
                except Exception as exc:
                    logger.warning("[CHAT] HF direct stream failed for %s: %s", active_model_id, exc)
                    err_ev = {
                        "type":"tool_error","tool":"model",
                        "model_id":active_model_id,
                        "model_label":hf_meta["label"],
                        "label":_provider_failure_label(hf_meta["label"], exc),
                    }
                    tool_log.append(err_ev)
                    yield f"data: {json.dumps(err_ev)}\n\n"
                    candidates = await _model_fallback_candidates(active_model_id, tok_in, uid)
                    fallback_model_id = next((c for c in candidates if c != active_model_id), active_model_id)

            if not streamed_direct:
                llm_result = await _llm_text_with_fallback(
                    messages, fallback_model_id,
                    max_tokens=3072 if image_inputs else 2048,
                    temperature=0.25 if image_inputs else 0.7,
                    user_id=uid)
                for ev in llm_result.get("events", []):
                    tool_log.append(ev)
                    yield f"data: {json.dumps(ev)}\n\n"
                final_model_id = llm_result.get("model_id") or active_model_id
                full_response = llm_result.get("content") or ""
                if not full_response.strip():
                    logger.warning("[CHAT] normal response was empty; recovering with fallback model")
                    full_response = await _recover_empty_reply(messages, final_model_id)
                async for ev in _stream_text(full_response):
                    yield f"data: {ev}\n\n"

            tok_out = _count_tokens(full_response); elapsed = int((time.time()-t0)*1000)
            now = _utcnow(); mid = _new_id()
            await db_execute(
                "INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,created_at) VALUES(?,?,?,'user',?,?,?)",
                (_new_id(), sid, uid, req.message, final_model_id, now))
            await db_execute(
                "INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,latency_ms,tool_calls_json,mode,created_at)"
                " VALUES(?,?,?,'assistant',?,?,?,?,'normal',?)",
                (mid, sid, uid, full_response, final_model_id, elapsed, json.dumps(tool_log), now))
            await db_execute("UPDATE chat_sessions SET model_id=?,last_message_at=?,turn_count=turn_count+1,updated_at=? WHERE id=?",
                             (final_model_id, now, now, sid))
            if skills:
                await _log_skill_usage(skills, skill_matches, uid, sid, mid)
            session_row = await db_fetchone("SELECT turn_count FROM chat_sessions WHERE id=?",(sid,))
            non_tool_count = await db_count("SELECT COUNT(*) as c FROM chat_history WHERE session_id=? AND role='user' AND mode != 'tool'",(sid,))
            if session_row and non_tool_count <= 1:
                asyncio.create_task(_auto_title(sid, req.message, full_response))
            context = await _context_window_status(sid, uid, final_model_id)
            yield f"data: {json.dumps({'type':'done','message_id':mid,'content':full_response,'latency_ms':elapsed,'tokens':{'input':tok_in,'output':tok_out},'mode':'normal','context':context,'model_id':final_model_id,'model_label':await _model_display_name(final_model_id)})}\n\n"
            if len(req.message.strip()) >= 30:
                asyncio.create_task(_auto_extract_memories(uid, req.message, full_response))
            asyncio.create_task(_maybe_summarize(sid, uid, final_model_id))
        except Exception as e:
            logger.error("[CHAT] stream error: %s", traceback.format_exc())
            yield f"data: {json.dumps({'type':'error','message':str(e)})}\n\n"

    return StreamingResponse(_gen(), media_type="text/event-stream",
                             headers={"Cache-Control":"no-cache","X-Accel-Buffering":"no","Connection":"keep-alive"})

@app.get("/chat/stream")
async def chat_stream_get(
    message: str, model_id: str = "llama-3.3-70b-versatile", use_rag: bool = True,
    session_id: Optional[str] = None, force_thinking: Optional[bool] = None,
    web_search: bool = False, plan_mode: bool = False,
    skill_ids: str = "", tools_json: str = "", image_url: Optional[str] = None,
    user: Dict = Depends(_get_current_user)) -> StreamingResponse:
    tools = _safe_json_loads(tools_json, {}) if tools_json else {}
    if not isinstance(tools, dict):
        tools = {}
    parsed_skill_ids = [s.strip() for s in (skill_ids or "").split(",") if s.strip()]
    if parsed_skill_ids and "skill_ids" not in tools:
        tools["skill_ids"] = parsed_skill_ids
    req = ChatReq(message=message, model_id=model_id, use_rag=use_rag,
                  session_id=session_id, force_thinking=force_thinking,
                  web_search=web_search, plan_mode=plan_mode,
                  skill_ids=parsed_skill_ids, tools=tools, image_url=image_url)
    return await chat_stream_post(req, user)

@app.post("/chat/message")
async def chat_message(req: ChatReq, background: BackgroundTasks, user: Dict = Depends(_get_current_user)):
    await _check_rate_limit(user, "messages_per_day")
    uid = user.get("id") or user.get("sub","")
    requested_model_id = _canonical_model_id(req.model_id)
    active_model_id, _ = await _auto_route_model(requested_model_id, req.message, uid)
    sid = await _get_or_create_session(uid, req.session_id, active_model_id)
    t0 = time.time()
    tool_opts = _chat_tool_options(req)
    use_rag = bool(tool_opts.get("rag", req.use_rag))
    use_web = bool(tool_opts.get("web_search", req.web_search))
    plan_mode = bool(tool_opts.get("plan_mode", req.plan_mode))
    skills, skill_matches = await _select_skills_for_message(req.message, tool_opts.get("skill_ids", req.skill_ids) or [])
    web_results = await _web_search(req.message, 6) if use_web else None
    tool_log: List[Dict[str, Any]] = []
    identity_reply = await _identity_recall_reply(uid, req.message, user)
    if identity_reply:
        elapsed = int((time.time()-t0)*1000); now = _utcnow(); mid = _new_id()
        tool_log.append({"type":"tool_result","tool":"identity","label":"Recovered your name from account context","count":1})
        await db_execute("INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,created_at) VALUES(?,?,?,'user',?,?,?)",
                         (_new_id(), sid, uid, req.message, active_model_id, now))
        await db_execute("INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,latency_ms,tool_calls_json,mode,created_at) VALUES(?,?,?,'assistant',?,?,?,?,'normal',?)",
                         (mid, sid, uid, identity_reply, active_model_id, elapsed, json.dumps(tool_log), now))
        await db_execute("UPDATE chat_sessions SET model_id=?,last_message_at=?,turn_count=turn_count+1,updated_at=? WHERE id=?",(active_model_id,now,now,sid))
        context = await _context_window_status(sid, uid, active_model_id)
        return {"message_id":mid,"content":identity_reply,"session_id":sid,"latency_ms":elapsed,"context":context,"model_id":active_model_id,"model_label":await _model_display_name(active_model_id)}

    missing_context_reply = await _verified_context_missing_reply(uid, req.message, user)
    if missing_context_reply:
        elapsed = int((time.time()-t0)*1000); now = _utcnow(); mid = _new_id()
        tool_log.append({"type":"tool_result","tool":"evidence_guard","label":"No verified context found for that private detail","count":0})
        await db_execute("INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,created_at) VALUES(?,?,?,'user',?,?,?)",
                         (_new_id(), sid, uid, req.message, active_model_id, now))
        await db_execute("INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,latency_ms,tool_calls_json,mode,created_at) VALUES(?,?,?,'assistant',?,?,?,?,'guard',?)",
                         (mid, sid, uid, missing_context_reply, active_model_id, elapsed, json.dumps(tool_log), now))
        await db_execute("UPDATE chat_sessions SET model_id=?,last_message_at=?,turn_count=turn_count+1,updated_at=? WHERE id=?",(active_model_id,now,now,sid))
        context = await _context_window_status(sid, uid, active_model_id)
        return {"message_id":mid,"content":missing_context_reply,"session_id":sid,"latency_ms":elapsed,"context":context,"model_id":active_model_id,"model_label":await _model_display_name(active_model_id)}

    effective_message, _ = await _build_image_context(req, uid, active_model_id, tool_log)
    messages = await _build_context(sid, uid, effective_message, use_rag, active_model_id,
                                    web_results, plan_mode, _format_skill_context(skills),
                                    _format_mcp_context(tool_opts.get("mcp_tools", [])))
    llm_result = await _llm_text_with_fallback(messages, active_model_id, max_tokens=4096, user_id=uid)
    reply = llm_result.get("content", "")
    final_model_id = llm_result.get("model_id") or active_model_id
    elapsed = int((time.time()-t0)*1000); now = _utcnow(); mid = _new_id()
    await db_execute("INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,created_at) VALUES(?,?,?,'user',?,?,?)",
                     (_new_id(), sid, uid, req.message, final_model_id, now))
    tool_log = tool_log + [{"type":"skill_loaded","skill_id":s["id"],"name":s["name"],"matched_by":skill_matches[i] if i < len(skill_matches) else ""} for i,s in enumerate(skills)]
    tool_log.extend(llm_result.get("events", []))
    await db_execute("INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,latency_ms,tool_calls_json,mode,created_at) VALUES(?,?,?,'assistant',?,?,?,?,'normal',?)",
                     (mid, sid, uid, reply, final_model_id, elapsed, json.dumps(tool_log), now))
    if skills:
        await _log_skill_usage(skills, skill_matches, uid, sid, mid)
    await db_execute("UPDATE chat_sessions SET model_id=?,last_message_at=?,turn_count=turn_count+1,updated_at=? WHERE id=?",(final_model_id,now,now,sid))
    context = await _context_window_status(sid, uid, final_model_id)
    background.add_task(_auto_extract_memories, uid, req.message, reply)
    background.add_task(_maybe_summarize, sid, uid, final_model_id)
    return {"message_id":mid,"content":reply,"session_id":sid,"latency_ms":elapsed,"context":context,"model_id":final_model_id,"model_label":await _model_display_name(final_model_id)}

@app.post("/chat/messages/{mid}/feedback")
async def message_feedback(mid: str, body: dict, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    await db_execute(
        "INSERT INTO message_feedback(id,message_id,user_id,rating,feedback,created_at) VALUES(?,?,?,?,?,?)"
        " ON CONFLICT(message_id,user_id) DO UPDATE SET rating=excluded.rating,feedback=excluded.feedback",
        (_new_id(), mid, uid, body.get("rating",0), body.get("feedback",""), _utcnow()))
    return {"ok":True}

@app.post("/chat/feedback")
async def message_feedback_alias(body: dict, user: Dict = Depends(_get_current_user)):
    mid = body.get("message_id","")
    if not mid: raise HTTPException(400, "message_id required")
    return await message_feedback(mid, body, user)

@app.delete("/chat/messages/{mid}")
async def delete_message(mid: str, user: Dict = Depends(_get_current_user)):
    """Soft-delete (hide) a single message by id."""
    uid = user.get("id") or user.get("sub","")
    row = await db_fetchone("SELECT id FROM chat_history WHERE id=? AND user_id=?", (mid, uid))
    if not row: raise HTTPException(404, "Message not found")
    await db_execute("UPDATE chat_history SET is_hidden=1,updated_at=? WHERE id=?",
                     (_utcnow(), mid)) if "updated_at" in row else \
    await db_execute("UPDATE chat_history SET is_hidden=1 WHERE id=?", (mid,))
    return {"ok": True}

@app.patch("/chat/messages/{mid}")
async def edit_message(mid: str, body: dict, user: Dict = Depends(_get_current_user)):
    """Edit the content of a user message and record the change in message_edits."""
    uid = user.get("id") or user.get("sub","")
    row = await db_fetchone(
        "SELECT id,content,session_id FROM chat_history WHERE id=? AND user_id=? AND role='user'", (mid, uid))
    if not row: raise HTTPException(404, "Message not found or not a user message")
    new_content = (body.get("content") or "").strip()
    if not new_content: raise HTTPException(400, "content required")
    now = _utcnow()
    await db_execute(
        "INSERT INTO message_edits(id,message_id,old_content,new_content,edited_at) VALUES(?,?,?,?,?)",
        (_new_id(), mid, row["content"], new_content, now))
    await db_execute(
        "UPDATE chat_history SET content=?,edit_count=edit_count+1 WHERE id=?",
        (new_content, mid))
    return {"ok": True, "message_id": mid, "content": new_content}

@app.post("/sessions/{sid}/clear")
async def clear_session_messages(sid: str, user: Dict = Depends(_get_current_user)):
    """Hide all messages in a session (keeps session metadata intact)."""
    uid = user.get("id") or user.get("sub","")
    s = await db_fetchone("SELECT id FROM chat_sessions WHERE id=? AND user_id=?", (sid, uid))
    if not s: raise HTTPException(404, "Session not found")
    await db_execute("UPDATE chat_history SET is_hidden=1 WHERE session_id=?", (sid,))
    await db_execute("UPDATE chat_sessions SET turn_count=0,updated_at=? WHERE id=?", (_utcnow(), sid))
    return {"ok": True}

@app.post("/chat/regenerate/{mid}")
async def regenerate_message(mid: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    msg = await db_fetchone("SELECT * FROM chat_history WHERE id=? AND user_id=? AND role='assistant'",(mid,uid))
    if not msg: raise HTTPException(404)
    await db_execute("UPDATE chat_history SET is_hidden=1 WHERE id=?",(mid,))
    session = await db_fetchone("SELECT * FROM chat_sessions WHERE id=?",(msg["session_id"],))
    t0 = time.time()
    messages = await _build_context(msg["session_id"], uid, "", True, session["model_id"])
    reply = await _llm_text(messages, session["model_id"], max_tokens=4096, temperature=0.8)
    if not reply.strip():
        reply = await _recover_empty_reply(messages, session["model_id"])
    now = _utcnow(); new_mid = _new_id()
    await db_execute(
        "INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,latency_ms,parent_msg_id,created_at)"
        " VALUES(?,?,?,'assistant',?,?,?,?,?)",
        (new_mid, msg["session_id"], uid, reply, session["model_id"],
         int((time.time()-t0)*1000), mid, now))
    return {"message_id":new_mid,"content":reply}

@app.post("/chat/regenerate")
async def regenerate_message_alias(body: dict, user: Dict = Depends(_get_current_user)):
    mid = body.get("message_id","")
    if not mid: raise HTTPException(400, "message_id required")
    return await regenerate_message(mid, user)

@app.post("/chat/multi-model")
async def multi_model_chat(req: MultiModelChatReq, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    allowed = set(await _allowed_model_ids_for_user(uid))
    model_ids = []
    for m in req.model_ids:
        mid = _canonical_model_id(m.strip()) if m and m.strip() else ""
        if mid and mid in allowed and mid not in model_ids:
            model_ids.append(mid)
    model_ids = model_ids[:6]
    if len(model_ids) < 2:
        raise HTTPException(400, "Select at least two enabled models")
    for _ in model_ids:
        await _check_rate_limit(user, "messages_per_day")
    sid = await _get_or_create_session(uid, req.session_id, model_ids[0])
    async def _call_one(mid: str) -> Dict:
        t0 = time.time()
        requested_mid = _canonical_model_id(mid)
        try:
            active_mid, route_reason = await _auto_route_model(requested_mid, req.message, uid)
            messages = await _build_context(sid, uid, req.message, req.use_rag, active_mid)
            llm_result = await _llm_text_with_fallback(messages, active_mid, max_tokens=2048, temperature=0.7, user_id=uid)
            final_mid = llm_result.get("model_id") or active_mid
            events = llm_result.get("events", [])
            if route_reason == "long_input" and active_mid != requested_mid:
                events = [{
                    "type":"model_switch", "from_model_id":requested_mid,
                    "model_id":active_mid, "model_label":await _model_display_name(active_mid),
                    "selected":False,
                    "label":f"Long input detected. Routed this run to {await _model_display_name(active_mid)}."
                }] + events
            return {
                "requested_model_id":requested_mid,
                "model_id":final_mid,
                "model_name":llm_result.get("provider_model") or final_mid,
                "model_label":llm_result.get("model_label") or await _model_display_name(final_mid),
                "response":llm_result.get("content", ""),
                "events":events,
                "latency_ms":int((time.time()-t0)*1000),
                "error":None if llm_result.get("content") else "empty response"
            }
        except Exception as e:
            return {"requested_model_id":requested_mid,"model_id":requested_mid,
                    "model_name":requested_mid,"response":None,"error":str(e),
                    "latency_ms":int((time.time()-t0)*1000)}
    results = await asyncio.gather(*[_call_one(m) for m in model_ids])
    return {"session_id":sid,"message":req.message,"responses":list(results),"charged_messages":len(model_ids)}

# ══════════════════════════════════════════════════════════════════════════════
# §25  WEBSOCKET
# ══════════════════════════════════════════════════════════════════════════════

@app.websocket("/ws/{user_id}")
async def websocket_endpoint(user_id: str, ws: WebSocket) -> None:
    await ws.accept()
    try:
        auth_msg = await asyncio.wait_for(ws.receive_json(), timeout=15.0)
    except asyncio.TimeoutError:
        await ws.close(code=4001, reason="Auth timeout"); return
    token = auth_msg.get("token","")
    payload = _decode_jwt(token) if token else None
    if not payload or payload.get("sub") != user_id:
        await ws.send_json({"type":"error","message":"Unauthorized"})
        await ws.close(code=4003, reason="Unauthorized"); return
    ws_manager._connections[user_id].add(ws)
    await ws.send_json({"type":"connected","user_id":user_id})
    try:
        while True:
            msg = await ws.receive_json(); mtype = msg.get("type","")
            if mtype == "ping": await ws.send_json({"type":"pong","ts":_utcnow()})
            elif mtype == "typing": await ws.send_json({"type":"typing_ack"})
    except WebSocketDisconnect: pass
    except Exception as e: logger.warning("[WS] %s error: %s", user_id[:8], e)
    finally: ws_manager.disconnect(user_id, ws)

# ══════════════════════════════════════════════════════════════════════════════
# §26  DOCUMENTS / RAG
# ══════════════════════════════════════════════════════════════════════════════

@app.post("/documents/upload")
async def upload_document(file: UploadFile = File(...), background_tasks: BackgroundTasks = None,
                           user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    await _check_rate_limit(user, "documents")
    data = await file.read()
    if len(data) > MAX_FILE_SIZE: raise HTTPException(413, "File too large (max 50MB)")
    doc_id = _new_id(); fname = f"{doc_id}_{re.sub(r'[^a-zA-Z0-9._-]','_',file.filename or 'file')}"
    path = UPLOADS_DIR / fname; path.write_bytes(data)
    now = _utcnow()
    await db_execute(
        "INSERT INTO documents(id,user_id,filename,original_name,file_type,file_size_bytes,uploaded_at)"
        " VALUES(?,?,?,?,?,?,?)",
        (doc_id, uid, fname, file.filename or "file", file.content_type or "", len(data), now))
    if background_tasks: background_tasks.add_task(_index_document, doc_id, uid, path)
    else: asyncio.create_task(_index_document(doc_id, uid, path))
    return await db_fetchone("SELECT * FROM documents WHERE id=?", (doc_id,))

@app.get("/documents")
async def list_documents(user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    docs = await db_fetchall(
        "SELECT id,original_name,file_type,file_size_bytes,chunk_count,is_indexed,index_error,uploaded_at"
        " FROM documents WHERE user_id=? ORDER BY uploaded_at DESC", (uid,))
    return {"documents":docs}

@app.delete("/documents/{doc_id}")
async def delete_document(doc_id: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    doc = await db_fetchone("SELECT * FROM documents WHERE id=? AND user_id=?", (doc_id, uid))
    if not doc: raise HTTPException(404)
    path = UPLOADS_DIR / doc["filename"]
    if path.exists(): path.unlink()
    col = _get_collection(uid)
    if col:
        try: col.delete(where={"doc_id":doc_id})
        except Exception: pass
    await db_execute("DELETE FROM documents WHERE id=?", (doc_id,))
    return {"ok":True}

@app.post("/rag/search")
async def rag_search_endpoint(body: dict, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    chunks = await _rag_search(uid, body.get("query",""), body.get("k", TOP_K_RETRIEVAL))
    return {"chunks":chunks}

@app.post("/web/search")
async def web_search_endpoint(body: dict, user: Dict = Depends(_get_current_user)):
    results = await _web_search(body.get("query", ""), body.get("count", 5))
    return {"results": results}

@app.get("/skills")
async def list_skills(user: Dict = Depends(_get_current_user)):
    rows = await db_fetchall(
        "SELECT id,name,description,activation_keywords,is_enabled,created_at,updated_at "
        "FROM skills WHERE is_enabled=1 ORDER BY name")
    return {"skills": [_normalise_skill_row(r) for r in rows]}

@app.get("/skills/{skill_id}")
async def get_skill(skill_id: str, user: Dict = Depends(_get_current_user)):
    row = await db_fetchone(
        "SELECT id,name,description,activation_keywords,instructions,resources_json,"
        "script_metadata_json,is_enabled,created_at,updated_at FROM skills WHERE id=? AND is_enabled=1",
        (skill_id,))
    if not row:
        raise HTTPException(404, "Skill not found")
    return _normalise_skill_row(row)

# ══════════════════════════════════════════════════════════════════════════════
# §27  MEMORY
# ══════════════════════════════════════════════════════════════════════════════

@app.get("/memories")
async def list_memories(user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    mems = await db_fetchall(
        "SELECT id,key,value,source,confidence,reinforcement_count,last_reinforced,created_at"
        " FROM memories WHERE user_id=? AND is_active=1 ORDER BY confidence DESC,last_reinforced DESC", (uid,))
    return {"memories":mems}

@app.post("/memories")
async def create_memory(body: MemoryCreate, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    count = await db_count("SELECT COUNT(*) as c FROM memories WHERE user_id=? AND is_active=1",(uid,))
    if count >= MAX_MEMORIES: raise HTTPException(400, "Memory limit reached")
    await _upsert_memory(uid, body.key, body.value, "manual", 1.0)
    return await db_fetchone("SELECT * FROM memories WHERE user_id=? AND key=?", (uid, body.key))

@app.delete("/memories/{mid}")
async def delete_memory(mid: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    await db_execute("UPDATE memories SET is_active=0 WHERE id=? AND user_id=?",(mid,uid))
    return {"ok":True}

@app.delete("/memories")
async def clear_memories(user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    await db_execute("UPDATE memories SET is_active=0 WHERE user_id=?",(uid,))
    return {"ok":True}

# ══════════════════════════════════════════════════════════════════════════════
# §28  CONNECTORS — OAUTH FLOWS (FIXED FULL PKCE)
# ══════════════════════════════════════════════════════════════════════════════

@app.get("/connectors/oauth/{app_id}/init")
async def oauth_init(app_id: str, request: Request,
                     state: str = "", token: str = ""):
    """Step 1: Build the correct OAuth URL and redirect user's browser to it."""
    providers = _oauth_providers()
    prov = providers.get(app_id)
    redirect_uri = _oauth_redirect_uri(app_id, request)

    if not prov:
        return HTMLResponse(_OAUTH_ERROR_HTML.format(
            message=f"Unknown connector: {app_id}",
            detail="This connector is not supported."), status_code=404)

    client_id = prov.get("client_id","")
    if not client_id:
        # Determine env var names for setup page
        app_upper = app_id.upper().replace("-","_")
        id_var = f"{app_upper}_CLIENT_ID"
        secret_var = f"{app_upper}_CLIENT_SECRET"
        if app_id in ("gmail","google_calendar","google_drive","google_sheets","bigquery","google_meet"):
            id_var, secret_var = "GOOGLE_CLIENT_ID", "GOOGLE_CLIENT_SECRET"
        elif app_id in ("outlook","onedrive","microsoft_teams","power_bi"):
            id_var, secret_var = "MICROSOFT_CLIENT_ID", "MICROSOFT_CLIENT_SECRET"
        elif app_id == "jira":
            id_var, secret_var = "ATLASSIAN_CLIENT_ID", "ATLASSIAN_CLIENT_SECRET"
        return HTMLResponse(_OAUTH_SETUP_HTML.format(
            name=app_id.replace("_"," ").title(), id_var=id_var, secret_var=secret_var,
            redirect_uri=redirect_uri))

    user_id = _user_id_from_token_str(token)
    if not user_id:
        return HTMLResponse(_OAUTH_ERROR_HTML.format(
            message="Authentication required",
            detail="Please log in to JAZZ first and try connecting again."), status_code=401)

    client_state = state or ""

    # PKCE if required
    code_verifier = None
    extra_params: Dict = {}
    if prov.get("pkce"):
        code_verifier, code_challenge = _pkce_pair()
        extra_params["code_challenge"] = code_challenge
        extra_params["code_challenge_method"] = "S256"

    # Use a signed OAuth state for the provider, while preserving the
    # frontend-generated state key for completion polling.
    state = _make_oauth_state_token(user_id, app_id, code_verifier or "", redirect_uri, client_state)

    # Clean expired pending states (> 30 min)
    expired = [k for k,v in _pending_oauth.items() if time.time()-v.get("ts",0) > 1800]
    for k in expired: _pending_oauth.pop(k,None); _completed_oauth.pop(k,None)

    # Save state to DB FIRST (primary source — survives server restarts & multi-worker)
    exp_at = (datetime.now(timezone.utc)+timedelta(minutes=30)).isoformat()
    await db_execute(
        "INSERT OR REPLACE INTO oauth_states(id,user_id,connector_type,state,code_verifier,redirect_uri,expires_at,created_at)"
        " VALUES(?,?,?,?,?,?,?,?)",
        (_new_id(), user_id, app_id, state, code_verifier or "", redirect_uri, exp_at, _utcnow()))

    # Also cache in memory for fast lookup
    _pending_oauth[state] = {
        "user_id": user_id, "connector_type": app_id,
        "code_verifier": code_verifier, "redirect_uri": redirect_uri,
        "client_state": client_state, "ts": time.time()
    }

    # Build authorize URL — CRITICAL: include response_type=code
    params: Dict = {
        "client_id":     client_id,
        "redirect_uri":  redirect_uri,
        "response_type": "code",          # ← Fixed: was missing in some connectors
        "state":         state,
    }
    scope = prov.get("scope","")
    if scope: params["scope"] = scope
    params.update(prov.get("extra_auth_params",{}))
    params.update(extra_params)

    auth_url = prov["auth_url"] + "?" + urllib.parse.urlencode(params)
    logger.info("[OAUTH] Redirecting to %s for connector %s", prov["auth_url"], app_id)
    return RedirectResponse(auth_url, status_code=302)

@app.get("/connectors/oauth/{app_id}/callback")
async def oauth_callback(app_id: str, request: Request,
                          code: Optional[str] = None, state: Optional[str] = None,
                          error: Optional[str] = None, error_description: Optional[str] = None):
    """Step 2: Exchange code for access token, save to DB."""
    app_name = app_id.replace("_"," ").title()

    if error:
        return HTMLResponse(_OAUTH_ERROR_HTML.format(
            message=f"Authorization denied for {app_name}",
            detail=error_description or error))

    if not code or not state:
        return HTMLResponse(_OAUTH_ERROR_HTML.format(
            message="Invalid callback", detail="Missing code or state parameter."))

    state_payload = _decode_oauth_state_token(state) or {}
    state_client_key = state_payload.get("client_state") or ""

    # Check DB FIRST (survives restarts & multi-worker), then fall back to memory
    db_state = await db_fetchone(
        "SELECT * FROM oauth_states WHERE state=?", (state,))
    if db_state:
        try:
            exp_dt = datetime.fromisoformat(db_state["expires_at"])
            now_dt = datetime.now(timezone.utc)
            if exp_dt.tzinfo is None: exp_dt = exp_dt.replace(tzinfo=timezone.utc)
            valid = (exp_dt + timedelta(minutes=10)) >= now_dt
        except Exception:
            valid = True
        if valid:
            pending = {"user_id":db_state["user_id"],"connector_type":db_state["connector_type"],
                       "code_verifier":db_state["code_verifier"] or None,
                       "redirect_uri":db_state["redirect_uri"] or "",
                       "client_state":state_client_key}
        else:
            pending = _pending_oauth.pop(state, None)
    else:
        pending = _pending_oauth.pop(state, None)
    if not pending and state_payload:
        pending = {
            "user_id": state_payload.get("sub", ""),
            "connector_type": state_payload.get("connector_type", app_id),
            "code_verifier": state_payload.get("code_verifier") or None,
            "redirect_uri": state_payload.get("redirect_uri") or "",
            "client_state": state_payload.get("client_state") or "",
        }
    if not pending:
        fallback_state = await db_fetchone(
            "SELECT * FROM oauth_states WHERE connector_type=? AND expires_at>? "
            "ORDER BY created_at DESC LIMIT 1",
            (app_id, _utcnow()))
        if fallback_state:
            logger.warning("[OAUTH] Recovering %s callback with latest unexpired state fallback", app_id)
            fallback_payload = _decode_oauth_state_token(fallback_state["state"]) or {}
            pending = {
                "user_id": fallback_state["user_id"],
                "connector_type": fallback_state["connector_type"],
                "code_verifier": fallback_state["code_verifier"] or fallback_payload.get("code_verifier") or None,
                "redirect_uri": fallback_state["redirect_uri"] or fallback_payload.get("redirect_uri") or "",
                "client_state": fallback_payload.get("client_state") or state,
            }
    if not pending:
        logger.warning("[OAUTH] State not found or expired: %s (app_id=%s)", state[:12]+"...", app_id)
        return HTMLResponse(_OAUTH_ERROR_HTML.format(
            message="OAuth state expired or invalid",
            detail="Please try connecting again from JAZZ. If this keeps happening, check that APP_BASE_URL is set correctly in your .env file."))

    user_id        = pending["user_id"]
    connector_type = pending["connector_type"]
    code_verifier  = pending.get("code_verifier")
    app_name       = connector_type.replace("_"," ").title()

    providers = _oauth_providers()
    prov = providers.get(connector_type, {})
    redirect_uri = pending.get("redirect_uri") or _oauth_redirect_uri(connector_type, request)

    token_data = {
        "grant_type":   "authorization_code",
        "code":         code,
        "redirect_uri": redirect_uri,
        "client_id":    prov.get("client_id",""),
        "client_secret": prov.get("client_secret",""),
    }
    if code_verifier: token_data["code_verifier"] = code_verifier

    basic_auth = None
    if prov.get("token_auth_type") == "basic":
        basic_auth = (prov["client_id"], prov["client_secret"])
        token_data.pop("client_id",None); token_data.pop("client_secret",None)

    try:
        tok_resp = await _oauth_post_form(prov["token_url"], token_data, basic_auth=basic_auth)
    except Exception as exc:
        return HTMLResponse(_OAUTH_ERROR_HTML.format(
            message="Token exchange failed", detail=str(exc)))

    if "error" in tok_resp:
        return HTMLResponse(_OAUTH_ERROR_HTML.format(
            message=f"Token error: {tok_resp.get('error')}",
            detail=tok_resp.get("error_description","")))

    # Handle Slack nested token format
    access_token = (tok_resp.get("access_token") or
                    tok_resp.get("authed_user",{}).get("access_token",""))
    if not access_token:
        return HTMLResponse(_OAUTH_ERROR_HTML.format(
            message="No access token returned",
            detail=str(tok_resp)[:300]))

    conn_data = {
        "access_token":  access_token,
        "refresh_token": tok_resp.get("refresh_token",""),
        "expires_in":    tok_resp.get("expires_in",0),
        "token_type":    tok_resp.get("token_type","Bearer"),
        "scope":         tok_resp.get("scope",""),
        "instance_url":  tok_resp.get("instance_url",""),  # for Salesforce
        "raw":           {k:v for k,v in tok_resp.items() if k not in ("access_token","refresh_token")},
    }
    encrypted = _encrypt(conn_data)
    now_iso   = _utcnow()

    # Save to connectors table (legacy path, always used)
    await db_execute(
        "INSERT INTO connectors(id,user_id,name,connector_type,encrypted_creds,is_active,created_at,updated_at)"
        " VALUES(?,?,?,?,?,1,?,?) "
        "ON CONFLICT(user_id,connector_type) DO UPDATE SET encrypted_creds=excluded.encrypted_creds,"
        "is_active=1,updated_at=excluded.updated_at",
        (_new_id(), user_id, app_name, connector_type, encrypted, now_iso, now_iso))

    # Also save to smart_connectors (primary lookup path used by _get_connector_creds_record)
    await db_execute(
        "INSERT INTO smart_connectors(id,user_id,app_name,connector_type,encrypted_conn_data,"
        "status,created_at,updated_at) VALUES(?,?,?,?,?,'active',?,?) "
        "ON CONFLICT(user_id,connector_type) DO UPDATE SET encrypted_conn_data=excluded.encrypted_conn_data,"
        "status='active',updated_at=excluded.updated_at",
        (_new_id(), user_id, app_name, connector_type, encrypted, now_iso, now_iso))

    # Mark connector as enabled in platform_connectors
    await db_execute(
        "UPDATE platform_connectors SET is_enabled=1,updated_at=? WHERE connector_type=?",
        (now_iso, connector_type))

    # Clean up DB state
    await db_execute("DELETE FROM oauth_states WHERE state=?", (state,))
    logger.info("[OAUTH] ✅ %s connected for user %s", connector_type, user_id)

    _completed_oauth[state] = {"connected":True,"connector_type":connector_type,"app_name":app_name}
    if pending.get("client_state"):
        _completed_oauth[pending["client_state"]] = {"connected":True,"connector_type":connector_type,"app_name":app_name}
    return HTMLResponse(_OAUTH_SUCCESS_HTML.format(name=app_name))

@app.get("/connectors/oauth/{app_id}/status")
async def oauth_status(app_id: str, state: str = "", user: Dict = Depends(_get_current_user)):
    """Step 3: Frontend polls this until {connected: true}."""
    uid = user.get("id") or user.get("sub","")
    if state and state in _completed_oauth:
        result = _completed_oauth.pop(state, {})
        if result.get("connector_type") == app_id or app_id in _GOOGLE_CONNECTOR_TYPES:
            return JSONResponse({"connected":True,"connector_type":result.get("connector_type",app_id),"app_name":result.get("app_name",app_id)})
    # DB fallback: check if connector now exists for this user (survives restarts)
    row = await db_fetchone(
        "SELECT id,connector_type,updated_at FROM connectors WHERE user_id=? AND connector_type=? AND is_active=1",
        (uid, app_id))
    if not row:
        # Also check smart_connectors for Google-shared connectors
        if app_id in _GOOGLE_CONNECTOR_TYPES:
            row = await db_fetchone(
                "SELECT id,connector_type,updated_at FROM smart_connectors WHERE user_id=? AND connector_type=? AND status='active'",
                (uid, app_id))
    if row:
        # If the connector exists, the authorization did complete. Returning true
        # here keeps the waiting modal in sync even if the callback was recovered
        # through the DB fallback or the browser missed the in-memory completion flag.
        return JSONResponse({"connected":True,"connector_type":app_id})
    return JSONResponse({"connected":False})

# ── Legacy + New Connector CRUD ────────────────────────────────────────────────
@app.get("/connectors")
async def list_connectors(user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    # Merge connectors + smart_connectors (deduplicated by connector_type)
    rows = await db_fetchall(
        "SELECT id,name,connector_type,is_active,last_tested_at,last_test_ok,created_at,'legacy' as source"
        " FROM connectors WHERE user_id=? AND is_active=1 ORDER BY created_at DESC", (uid,))
    smart_rows = await db_fetchall(
        "SELECT id,app_name as name,connector_type,1 as is_active,last_tested_at,last_test_ok,created_at,'smart' as source"
        " FROM smart_connectors WHERE user_id=? AND status='active' ORDER BY created_at DESC", (uid,))
    seen = {r["connector_type"] for r in rows}
    for sr in smart_rows:
        if sr["connector_type"] not in seen:
            rows.append(sr)
            seen.add(sr["connector_type"])
    # Annotate with OAuth providers list
    oauth_types = set(_oauth_providers().keys())
    for r in rows:
        r["oauth"] = r["connector_type"] in oauth_types
    return {"connectors":rows, "connected_types": list(seen)}

@app.get("/connectors/connected")
async def list_connected_types(user: Dict = Depends(_get_current_user)):
    """Returns a list of connector_type strings the user has connected — used by the chat UI."""
    uid = user.get("id") or user.get("sub","")
    legacy = await db_fetchall(
        "SELECT connector_type FROM connectors WHERE user_id=? AND is_active=1", (uid,))
    smart = await db_fetchall(
        "SELECT connector_type FROM smart_connectors WHERE user_id=? AND status='active'", (uid,))
    types = list({r["connector_type"] for r in legacy + smart})
    return {"connected": types}

@app.post("/connectors/{connector_type}/action")
async def connector_action(connector_type: str, body: ConnectorActionReq,
                            user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    result = await _connector_api_call(uid, connector_type, body.action, body.params)
    return result

@app.post("/connectors/{cid}/test")
async def test_connector(cid: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    row = await db_fetchone("SELECT * FROM connectors WHERE id=? AND user_id=?", (cid, uid))
    if not row: raise HTTPException(404, "Connector not found")
    try:
        creds = _decrypt(row["encrypted_creds"])
        success = bool(creds.get("access_token","") or creds.get("api_key",""))
        err = "" if success else "No credentials found"
        await db_execute("UPDATE connectors SET last_tested_at=?,last_test_ok=?,test_error=? WHERE id=?",
                         (_utcnow(), 1 if success else 0, err, cid))
        return {"success":success,"error":err}
    except Exception as e:
        await db_execute("UPDATE connectors SET last_tested_at=?,last_test_ok=0,test_error=? WHERE id=?",
                         (_utcnow(), str(e)[:300], cid))
        return {"success":False,"error":str(e)}

@app.delete("/connectors/{connector_type_or_id}")
async def delete_connector(connector_type_or_id: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    # Try by ID first, then by type
    row = await db_fetchone("SELECT id FROM connectors WHERE id=? AND user_id=?", (connector_type_or_id, uid))
    if not row:
        row = await db_fetchone("SELECT id FROM connectors WHERE connector_type=? AND user_id=?", (connector_type_or_id, uid))
    if not row: raise HTTPException(404, "Connector not found")
    await db_execute("DELETE FROM connectors WHERE id=?", (row["id"],))
    return {"ok":True}

# ── Smart Connectors ───────────────────────────────────────────────────────────
@app.get("/smart-connectors")
async def list_smart_connectors(user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    smart_rows = await db_fetchall(
        "SELECT id,app_name,connector_type,status,last_tested_at,last_test_ok,created_at"
        " FROM smart_connectors WHERE user_id=? ORDER BY created_at DESC", (uid,))
    legacy_rows = await db_fetchall(
        "SELECT id,name,connector_type,is_active,last_tested_at,last_test_ok,created_at"
        " FROM connectors WHERE user_id=? ORDER BY created_at DESC", (uid,))

    merged: List[Dict[str, Any]] = []
    seen_types = set()
    for r in smart_rows:
        merged.append(dict(r))
        seen_types.add(r["connector_type"])

    for r in legacy_rows:
        if r["connector_type"] in seen_types:
            continue
        merged.append({
            "id": r["id"],
            "app_name": r.get("name") or r["connector_type"],
            "connector_type": r["connector_type"],
            "status": "active" if r.get("is_active") else "inactive",
            "last_tested_at": r.get("last_tested_at"),
            "last_test_ok": r.get("last_test_ok"),
            "created_at": r.get("created_at"),
            "source": "legacy_oauth",
        })
    return {"connectors":merged}

@app.post("/smart-connectors")
async def create_smart_connector(body: dict, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    await _check_rate_limit(user, "connectors")
    cid = _new_id(); now = _utcnow()
    conn_data = body.get("conn_data",{})
    encrypted = _encrypt(conn_data)
    ctype = body.get("connector_type","custom")
    await db_execute(
        "INSERT INTO smart_connectors(id,user_id,app_name,connector_type,encrypted_conn_data,status,created_at,updated_at)"
        " VALUES(?,?,?,?,?,?,?,?)",
        (cid, uid, body.get("app_name",ctype), ctype, encrypted, "active", now, now))
    return {"id":cid,"status":"active"}

@app.put("/smart-connectors/{cid}")
async def update_smart_connector(cid: str, body: dict, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    now = _utcnow()
    status = body.get("status")
    conn_data = body.get("conn_data")

    smart = await db_fetchone("SELECT id FROM smart_connectors WHERE id=? AND user_id=?", (cid, uid))
    if smart:
        sets, vals = [], []
        if conn_data is not None:
            sets.append("encrypted_conn_data=?"); vals.append(_encrypt(conn_data))
        if status is not None:
            sets.append("status=?"); vals.append(str(status))
        sets.append("updated_at=?"); vals.append(now)
        vals.append(cid)
        await db_execute(f"UPDATE smart_connectors SET {', '.join(sets)} WHERE id=?", tuple(vals))
        return {"ok":True}

    legacy = await db_fetchone("SELECT id FROM connectors WHERE id=? AND user_id=?", (cid, uid))
    if not legacy:
        raise HTTPException(404, "Connector not found")

    sets, vals = [], []
    if conn_data is not None:
        sets.append("encrypted_creds=?"); vals.append(_encrypt(conn_data))
    if status is not None:
        sets.append("is_active=?"); vals.append(1 if str(status).lower() == "active" else 0)
    sets.append("updated_at=?"); vals.append(now)
    vals.append(cid)
    await db_execute(f"UPDATE connectors SET {', '.join(sets)} WHERE id=?", tuple(vals))
    return {"ok":True}

@app.post("/smart-connectors/{cid}/test")
async def test_smart_connector(cid: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    now = _utcnow()

    smart = await db_fetchone(
        "SELECT id,encrypted_conn_data FROM smart_connectors WHERE id=? AND user_id=?", (cid, uid))
    if smart:
        try:
            creds = _decrypt(smart["encrypted_conn_data"])
            success = bool(
                creds.get("access_token","") or creds.get("api_key","") or creds.get("webhook_url","")
            )
            err = "" if success else "No credentials found"
            await db_execute(
                "UPDATE smart_connectors SET last_tested_at=?,last_test_ok=?,test_error=?,updated_at=? WHERE id=?",
                (now, 1 if success else 0, err, now, cid))
            return {"success":success,"error":err}
        except Exception as e:
            await db_execute(
                "UPDATE smart_connectors SET last_tested_at=?,last_test_ok=0,test_error=?,updated_at=? WHERE id=?",
                (now, str(e)[:300], now, cid))
            return {"success":False,"error":str(e)}

    legacy = await db_fetchone("SELECT id,encrypted_creds FROM connectors WHERE id=? AND user_id=?", (cid, uid))
    if not legacy:
        raise HTTPException(404, "Connector not found")
    try:
        creds = _decrypt(legacy["encrypted_creds"])
        success = bool(creds.get("access_token","") or creds.get("api_key",""))
        err = "" if success else "No credentials found"
        await db_execute("UPDATE connectors SET last_tested_at=?,last_test_ok=?,test_error=? WHERE id=?",
                         (now, 1 if success else 0, err, cid))
        return {"success":success,"error":err}
    except Exception as e:
        await db_execute("UPDATE connectors SET last_tested_at=?,last_test_ok=0,test_error=? WHERE id=?",
                         (now, str(e)[:300], cid))
        return {"success":False,"error":str(e)}

@app.delete("/smart-connectors/{cid}")
async def delete_smart_connector(cid: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    smart = await db_fetchone("SELECT id FROM smart_connectors WHERE id=? AND user_id=?", (cid, uid))
    if smart:
        await db_execute("DELETE FROM smart_connectors WHERE id=? AND user_id=?",(cid,uid))
        return {"ok":True}
    legacy = await db_fetchone("SELECT id FROM connectors WHERE id=? AND user_id=?", (cid, uid))
    if legacy:
        await db_execute("DELETE FROM connectors WHERE id=? AND user_id=?",(cid,uid))
        return {"ok":True}
    by_type_smart = await db_fetchone("SELECT id FROM smart_connectors WHERE connector_type=? AND user_id=?", (cid, uid))
    if by_type_smart:
        await db_execute("DELETE FROM smart_connectors WHERE id=? AND user_id=?", (by_type_smart["id"], uid))
        return {"ok":True}
    by_type_legacy = await db_fetchone("SELECT id FROM connectors WHERE connector_type=? AND user_id=?", (cid, uid))
    if by_type_legacy:
        await db_execute("DELETE FROM connectors WHERE id=? AND user_id=?", (by_type_legacy["id"], uid))
        return {"ok":True}
    raise HTTPException(404, "Connector not found")
    return {"ok":True}

# ══════════════════════════════════════════════════════════════════════════════
# §29  API KEYS
# ══════════════════════════════════════════════════════════════════════════════

@app.get("/api-keys")
async def list_api_keys(user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    keys = await db_fetchall(
        "SELECT id,name,key_prefix,scopes_json,rate_limit_rpm,is_active,last_used_at,usage_count,expires_at,created_at"
        " FROM api_keys WHERE user_id=? ORDER BY created_at DESC", (uid,))
    return {"api_keys":keys}

@app.post("/api-keys")
async def create_api_key(req: ApiKeyCreate, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    await _check_rate_limit(user, "api_keys")
    raw = f"jz_{secrets.token_urlsafe(32)}"; h = _hash_token(raw); kid = _new_id()
    exp = None
    if req.expires_days:
        exp = (datetime.now(timezone.utc)+timedelta(days=req.expires_days)).isoformat()
    await db_execute(
        "INSERT INTO api_keys(id,user_id,name,key_hash,key_prefix,scopes_json,rate_limit_rpm,expires_at,is_active,created_at)"
        " VALUES(?,?,?,?,?,?,?,?,1,?)",
        (kid, uid, req.name, h, raw[:12], req.scopes_json, req.rate_limit_rpm, exp, _utcnow()))
    return {"id":kid,"key":raw,"prefix":raw[:12],"name":req.name}

@app.delete("/api-keys/{kid}")
async def delete_api_key(kid: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    await db_execute("DELETE FROM api_keys WHERE id=? AND user_id=?",(kid,uid))
    return {"ok":True}

# ══════════════════════════════════════════════════════════════════════════════
# §30  AGENT JOBS
# ══════════════════════════════════════════════════════════════════════════════

@app.get("/agent/jobs")
@app.get("/agent-jobs")
async def list_agent_jobs(user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    jobs = await db_fetchall("SELECT * FROM agent_jobs WHERE user_id=? ORDER BY created_at DESC",(uid,))
    return {"jobs":jobs}

@app.post("/agent/jobs")
@app.post("/agent-jobs")
async def create_agent_job(req: AgentJobCreate, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    await _check_rate_limit(user, "agent_jobs")
    jid = _new_id(); now = _utcnow()
    await db_execute(
        "INSERT INTO agent_jobs(id,user_id,name,description,job_type,trigger_json,"
        "prompt_template,tools_json,model_id,enabled,max_retries,timeout_seconds,metadata_json,created_at,updated_at)"
        " VALUES(?,?,?,?,?,?,?,?,?,1,?,?,'{}',?,?)",
        (jid, uid, req.name, req.description, req.job_type, req.trigger_json,
         req.prompt_template, req.tools_json, req.model_id, req.max_retries, req.timeout_seconds, now, now))
    job = await db_fetchone("SELECT * FROM agent_jobs WHERE id=?",(jid,))
    _schedule_job(job)
    return job

@app.put("/agent/jobs/{jid}")
@app.put("/agent-jobs/{jid}")
async def update_agent_job(jid: str, req: AgentJobUpdate, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    job = await db_fetchone("SELECT id FROM agent_jobs WHERE id=? AND user_id=?",(jid,uid))
    if not job: raise HTTPException(404)
    now = _utcnow()
    fields: Dict = {}
    if req.name            is not None: fields["name"]            = req.name
    if req.prompt_template is not None: fields["prompt_template"] = req.prompt_template
    if req.trigger_json    is not None: fields["trigger_json"]    = req.trigger_json
    if req.tools_json      is not None: fields["tools_json"]      = req.tools_json
    if req.enabled         is not None: fields["enabled"]         = int(req.enabled)
    if req.model_id        is not None: fields["model_id"]        = req.model_id
    fields["updated_at"] = now
    set_clause = ", ".join(f"{k}=?" for k in fields)
    await db_execute(f"UPDATE agent_jobs SET {set_clause} WHERE id=?",tuple(fields.values())+(jid,))
    job = await db_fetchone("SELECT * FROM agent_jobs WHERE id=?",(jid,))
    if job["enabled"]: _schedule_job(job)
    else:
        try: _scheduler.remove_job(jid)
        except Exception: pass
    return {"ok":True}

@app.delete("/agent/jobs/{jid}")
@app.delete("/agent-jobs/{jid}")
async def delete_agent_job(jid: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    await db_execute("DELETE FROM agent_jobs WHERE id=? AND user_id=?",(jid,uid))
    try: _scheduler.remove_job(jid)
    except Exception: pass
    return {"ok":True}

@app.post("/agent/jobs/{jid}/run")
@app.post("/agent-jobs/{jid}/run")
async def run_agent_job_now(jid: str, background: BackgroundTasks, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    job = await db_fetchone("SELECT id FROM agent_jobs WHERE id=? AND user_id=?",(jid,uid))
    if not job: raise HTTPException(404)
    background.add_task(_run_agent_job, jid)
    return {"ok":True,"message":"Job triggered"}

@app.get("/agent/jobs/{jid}/logs")
@app.get("/agent-jobs/{jid}/logs")
async def get_agent_job_logs(jid: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    job = await db_fetchone("SELECT id FROM agent_jobs WHERE id=? AND user_id=?",(jid,uid))
    if not job: raise HTTPException(404)
    logs = await db_fetchall(
        "SELECT * FROM agent_job_logs WHERE job_id=? ORDER BY started_at DESC LIMIT 50",(jid,))
    return {"logs":logs}

@app.post("/agent/run")
async def agent_run_inline(req: AgentRunReq, user: Dict = Depends(_get_current_user)):
    await _check_rate_limit(user, "messages_per_day")
    t0 = time.time()
    answer, tool_log = await _agent_loop(req.prompt, user, req.tools, req.model_id, max_steps=12)
    return {"answer":answer,"tool_calls":tool_log,"steps":len(tool_log),"latency_ms":int((time.time()-t0)*1000)}

@app.post("/agent/run/stream")
async def agent_run_stream(req: AgentRunReq, user: Dict = Depends(_get_current_user)) -> StreamingResponse:
    await _check_rate_limit(user, "messages_per_day")

    async def _gen():
        t0 = time.time()
        client, model_name = await _get_model_client(req.model_id)
        messages = [{"role":"system","content":_AGENT_SYSTEM},{"role":"user","content":req.prompt}]
        try:
            for step in range(12):
                try:
                    resp = await asyncio.get_running_loop().run_in_executor(
                        _executor, lambda c=client, m=model_name: c.chat.completions.create(
                            model=m, messages=messages, max_tokens=800, temperature=0.2))
                    raw = resp.choices[0].message.content.strip()
                    clean = re.sub(r"```json|```", "", raw).strip()
                    step_json = json.loads(clean)
                except Exception as e:
                    yield f"data: {json.dumps({'type':'error','message':f'Agent parse error step {step}: {str(e)}'})}\n\n"
                    break
                thought = step_json.get("thought", "")
                action  = step_json.get("action", "FINISH")
                args    = step_json.get("args", {})
                messages.append({"role":"assistant","content":raw})
                yield f"data: {json.dumps({'type':'thought','step':step,'thought':thought,'action':action})}\n\n"
                if action == "FINISH":
                    answer = step_json.get("final_answer", "Task complete.")
                    yield f"data: {json.dumps({'type':'done','answer':answer,'latency_ms':int((time.time()-t0)*1000)})}\n\n"
                    return
                t_step = time.time()
                obs = await _run_tool(action, args, user)
                success = not str(obs).startswith("Tool error")
                yield f"data: {json.dumps({'type':'tool_result','step':step,'tool':action,'success':success,'latency_ms':int((time.time()-t_step)*1000),'output':str(obs)[:300]})}\n\n"
                messages.append({"role":"user","content":f"[Observation]: {obs}"})
            yield f"data: {json.dumps({'type':'done','answer':'Agent reached max steps.','latency_ms':int((time.time()-t0)*1000)})}\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'type':'error','message':str(e)})}\n\n"

    return StreamingResponse(_gen(), media_type="text/event-stream",
                             headers={"Cache-Control":"no-cache","X-Accel-Buffering":"no","Connection":"keep-alive"})

# ══════════════════════════════════════════════════════════════════════════════
# §31  SYSTEM ACCESS
# ══════════════════════════════════════════════════════════════════════════════

@app.post("/system/shell")
async def run_shell_cmd(req: ShellReq, user: Dict = Depends(_get_current_user)):
    command = (req.command or req.cmd or "").strip()
    if not command:
        raise HTTPException(400, "command is required")
    result = await asyncio.get_running_loop().run_in_executor(
        _executor, lambda: _run_shell_sync(command, req.timeout, req.cwd))
    return result

@app.post("/system/code")
async def run_code(req: CodeRunReq, user: Dict = Depends(_get_current_user)):
    await _check_rate_limit(user, "code_runs_per_day")
    uid = user.get("id") or user.get("sub","")
    result = await asyncio.get_running_loop().run_in_executor(
        _executor, lambda: _run_code_sync(req.code, req.language, req.cwd))
    if req.save_log:
        await db_execute(
            "INSERT INTO code_run_logs(id,user_id,language,cwd,code,stdout,stderr,exit_code,duration_ms,created_at)"
            " VALUES(?,?,?,?,?,?,?,?,?,?)",
            (_new_id(), uid, req.language, req.cwd or "", req.code[:200000],
             (result.get("stdout") or "")[:CODE_EXEC_MAX_OUT],
             (result.get("stderr") or "")[:CODE_EXEC_MAX_OUT],
             int(result.get("exit_code") or 0), int(result.get("duration_ms") or 0), _utcnow()))
    return result

@app.post("/code/run")
async def run_code_legacy(req: CodeRunReq, user: Dict = Depends(_get_current_user)):
    return await run_code(req, user)

@app.get("/system/info")
async def system_info_ep(user: Dict = Depends(_get_current_user)):
    return await asyncio.get_running_loop().run_in_executor(
        _executor, lambda: {"os":platform.system(),"version":platform.version(),
                            "python":sys.version,"cwd":os.getcwd(),"hostname":platform.node()})

@app.get("/code/history")
async def code_history(limit: int = 30, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    limit = max(1, min(100, int(limit or 30)))
    rows = await db_fetchall(
        "SELECT id,language,cwd,code,stdout,stderr,exit_code,duration_ms,created_at"
        " FROM code_run_logs WHERE user_id=? ORDER BY created_at DESC LIMIT ?",
        (uid, limit))
    for r in rows:
        r["code_preview"] = (r.get("code") or "")[:600]
        r["stdout_preview"] = (r.get("stdout") or "")[:600]
        r["stderr_preview"] = (r.get("stderr") or "")[:600]
    return {"runs": rows}

@app.delete("/code/history")
async def clear_code_history(user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    await db_execute("DELETE FROM code_run_logs WHERE user_id=?", (uid,))
    return {"ok": True}

# ══════════════════════════════════════════════════════════════════════════════
# §32  WEBSITES
# ══════════════════════════════════════════════════════════════════════════════

@app.post("/websites/build")
async def build_website_ep(req: WebsiteCreate, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    await _check_rate_limit(user, "websites")
    html = await _build_website(
        req.description, req.title, req.style, req.model_id,
        req.pages, req.color_palette or "", req.extra_instructions or "",
        req.site_type or "", req.audience or "", req.features or [],
        req.theme_mode or "auto", req.interactivity or "rich", req.seo_keywords or "")
    wid = _new_id(); now = _utcnow(); fname = f"{wid}.html"
    fpath = SITES_DIR / fname; fpath.write_text(html, encoding="utf-8")
    prompt_used = json.dumps({
        "description": req.description, "pages": req.pages or [],
        "color_palette": req.color_palette or "", "extra_instructions": req.extra_instructions or "",
        "site_type": req.site_type or "", "audience": req.audience or "",
        "features": req.features or [], "theme_mode": req.theme_mode or "auto",
        "interactivity": req.interactivity or "rich", "seo_keywords": req.seo_keywords or "",
    })
    await db_execute(
        "INSERT INTO websites(id,user_id,title,description,filename,style,prompt_used,html_size_bytes,created_at,updated_at)"
        " VALUES(?,?,?,?,?,?,?,?,?,?)",
        (wid, uid, req.title, req.description, fname, req.style, prompt_used, len(html.encode()), now, now))
    return {"id":wid,"title":req.title,"preview_url":f"/preview/{fname}","html_size_bytes":len(html.encode())}

@app.post("/websites/generate")  # legacy alias
async def generate_website(body: dict, user: Dict = Depends(_get_current_user)):
    req = WebsiteCreate(title=body.get("title","My Website"),
                        description=body.get("prompt","A modern landing page"),
                        style=body.get("style","modern"),
                        model_id=body.get("model_id","llama-3.3-70b-versatile"))
    return await build_website_ep(req, user)

@app.get("/websites")
async def list_websites(user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    sites = await db_fetchall(
        "SELECT id,title,description,filename,style,prompt_used,html_size_bytes,view_count,created_at,updated_at"
        " FROM websites WHERE user_id=? ORDER BY created_at DESC", (uid,))
    for r in sites:
        r["preview_url"] = f"/preview/{r['filename']}"
        try: r["prompt"] = json.loads(r.get("prompt_used") or "{}")
        except Exception: r["prompt"] = {}
    return {"websites":sites}

@app.get("/websites/{wid}")
async def get_website(wid: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    site = await db_fetchone("SELECT * FROM websites WHERE id=? AND user_id=?", (wid, uid))
    if not site:
        raise HTTPException(404, "Website not found")
    path = SITES_DIR / site["filename"]
    html = path.read_text(encoding="utf-8", errors="replace") if path.exists() else ""
    return {**dict(site), "html": html, "preview_url": f"/preview/{site['filename']}"}

@app.put("/websites/{wid}")
async def update_website(wid: str, body: dict, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    site = await db_fetchone("SELECT * FROM websites WHERE id=? AND user_id=?", (wid, uid))
    if not site:
        raise HTTPException(404, "Website not found")
    title = str(body.get("title") or site["title"])[:100]
    description = str(body.get("description") or site["description"])
    style = str(body.get("style") or site["style"])[:60]
    html = body.get("html")
    now = _utcnow()
    if html is not None:
        html = str(html)
        path = SITES_DIR / site["filename"]
        path.write_text(html, encoding="utf-8")
        size = len(html.encode())
    else:
        size = site.get("html_size_bytes") or 0
    await db_execute(
        "UPDATE websites SET title=?,description=?,style=?,html_size_bytes=?,updated_at=? WHERE id=? AND user_id=?",
        (title, description, style, size, now, wid, uid))
    return {"ok": True, "id": wid, "preview_url": f"/preview/{site['filename']}"}

@app.post("/websites/{wid}/iterate")
async def iterate_website(wid: str, req: WebsiteIterateReq, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    await _check_rate_limit(user, "websites")
    site = await db_fetchone("SELECT * FROM websites WHERE id=? AND user_id=?", (wid, uid))
    if not site:
        raise HTTPException(404, "Website not found")
    src = SITES_DIR / site["filename"]
    existing = src.read_text(encoding="utf-8", errors="replace") if src.exists() else ""
    html = await _iterate_website(existing, req.change_request, site["title"], req.model_id, req.preserve_brand)
    now = _utcnow()
    if req.save_as_new:
        new_id = _new_id(); fname = f"{new_id}.html"
        (SITES_DIR / fname).write_text(html, encoding="utf-8")
        prompt_used = json.dumps({
            "source_website_id": wid,
            "iteration_request": req.change_request,
            "preserve_brand": req.preserve_brand,
        })
        title = f"{site['title']} Iteration"[:100]
        await db_execute(
            "INSERT INTO websites(id,user_id,title,description,filename,style,prompt_used,html_size_bytes,created_at,updated_at)"
            " VALUES(?,?,?,?,?,?,?,?,?,?)",
            (new_id, uid, title, site["description"], fname, site["style"], prompt_used, len(html.encode()), now, now))
        return {"ok": True, "id": new_id, "title": title, "preview_url": f"/preview/{fname}", "html_size_bytes": len(html.encode())}
    src.write_text(html, encoding="utf-8")
    try:
        prompt_meta = json.loads(site.get("prompt_used") or "{}")
    except Exception:
        prompt_meta = {}
    prompt_meta["last_iteration"] = {"request": req.change_request, "at": now}
    await db_execute(
        "UPDATE websites SET prompt_used=?,html_size_bytes=?,updated_at=? WHERE id=? AND user_id=?",
        (json.dumps(prompt_meta), len(html.encode()), now, wid, uid))
    return {"ok": True, "id": wid, "preview_url": f"/preview/{site['filename']}", "html_size_bytes": len(html.encode())}

@app.post("/websites/{wid}/duplicate")
async def duplicate_website(wid: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    site = await db_fetchone("SELECT * FROM websites WHERE id=? AND user_id=?", (wid, uid))
    if not site:
        raise HTTPException(404, "Website not found")
    src = SITES_DIR / site["filename"]
    html = src.read_text(encoding="utf-8", errors="replace") if src.exists() else ""
    new_id = _new_id(); fname = f"{new_id}.html"; now = _utcnow()
    (SITES_DIR / fname).write_text(html, encoding="utf-8")
    await db_execute(
        "INSERT INTO websites(id,user_id,title,description,filename,style,prompt_used,html_size_bytes,created_at,updated_at)"
        " VALUES(?,?,?,?,?,?,?,?,?,?)",
        (new_id, uid, f"{site['title']} Copy"[:100], site["description"], fname,
         site["style"], site.get("prompt_used") or "", len(html.encode()), now, now))
    return {"id": new_id, "title": f"{site['title']} Copy", "preview_url": f"/preview/{fname}"}

@app.get("/websites/{wid}/download")
async def download_website(wid: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    site = await db_fetchone("SELECT * FROM websites WHERE id=? AND user_id=?", (wid, uid))
    if not site:
        raise HTTPException(404, "Website not found")
    path = SITES_DIR / site["filename"]
    if not path.exists():
        raise HTTPException(404, "Website file missing")
    safe_title = re.sub(r"[^A-Za-z0-9._-]+", "_", site["title"]).strip("_") or "website"
    return FileResponse(str(path), media_type="text/html", filename=f"{safe_title}.html")

@app.delete("/websites/{wid}")
async def delete_website(wid: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    site = await db_fetchone("SELECT * FROM websites WHERE id=? AND user_id=?",(wid,uid))
    if not site: raise HTTPException(404)
    p = SITES_DIR / site["filename"]
    if p.exists(): p.unlink()
    await db_execute("DELETE FROM websites WHERE id=?",(wid,))
    return {"ok":True}

# ══════════════════════════════════════════════════════════════════════════════
# §33  VOICE / LIVEKIT
# ══════════════════════════════════════════════════════════════════════════════

@app.post("/voice/token")
async def get_voice_token(req: VoiceTokenReq, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    identity = req.identity or uid
    name = user.get("email", identity) or identity
    token = _lk_token(identity, name, req.room_name, req.ttl_seconds)
    return {"token":token,"room":req.room_name,"url":LIVEKIT_URL,"identity":identity}

@app.get("/voice/livekit/status")
async def livekit_voice_status(user: Dict = Depends(_get_current_user)):
    return {
        "available": _lk_available(),
        "configured": bool(LIVEKIT_URL and LIVEKIT_API_KEY and LIVEKIT_API_SECRET),
        "has_sdk": HAS_LIVEKIT,
        "url": LIVEKIT_URL,
        "mode": "livekit",
        "agent_required": True,
        "service": "JAZZ Live Voice",
    }

@app.post("/voice/livekit/session")
async def create_livekit_voice_session(req: LiveVoiceSessionReq, user: Dict = Depends(_get_current_user)):
    if not _lk_available():
        raise HTTPException(503, "LiveKit is not configured")
    uid = user.get("id") or user.get("sub","")
    base_room = req.room_name or f"jazz-live-{uid[:8]}-{secrets.token_hex(4)}"
    room = re.sub(r"[^a-zA-Z0-9_-]+", "-", base_room).strip("-")[:96] or f"jazz-live-{secrets.token_hex(4)}"
    identity = f"user-{uid[:18]}-{secrets.token_hex(3)}"
    name = user.get("full_name") or user.get("email") or "JAZZ User"
    token = _lk_token(identity, name, room, req.ttl_seconds)
    sid = _new_id(); now = _utcnow()
    await db_execute(
        "INSERT INTO voice_sessions(id,user_id,room_name,livekit_token,language,status,started_at)"
        " VALUES(?,?,?,?,?,'active',?)",
        (sid, uid, room, token, req.language or "auto", now))
    return {
        "session_id": sid,
        "room": room,
        "token": token,
        "url": LIVEKIT_URL,
        "identity": identity,
        "configured": True,
        "agent_required": True,
    }

@app.post("/voice/livekit/session/{session_id}/end")
async def end_livekit_voice_session(session_id: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    row = await db_fetchone("SELECT id FROM voice_sessions WHERE id=? AND user_id=?", (session_id, uid))
    if not row:
        raise HTTPException(404, "Voice session not found")
    await db_execute(
        "UPDATE voice_sessions SET status='ended', ended_at=? WHERE id=? AND user_id=?",
        (_utcnow(), session_id, uid))
    return {"ok": True}

@app.get("/voice/livekit/sessions")
async def list_livekit_voice_sessions(limit: int = 20, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    limit = max(1, min(100, int(limit or 20)))
    rows = await db_fetchall(
        "SELECT id,room_name,language,status,started_at,ended_at FROM voice_sessions"
        " WHERE user_id=? ORDER BY started_at DESC LIMIT ?",
        (uid, limit))
    return {"sessions": rows}

@app.post("/voice/transcribe")
async def transcribe_audio(file: UploadFile = File(...), user: Dict = Depends(_get_current_user)):
    await _check_rate_limit(user, "messages_per_day")
    content = await file.read()
    if len(content) > 25*1024*1024: raise HTTPException(413, "Audio file too large (max 25MB)")
    transcript = await _stt_transcribe(content, file.content_type or "audio/webm")
    return {"transcript":transcript,"chars":len(transcript)}

@app.post("/voice/chat")
async def voice_chat(req: VoiceChatReq, background: BackgroundTasks, user: Dict = Depends(_get_current_user)):
    await _check_rate_limit(user, "messages_per_day")
    uid = user.get("id") or user.get("sub","")
    sid = await _get_or_create_session(uid, req.session_id, req.model_id)
    messages = await _build_context(sid, uid, req.transcript, True, req.model_id)
    reply = await _llm_text(messages, req.model_id, max_tokens=2048)
    mid = _new_id(); now = _utcnow()
    await db_execute("INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,created_at) VALUES(?,?,?,'user',?,?,?)",
                     (_new_id(), sid, uid, req.transcript, req.model_id, now))
    await db_execute("INSERT INTO chat_history(id,session_id,user_id,role,content,model_used,mode,created_at) VALUES(?,?,?,'assistant',?,?,'voice',?)",
                     (mid, sid, uid, reply, req.model_id, now))
    await db_execute("UPDATE chat_sessions SET last_message_at=?,turn_count=turn_count+1,updated_at=? WHERE id=?",(now,now,sid))
    response: Dict = {"session_id":sid,"message_id":mid,"response":reply}
    if req.return_audio:
        try:
            client = _groq_client()
            tts_resp = await asyncio.get_running_loop().run_in_executor(
                None, lambda: client.audio.speech.create(
                    model="playai-tts", voice=req.voice, input=reply[:4096], response_format="mp3"))
            response["audio_base64"] = base64.b64encode(tts_resp.content).decode()
            response["audio_format"] = "mp3"
        except Exception as e:
            response["tts_error"] = str(e)
    return response

@app.post("/voice/tts")
async def text_to_speech(body: dict, user: Dict = Depends(_get_current_user)):
    text = body.get("text","")[:4096]
    if not text.strip(): raise HTTPException(400, "text is required")
    try:
        client = _groq_client()
        tts_resp = await asyncio.get_running_loop().run_in_executor(
            None, lambda: client.audio.speech.create(
                model="playai-tts", voice=body.get("voice",TTS_VOICE), input=text, response_format="mp3"))
        return StreamingResponse(io.BytesIO(tts_resp.content), media_type="audio/mpeg",
                                 headers={"Content-Disposition":"inline; filename=speech.mp3"})
    except Exception as e:
        raise HTTPException(500, f"TTS failed: {e}")

@app.get("/voice/tts")
async def text_to_speech_get(text: str = "", voice: str = TTS_VOICE, user: Dict = Depends(_get_current_user)):
    if not text.strip(): raise HTTPException(400, "text is required")
    return await text_to_speech({"text": text, "voice": voice}, user)

# ══════════════════════════════════════════════════════════════════════════════
# §34  NOTIFICATIONS
# ══════════════════════════════════════════════════════════════════════════════

@app.get("/notifications")
async def list_notifications(user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    notes = await db_fetchall(
        "SELECT * FROM notifications WHERE (user_id=? OR is_broadcast=1)"
        " AND (expires_at IS NULL OR expires_at>?) ORDER BY created_at DESC LIMIT 50",
        (uid, _utcnow()))
    return {"notifications":notes,"unread":sum(1 for n in notes if not n["is_read"])}

@app.post("/notifications/{nid}/read")
async def mark_notification_read(nid: str, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    await db_execute("UPDATE notifications SET is_read=1 WHERE id=? AND user_id=?",(nid,uid))
    return {"ok":True}

@app.post("/notifications/read-all")
async def mark_all_read(user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    await db_execute("UPDATE notifications SET is_read=1 WHERE user_id=?",(uid,))
    return {"ok":True}

# ══════════════════════════════════════════════════════════════════════════════
# §35  USER SETTINGS & USAGE
# ══════════════════════════════════════════════════════════════════════════════

@app.get("/settings")
async def get_settings(user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    u = await db_fetchone(
        "SELECT id,email,full_name,role,subscription,memory_enabled,timezone,created_at FROM users WHERE id=?", (uid,))
    return u or {}

@app.get("/features/sidebar")
async def get_sidebar_features(user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    return await _sidebar_feature_access_for_user(uid)

@app.get("/usage")
async def get_usage(user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    sub = user.get("subscription","free")
    limits = await _subscription_limits_for_plan(sub)
    today = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    async def _cnt(res):
        row = await db_fetchone(
            "SELECT count FROM rate_limits WHERE user_id=? AND resource=? AND window_key=?", (uid,res,today))
        return row["count"] if row else 0
    return {
        "subscription": sub, "limits": limits,
        "usage": {
            "messages_today":    await _cnt("messages_per_day"),
            "code_runs_today":   await _cnt("code_runs_per_day"),
            "documents":         await db_count("SELECT COUNT(*) as c FROM documents WHERE user_id=?",(uid,)),
            "memories":          await db_count("SELECT COUNT(*) as c FROM memories WHERE user_id=? AND is_active=1",(uid,)),
            "agent_jobs":        await db_count("SELECT COUNT(*) as c FROM agent_jobs WHERE user_id=?",(uid,)),
            "connectors":        await db_count("SELECT COUNT(*) as c FROM connectors WHERE user_id=? AND is_active=1",(uid,)),
        }
    }

def _billing_plan_amount(tier: str) -> int:
    return {
        "pro": RAZORPAY_PRO_MONTHLY,
        "premium": RAZORPAY_PREMIUM_MONTHLY,
        "enterprise": RAZORPAY_ENT_MONTHLY,
    }.get(tier, 0)

def _billing_plan_rows(current: str = "free") -> List[Dict[str, Any]]:
    return [
        {
            "tier": "free", "name": "Free", "amount": 0, "currency": RAZORPAY_CURRENCY,
            "display": "Free", "current": current == "free",
            "features": ["Starter chat usage", "Basic file analysis", "Limited documents"],
        },
        {
            "tier": "pro", "name": "Pro", "amount": _billing_plan_amount("pro"), "currency": RAZORPAY_CURRENCY,
            "display": f"{RAZORPAY_CURRENCY} {RAZORPAY_PRO_MONTHLY / 100:,.0f}/mo", "current": current == "pro",
            "features": ["More daily messages", "Larger uploads and documents", "Workspace tools"],
        },
        {
            "tier": "premium", "name": "Premium", "amount": _billing_plan_amount("premium"), "currency": RAZORPAY_CURRENCY,
            "display": f"{RAZORPAY_CURRENCY} {RAZORPAY_PREMIUM_MONTHLY / 100:,.0f}/mo", "current": current == "premium",
            "features": ["Unlimited-style chat", "Priority model access", "Advanced files, RAG, and agents"],
        },
        {
            "tier": "enterprise", "name": "Enterprise", "amount": _billing_plan_amount("enterprise"), "currency": RAZORPAY_CURRENCY,
            "display": f"{RAZORPAY_CURRENCY} {RAZORPAY_ENT_MONTHLY / 100:,.0f}/mo", "current": current == "enterprise",
            "features": ["Unlimited-style limits", "Admin controls", "Advanced connectors and models"],
        },
    ]

def _razorpay_post(path: str, payload: Dict[str, Any]) -> Dict[str, Any]:
    url = "https://api.razorpay.com/v1" + path
    raw = json.dumps(payload).encode("utf-8")
    auth = base64.b64encode(f"{RAZORPAY_KEY_ID}:{RAZORPAY_KEY_SECRET}".encode()).decode()
    req = urllib.request.Request(
        url,
        data=raw,
        headers={"Authorization": f"Basic {auth}", "Content-Type": "application/json"},
        method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=20) as resp:
            return json.loads(resp.read().decode("utf-8"))
    except urllib.error.HTTPError as e:
        detail = e.read().decode("utf-8", "replace")
        raise HTTPException(502, f"Razorpay order failed: {detail[:500]}")

@app.get("/billing/plans")
async def billing_plans(user: Dict = Depends(_get_current_user)):
    return {
        "provider": "razorpay",
        "enabled": bool(RAZORPAY_KEY_ID and RAZORPAY_KEY_SECRET),
        "key_id": RAZORPAY_KEY_ID,
        "currency": RAZORPAY_CURRENCY,
        "current": user.get("subscription", "free"),
        "plans": _billing_plan_rows(user.get("subscription", "free")),
    }

@app.post("/billing/razorpay/order")
async def billing_create_razorpay_order(req: BillingOrderReq, user: Dict = Depends(_get_current_user)):
    if not (RAZORPAY_KEY_ID and RAZORPAY_KEY_SECRET):
        raise HTTPException(400, "Razorpay is not configured. Add RAZORPAY_KEY_ID and RAZORPAY_KEY_SECRET.")
    uid = user.get("id") or user.get("sub", "")
    amount = _billing_plan_amount(req.tier)
    if amount <= 0:
        raise HTTPException(400, "Invalid billing amount")
    now = _utcnow()
    receipt = f"jz_{uid[:8]}_{int(time.time())}_{secrets.token_hex(3)}"[:40]
    payload = {
        "amount": amount,
        "currency": RAZORPAY_CURRENCY,
        "receipt": receipt,
        "notes": {"user_id": uid[:36], "tier": req.tier, "email": user.get("email", "")[:128]},
    }
    order = await asyncio.get_running_loop().run_in_executor(_executor, _razorpay_post, "/orders", payload)
    await db_execute(
        "INSERT INTO billing_orders(id,user_id,provider,tier,amount,currency,provider_order_id,receipt,status,raw_json,created_at,updated_at)"
        " VALUES(?,?,?,?,?,?,?,?,?,?,?,?)",
        (_new_id(), uid, "razorpay", req.tier, amount, RAZORPAY_CURRENCY, order.get("id"), receipt,
         order.get("status", "created"), json.dumps(order), now, now),
    )
    return {
        "key_id": RAZORPAY_KEY_ID,
        "order_id": order.get("id"),
        "amount": amount,
        "currency": RAZORPAY_CURRENCY,
        "tier": req.tier,
        "name": "JAZZ AI",
        "description": f"{req.tier.title()} plan",
        "prefill": {"name": user.get("full_name", ""), "email": user.get("email", "")},
    }

@app.post("/billing/razorpay/verify")
async def billing_verify_razorpay_payment(req: RazorpayVerifyReq, user: Dict = Depends(_get_current_user)):
    if not RAZORPAY_KEY_SECRET:
        raise HTTPException(400, "Razorpay is not configured")
    uid = user.get("id") or user.get("sub", "")
    order = await db_fetchone(
        "SELECT * FROM billing_orders WHERE provider_order_id=? AND user_id=?",
        (req.razorpay_order_id, uid),
    )
    if not order:
        raise HTTPException(404, "Billing order not found")
    digest = hmac.new(
        RAZORPAY_KEY_SECRET.encode("utf-8"),
        f"{req.razorpay_order_id}|{req.razorpay_payment_id}".encode("utf-8"),
        hashlib.sha256,
    ).hexdigest()
    if not hmac.compare_digest(digest, req.razorpay_signature):
        raise HTTPException(400, "Invalid Razorpay signature")
    now = _utcnow()
    tier = order["tier"]
    expires_at = (datetime.now(timezone.utc) + timedelta(days=30)).isoformat()
    await db_execute(
        "UPDATE billing_orders SET provider_payment_id=?,status='paid',updated_at=? WHERE id=?",
        (req.razorpay_payment_id, now, order["id"]),
    )
    await db_execute(
        "UPDATE users SET subscription=?,subscription_expires_at=?,updated_at=? WHERE id=?",
        (tier, expires_at, now, uid),
    )
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), uid, uid, "razorpay_subscription_payment",
         json.dumps({"tier": tier, "order_id": req.razorpay_order_id, "payment_id": req.razorpay_payment_id,
                     "subscription_expires_at": expires_at}), now),
    )
    return {"ok": True, "tier": tier, "payment_id": req.razorpay_payment_id, "subscription_expires_at": expires_at}

@app.post("/coupons/redeem")
async def redeem_coupon(body: dict, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    code = body.get("code","").strip()
    coupon = await db_fetchone(
        "SELECT * FROM coupons WHERE code=? AND is_active=1 AND (max_uses<=0 OR current_uses<max_uses)"
        " AND (expires_at IS NULL OR expires_at>?)", (code, _utcnow()))
    if not coupon: raise HTTPException(400, "Invalid or expired coupon code")
    existing = await db_fetchone("SELECT id FROM coupon_redemptions WHERE coupon_id=? AND user_id=?",(coupon["id"],uid))
    if existing: raise HTTPException(400, "Coupon already used")
    exp = (datetime.now(timezone.utc)+timedelta(days=coupon["duration_days"])).isoformat()
    now = _utcnow()
    async with db_transaction():
        await db_execute("UPDATE users SET subscription=?,subscription_expires_at=?,updated_at=? WHERE id=?",
                         (coupon["grants_tier"],exp,now,uid))
        await db_execute("UPDATE coupons SET current_uses=current_uses+1 WHERE id=?",(coupon["id"],))
        await db_execute("INSERT INTO coupon_redemptions(id,coupon_id,user_id,redeemed_at) VALUES(?,?,?,?)",
                         (_new_id(),coupon["id"],uid,now))
    return {"ok":True,"granted_tier":coupon["grants_tier"],"expires_at":exp}

@app.get("/search")
async def search(q: str, scope: str = "messages", limit: int = 20,
                 user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    if len(q.strip()) < 2: raise HTTPException(400, "Query must be at least 2 characters")
    pattern = f"%{q}%"; results: Dict = {}
    if scope in ("messages","all"):
        rows = await db_fetchall(
            "SELECT ch.id,ch.role,ch.content,ch.created_at,cs.id as session_id,cs.title as session_title"
            " FROM chat_history ch JOIN chat_sessions cs ON ch.session_id=cs.id"
            " WHERE ch.user_id=? AND ch.is_hidden=0 AND ch.content LIKE ? ORDER BY ch.created_at DESC LIMIT ?",
            (uid, pattern, limit))
        results["messages"] = rows
    if scope in ("memories","all"):
        results["memories"] = await db_fetchall(
            "SELECT id,key,value,source,confidence FROM memories"
            " WHERE user_id=? AND is_active=1 AND (key LIKE ? OR value LIKE ?) LIMIT ?",
            (uid, pattern, pattern, limit))
    return {"query":q,"scope":scope,"results":results}

@app.get("/gdpr/delete-account")
async def gdpr_delete(confirm: str = "", user: Dict = Depends(_get_current_user)):
    if confirm.replace("+"," ") != "DELETE MY ACCOUNT":
        raise HTTPException(400, "Confirmation text mismatch")
    uid = user.get("id") or user.get("sub","")
    await db_execute("DELETE FROM users WHERE id=?",(uid,))
    return {"ok":True,"message":"Account deleted"}

# ══════════════════════════════════════════════════════════════════════════════
# §36  MCP SERVERS
# ══════════════════════════════════════════════════════════════════════════════

@app.get("/mcp-servers")
async def list_mcp_servers(user: Dict = Depends(_get_current_user)):
    rows = await db_fetchall(
        "SELECT id,name,description,server_url,tools_json,ping_ok,last_pinged_at"
        " FROM mcp_servers WHERE is_active=1 AND is_public=1 ORDER BY name")
    return {"servers":rows}

@app.post("/mcp/call")
async def call_mcp_tool(req: MCPToolCallReq, user: Dict = Depends(_get_current_user)):
    row = await db_fetchone("SELECT * FROM mcp_servers WHERE id=? AND is_active=1",(req.server_id,))
    if not row: raise HTTPException(404, "MCP server not found")
    if not row["is_public"] and user.get("role") != "admin": raise HTTPException(403)
    try:
        result = await _mcp_call_tool(req.server_id, req.tool_name, req.tool_args)
        return {"success":True,"server":row["name"],"tool":req.tool_name,"result":result}
    except Exception as e:
        raise HTTPException(502, f"MCP tool call failed: {e}")

@app.get("/mcp-servers/{server_id}/tools")
async def list_mcp_server_tools(server_id: str, user: Dict = Depends(_get_current_user)):
    row = await db_fetchone("SELECT * FROM mcp_servers WHERE id=? AND is_active=1",(server_id,))
    if not row: raise HTTPException(404)
    if not row["is_public"] and user.get("role") != "admin": raise HTTPException(403)
    tools = await _mcp_list_tools(server_id)
    return {"server":row["name"],"tools":tools}

# ══════════════════════════════════════════════════════════════════════════════
# §37  APACHE LIVY ROUTES
# ══════════════════════════════════════════════════════════════════════════════

@app.get("/livy/health")
async def livy_health(user: Dict = Depends(_get_current_user)):
    try:
        result = await _livy_request("GET","/version")
        return {"available":True,"url":LIVY_URL,"version":result}
    except Exception as e:
        return {"available":False,"url":LIVY_URL,"error":str(e)}

@app.post("/livy/sessions")
async def livy_create_session_ep(req: LivySessionCreate, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    try:
        result = await _livy_request("POST","/sessions",{"kind":req.kind,"name":req.name})
        livy_id = result.get("id")
        sid = _new_id(); now = _utcnow()
        await db_execute(
            "INSERT INTO livy_sessions(id,user_id,livy_id,livy_url,kind,state,created_at,updated_at)"
            " VALUES(?,?,?,?,?,?,?,?)",
            (sid, uid, livy_id, LIVY_URL, req.kind, result.get("state","starting"), now, now))
        return {"id":sid,"livy_id":livy_id,"state":result.get("state"),"kind":req.kind}
    except Exception as e: raise HTTPException(502, f"Livy session creation failed: {e}")

@app.post("/livy/sessions/{livy_session_id}/statements")
async def livy_submit_statement(livy_session_id: int, req: LivyStatementCreate, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    try:
        result = await _livy_request("POST",f"/sessions/{livy_session_id}/statements",{"code":req.code,"kind":req.kind})
        return {"statement_id":result.get("id"),"state":result.get("state")}
    except Exception as e: raise HTTPException(502, f"Livy statement failed: {e}")

@app.get("/livy/sessions/{livy_session_id}/statements/{statement_id}")
async def livy_get_statement(livy_session_id: int, statement_id: int, user: Dict = Depends(_get_current_user)):
    try:
        result = await _livy_request("GET",f"/sessions/{livy_session_id}/statements/{statement_id}")
        output = result.get("output",{})
        return {"statement_id":statement_id,"state":result.get("state"),"progress":result.get("progress",0.0),
                "output":output,"text":output.get("text","") if output.get("status")=="ok" else "",
                "error":output.get("evalue","") if output.get("status")=="error" else ""}
    except Exception as e: raise HTTPException(502, f"Livy error: {e}")

@app.post("/livy/batches")
async def livy_submit_batch_ep(req: LivyBatchCreate, user: Dict = Depends(_get_current_user)):
    try:
        body: Dict = {"file":req.file,"args":req.args,"conf":req.conf}
        if req.class_name: body["className"] = req.class_name
        result = await _livy_request("POST","/batches",body)
        return {"batch_id":result.get("id"),"state":result.get("state"),"app_id":result.get("appId")}
    except Exception as e: raise HTTPException(502, f"Livy batch failed: {e}")

@app.get("/livy/batches/{batch_id}")
async def livy_get_batch(batch_id: int, user: Dict = Depends(_get_current_user)):
    try: return await _livy_request("GET",f"/batches/{batch_id}")
    except Exception as e: raise HTTPException(502, f"Livy error: {e}")

# ══════════════════════════════════════════════════════════════════════════════
# §38  ADMIN ROUTES
# ══════════════════════════════════════════════════════════════════════════════

@app.get("/admin/stats")
async def admin_stats(user: Dict = Depends(_require_admin)):
    uptime = (datetime.now(timezone.utc)-_SERVER_START).total_seconds()
    user_stats = {
        "total":      await db_count("SELECT COUNT(*) as c FROM users"),
        "active":     await db_count("SELECT COUNT(*) as c FROM users WHERE is_active=1"),
        "free":       await db_count("SELECT COUNT(*) as c FROM users WHERE COALESCE(subscription,'free')='free'"),
        "pro":        await db_count("SELECT COUNT(*) as c FROM users WHERE subscription='pro'"),
        "premium":    await db_count("SELECT COUNT(*) as c FROM users WHERE subscription='premium'"),
        "enterprise": await db_count("SELECT COUNT(*) as c FROM users WHERE subscription='enterprise'"),
    }
    content_stats = {
        "messages":   await db_count("SELECT COUNT(*) as c FROM chat_history"),
        "documents":  await db_count("SELECT COUNT(*) as c FROM documents"),
        "sessions":   await db_count("SELECT COUNT(*) as c FROM chat_sessions"),
        "agent_jobs": await db_count("SELECT COUNT(*) as c FROM agent_jobs"),
        "connectors": await db_count("SELECT COUNT(*) as c FROM connectors"),
        "websites":   await db_count("SELECT COUNT(*) as c FROM websites"),
        "ai_models":  await db_count("SELECT COUNT(*) as c FROM ai_models WHERE is_active=1"),
        "mcp_servers":await db_count("SELECT COUNT(*) as c FROM mcp_servers WHERE is_active=1"),
        "skills":     await db_count("SELECT COUNT(*) as c FROM skills WHERE is_enabled=1"),
    }
    tier_stats = {
        "free": user_stats["free"],
        "pro": user_stats["pro"],
        "premium": user_stats["premium"],
        "enterprise": user_stats["enterprise"],
    }
    return {
        "users": user_stats["total"],
        "active_users": user_stats["active"],
        "user_stats": user_stats,
        "by_tier": tier_stats,
        "plans": tier_stats,
        "content": content_stats,
        **content_stats,
        "system": {
            "uptime_seconds": int(uptime), "db_path": DB_PATH,
            "chroma": HAS_CHROMA and _chroma_client is not None,
            "fitz": HAS_FITZ, "docx": HAS_DOCX, "livekit": _lk_available(),
        },
        "groq_configured": bool(GROQ_API_KEY),
    }

@app.get("/admin/users")
async def admin_list_users(limit: int = 100, offset: int = 0, user: Dict = Depends(_require_admin)):
    users = await db_fetchall(
        "SELECT id,email,full_name,role,subscription,is_active,is_verified,last_login_at,created_at"
        " FROM users ORDER BY created_at DESC LIMIT ? OFFSET ?", (limit, offset))
    total = await db_count("SELECT COUNT(*) as c FROM users")
    return {"users":users,"total":total}

@app.patch("/admin/users/{uid}")
async def admin_update_user(uid: str, body: dict, admin: Dict = Depends(_require_admin)):
    allowed = {"subscription","role","is_active","is_verified","subscription_expires_at"}
    fields  = {k:v for k,v in body.items() if k in allowed}
    if not fields: return {"ok":True}
    fields["updated_at"] = _utcnow()
    set_clause = ", ".join(f"{k}=?" for k in fields)
    await db_execute(f"UPDATE users SET {set_clause} WHERE id=?",tuple(fields.values())+(uid,))
    await db_execute("INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
                     (_new_id(), admin.get("id") or admin.get("sub",""), uid, "admin_update_user",
                      json.dumps(fields), _utcnow()))
    return {"ok":True}

@app.put("/admin/subscription")
async def admin_update_subscription(req: SubscriptionUpdate, admin: Dict = Depends(_require_admin)):
    row = await db_fetchone(
        "SELECT id,email,subscription,subscription_expires_at FROM users WHERE id=?",
        (req.user_id,))
    if not row:
        raise HTTPException(404, "User not found")
    now = _utcnow()
    admin_uid = admin.get("id") or admin.get("sub","")
    previous = {
        "subscription": row.get("subscription") or "free",
        "subscription_expires_at": row.get("subscription_expires_at"),
    }
    if req.tier == "free":
        await db_execute(
            "UPDATE users SET subscription=?,subscription_expires_at=NULL,updated_at=? WHERE id=?",
            (req.tier, now, req.user_id))
        expires_at = None
    else:
        await db_execute(
            "UPDATE users SET subscription=?,updated_at=? WHERE id=?",
            (req.tier, now, req.user_id))
        expires_at = row.get("subscription_expires_at")
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), admin_uid, req.user_id, "admin_subscription_update",
         json.dumps({"previous": previous, "tier": req.tier, "subscription_expires_at": expires_at}), now))
    return {"ok": True, "user_id": req.user_id, "tier": req.tier, "subscription_expires_at": expires_at}

@app.get("/admin/sidebar-features")
async def admin_list_sidebar_features(admin: Dict = Depends(_require_admin)):
    rows = await _sidebar_feature_rows()
    return {"features": rows, "total": len(rows)}

@app.patch("/admin/sidebar-features/{feature_key}")
async def admin_update_sidebar_feature(feature_key: str, body: dict,
                                       admin: Dict = Depends(_require_admin)):
    feature_key = (feature_key or "").strip()
    known_keys = await _sidebar_feature_key_set()
    if feature_key not in known_keys:
        raise HTTPException(404, "Sidebar feature not found")
    is_enabled = body.get("is_enabled")
    if is_enabled is None:
        return {"ok": True}
    await db_execute(
        "UPDATE sidebar_features SET is_enabled=?,updated_at=? WHERE feature_key=?",
        (1 if bool(is_enabled) else 0, _utcnow(), feature_key))
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), admin.get("id") or admin.get("sub",""), feature_key, "sidebar_feature_toggle",
         json.dumps({"is_enabled": bool(is_enabled)}), _utcnow()))
    return {"ok": True, "feature_key": feature_key, "is_enabled": bool(is_enabled)}

@app.post("/admin/sidebar-features/bulk")
async def admin_bulk_sidebar_features(body: dict, admin: Dict = Depends(_require_admin)):
    is_enabled = bool(body.get("is_enabled"))
    now = _utcnow()
    await db_execute(
        "UPDATE sidebar_features SET is_enabled=?,updated_at=?",
        (1 if is_enabled else 0, now))
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), admin.get("id") or admin.get("sub",""), "sidebar_features",
         "sidebar_features_bulk_toggle", json.dumps({"is_enabled": is_enabled}), now))
    return {"ok": True, "is_enabled": is_enabled}

@app.get("/admin/users/{uid}/features")
async def admin_get_user_sidebar_features(uid: str, admin: Dict = Depends(_require_admin)):
    user_row = await db_fetchone(
        "SELECT id,email,full_name,role,subscription FROM users WHERE id=?", (uid,))
    if not user_row:
        raise HTTPException(404, "User not found")
    access = await _sidebar_feature_access_for_user(uid)
    return {"user": user_row, **access}

@app.patch("/admin/users/{uid}/features")
async def admin_update_user_sidebar_features(uid: str, body: dict,
                                             admin: Dict = Depends(_require_admin)):
    user_row = await db_fetchone("SELECT id FROM users WHERE id=?", (uid,))
    if not user_row:
        raise HTTPException(404, "User not found")
    overrides = body.get("overrides", {})
    if not isinstance(overrides, dict):
        raise HTTPException(400, "overrides must be an object")
    now = _utcnow()
    admin_uid = admin.get("id") or admin.get("sub","")
    changed: Dict[str, Any] = {}
    known_keys = await _sidebar_feature_key_set()
    for key, value in overrides.items():
        key = str(key).strip()
        if key not in known_keys:
            continue
        if value is None:
            await db_execute(
                "DELETE FROM user_sidebar_feature_access WHERE user_id=? AND feature_key=?",
                (uid, key))
            changed[key] = None
        else:
            enabled = 1 if bool(value) else 0
            await db_execute(
                "INSERT INTO user_sidebar_feature_access"
                "(id,user_id,feature_key,is_enabled,updated_by,updated_at) VALUES(?,?,?,?,?,?)"
                " ON CONFLICT(user_id,feature_key) DO UPDATE SET "
                "is_enabled=excluded.is_enabled,updated_by=excluded.updated_by,updated_at=excluded.updated_at",
                (_new_id(), uid, key, enabled, admin_uid, now))
            changed[key] = bool(value)
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), admin_uid, uid, "user_sidebar_features_update",
         json.dumps(changed), now))
    access = await _sidebar_feature_access_for_user(uid)
    return {"ok": True, "changed": changed, **access}

@app.get("/admin/users/{uid}/models")
async def admin_get_user_model_access(uid: str, admin: Dict = Depends(_require_admin)):
    user_row = await db_fetchone(
        "SELECT id,email,full_name,role,subscription FROM users WHERE id=?", (uid,))
    if not user_row:
        raise HTTPException(404, "User not found")
    access = await _model_access_for_user(uid)
    return {"user": user_row, **access}

@app.patch("/admin/users/{uid}/models")
async def admin_update_user_model_access(uid: str, body: dict,
                                         admin: Dict = Depends(_require_admin)):
    user_row = await db_fetchone("SELECT id FROM users WHERE id=?", (uid,))
    if not user_row:
        raise HTTPException(404, "User not found")
    overrides = body.get("overrides", {})
    if not isinstance(overrides, dict):
        raise HTTPException(400, "overrides must be an object")
    known = {_canonical_model_id(r["id"]) for r in await _model_access_rows(include_inactive_db=True)}
    now = _utcnow()
    admin_uid = admin.get("id") or admin.get("sub","")
    changed: Dict[str, Any] = {}
    for key, value in overrides.items():
        model_id = _canonical_model_id(str(key).strip())
        if model_id not in known:
            continue
        if value is None:
            await db_execute(
                "DELETE FROM user_model_access WHERE user_id=? AND model_id=?",
                (uid, model_id))
            changed[model_id] = None
        else:
            enabled = 1 if bool(value) else 0
            await db_execute(
                "INSERT INTO user_model_access"
                "(id,user_id,model_id,is_enabled,updated_by,updated_at) VALUES(?,?,?,?,?,?)"
                " ON CONFLICT(user_id,model_id) DO UPDATE SET "
                "is_enabled=excluded.is_enabled,updated_by=excluded.updated_by,updated_at=excluded.updated_at",
                (_new_id(), uid, model_id, enabled, admin_uid, now))
            changed[model_id] = bool(value)
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), admin_uid, uid, "user_model_access_update",
         json.dumps(changed), now))
    access = await _model_access_for_user(uid)
    return {"ok": True, "changed": changed, **access}

@app.delete("/admin/users/{uid}")
async def admin_delete_user(uid: str, admin: Dict = Depends(_require_admin)):
    admin_uid = admin.get("id") or admin.get("sub","")
    if uid == admin_uid: raise HTTPException(400, "Cannot delete yourself")
    await db_execute("DELETE FROM users WHERE id=?",(uid,))
    return {"ok":True}

@app.post("/admin/users/{uid}/impersonate")
async def admin_impersonate(uid: str, admin: Dict = Depends(_require_admin)):
    user = await db_fetchone("SELECT * FROM users WHERE id=?",(uid,))
    if not user: raise HTTPException(404)
    token = _make_access_token(uid, user["role"], user["email"], user["subscription"])
    await db_execute("INSERT INTO audit_log(id,actor_id,target_id,action,created_at) VALUES(?,?,?,?,?)",
                     (_new_id(), admin.get("id") or admin.get("sub",""), uid, "impersonate", _utcnow()))
    return {"access_token":token,"token_type":"bearer","warning":"1-hour impersonation token"}

@app.get("/admin/audit-log")
async def admin_audit_log(limit: int = 100, admin: Dict = Depends(_require_admin)):
    logs = await db_fetchall("SELECT * FROM audit_log ORDER BY created_at DESC LIMIT ?",(limit,))
    return {"logs":logs}

@app.get("/admin/connector-logs")
async def admin_connector_logs(
    connector_type: str = "",
    user: str = "",
    action: str = "",
    status: str = "",
    date_from: str = "",
    date_to: str = "",
    limit: int = 100,
    admin: Dict = Depends(_require_admin)
):
    limit = max(1, min(500, limit))
    where, vals = [], []
    if connector_type:
        where.append("cl.connector_type=?"); vals.append(connector_type)
    if user:
        where.append("(cl.user_id=? OR u.email LIKE ?)"); vals.extend([user, f"%{user}%"])
    if action:
        where.append("cl.action LIKE ?"); vals.append(f"%{action}%")
    if status in ("success", "failure"):
        where.append("cl.status=?"); vals.append(status)
    if date_from:
        where.append("date(cl.created_at) >= date(?)"); vals.append(date_from)
    if date_to:
        where.append("date(cl.created_at) <= date(?)"); vals.append(date_to)
    sql_where = (" WHERE " + " AND ".join(where)) if where else ""
    rows = await db_fetchall(
        "SELECT cl.*, u.email FROM connector_logs cl LEFT JOIN users u ON u.id=cl.user_id"
        f"{sql_where} ORDER BY cl.created_at DESC LIMIT ?",
        tuple(vals + [limit]))
    return {"logs": rows, "total": len(rows)}

@app.post("/admin/coupons")
async def admin_create_coupon(req: CouponCreate, admin: Dict = Depends(_require_admin)):
    cid = _new_id(); now = _utcnow()
    code = req.code or secrets.token_hex(4).upper()
    exp = None
    if req.expires_days: exp = (datetime.now(timezone.utc)+timedelta(days=req.expires_days)).isoformat()
    admin_uid = admin.get("id") or admin.get("sub","")
    await db_execute(
        "INSERT INTO coupons(id,code,grants_tier,duration_days,max_uses,is_active,created_by,expires_at,created_at)"
        " VALUES(?,?,?,?,?,1,?,?,?)",
        (cid, code, req.grants_tier, req.duration_days, req.max_uses, admin_uid, exp, now))
    return {"id":cid,"code":code}

@app.get("/admin/coupons")
async def admin_list_coupons(admin: Dict = Depends(_require_admin)):
    coupons = await db_fetchall("SELECT * FROM coupons ORDER BY created_at DESC")
    return {"coupons":coupons}

@app.get("/admin/skills/markdown-template")
async def admin_skill_markdown_template(admin: Dict = Depends(_require_admin)):
    return {"template": _SKILL_MD_TEMPLATE}

@app.post("/admin/skills/import-markdown")
async def admin_import_skill_markdown(body: dict, admin: Dict = Depends(_require_admin)):
    try:
        parsed = _parse_skill_markdown(body.get("markdown", ""), body.get("name", ""))
    except ValueError as exc:
        raise HTTPException(400, str(exc))
    now = _utcnow()
    admin_uid = admin.get("id") or admin.get("sub","")
    existing = await db_fetchone("SELECT id FROM skills WHERE name=?", (parsed["name"],))
    if existing:
        await db_execute(
            "UPDATE skills SET description=?,activation_keywords=?,instructions=?,resources_json=?,"
            "script_metadata_json=?,is_enabled=?,updated_at=? WHERE id=?",
            (parsed["description"], json.dumps(parsed["activation_keywords"]), parsed["instructions"],
             json.dumps(parsed["resources"]), json.dumps(parsed["script_metadata"]),
             int(body.get("is_enabled", True)), now, existing["id"]))
        return {"id": existing["id"], "name": parsed["name"], "updated": True}
    sid = _new_id()
    await db_execute(
        "INSERT INTO skills(id,name,description,activation_keywords,instructions,resources_json,"
        "script_metadata_json,is_enabled,created_by,created_at,updated_at) VALUES(?,?,?,?,?,?,?,?,?,?,?)",
        (sid, parsed["name"], parsed["description"], json.dumps(parsed["activation_keywords"]),
         parsed["instructions"], json.dumps(parsed["resources"]), json.dumps(parsed["script_metadata"]),
         int(body.get("is_enabled", True)), admin_uid, now, now))
    return {"id": sid, "name": parsed["name"], "updated": False}

@app.get("/admin/skills")
async def admin_list_skills(admin: Dict = Depends(_require_admin)):
    rows = await db_fetchall("SELECT * FROM skills ORDER BY is_enabled DESC,name")
    return {"skills": [_normalise_skill_row(r) for r in rows], "total": len(rows)}

@app.post("/admin/skills")
async def admin_create_skill(req: SkillCreate, admin: Dict = Depends(_require_admin)):
    sid = _new_id(); now = _utcnow()
    await db_execute(
        "INSERT INTO skills(id,name,description,activation_keywords,instructions,resources_json,"
        "script_metadata_json,is_enabled,created_by,created_at,updated_at) VALUES(?,?,?,?,?,?,?,?,?,?,?)",
        (sid, req.name.strip(), req.description, json.dumps(req.activation_keywords),
         req.instructions, json.dumps(req.resources), json.dumps(req.script_metadata),
         int(req.is_enabled), admin.get("id") or admin.get("sub",""), now, now))
    return {"id": sid, "name": req.name}

@app.put("/admin/skills/{skill_id}")
async def admin_update_skill(skill_id: str, req: SkillUpdate, admin: Dict = Depends(_require_admin)):
    row = await db_fetchone("SELECT id FROM skills WHERE id=?", (skill_id,))
    if not row:
        raise HTTPException(404, "Skill not found")
    sets, vals = [], []
    if req.name is not None: sets.append("name=?"); vals.append(req.name.strip())
    if req.description is not None: sets.append("description=?"); vals.append(req.description)
    if req.activation_keywords is not None: sets.append("activation_keywords=?"); vals.append(json.dumps(req.activation_keywords))
    if req.instructions is not None: sets.append("instructions=?"); vals.append(req.instructions)
    if req.resources is not None: sets.append("resources_json=?"); vals.append(json.dumps(req.resources))
    if req.script_metadata is not None: sets.append("script_metadata_json=?"); vals.append(json.dumps(req.script_metadata))
    if req.is_enabled is not None: sets.append("is_enabled=?"); vals.append(int(req.is_enabled))
    if not sets:
        return {"ok": True}
    sets.append("updated_at=?"); vals.append(_utcnow()); vals.append(skill_id)
    await db_execute(f"UPDATE skills SET {', '.join(sets)} WHERE id=?", tuple(vals))
    return {"ok": True}

@app.delete("/admin/skills/{skill_id}")
async def admin_delete_skill(skill_id: str, admin: Dict = Depends(_require_admin)):
    await db_execute("DELETE FROM skills WHERE id=?", (skill_id,))
    return {"ok": True}

@app.post("/admin/skills/{skill_id}/test")
async def admin_test_skill(skill_id: str, body: dict, admin: Dict = Depends(_require_admin)):
    row = await db_fetchone("SELECT * FROM skills WHERE id=?", (skill_id,))
    if not row:
        raise HTTPException(404, "Skill not found")
    skill = _normalise_skill_row(row)
    prompt = body.get("message") or "Use this skill on a short example."
    messages = [
        {"role":"system","content":_format_skill_context([skill])},
        {"role":"user","content":prompt},
    ]
    answer = await _llm_text(messages, body.get("model_id","llama-3.3-70b-versatile"), max_tokens=1000)
    return {"skill": skill, "prompt": prompt, "answer": answer}

@app.post("/admin/broadcast")
async def admin_broadcast(body: dict, admin: Dict = Depends(_require_admin)):
    nid = _new_id(); now = _utcnow()
    await db_execute(
        "INSERT INTO notifications(id,user_id,title,message,type,is_broadcast,is_read,created_at)"
        " VALUES(?,NULL,?,?,?,1,0,?)",
        (nid, body.get("title",""), body.get("message",""), body.get("type","info"), now))
    await ws_manager.broadcast({"type":"broadcast","title":body.get("title",""),"message":body.get("message","")})
    return {"ok":True,"id":nid}

def _model_slug(value: str) -> str:
    s = re.sub(r"[^a-zA-Z0-9]+", "-", str(value or "").strip().lower()).strip("-")
    return (s or f"model-{_new_id()[:8]}")[:120]

def _balanced_model_dump_blocks(src: str) -> List[str]:
    blocks: List[str] = []
    stack: List[Tuple[str, int]] = []
    quote = ""
    escape = False
    pairs = {"[": "]", "{": "}"}
    for i, ch in enumerate(src or ""):
        if quote:
            if escape:
                escape = False
            elif ch == "\\":
                escape = True
            elif ch == quote:
                quote = ""
            continue
        if ch in ("'", '"'):
            quote = ch
            continue
        if ch in pairs:
            stack.append((ch, i))
        elif ch in ("]", "}") and stack:
            open_ch, start = stack[-1]
            if pairs[open_ch] != ch:
                stack.clear()
                continue
            stack.pop()
            if not stack:
                block = src[start:i + 1].strip()
                if len(block) > 2:
                    blocks.append(block)
    return blocks

def _model_dump_candidates(raw: str) -> List[str]:
    src = (raw or "").strip()
    fenced = re.findall(r"```(?:json|js|javascript|python)?\s*([\s\S]*?)```", src, flags=re.I)
    candidates = [x.strip() for x in fenced if x.strip()] + [src]
    out: List[str] = []
    for cand in candidates:
        c = cand.strip().rstrip(";")
        c = re.sub(r"^\s*export\s+default\s+", "", c, flags=re.I)
        c = re.sub(r"^\s*module\.exports\s*=\s*", "", c, flags=re.I)
        c = re.sub(r"^\s*(?:const|let|var)\s+[A-Za-z_$][\w$]*\s*=\s*", "", c, flags=re.I)
        c = re.sub(r"^\s*[A-Za-z_$][\w$\.]*\s*=\s*", "", c)
        out.append(c.strip().rstrip(";"))
        starts = [(c.find("["), "["), (c.find("{"), "{")]
        starts = [(idx, ch) for idx, ch in starts if idx >= 0]
        if starts:
            start, ch = min(starts, key=lambda x: x[0])
            end = c.rfind("]" if ch == "[" else "}")
            if end > start:
                out.append(c[start:end + 1].strip())
        out.extend(_balanced_model_dump_blocks(c))
    seen: Set[str] = set()
    return [x for x in out if x and not (x in seen or seen.add(x))]

def _parse_model_dump(raw: str) -> Any:
    last_error = "Could not parse model dump"
    for cand in _model_dump_candidates(raw):
        for text in (cand, re.sub(r"([{\[,]\s*)([A-Za-z_][\w-]*)\s*:", r'\1"\2":', cand)):
            try:
                return json.loads(text)
            except Exception as exc:
                last_error = str(exc)
            literal = re.sub(r"\btrue\b", "True", text)
            literal = re.sub(r"\bfalse\b", "False", literal)
            literal = re.sub(r"\bnull\b", "None", literal)
            try:
                return ast.literal_eval(literal)
            except Exception as exc:
                last_error = str(exc)
    raise ValueError(last_error)

def _model_dump_items(data: Any) -> List[Dict[str, Any]]:
    if isinstance(data, list):
        return [x for x in data if isinstance(x, dict)]
    if not isinstance(data, dict):
        return []
    for key in ("models", "db_models", "builtin_models", "items", "data", "payload"):
        value = data.get(key)
        if isinstance(value, list):
            return [x for x in value if isinstance(x, dict)]
        if isinstance(value, dict):
            nested = _model_dump_items(value)
            if nested:
                return nested
    modelish = {"id", "model_id", "model_name", "model", "name", "label", "display_name", "provider_model"}
    if modelish.intersection(data.keys()):
        return [data]
    if data and all(isinstance(v, dict) for v in data.values()):
        return [dict(v, id=str(k)) for k, v in data.items()]
    return []

def _dump_str(item: Dict[str, Any], *names: str, default: str = "") -> str:
    for name in names:
        value = item.get(name)
        if value is not None and not isinstance(value, (dict, list)):
            text = str(value).strip()
            if text:
                return text
    return default

def _dump_bool(item: Dict[str, Any], names: Tuple[str, ...], default: bool) -> bool:
    for name in names:
        if name not in item:
            continue
        value = item.get(name)
        if isinstance(value, bool):
            return value
        if isinstance(value, (int, float)):
            return bool(value)
        text = str(value).strip().lower()
        if text in {"1", "true", "yes", "on", "enabled", "active"}:
            return True
        if text in {"0", "false", "no", "off", "disabled", "inactive"}:
            return False
    return default

def _dump_int(item: Dict[str, Any], names: Tuple[str, ...], default: int, lo: int, hi: int) -> int:
    for name in names:
        if name not in item:
            continue
        try:
            return max(lo, min(hi, int(float(str(item.get(name)).replace(",", "")))))
        except Exception:
            pass
    return default

def _normalize_imported_model(item: Dict[str, Any]) -> Dict[str, Any]:
    model_name = _dump_str(item, "model_name", "provider_model", "model", "modelId", "model_id")
    name = _dump_str(item, "name", "display_name", "label", default=model_name)
    if not model_name and name:
        model_name = name
    if not model_name:
        raise ValueError("model_name/model is required")
    raw_base = _dump_str(item, "base_url", "api_base", "api_base_url", "endpoint", "invoke_url", "url")
    base_url = re.sub(r"/chat/completions/?$", "", raw_base.strip())
    provider = _dump_str(item, "provider", "source", "api_provider", default="").lower()
    if provider in {"hf", "hugging-face", "featherless", "featherless-ai"}:
        provider = "huggingface"
    if not provider:
        hay = f"{base_url} {model_name}".lower()
        if "nvidia.com" in hay or "moonshotai/kimi" in hay:
            provider = "nvidia"
        elif "huggingface" in hay or ":featherless-ai" in hay or ":fireworks-ai" in hay or ":together" in hay:
            provider = "huggingface"
        elif "openrouter" in hay:
            provider = "openrouter"
        elif "groq" in hay:
            provider = "groq"
        else:
            provider = "custom"
    if not base_url:
        base_url = _PROVIDER_DEFAULTS.get(provider, "")
    mid = _dump_str(item, "id", "model_id", "slug", default="")
    mid = _model_slug(mid or name or model_name)
    if mid in BUILTIN_MODELS:
        mid = f"{mid}-custom"
    api_key = _dump_str(item, "api_key", "key", "token", "authorization", default="")
    if api_key.startswith("$") or "os.environ" in api_key or api_key.lower().startswith("bearer $"):
        api_key = ""
    tags = item.get("tags") or item.get("tags_json") or []
    if isinstance(tags, str):
        try:
            tags = json.loads(tags)
        except Exception:
            tags = [x.strip() for x in tags.split(",") if x.strip()]
    if not isinstance(tags, list):
        tags = []
    tags = [str(t).strip() for t in tags if str(t).strip()]
    for tag in ("imported", provider):
        if tag and tag not in tags:
            tags.append(tag)
    return {
        "id": mid,
        "name": name[:100] or mid,
        "provider": provider[:40] or "custom",
        "base_url": base_url[:500],
        "model_name": model_name[:200],
        "api_key": api_key,
        "is_active": _dump_bool(item, ("is_active", "enabled", "active", "access_enabled"), True),
        "is_default": _dump_bool(item, ("is_default", "default"), False),
        "is_fast": _dump_bool(item, ("is_fast", "fast"), False),
        "is_vision": _dump_bool(item, ("is_vision", "vision", "multimodal"), False),
        "is_code": _dump_bool(item, ("is_code", "code", "coder"), False),
        "context_length": _dump_int(item, ("context_length", "ctx", "context", "max_context", "max_context_tokens"), 32768, 512, 2_000_000),
        "max_output_tokens": _dump_int(item, ("max_output_tokens", "max_tokens", "output_tokens"), 4096, 128, 128_000),
        "temperature_default": max(0.0, min(2.0, float(item.get("temperature_default", item.get("temperature", 0.7)) or 0.7))),
        "description": _dump_str(item, "description", "desc", "notes", default="Imported from admin model dump.")[:1000],
        "tags_json": json.dumps(tags[:20]),
    }

async def _upsert_imported_model(spec: Dict[str, Any], admin_uid: str, dry_run: bool = False) -> Dict[str, Any]:
    now = _utcnow()
    existing = await db_fetchone(
        "SELECT id FROM ai_models WHERE id=? OR name=? OR model_name=? LIMIT 1",
        (spec["id"], spec["name"], spec["model_name"]))
    mid = existing["id"] if existing else spec["id"]
    action = "updated" if existing else "created"
    if dry_run:
        return {"id": mid, "name": spec["name"], "provider": spec["provider"], "model_name": spec["model_name"], "action": action}
    if spec["is_default"]:
        await db_execute("UPDATE ai_models SET is_default=0,updated_at=?", (now,))
    if spec["is_fast"]:
        await db_execute("UPDATE ai_models SET is_fast=0,updated_at=?", (now,))
    enc = _encrypt({"key": spec.get("api_key", "")})
    if existing:
        sets = [
            "name=?", "provider=?", "base_url=?", "model_name=?", "is_active=?", "is_default=?",
            "is_fast=?", "is_vision=?", "is_code=?", "context_length=?", "max_output_tokens=?",
            "temperature_default=?", "description=?", "tags_json=?", "updated_at=?",
        ]
        vals: List[Any] = [
            spec["name"], spec["provider"], spec["base_url"], spec["model_name"], int(spec["is_active"]),
            int(spec["is_default"]), int(spec["is_fast"]), int(spec["is_vision"]), int(spec["is_code"]),
            spec["context_length"], spec["max_output_tokens"], spec["temperature_default"],
            spec["description"], spec["tags_json"], now,
        ]
        if spec.get("api_key"):
            sets.insert(4, "encrypted_api_key=?")
            vals.insert(4, enc)
        vals.append(mid)
        await db_execute(f"UPDATE ai_models SET {','.join(sets)} WHERE id=?", tuple(vals))
    else:
        await db_execute(
            "INSERT INTO ai_models(id,name,provider,base_url,model_name,encrypted_api_key,"
            "is_active,is_default,is_fast,is_vision,is_code,context_length,max_output_tokens,"
            "temperature_default,description,tags_json,created_by,created_at,updated_at)"
            " VALUES(?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)",
            (mid, spec["name"], spec["provider"], spec["base_url"], spec["model_name"], enc,
             int(spec["is_active"]), int(spec["is_default"]), int(spec["is_fast"]), int(spec["is_vision"]),
             int(spec["is_code"]), spec["context_length"], spec["max_output_tokens"],
             spec["temperature_default"], spec["description"], spec["tags_json"], admin_uid, now, now))
    await db_execute(
        "INSERT INTO model_access(id,model_id,display_name,source,is_enabled,created_at,updated_at)"
        " VALUES(?,?,?,?,?,?,?)"
        " ON CONFLICT(model_id) DO UPDATE SET display_name=excluded.display_name,"
        "source='db',is_enabled=excluded.is_enabled,updated_at=excluded.updated_at",
        (_new_id(), mid, spec["name"], "db", int(spec["is_active"]), now, now))
    return {"id": mid, "name": spec["name"], "provider": spec["provider"], "model_name": spec["model_name"], "action": action}

@app.get("/admin/models")
async def admin_list_models(user: Dict = Depends(_require_admin)):
    models = await _model_access_rows(include_inactive_db=True)
    return {"models":models,"total":len(models)}

@app.post("/admin/models")
async def admin_create_model(req: AIModelCreate, admin: Dict = Depends(_require_admin)):
    mid = _new_id(); now = _utcnow()
    url = req.base_url or _PROVIDER_DEFAULTS.get(req.provider,"")
    enc = _encrypt({"key":req.api_key})
    admin_uid = admin.get("id") or admin.get("sub","")
    if req.is_default:
        await db_execute("UPDATE ai_models SET is_default=0,updated_at=?",(now,))
    if req.is_fast:
        await db_execute("UPDATE ai_models SET is_fast=0,updated_at=?",(now,))
    await db_execute(
        "INSERT INTO ai_models(id,name,provider,base_url,model_name,encrypted_api_key,"
        "is_active,is_default,is_fast,is_vision,is_code,context_length,max_output_tokens,"
        "temperature_default,description,tags_json,created_by,created_at,updated_at)"
        " VALUES(?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)",
        (mid, req.name, req.provider, url, req.model_name, enc,
         int(req.is_active), int(req.is_default), int(req.is_fast),
         int(req.is_vision), int(req.is_code), req.context_length, req.max_output_tokens,
         req.temperature_default, req.description, "[]", admin_uid, now, now))
    await db_execute(
        "INSERT INTO model_access(id,model_id,display_name,source,is_enabled,created_at,updated_at)"
        " VALUES(?,?,?,?,?,?,?)"
        " ON CONFLICT(model_id) DO UPDATE SET display_name=excluded.display_name,"
        "source='db',is_enabled=excluded.is_enabled,updated_at=excluded.updated_at",
        (_new_id(), mid, req.name, "db", int(req.is_active), now, now))
    return {"id":mid,"name":req.name}

@app.post("/admin/models/seed-defaults")
async def admin_seed_default_models(admin: Dict = Depends(_require_admin)):
    admin_uid = admin.get("id") or admin.get("sub","")
    before = await db_count("SELECT COUNT(*) as c FROM ai_models")
    await _seed_default_ai_models(admin_uid)
    after = await db_count("SELECT COUNT(*) as c FROM ai_models")
    return {"ok": True, "created": max(0, after - before), "total": after}

@app.post("/admin/models/import-dump")
async def admin_import_model_dump(req: AIModelDumpImport, admin: Dict = Depends(_require_admin)):
    admin_uid = admin.get("id") or admin.get("sub","")
    try:
        parsed = _parse_model_dump(req.raw)
    except ValueError as exc:
        raise HTTPException(400, f"Model dump parse failed: {exc}")
    items = _model_dump_items(parsed)
    if not items:
        raise HTTPException(400, "No model objects found in dump")
    if len(items) > 200:
        raise HTTPException(400, "Import up to 200 models at a time")
    imported: List[Dict[str, Any]] = []
    errors: List[Dict[str, Any]] = []
    for idx, item in enumerate(items):
        try:
            spec = _normalize_imported_model(item)
            imported.append(await _upsert_imported_model(spec, admin_uid, req.dry_run))
        except Exception as exc:
            errors.append({"index": idx, "error": str(exc)[:300]})
    if imported and not req.dry_run:
        await db_execute(
            "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
            (_new_id(), admin_uid, "models", "model_dump_import",
             json.dumps({"imported": len(imported), "errors": len(errors)}), _utcnow()))
    return {
        "ok": True,
        "dry_run": bool(req.dry_run),
        "imported": len(imported),
        "errors": errors,
        "models": imported,
    }

@app.put("/admin/models/{mid}")
async def admin_update_model(mid: str, req: AIModelUpdate, admin: Dict = Depends(_require_admin)):
    row = await db_fetchone("SELECT id FROM ai_models WHERE id=?",(mid,))
    if not row: raise HTTPException(404)
    now = _utcnow(); sets: List[str] = []; vals: List[Any] = []
    if req.name       is not None: sets.append("name=?"); vals.append(req.name)
    if req.provider   is not None: sets.append("provider=?"); vals.append(req.provider)
    if req.base_url   is not None: sets.append("base_url=?"); vals.append(req.base_url)
    if req.model_name is not None: sets.append("model_name=?"); vals.append(req.model_name)
    if req.api_key    is not None: sets.append("encrypted_api_key=?"); vals.append(_encrypt({"key":req.api_key}))
    if req.is_active  is not None: sets.append("is_active=?"); vals.append(int(req.is_active))
    if req.is_default is not None:
        if req.is_default: await db_execute("UPDATE ai_models SET is_default=0,updated_at=? WHERE id!=?",(now,mid))
        sets.append("is_default=?"); vals.append(int(req.is_default))
    if req.is_fast    is not None:
        if req.is_fast: await db_execute("UPDATE ai_models SET is_fast=0,updated_at=? WHERE id!=?",(now,mid))
        sets.append("is_fast=?"); vals.append(int(req.is_fast))
    if req.is_vision  is not None: sets.append("is_vision=?"); vals.append(int(req.is_vision))
    if req.is_code    is not None: sets.append("is_code=?"); vals.append(int(req.is_code))
    if req.context_length is not None: sets.append("context_length=?"); vals.append(int(req.context_length))
    if req.max_output_tokens is not None: sets.append("max_output_tokens=?"); vals.append(int(req.max_output_tokens))
    if req.temperature_default is not None: sets.append("temperature_default=?"); vals.append(float(req.temperature_default))
    if req.description is not None: sets.append("description=?"); vals.append(req.description)
    if not sets: return {"ok":True}
    sets.append("updated_at=?"); vals.append(now); vals.append(mid)
    await db_execute(f"UPDATE ai_models SET {','.join(sets)} WHERE id=?",tuple(vals))
    if req.name is not None or req.is_active is not None:
        row = await db_fetchone("SELECT id,name,is_active FROM ai_models WHERE id=?", (mid,))
        if row:
            await db_execute(
                "INSERT INTO model_access(id,model_id,display_name,source,is_enabled,created_at,updated_at)"
                " VALUES(?,?,?,?,?,?,?)"
                " ON CONFLICT(model_id) DO UPDATE SET display_name=excluded.display_name,"
                "source='db',is_enabled=excluded.is_enabled,updated_at=excluded.updated_at",
                (_new_id(), mid, row["name"], "db", int(row["is_active"]), now, now))
    return {"ok":True}

@app.patch("/admin/models/{mid}/access")
async def admin_update_model_access(mid: str, body: dict, admin: Dict = Depends(_require_admin)):
    model_id = _canonical_model_id(mid)
    is_enabled = bool(body.get("is_enabled"))
    rows = await _model_access_rows(include_inactive_db=True)
    match = next((r for r in rows if _canonical_model_id(r["id"]) == model_id), None)
    if not match:
        raise HTTPException(404, "Model not found")
    now = _utcnow()
    await db_execute(
        "INSERT INTO model_access(id,model_id,display_name,source,is_enabled,created_at,updated_at)"
        " VALUES(?,?,?,?,?,?,?)"
        " ON CONFLICT(model_id) DO UPDATE SET display_name=excluded.display_name,"
        "source=excluded.source,is_enabled=excluded.is_enabled,updated_at=excluded.updated_at",
        (_new_id(), model_id, match.get("name") or model_id,
         "builtin" if match.get("is_builtin") else "db", int(is_enabled), now, now))
    if not match.get("is_builtin"):
        await db_execute("UPDATE ai_models SET is_active=?,updated_at=? WHERE id=?",
                         (int(is_enabled), now, model_id))
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), admin.get("id") or admin.get("sub",""), model_id, "model_access_toggle",
         json.dumps({"is_enabled": is_enabled}), now))
    return {"ok": True, "model_id": model_id, "is_enabled": is_enabled}

@app.post("/admin/models/access/bulk")
async def admin_bulk_model_access(body: dict, admin: Dict = Depends(_require_admin)):
    is_enabled = bool(body.get("is_enabled"))
    now = _utcnow()
    rows = await _model_access_rows(include_inactive_db=True)
    for row in rows:
        mid = _canonical_model_id(row["id"])
        await db_execute(
            "INSERT INTO model_access(id,model_id,display_name,source,is_enabled,created_at,updated_at)"
            " VALUES(?,?,?,?,?,?,?)"
            " ON CONFLICT(model_id) DO UPDATE SET is_enabled=excluded.is_enabled,updated_at=excluded.updated_at",
            (_new_id(), mid, row.get("name") or mid,
             "builtin" if row.get("is_builtin") else "db", int(is_enabled), now, now))
        if not row.get("is_builtin"):
            await db_execute("UPDATE ai_models SET is_active=?,updated_at=? WHERE id=?",
                             (int(is_enabled), now, mid))
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), admin.get("id") or admin.get("sub",""), "models", "model_access_bulk_toggle",
         json.dumps({"is_enabled": is_enabled}), now))
    return {"ok": True, "updated": len(rows), "is_enabled": is_enabled}

@app.delete("/admin/models/{mid}")
async def admin_delete_model(mid: str, admin: Dict = Depends(_require_admin)):
    await db_execute("DELETE FROM ai_models WHERE id=?",(mid,))
    await db_execute("DELETE FROM model_access WHERE model_id=?", (mid,))
    await db_execute("DELETE FROM user_model_access WHERE model_id=?", (mid,))
    return {"ok":True}

@app.post("/admin/mcp-servers")
async def admin_create_mcp_server(req: MCPServerCreate, admin: Dict = Depends(_require_admin)):
    sid = _new_id(); now = _utcnow()
    enc = _encrypt(req.creds) if req.creds else ""
    admin_uid = admin.get("id") or admin.get("sub","")
    try:
        await db_execute(
            "INSERT INTO mcp_servers(id,name,description,server_url,auth_type,encrypted_creds,"
            "is_active,is_public,created_by,created_at,updated_at) VALUES(?,?,?,?,?,?,1,?,?,?,?)",
            (sid, req.name, req.description, req.server_url, req.auth_type, enc,
             int(req.is_public), admin_uid, now, now))
    except Exception:
        raise HTTPException(400, f"MCP server named '{req.name}' already exists")
    # Auto-discover tools
    try:
        tools = await _mcp_list_tools(sid)
        if tools:
            await db_execute("UPDATE mcp_servers SET tools_json=?,ping_ok=1,last_pinged_at=? WHERE id=?",
                             (json.dumps(tools), now, sid))
    except Exception: pass
    return {"id":sid,"name":req.name}

@app.get("/admin/mcp-servers")
async def admin_list_mcp_servers(admin: Dict = Depends(_require_admin)):
    servers = await db_fetchall(
        "SELECT id,name,description,server_url,auth_type,is_active,is_public,ping_ok,"
        "last_pinged_at,tools_json,created_at,updated_at FROM mcp_servers ORDER BY created_at DESC")
    for s in servers:
        try:
            s["tool_count"] = len(json.loads(s.get("tools_json") or "[]"))
        except Exception:
            s["tool_count"] = 0
    return {"servers": servers, "total": len(servers)}

@app.delete("/admin/mcp-servers/{sid}")
async def admin_delete_mcp_server(sid: str, admin: Dict = Depends(_require_admin)):
    await db_execute("DELETE FROM mcp_servers WHERE id=?",(sid,))
    return {"ok":True}

@app.post("/admin/mcp-servers/{sid}/ping")
async def admin_ping_mcp_server(sid: str, admin: Dict = Depends(_require_admin)):
    row = await db_fetchone("SELECT * FROM mcp_servers WHERE id=?",(sid,))
    if not row: raise HTTPException(404)
    now = _utcnow()
    tools = await _mcp_list_tools(sid)
    ok = len(tools) >= 0
    await db_execute("UPDATE mcp_servers SET ping_ok=?,last_pinged_at=?,tools_json=? WHERE id=?",
                     (int(ok), now, json.dumps(tools), sid))
    return {"success":ok,"tools":tools,"tool_count":len(tools)}

@app.post("/admin/mcp-servers/{sid}/tools/{tool_name}/test")
async def admin_test_mcp_tool(sid: str, tool_name: str, body: dict, admin: Dict = Depends(_require_admin)):
    args = body.get("tool_args", body.get("args", {})) or {}
    admin_uid = admin.get("id") or admin.get("sub","")
    try:
        result = await _mcp_call_tool(sid, tool_name, args)
        await db_execute(
            "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
            (_new_id(), admin_uid, sid, "mcp_tool_test",
             json.dumps({"tool":tool_name,"success":True}), _utcnow()))
        return {"success": True, "server_id": sid, "tool": tool_name, "result": result}
    except Exception as e:
        await db_execute(
            "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
            (_new_id(), admin_uid, sid, "mcp_tool_test",
             json.dumps({"tool":tool_name,"success":False,"error":str(e)[:500]}), _utcnow()))
        raise HTTPException(502, f"MCP tool call failed: {e}")

@app.post("/admin/chat/multi-model")
async def admin_multi_model_chat(req: MultiModelChatReq, admin: Dict = Depends(_require_admin)):
    return await multi_model_chat(req, admin)

@app.get("/admin/env")
async def admin_get_env(admin: Dict = Depends(_require_admin)):
    path = _env_file_path()
    pairs = _env_read_pairs()
    variables = [
        {
            "key": key,
            "value": value,
            "is_secret": bool(_ENV_SECRET_RE.search(key)),
            "is_runtime": key in _RUNTIME_ENV_KEYS,
        }
        for key, value in sorted(pairs.items())
    ]
    return {"path": str(path), "exists": path.exists(), "variables": variables, "total": len(variables)}

@app.post("/admin/env")
async def admin_create_env(req: EnvVarCreate, admin: Dict = Depends(_require_admin)):
    key = _validate_env_key(req.key)
    path = _env_write_value(key, req.value)
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), admin.get("id") or admin.get("sub",""), key, "env_set",
         json.dumps({"path": str(path), "runtime": key in _RUNTIME_ENV_KEYS}), _utcnow()))
    return {"ok": True, "key": key, "path": str(path), "runtime_updated": key in _RUNTIME_ENV_KEYS}

@app.post("/admin/env/reload")
async def admin_reload_env(admin: Dict = Depends(_require_admin)):
    _reload_runtime_env()
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), admin.get("id") or admin.get("sub",""), "env", "env_reload",
         json.dumps({"path": str(_env_file_path())}), _utcnow()))
    return {"ok": True, "path": str(_env_file_path())}

@app.put("/admin/env/{key}")
async def admin_update_env(key: str, req: EnvVarUpdate, admin: Dict = Depends(_require_admin)):
    key = _validate_env_key(key)
    path = _env_write_value(key, req.value)
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), admin.get("id") or admin.get("sub",""), key, "env_set",
         json.dumps({"path": str(path), "runtime": key in _RUNTIME_ENV_KEYS}), _utcnow()))
    return {"ok": True, "key": key, "path": str(path), "runtime_updated": key in _RUNTIME_ENV_KEYS}

@app.delete("/admin/env/{key}")
async def admin_delete_env(key: str, admin: Dict = Depends(_require_admin)):
    key = _validate_env_key(key)
    path = _env_delete_value(key)
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), admin.get("id") or admin.get("sub",""), key, "env_delete",
         json.dumps({"path": str(path), "runtime": key in _RUNTIME_ENV_KEYS}), _utcnow()))
    return {"ok": True, "key": key, "path": str(path), "runtime_updated": key in _RUNTIME_ENV_KEYS}

@app.get("/admin/system")
async def admin_system(admin: Dict = Depends(_require_admin)):
    search_keys = ["BRAVE_SEARCH_API_KEY", "TAVILY_API_KEY", "SERPAPI_API_KEY"]
    model_keys = ["GROQ_API_KEY", "HF_TOKEN"]
    oauth_keys = sorted({k for _, _, _, _, _, keys in _PLATFORM_CONNECTOR_CATALOG for k in keys})
    return {
        "frontend_path": str(FRONTEND_PATH),
        "db_path": DB_PATH,
        "env_path": str(_env_file_path()),
        "env_exists": _env_file_path().exists(),
        "app_base_url": APP_BASE_URL,
        "oauth_redirect_uris": _oauth_redirect_uris(),
        "google_oauth_redirect_uri": _oauth_redirect_uri("gmail"),
        "web_search": {
            "enabled": any(os.getenv(k, "").strip() for k in search_keys) or True,
            "providers": {k: bool(os.getenv(k, "").strip()) for k in search_keys},
            "fallback": "DuckDuckGo HTML fallback",
        },
        "models": {k: bool(os.getenv(k, "").strip()) for k in model_keys},
        "oauth": {k: bool(os.getenv(k, "").strip()) for k in oauth_keys},
        "services": {
            "chromadb": HAS_CHROMA and _chroma_client is not None,
            "scheduler": _scheduler is not None and _scheduler.running,
            "livekit": _lk_available(),
            "livy": _livy_reachable_quick(),
            "livy_configured": _livy_configured(),
        },
    }

@app.get("/admin/model-usage")
async def admin_model_usage(days: int = 14, limit: int = 30,
                            admin: Dict = Depends(_require_admin)):
    model_rates = {
        "llama-3.3-70b-versatile": (0.59, 0.79),
        "llama-3.1-8b-instant": (0.05, 0.08),
    }
    def _estimate_cost(mid: str, in_tokens: int, out_tokens: int) -> float:
        rates = model_rates.get(_canonical_model_id(mid), (0.0, 0.0))
        return round((in_tokens / 1_000_000 * rates[0]) + (out_tokens / 1_000_000 * rates[1]), 6)

    days = max(1, min(90, int(days or 14)))
    limit = max(1, min(100, int(limit or 30)))
    cutoff = (datetime.now(timezone.utc)-timedelta(days=days)).isoformat()
    rows = await db_fetchall(
        "SELECT COALESCE(NULLIF(model_used,''),'unknown') as model_id,"
        " COUNT(*) as responses,"
        " SUM(CASE WHEN tokens_input>0 THEN tokens_input ELSE 0 END) as input_tokens,"
        " SUM(CASE WHEN tokens_output>0 THEN tokens_output ELSE CAST((LENGTH(content)+3)/4 AS INTEGER) END) as output_tokens,"
        " AVG(CASE WHEN latency_ms>0 THEN latency_ms END) as avg_latency_ms,"
        " MAX(created_at) as last_used_at,"
        " SUM(CASE WHEN content LIKE 'I could not get a response%' OR content LIKE '%model API key%' OR content LIKE '%fallback failed%' THEN 1 ELSE 0 END) as error_count,"
        " SUM(CASE WHEN mode='tool' THEN 1 ELSE 0 END) as tool_responses,"
        " SUM(CASE WHEN mode='agent' THEN 1 ELSE 0 END) as agent_responses,"
        " SUM(CASE WHEN mode='thinking' THEN 1 ELSE 0 END) as thinking_responses"
        " FROM chat_history WHERE role='assistant' AND is_hidden=0 AND created_at>=?"
        " GROUP BY COALESCE(NULLIF(model_used,''),'unknown')"
        " ORDER BY responses DESC LIMIT ?",
        (cutoff, limit))
    for r in rows:
        r["model_label"] = await _model_display_name(r["model_id"])
        r["input_tokens"] = int(r.get("input_tokens") or 0)
        r["output_tokens"] = int(r.get("output_tokens") or 0)
        r["total_tokens"] = r["input_tokens"] + r["output_tokens"]
        r["avg_latency_ms"] = int(r.get("avg_latency_ms") or 0)
        r["error_count"] = int(r.get("error_count") or 0)
        r["cost_estimate_usd"] = _estimate_cost(r["model_id"], r["input_tokens"], r["output_tokens"])
    total_row = await db_fetchone(
        "SELECT COUNT(*) as responses,"
        " SUM(CASE WHEN tokens_input>0 THEN tokens_input ELSE 0 END) as input_tokens,"
        " SUM(CASE WHEN tokens_output>0 THEN tokens_output ELSE CAST((LENGTH(content)+3)/4 AS INTEGER) END) as output_tokens,"
        " AVG(CASE WHEN latency_ms>0 THEN latency_ms END) as avg_latency_ms,"
        " SUM(CASE WHEN content LIKE 'I could not get a response%' OR content LIKE '%model API key%' OR content LIKE '%fallback failed%' THEN 1 ELSE 0 END) as error_count"
        " FROM chat_history WHERE role='assistant' AND is_hidden=0 AND created_at>=?",
        (cutoff,))
    daily = await db_fetchall(
        "SELECT DATE(created_at) as date, COUNT(*) as responses,"
        " SUM(CASE WHEN tokens_output>0 THEN tokens_output ELSE CAST((LENGTH(content)+3)/4 AS INTEGER) END) as output_tokens"
        " FROM chat_history WHERE role='assistant' AND is_hidden=0 AND created_at>=?"
        " GROUP BY DATE(created_at) ORDER BY date",
        (cutoff,))
    totals = {
        "responses": int((total_row or {}).get("responses") or 0),
        "input_tokens": int((total_row or {}).get("input_tokens") or 0),
        "output_tokens": int((total_row or {}).get("output_tokens") or 0),
        "avg_latency_ms": int((total_row or {}).get("avg_latency_ms") or 0),
        "error_count": int((total_row or {}).get("error_count") or 0),
    }
    totals["total_tokens"] = totals["input_tokens"] + totals["output_tokens"]
    totals["cost_estimate_usd"] = round(sum(float(r.get("cost_estimate_usd") or 0) for r in rows), 6)
    return {"days": days, "models": rows, "daily": daily, "totals": totals}

@app.get("/admin/error-monitor")
async def admin_error_monitor(limit: int = 100, admin: Dict = Depends(_require_admin)):
    limit = max(1, min(300, int(limit or 100)))
    log_path = Path("jazz.log")
    lines: List[str] = []
    if log_path.exists():
        try:
            lines = log_path.read_text(encoding="utf-8", errors="replace").splitlines()[-3000:]
        except Exception:
            lines = []
    error_re = re.compile(r"\b(ERROR|CRITICAL|Traceback|Exception|failed|failure)\b", re.IGNORECASE)
    warning_re = re.compile(r"\b(WARNING|WARN)\b", re.IGNORECASE)
    entries: List[Dict[str, Any]] = []
    for idx, line in enumerate(lines, 1):
        if not error_re.search(line) and not warning_re.search(line):
            continue
        level = "error" if error_re.search(line) else "warning"
        ts = ""
        m = re.match(r"^(\d{4}-\d{2}-\d{2}[ T]\d{2}:\d{2}:\d{2}(?:,\d+)?)", line)
        if m:
            ts = m.group(1)
        entries.append({"level": level, "timestamp": ts, "line_no": idx, "message": line[:1200]})
    entries = entries[-limit:][::-1]
    connector_failures = await db_fetchall(
        "SELECT connector_type,action,error_msg,created_at FROM connector_logs"
        " WHERE status='failure' ORDER BY created_at DESC LIMIT 20")
    failed_jobs = await db_fetchall(
        "SELECT job_id,status,error_message,started_at,finished_at FROM agent_job_logs"
        " WHERE status='failed' ORDER BY started_at DESC LIMIT 20")
    return {
        "log_path": str(log_path),
        "errors": entries,
        "counts": {
            "errors": sum(1 for e in entries if e["level"] == "error"),
            "warnings": sum(1 for e in entries if e["level"] == "warning"),
            "connector_failures": len(connector_failures),
            "failed_jobs": len(failed_jobs),
        },
        "connector_failures": connector_failures,
        "failed_jobs": failed_jobs,
    }

@app.get("/admin/platform-controls")
async def admin_platform_controls(admin: Dict = Depends(_require_admin)):
    status = await _platform_status()
    db_path = Path(DB_PATH)
    backups = sorted(BACKUPS_DIR.glob("jazz_v14_backup_*.db"), key=lambda p: p.stat().st_mtime, reverse=True)
    return {
        **status,
        "database": {
            "path": str(db_path),
            "size_bytes": db_path.stat().st_size if db_path.exists() else 0,
            "latest_backup": str(backups[0]) if backups else "",
            "latest_backup_at": datetime.fromtimestamp(backups[0].stat().st_mtime, timezone.utc).isoformat() if backups else "",
        },
    }

@app.patch("/admin/platform-controls")
async def admin_update_platform_controls(body: dict, admin: Dict = Depends(_require_admin)):
    admin_uid = admin.get("id") or admin.get("sub", "")
    changed: Dict[str, Any] = {}
    maintenance = body.get("maintenance")
    if isinstance(maintenance, dict):
        if "enabled" in maintenance:
            await _setting_set("maintenance_enabled", bool(maintenance.get("enabled")), admin_uid)
            changed["maintenance_enabled"] = bool(maintenance.get("enabled"))
        if "message" in maintenance:
            msg = str(maintenance.get("message") or "")[:500]
            await _setting_set("maintenance_message", msg, admin_uid)
            changed["maintenance_message"] = msg
    banner = body.get("banner")
    if isinstance(banner, dict):
        if "enabled" in banner:
            await _setting_set("broadcast_banner_enabled", bool(banner.get("enabled")), admin_uid)
            changed["broadcast_banner_enabled"] = bool(banner.get("enabled"))
        for key, setting_key, max_len in (
            ("title", "broadcast_banner_title", 120),
            ("message", "broadcast_banner_message", 500),
            ("type", "broadcast_banner_type", 20),
        ):
            if key in banner:
                value = str(banner.get(key) or "")[:max_len]
                if key == "type" and value not in ("info", "warn", "err", "ok"):
                    value = "info"
                await _setting_set(setting_key, value, admin_uid)
                changed[setting_key] = value
    if changed:
        await db_execute(
            "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
            (_new_id(), admin_uid, "platform_controls", "platform_controls_update",
             json.dumps(changed), _utcnow()))
    return {"ok": True, "changed": changed, **(await _platform_status())}

@app.get("/admin/user-activity")
async def admin_user_activity(limit: int = 100, admin: Dict = Depends(_require_admin)):
    limit = max(1, min(500, int(limit or 100)))
    rows = await db_fetchall(
        "SELECT u.id,u.email,u.full_name,u.role,u.subscription,u.is_active,u.last_login_at,u.created_at,"
        " COALESCE(cs.chats,0) as chats,"
        " COALESCE(ch.messages,0) as messages,"
        " ch.last_message_at as last_message_at"
        " FROM users u"
        " LEFT JOIN (SELECT user_id,COUNT(*) as chats FROM chat_sessions GROUP BY user_id) cs ON cs.user_id=u.id"
        " LEFT JOIN (SELECT user_id,COUNT(*) as messages,MAX(created_at) as last_message_at FROM chat_history GROUP BY user_id) ch ON ch.user_id=u.id"
        " ORDER BY COALESCE(ch.last_message_at,u.last_login_at,u.created_at) DESC LIMIT ?",
        (limit,))
    model_rows = await db_fetchall(
        "SELECT user_id,GROUP_CONCAT(DISTINCT model_used) as models"
        " FROM chat_history WHERE role='assistant' AND COALESCE(model_used,'')!=''"
        " GROUP BY user_id")
    conn_rows = await db_fetchall(
        "SELECT user_id,GROUP_CONCAT(DISTINCT connector_type) as connectors FROM ("
        " SELECT user_id,connector_type FROM connectors WHERE is_active=1"
        " UNION SELECT user_id,connector_type FROM smart_connectors WHERE status='active'"
        ") GROUP BY user_id")
    models_by_user = {r["user_id"]: [x for x in (r.get("models") or "").split(",") if x] for r in model_rows}
    conns_by_user = {r["user_id"]: [x for x in (r.get("connectors") or "").split(",") if x] for r in conn_rows}
    for row in rows:
        row["models_used"] = models_by_user.get(row["id"], [])[:8]
        row["connectors_used"] = conns_by_user.get(row["id"], [])[:10]
        row["last_active_at"] = row.get("last_message_at") or row.get("last_login_at") or row.get("created_at")
    return {"users": rows, "total": len(rows)}

@app.get("/admin/connector-health")
async def admin_connector_health(admin: Dict = Depends(_require_admin)):
    platform_rows = await db_fetchall(
        "SELECT connector_type,display_name,category,is_enabled,setup_status,env_keys,updated_at"
        " FROM platform_connectors ORDER BY category,display_name")
    smart_rows = await db_fetchall(
        "SELECT connector_type,user_id,status,last_tested_at,last_test_ok,test_error,updated_at,encrypted_conn_data"
        " FROM smart_connectors")
    legacy_rows = await db_fetchall(
        "SELECT connector_type,user_id,is_active,last_tested_at,last_test_ok,test_error,updated_at,encrypted_creds"
        " FROM connectors")
    failures = await db_fetchall(
        "SELECT connector_type,action,error_msg,created_at FROM connector_logs"
        " WHERE status='failure' ORDER BY created_at DESC LIMIT 200")
    last_failure: Dict[str, Dict[str, Any]] = {}
    for f in failures:
        last_failure.setdefault(f["connector_type"], f)

    def _expired_count(rows: List[Dict], encrypted_key: str) -> int:
        expired = 0
        now_dt = datetime.now(timezone.utc)
        for row in rows:
            try:
                data = _decrypt(row.get(encrypted_key) or "")
            except Exception:
                continue
            expires_at = data.get("expires_at")
            expires_in = int(data.get("expires_in") or 0)
            exp_dt = None
            if expires_at:
                try:
                    exp_dt = datetime.fromisoformat(str(expires_at).replace("Z", "+00:00"))
                except Exception:
                    exp_dt = None
            elif expires_in and row.get("updated_at"):
                try:
                    base = datetime.fromisoformat(str(row["updated_at"]).replace("Z", "+00:00"))
                    if base.tzinfo is None:
                        base = base.replace(tzinfo=timezone.utc)
                    exp_dt = base + timedelta(seconds=expires_in)
                except Exception:
                    exp_dt = None
            if exp_dt:
                if exp_dt.tzinfo is None:
                    exp_dt = exp_dt.replace(tzinfo=timezone.utc)
                if exp_dt < now_dt:
                    expired += 1
        return expired

    by_type: Dict[str, Dict[str, Any]] = {}
    for row in smart_rows:
        item = by_type.setdefault(row["connector_type"], {"connected_users": set(), "last_tested_at": "", "test_failures": 0, "token_expired": 0})
        if row.get("status") == "active":
            item["connected_users"].add(row["user_id"])
        if row.get("last_tested_at") and row["last_tested_at"] > (item.get("last_tested_at") or ""):
            item["last_tested_at"] = row["last_tested_at"]
        if row.get("last_test_ok") == 0:
            item["test_failures"] += 1
    for row in legacy_rows:
        item = by_type.setdefault(row["connector_type"], {"connected_users": set(), "last_tested_at": "", "test_failures": 0, "token_expired": 0})
        if row.get("is_active"):
            item["connected_users"].add(row["user_id"])
        if row.get("last_tested_at") and row["last_tested_at"] > (item.get("last_tested_at") or ""):
            item["last_tested_at"] = row["last_tested_at"]
        if row.get("last_test_ok") == 0:
            item["test_failures"] += 1
    for ctype in set([r["connector_type"] for r in smart_rows] + [r["connector_type"] for r in legacy_rows]):
        smart_for_type = [r for r in smart_rows if r["connector_type"] == ctype]
        legacy_for_type = [r for r in legacy_rows if r["connector_type"] == ctype]
        item = by_type.setdefault(ctype, {"connected_users": set(), "last_tested_at": "", "test_failures": 0, "token_expired": 0})
        item["token_expired"] = _expired_count(smart_for_type, "encrypted_conn_data") + _expired_count(legacy_for_type, "encrypted_creds")

    out = []
    seen = set()
    for p in platform_rows:
        ctype = p["connector_type"]; seen.add(ctype)
        h = by_type.get(ctype, {"connected_users": set(), "last_tested_at": "", "test_failures": 0, "token_expired": 0})
        lf = last_failure.get(ctype, {})
        out.append({
            **dict(p),
            "connected_users": len(h["connected_users"]),
            "last_tested_at": h.get("last_tested_at") or "",
            "test_failures": int(h.get("test_failures") or 0),
            "token_expired": int(h.get("token_expired") or 0),
            "last_error": lf.get("error_msg", ""),
            "last_error_at": lf.get("created_at", ""),
        })
    for ctype, h in by_type.items():
        if ctype in seen:
            continue
        out.append({
            "connector_type": ctype, "display_name": ctype.replace("_", " ").title(),
            "category": "Custom", "is_enabled": 1, "setup_status": "ready", "env_keys": "[]",
            "connected_users": len(h["connected_users"]), "last_tested_at": h.get("last_tested_at") or "",
            "test_failures": int(h.get("test_failures") or 0), "token_expired": int(h.get("token_expired") or 0),
            "last_error": (last_failure.get(ctype) or {}).get("error_msg", ""),
            "last_error_at": (last_failure.get(ctype) or {}).get("created_at", ""),
        })
    counts = {
        "connectors": len(out),
        "enabled": sum(1 for r in out if r.get("is_enabled")),
        "connected_users": sum(int(r.get("connected_users") or 0) for r in out),
        "token_expired": sum(int(r.get("token_expired") or 0) for r in out),
        "test_failures": sum(int(r.get("test_failures") or 0) for r in out),
    }
    return {"connectors": out, "counts": counts}

@app.get("/admin/role-permissions")
async def admin_role_permissions(admin: Dict = Depends(_require_admin)):
    features = await _sidebar_feature_rows()
    models = await _model_access_rows(include_inactive_db=True)
    limits_by_plan = await _subscription_limits_all()
    feature_access = await db_fetchall("SELECT plan,feature_key,is_enabled FROM plan_sidebar_feature_access")
    model_access = await db_fetchall("SELECT plan,model_id,is_enabled FROM plan_model_access")
    f_map = {(r["plan"], r["feature_key"]): bool(r["is_enabled"]) for r in feature_access}
    m_map = {(r["plan"], r["model_id"]): bool(r["is_enabled"]) for r in model_access}
    plans = []
    for plan in PLAN_TIERS:
        plans.append({
            "plan": plan,
            "limits": limits_by_plan.get(plan, {}),
            "features": [{**dict(f), "plan_enabled": f_map.get((plan, f["feature_key"]), True)} for f in features],
            "models": [{**dict(m), "plan_enabled": m_map.get((plan, m["id"]), True)} for m in models],
        })
    return {"plans": plans}

@app.patch("/admin/role-permissions/{plan}")
async def admin_update_role_permissions(plan: str, body: dict, admin: Dict = Depends(_require_admin)):
    if plan not in PLAN_TIERS:
        raise HTTPException(404, "Plan not found")
    admin_uid = admin.get("id") or admin.get("sub", "")
    now = _utcnow()
    changed: Dict[str, Any] = {"features": {}, "models": {}}
    features = body.get("features", {})
    if isinstance(features, dict):
        known_feature_keys = await _sidebar_feature_key_set()
        for key, value in features.items():
            key = str(key).strip()
            if key not in known_feature_keys:
                continue
            if value is None:
                await db_execute("DELETE FROM plan_sidebar_feature_access WHERE plan=? AND feature_key=?", (plan, key))
                changed["features"][key] = None
            else:
                await db_execute(
                    "INSERT INTO plan_sidebar_feature_access(id,plan,feature_key,is_enabled,updated_by,updated_at)"
                    " VALUES(?,?,?,?,?,?) ON CONFLICT(plan,feature_key) DO UPDATE SET "
                    "is_enabled=excluded.is_enabled,updated_by=excluded.updated_by,updated_at=excluded.updated_at",
                    (_new_id(), plan, key, 1 if bool(value) else 0, admin_uid, now))
                changed["features"][key] = bool(value)
    models = body.get("models", {})
    if isinstance(models, dict):
        known = {_canonical_model_id(r["id"]) for r in await _model_access_rows(include_inactive_db=True)}
        for key, value in models.items():
            mid = _canonical_model_id(str(key).strip())
            if mid not in known:
                continue
            if value is None:
                await db_execute("DELETE FROM plan_model_access WHERE plan=? AND model_id=?", (plan, mid))
                changed["models"][mid] = None
            else:
                await db_execute(
                    "INSERT INTO plan_model_access(id,plan,model_id,is_enabled,updated_by,updated_at)"
                    " VALUES(?,?,?,?,?,?) ON CONFLICT(plan,model_id) DO UPDATE SET "
                    "is_enabled=excluded.is_enabled,updated_by=excluded.updated_by,updated_at=excluded.updated_at",
                    (_new_id(), plan, mid, 1 if bool(value) else 0, admin_uid, now))
                changed["models"][mid] = bool(value)
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), admin_uid, plan, "role_permissions_update", json.dumps(changed), now))
    return {"ok": True, "plan": plan, "changed": changed}

@app.get("/admin/db-backup/download")
async def admin_download_db_backup(admin: Dict = Depends(_require_admin)):
    src = Path(DB_PATH)
    if not src.exists():
        raise HTTPException(404, "Database file not found")
    stamp = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
    dest = BACKUPS_DIR / f"jazz_v14_backup_{stamp}.db"
    shutil.copy2(src, dest)
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), admin.get("id") or admin.get("sub",""), "database", "db_backup",
         json.dumps({"path": str(dest), "size_bytes": dest.stat().st_size}), _utcnow()))
    return FileResponse(str(dest), media_type="application/octet-stream", filename=dest.name)

def _admin_shell(cmd: str, timeout: int = 8) -> Dict[str, Any]:
    try:
        if os.name == "nt":
            proc = subprocess.run(cmd, shell=True, capture_output=True, text=True, timeout=timeout)
        else:
            proc = subprocess.run(["bash", "-lc", cmd], capture_output=True, text=True, timeout=timeout)
        return {
            "ok": proc.returncode == 0,
            "code": proc.returncode,
            "stdout": (proc.stdout or "").strip()[-5000:],
            "stderr": (proc.stderr or "").strip()[-2000:],
        }
    except subprocess.TimeoutExpired as exc:
        return {"ok": False, "code": 124, "stdout": (exc.stdout or "")[-2000:] if exc.stdout else "", "stderr": "Command timed out"}
    except Exception as exc:
        return {"ok": False, "code": -1, "stdout": "", "stderr": str(exc)}

def _system_jobs_status_sync() -> Dict[str, Any]:
    docker_ok = bool(shutil.which("docker"))
    ngrok_ok = bool(shutil.which("ngrok"))
    n8n_rows: List[Dict[str, str]] = []
    if docker_ok:
        n8n = _admin_shell("docker ps -a --format '{{.Names}}|{{.Status}}|{{.Ports}}' | grep -E '^(jazzai-n8n|n8n)\\|' || true", 6)
        for line in (n8n.get("stdout") or "").splitlines():
            parts = line.split("|", 2)
            if len(parts) == 3:
                n8n_rows.append({"name": parts[0], "status": parts[1], "ports": parts[2]})
    ngrok = _admin_shell("pgrep -af ngrok || true", 5) if os.name != "nt" else _admin_shell("tasklist | findstr ngrok", 5)
    server_dir = str(Path(__file__).resolve().parent)
    return {
        "backend": {
            "pid": os.getpid(),
            "port": os.getenv("PORT", "8000"),
            "cwd": os.getcwd(),
            "server_dir": server_dir,
            "uptime_seconds": int((datetime.now(timezone.utc) - _SERVER_START).total_seconds()),
        },
        "scheduler": {
            "running": bool(_scheduler and _scheduler.running),
            "jobs": len(_scheduler.get_jobs()) if _scheduler and _scheduler.running else 0,
        },
        "docker": {"available": docker_ok},
        "n8n": {
            "available": docker_ok,
            "containers": n8n_rows,
            "url": "http://45.79.124.28:5678/",
            "running": any("Up " in r.get("status", "") for r in n8n_rows),
        },
        "ngrok": {
            "available": ngrok_ok,
            "running": bool((ngrok.get("stdout") or "").strip()),
            "processes": (ngrok.get("stdout") or "").splitlines()[:8],
        },
    }

def _restart_backend_later(port: str) -> None:
    if os.name == "nt":
        return
    server_dir = shlex.quote(str(Path(__file__).resolve().parent))
    py = "/root/jazzai/venv/bin/python" if Path("/root/jazzai/venv/bin/python").exists() else sys.executable
    log_path = "server.out"
    cmd = (
        f"cd {server_dir}; "
        f"(sleep 1; kill -TERM {os.getpid()} 2>/dev/null; sleep 1; "
        f"nohup env PORT={shlex.quote(str(port))} {shlex.quote(py)} server14.py > {shlex.quote(log_path)} 2>&1 &) "
        f"> /dev/null 2>&1 &"
    )
    subprocess.Popen(["bash", "-lc", cmd], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)

async def _system_job_action(action: str, body: dict, admin_uid: str,
                             background: BackgroundTasks) -> Dict[str, Any]:
    action = (action or "").strip()
    result: Dict[str, Any]
    if action == "reload_env":
        _reload_runtime_env()
        result = {"ok": True, "stdout": "Runtime environment reloaded.", "stderr": ""}
    elif action == "backend_restart":
        port = str(body.get("port") or os.getenv("PORT", "8001"))
        background.add_task(_restart_backend_later, port)
        result = {"ok": True, "stdout": f"Backend restart scheduled on port {port}.", "stderr": ""}
    elif action == "n8n_start":
        cmd = (
            "docker start jazzai-n8n 2>/dev/null || docker start n8n 2>/dev/null || "
            "docker run -d --name jazzai-n8n -p 5678:5678 "
            "-e N8N_HOST=45.79.124.28 -e N8N_PORT=5678 -e N8N_PROTOCOL=http "
            "-e N8N_SECURE_COOKIE=false "
            "-e N8N_EDITOR_BASE_URL=http://45.79.124.28:5678/ "
            "-e WEBHOOK_URL=http://45.79.124.28:5678/ "
            "-v n8n_data:/home/node/.n8n n8nio/n8n:latest"
        )
        result = await asyncio.get_running_loop().run_in_executor(_executor, lambda: _admin_shell(cmd, 40))
    elif action == "n8n_stop":
        result = await asyncio.get_running_loop().run_in_executor(
            _executor, lambda: _admin_shell("docker stop jazzai-n8n 2>/dev/null || docker stop n8n 2>/dev/null", 15))
    elif action == "ngrok_start_8000":
        domain = str(body.get("domain") or await _setting_get("ngrok_domain", "imperceptibly-hymnlike-leesa.ngrok-free.dev")).strip()
        domain_arg = f" --domain={shlex.quote(domain)}" if domain else ""
        result = await asyncio.get_running_loop().run_in_executor(
            _executor, lambda: _admin_shell(f"pgrep -af 'ngrok http' >/dev/null || nohup ngrok http 8000{domain_arg} > ngrok.log 2>&1 &", 8))
    elif action == "ngrok_stop":
        result = await asyncio.get_running_loop().run_in_executor(_executor, lambda: _admin_shell("pkill ngrok || true", 8))
    else:
        raise HTTPException(400, "Unknown system job action")
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), admin_uid, "system_jobs", "system_job_action",
         json.dumps({"action": action, "ok": bool(result.get("ok")), "code": result.get("code")}), _utcnow()))
    return {"action": action, "result": result, "status": await asyncio.get_running_loop().run_in_executor(_executor, _system_jobs_status_sync)}

@app.get("/admin/system-jobs")
async def admin_system_jobs(admin: Dict = Depends(_require_admin)):
    status = await asyncio.get_running_loop().run_in_executor(_executor, _system_jobs_status_sync)
    return {"status": status, "actions": ["reload_env", "backend_restart", "n8n_start", "n8n_stop", "ngrok_start_8000", "ngrok_stop"]}

@app.post("/admin/system-jobs/{action}")
async def admin_run_system_job(action: str, background: BackgroundTasks,
                               body: dict = None,
                               admin: Dict = Depends(_require_admin)):
    return await _system_job_action(action, body or {}, admin.get("id") or admin.get("sub",""), background)

@app.get("/admin/rate-limits")
async def admin_rate_limits(admin: Dict = Depends(_require_admin)):
    limits = await _subscription_limits_all()
    return {
        "resources": [{"key": k, "label": label, "description": desc} for k, label, desc in RATE_LIMIT_RESOURCES],
        "plans": [{"plan": plan, "limits": limits.get(plan, {})} for plan in PLAN_TIERS],
    }

@app.patch("/admin/rate-limits")
async def admin_update_rate_limits(body: dict, admin: Dict = Depends(_require_admin)):
    current = await _subscription_limits_all()
    incoming = body.get("limits", {})
    if not isinstance(incoming, dict):
        raise HTTPException(400, "limits object required")
    valid_keys = {k for k, _, _ in RATE_LIMIT_RESOURCES}
    changed: Dict[str, Dict[str, int]] = {}
    for plan, values in incoming.items():
        if plan not in PLAN_TIERS or not isinstance(values, dict):
            continue
        changed[plan] = {}
        for key, value in values.items():
            if key not in valid_keys:
                continue
            try:
                n = max(-1, min(1_000_000, int(value)))
            except Exception:
                raise HTTPException(400, f"{plan}.{key} must be an integer")
            current.setdefault(plan, {})[key] = n
            changed[plan][key] = n
    admin_uid = admin.get("id") or admin.get("sub", "")
    await _setting_set("subscription_limits", current, admin_uid)
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), admin_uid, "rate_limits", "rate_limits_update", json.dumps(changed), _utcnow()))
    return {"ok": True, "changed": changed, "plans": [{"plan": p, "limits": current.get(p, {})} for p in PLAN_TIERS]}

@app.get("/admin/security")
async def admin_security(admin: Dict = Depends(_require_admin)):
    now = _utcnow()
    rows = await db_fetchall(
        "SELECT u.id,u.email,u.full_name,u.role,u.subscription,u.is_active,u.is_verified,u.last_login_at,u.created_at,"
        " COALESCE(rt.active_tokens,0) as active_tokens,"
        " COALESCE(ak.active_api_keys,0) as active_api_keys"
        " FROM users u"
        " LEFT JOIN (SELECT user_id,COUNT(*) as active_tokens FROM refresh_tokens WHERE revoked=0 AND expires_at>? GROUP BY user_id) rt ON rt.user_id=u.id"
        " LEFT JOIN (SELECT user_id,COUNT(*) as active_api_keys FROM api_keys WHERE is_active=1 GROUP BY user_id) ak ON ak.user_id=u.id"
        " ORDER BY COALESCE(u.last_login_at,u.created_at) DESC",
        (now,))
    recent = await db_fetchall(
        "SELECT actor_id,target_id,action,detail_json,created_at FROM audit_log"
        " WHERE action LIKE 'security_%' OR action IN('admin_update_user','impersonate')"
        " ORDER BY created_at DESC LIMIT 40")
    return {"users": rows, "recent": recent, "totals": {
        "users": len(rows),
        "active_tokens": sum(int(r.get("active_tokens") or 0) for r in rows),
        "active_api_keys": sum(int(r.get("active_api_keys") or 0) for r in rows),
        "disabled_users": sum(1 for r in rows if not r.get("is_active")),
    }}

@app.post("/admin/security/users/{uid}/{action}")
async def admin_security_user_action(uid: str, action: str, body: dict = None,
                                     admin: Dict = Depends(_require_admin)):
    body = body or {}
    admin_uid = admin.get("id") or admin.get("sub", "")
    if uid == admin_uid and action in ("disable", "delete"):
        raise HTTPException(400, "Cannot disable or delete yourself")
    row = await db_fetchone("SELECT id,email FROM users WHERE id=?", (uid,))
    if not row:
        raise HTTPException(404, "User not found")
    now = _utcnow()
    detail: Dict[str, Any] = {}
    if action == "revoke_tokens":
        await db_execute("UPDATE refresh_tokens SET revoked=1 WHERE user_id=? AND revoked=0", (uid,))
    elif action == "reset_password":
        pw = str(body.get("password") or "").strip()
        _validated_password(pw)
        new_hash = await _hash_pw_async(pw)
        await db_execute("UPDATE users SET password_hash=?,updated_at=? WHERE id=?", (new_hash, now, uid))
        await db_execute("UPDATE refresh_tokens SET revoked=1 WHERE user_id=? AND revoked=0", (uid,))
        detail["password_changed"] = True
    elif action == "disable":
        await db_execute("UPDATE users SET is_active=0,updated_at=? WHERE id=?", (now, uid))
        await db_execute("UPDATE refresh_tokens SET revoked=1 WHERE user_id=? AND revoked=0", (uid,))
    elif action == "enable":
        await db_execute("UPDATE users SET is_active=1,updated_at=? WHERE id=?", (now, uid))
    elif action == "verify":
        await db_execute("UPDATE users SET is_verified=1,updated_at=? WHERE id=?", (now, uid))
    elif action == "unverify":
        await db_execute("UPDATE users SET is_verified=0,updated_at=? WHERE id=?", (now, uid))
    else:
        raise HTTPException(400, "Unknown security action")
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), admin_uid, uid, f"security_{action}", json.dumps(detail), now))
    return {"ok": True, "user_id": uid, "action": action}

@app.get("/admin/prompt-policies")
async def admin_get_prompt_policies(admin: Dict = Depends(_require_admin)):
    return {"policies": await _prompt_policies(), "note": "Stored for admin governance. Chat runtime is intentionally not modified here."}

@app.patch("/admin/prompt-policies")
async def admin_update_prompt_policies(body: dict, admin: Dict = Depends(_require_admin)):
    allowed = _prompt_policy_defaults()
    current = await _prompt_policies()
    changed: Dict[str, Any] = {}
    for key in allowed:
        if key == "chat_runtime_applied":
            continue
        if key in body:
            value = str(body.get(key) or "")[:8000]
            current[key] = value
            changed[key] = value
    current["chat_runtime_applied"] = False
    admin_uid = admin.get("id") or admin.get("sub", "")
    await _setting_set("prompt_policies", current, admin_uid)
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), admin_uid, "prompt_policies", "prompt_policies_update",
         json.dumps({"changed": list(changed.keys())}), _utcnow()))
    return {"ok": True, "policies": current}

@app.get("/admin/export/{kind}")
async def admin_export_csv(kind: str, admin: Dict = Depends(_require_admin)):
    kind = (kind or "").strip().lower()
    specs = {
        "users": (
            "SELECT id,email,full_name,role,subscription,is_active,is_verified,last_login_at,created_at FROM users ORDER BY created_at DESC",
            "jazz_users.csv",
        ),
        "websites": (
            "SELECT w.id,u.email,w.title,w.style,w.html_size_bytes,w.view_count,w.created_at,w.updated_at FROM websites w LEFT JOIN users u ON u.id=w.user_id ORDER BY w.created_at DESC",
            "jazz_websites.csv",
        ),
        "code-runs": (
            "SELECT cr.id,u.email,cr.language,cr.exit_code,cr.duration_ms,cr.created_at FROM code_run_logs cr LEFT JOIN users u ON u.id=cr.user_id ORDER BY cr.created_at DESC LIMIT 5000",
            "jazz_code_runs.csv",
        ),
        "connectors": (
            "SELECT connector_type,user_id,is_active,last_tested_at,last_test_ok,created_at FROM connectors ORDER BY created_at DESC",
            "jazz_connectors.csv",
        ),
    }
    if kind not in specs:
        raise HTTPException(404, "Unknown export kind")
    sql, filename = specs[kind]
    rows = await db_fetchall(sql)
    buf = io.StringIO()
    if rows:
        writer = csv.DictWriter(buf, fieldnames=list(rows[0].keys()))
        writer.writeheader()
        writer.writerows([dict(r) for r in rows])
    else:
        buf.write("")
    data = buf.getvalue().encode("utf-8")
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,detail_json,created_at) VALUES(?,?,?,?,?,?)",
        (_new_id(), admin.get("id") or admin.get("sub",""), kind, "admin_export_csv",
         json.dumps({"rows": len(rows), "filename": filename}), _utcnow()))
    return StreamingResponse(io.BytesIO(data), media_type="text/csv",
                             headers={"Content-Disposition": f'attachment; filename="{filename}"'})

@app.post("/admin/web-search/test")
async def admin_test_web_search(body: dict, admin: Dict = Depends(_require_admin)):
    q = body.get("query") or "OpenAI news"
    results = await _web_search(q, body.get("count", 3))
    return {"query": q, "results": results}

@app.get("/admin/metrics")
async def admin_metrics(days: int = 14, admin: Dict = Depends(_require_admin)):
    cutoff = (datetime.now(timezone.utc)-timedelta(days=days)).isoformat()
    msgs = await db_fetchall(
        "SELECT DATE(created_at) as date,COUNT(*) as count FROM chat_history"
        " WHERE created_at>=? GROUP BY DATE(created_at) ORDER BY date", (cutoff,))
    users = await db_fetchall(
        "SELECT DATE(created_at) as date,COUNT(*) as count FROM users"
        " WHERE created_at>=? GROUP BY DATE(created_at) ORDER BY date", (cutoff,))
    return {"messages":msgs,"new_users":users,"days":days}

@app.get("/logs")
async def get_logs(admin: Dict = Depends(_require_admin)):
    try: return {"logs":Path("jazz.log").read_text(encoding="utf-8").splitlines()[-500:]}
    except FileNotFoundError: return {"logs":[]}

# ══════════════════════════════════════════════════════════════════════════════
# §39  HEALTH & INFO
# ══════════════════════════════════════════════════════════════════════════════

@app.get("/health")
async def health():
    db_ok = False
    try: await db_fetchone("SELECT 1"); db_ok = True
    except Exception: pass
    uptime = int((datetime.now(timezone.utc)-_SERVER_START).total_seconds())
    return {
        "status": "healthy" if db_ok else "degraded",
        "version": "14.0",
        "uptime_s": uptime,
        "services": {
            "sqlite":    db_ok,
            "chromadb":  HAS_CHROMA and _chroma_client is not None,
            "scheduler": _scheduler is not None and _scheduler.running,
            "livekit":   _lk_available(),
            "livy":      _livy_reachable_quick(),
            "livy_configured": _livy_configured(),
        },
        "features": [
            "streaming_sse","multi_model","agent_loop","cron_jobs","branching",
            "vision","rag","hybrid_search","reranking","rolling_summaries",
            "tiered_memory","auto_memory_extract","auto_thinking","website_builder",
            "code_execution","shell_access","file_ops","latex_compile","latex_pdf","voice_stt_tts","livekit_voice",
            "full_oauth_pkce","25_connectors","mcp_servers","livy_spark",
            "slash_commands","api_keys","gdpr","audit_log","impersonation",
            "coupons","notifications","broadcast","multi_model_compare",
        ],
        "oauth_connectors": list(_oauth_providers().keys()),
    }

@app.get("/models")
async def list_models(user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    allowed = set(await _allowed_model_ids_for_user(uid))
    builtin = [{
        "id":k, "label":v["label"], "provider":v.get("provider","groq"),
        "ctx":v["ctx"], "fast":v["fast"], "censored":v.get("censored", True),
        "protected":v.get("protected", False), "source":"builtin",
    } for k,v in BUILTIN_MODELS.items() if k in allowed]
    custom = await db_fetchall(
        "SELECT id,name as label,provider,model_name,is_fast,is_vision,is_code,context_length as ctx,description,tags_json"
        " FROM ai_models WHERE is_active=1 ORDER BY is_default DESC,name")
    custom = [m for m in custom if m["id"] in allowed and not _model_is_internal(m)]
    for m in custom:
        tags = [str(t).lower() for t in _safe_json_loads(m.get("tags_json"), [])]
        haystack = " ".join(str(m.get(k, "")) for k in ("id","label","model_name","description")).lower()
        m["source"] = "db"
        m["tags"] = tags
        m["censored"] = not ("uncensored" in tags or "dolphin" in haystack)
        m["protected"] = "protected" in tags
    return {"models":builtin,"custom_models":custom}

@app.get("/model-info")
async def model_info(user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub","")
    allowed = set(await _allowed_model_ids_for_user(uid))
    builtin_models = [{
        "id":k, "label":v["label"], "ctx":v["ctx"], "fast":v["fast"],
        "provider":v.get("provider","groq"), "censored":v.get("censored", True),
        "protected":v.get("protected", False),
    } for k,v in BUILTIN_MODELS.items() if k in allowed]
    db_models = await db_fetchall(
        "SELECT id,name,provider,base_url,model_name,is_active,is_default,is_fast,is_vision,is_code,"
        "context_length,max_output_tokens,temperature_default,description,tags_json"
        " FROM ai_models WHERE is_active=1 ORDER BY is_default DESC,name")
    db_models = [m for m in db_models if m["id"] in allowed and not _model_is_internal(m)]
    for m in db_models:
        tags = [str(t).lower() for t in _safe_json_loads(m.get("tags_json"), [])]
        haystack = " ".join(str(m.get(k, "")) for k in ("id","name","model_name","description")).lower()
        m["tags"] = tags
        m["censored"] = not ("uncensored" in tags or "dolphin" in haystack)
        m["protected"] = "protected" in tags
    return {"builtin_models":builtin_models,"db_models":db_models}

@app.get("/favicon.ico", include_in_schema=False)
async def favicon(): return JSONResponse(status_code=204, content=None)

@app.post("/upload/image")
async def upload_image(file: UploadFile = File(...), user: Dict = Depends(_get_current_user)):
    data = await file.read()
    if len(data) > MAX_IMAGE_SIZE: raise HTTPException(413, "Image too large (max 20MB)")
    ct = file.content_type or "image/jpeg"
    b64 = base64.b64encode(data).decode()
    return {"image_b64":b64,"mime":ct,"size":len(data),"url":f"data:{ct};base64,{b64}"}

@app.post("/upload/file")
async def upload_file_generic(file: UploadFile = File(...), user: Dict = Depends(_get_current_user)):
    data = await file.read()
    if len(data) > MAX_FILE_SIZE: raise HTTPException(413)
    suf = Path(file.filename or "file").suffix.lower()
    fname = f"{_new_id()}_{re.sub(r'[^a-zA-Z0-9._-]','_',file.filename or 'file')}"
    path = UPLOADS_DIR / fname; path.write_bytes(data)
    text, _ = await asyncio.get_running_loop().run_in_executor(
        _executor, _extract_text_sync, path, file.content_type or "")
    return {"filename":fname,"original":file.filename,"size":len(data),"text_preview":text[:12000],"text_chars":len(text or "")}

def _safe_generated_filename(name: str, ext: str) -> str:
    base = re.sub(r"[^A-Za-z0-9._ -]+", "_", name or "jazz-generated").strip(" ._-")
    base = re.sub(r"\s+", "_", base)[:80] or "jazz-generated"
    if not base.lower().endswith("." + ext):
        base = re.sub(r"\.[A-Za-z0-9]{1,8}$", "", base) + "." + ext
    return base

_GENERATED_EXT_ALIASES = {
    "txt": "txt", "text": "txt",
    "md": "md", "markdown": "md",
    "csv": "csv",
    "xlsx": "xlsx", "excel": "xlsx",
    "pdf": "pdf",
    "tex": "tex", "latex": "tex",
    "docx": "docx", "word": "docx",
    "pptx": "pptx", "ppt": "pptx", "powerpoint": "pptx", "slides": "pptx", "presentation": "pptx",
    "html": "html",
    "json": "json",
    "zip": "zip", "archive": "zip",
}

_GENERATED_MIME = {
    "txt": "text/plain",
    "md": "text/markdown",
    "csv": "text/csv",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pdf": "application/pdf",
    "tex": "application/x-tex",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "html": "text/html",
    "json": "application/json",
    "zip": "application/zip",
}

def _generated_ext(kind: str) -> str:
    return _GENERATED_EXT_ALIASES.get((kind or "txt").lower(), "txt")

def _infer_generated_filename(prompt: str, content: str, ext: str) -> str:
    text = f"{prompt}\n{content[:500]}"
    m = re.search(r"(?:called|named|filename|file name|as)\s+[`\"']?([A-Za-z0-9][A-Za-z0-9._ -]{0,80}?\.(?:pdf|tex|latex|xlsx|csv|txt|md|markdown|docx|pptx|html|json|zip))", text, re.I)
    if not m:
        m = re.search(
            r"(?:called|named|filename|file name|as)\s+[`\"']?([A-Za-z0-9][A-Za-z0-9._ -]{1,80}?)(?=$|[`\"',.!?:;]|\s+(?:with|for|from|about|containing|using|and)\b)",
            text,
            re.I,
        )
    if m:
        return _safe_generated_filename(m.group(1), ext)
    words = re.findall(r"[A-Za-z0-9]+", prompt)[:8]
    return _safe_generated_filename("_".join(words) or "jazz-generated", ext)

def _rows_from_generated_text(text: str) -> List[List[str]]:
    lines = [ln.strip() for ln in (text or "").splitlines() if ln.strip()]
    table_lines = [ln for ln in lines if ln.startswith("|") and ln.endswith("|") and ln.count("|") >= 2]
    if table_lines:
        rows: List[List[str]] = []
        for ln in table_lines:
            cells = [c.strip() for c in ln.strip("|").split("|")]
            if cells and all(re.fullmatch(r":?-{2,}:?", c or "") for c in cells):
                continue
            rows.append(cells)
        if rows:
            return rows
    csv_lines = [ln for ln in lines if "," in ln]
    if csv_lines:
        try:
            rows = list(csv.reader(io.StringIO("\n".join(csv_lines))))
            if rows:
                return [[str(c).strip() for c in row] for row in rows]
        except Exception:
            pass
    bullets = [re.sub(r"^[-*•\d.)\s]+", "", ln).strip() for ln in lines]
    bullets = [b for b in bullets if b]
    return [["Item"], *[[b] for b in bullets[:500]]] if bullets else [["Content"], [text[:32000]]]

def _clean_generated_content(content: str) -> str:
    text = (content or "").strip()
    m = re.fullmatch(r"```(?:[A-Za-z0-9_-]+)?\s*\n([\s\S]*?)\n```\s*", text)
    return m.group(1).strip() if m else text

_LATEX_HINT_RE = re.compile(
    r"\\(?:documentclass|begin\{document\}|usepackage|section\{|subsection\{|chapter\{|title\{|author\{|begin\{(?:equation|align|tabular|figure|enumerate|itemize))",
    re.IGNORECASE,
)

def _extract_latex_source(content: str) -> str:
    text = _clean_generated_content(content)
    fenced = re.search(r"```(?:latex|tex)\s*\n([\s\S]*?)\n```", text, re.IGNORECASE)
    if fenced:
        text = fenced.group(1).strip()
    text = text.replace("<\\/JAZZ_PASTED_CONTENT>", "</JAZZ_PASTED_CONTENT>")
    pasted = re.search(r"<JAZZ_PASTED_CONTENT\b[^>]*>\s*([\s\S]*?)\s*</JAZZ_PASTED_CONTENT>", text, re.IGNORECASE)
    if pasted and _LATEX_HINT_RE.search(pasted.group(1)):
        text = pasted.group(1).strip()
    starts = [i for i in (text.find("\\documentclass"), text.find("\\begin{document}")) if i >= 0]
    if starts:
        text = text[min(starts):].strip()
    return text

def _looks_like_latex(content: str) -> bool:
    return bool(_LATEX_HINT_RE.search(_extract_latex_source(content)))

def _latex_analysis(source: str) -> Dict[str, Any]:
    src = source or ""
    packages = re.findall(r"\\usepackage(?:\[[^\]]*\])?\{([^}]+)\}", src)
    pkg_list = []
    for p in packages:
        pkg_list.extend([x.strip() for x in p.split(",") if x.strip()])
    sections = re.findall(r"\\(?:chapter|section|subsection|subsubsection)\*?\{([^}]+)\}", src)
    title = ""
    m_title = re.search(r"\\title\{([^}]+)\}", src)
    if m_title:
        title = re.sub(r"\s+", " ", m_title.group(1)).strip()
    cls = ""
    m_cls = re.search(r"\\documentclass(?:\[[^\]]*\])?\{([^}]+)\}", src)
    if m_cls:
        cls = m_cls.group(1).strip()
    equations = len(re.findall(r"\\begin\{(?:equation|align|gather|multline)\}|\$\$|\\\[", src))
    return {
        "document_class": cls or "article",
        "title": title,
        "packages": sorted(set(pkg_list)),
        "sections": [re.sub(r"\s+", " ", s).strip() for s in sections[:40]],
        "equations": equations,
        "has_document": bool(re.search(r"\\begin\{document\}", src)),
        "chars": len(src),
        "lines": len(src.splitlines()),
    }

def _normalise_latex_document(source: str) -> str:
    src = _extract_latex_source(source).strip()
    if not src:
        src = "Empty LaTeX document."
    if re.search(r"\\begin\{document\}", src):
        return src
    return (
        "\\documentclass[11pt]{article}\n"
        "\\usepackage[utf8]{inputenc}\n"
        "\\usepackage{amsmath,amssymb,geometry,hyperref}\n"
        "\\geometry{margin=1in}\n"
        "\\begin{document}\n"
        f"{src}\n"
        "\\end{document}\n"
    )

def _latex_plain_preview(source: str) -> str:
    src = _normalise_latex_document(source)
    body = re.sub(r"[\s\S]*?\\begin\{document\}", "", src, count=1)
    body = re.sub(r"\\end\{document\}[\s\S]*", "", body, count=1)
    body = re.sub(r"\\(chapter|section|subsection|subsubsection)\*?\{([^}]*)\}", r"\n\2\n", body)
    body = re.sub(r"\\(title|author|date)\{([^}]*)\}", r"\2", body)
    body = re.sub(r"\\begin\{[^}]+\}|\\end\{[^}]+\}", "\n", body)
    body = re.sub(r"\\[a-zA-Z]+\*?(?:\[[^\]]*\])?(?:\{([^{}]*)\})?", lambda m: m.group(1) or "", body)
    body = body.replace("~", " ").replace("\\%", "%").replace("\\&", "&").replace("\\_", "_")
    body = re.sub(r"[{}]", "", body)
    return re.sub(r"\n{3,}", "\n\n", body).strip() or src[:8000]

def _latex_engine() -> Optional[str]:
    for exe in ("tectonic", "pdflatex", "xelatex"):
        found = shutil.which(exe)
        if found:
            return found
    return None

def _run_latex_engine(source: str, timeout: int = 35) -> Tuple[Optional[bytes], str, str]:
    engine = _latex_engine()
    if not engine:
        return None, "none", "No LaTeX engine found on this server. Used readable PDF fallback."
    with tempfile.TemporaryDirectory(prefix="jazz_latex_") as tmp:
        tmp_path = Path(tmp)
        tex_path = tmp_path / "main.tex"
        tex_path.write_text(_normalise_latex_document(source), encoding="utf-8")
        exe = Path(engine).name.lower()
        if exe == "tectonic":
            cmd = [engine, "--keep-logs", "--synctex=0", "--outdir", str(tmp_path), str(tex_path)]
        else:
            cmd = [engine, "-interaction=nonstopmode", "-halt-on-error", "-file-line-error", "-no-shell-escape", "main.tex"]
        env = os.environ.copy()
        env.update({"openin_any": "p", "openout_any": "p", "TEXMFOUTPUT": str(tmp_path)})
        try:
            proc = subprocess.run(cmd, cwd=str(tmp_path), env=env, capture_output=True, text=True, timeout=timeout)
        except subprocess.TimeoutExpired:
            return None, exe, f"LaTeX compile timed out after {timeout}s."
        log = ((proc.stdout or "") + "\n" + (proc.stderr or "")).strip()
        pdf = tmp_path / "main.pdf"
        if proc.returncode == 0 and pdf.exists():
            return pdf.read_bytes(), exe, log[-6000:]
        return None, exe, log[-6000:] or f"{exe} exited with code {proc.returncode}"

def _write_latex_pdf(path: Path, title: str, source: str) -> Dict[str, Any]:
    doc = _normalise_latex_document(source)
    analysis = _latex_analysis(doc)
    pdf_bytes, engine, log = _run_latex_engine(doc)
    if pdf_bytes:
        path.write_bytes(pdf_bytes)
        return {"compiled": True, "engine": engine, "compile_log": log, "analysis": analysis}
    fallback = (
        f"{title}\n\n"
        "LaTeX interpreted preview\n"
        f"Engine: {engine}\n"
        f"Compile note: {log}\n\n"
        f"Document class: {analysis.get('document_class')}\n"
        f"Packages: {', '.join(analysis.get('packages') or []) or 'none'}\n"
        f"Sections: {', '.join(analysis.get('sections') or []) or 'none'}\n"
        f"Equations: {analysis.get('equations', 0)}\n\n"
        f"{_latex_plain_preview(doc)}"
    )
    _write_generated_pdf(path, title, fallback[:60000])
    return {"compiled": False, "engine": engine, "compile_log": log, "analysis": analysis}

def _html_from_generated_text(title: str, content: str) -> str:
    clean = _clean_generated_content(content)
    if re.search(r"<!doctype\s+html|<html[\s>]", clean, re.I):
        return clean
    rows = _rows_from_generated_text(clean) if ("|" in clean or "," in clean) else []
    body_parts = [f"<h1>{html_lib.escape(title)}</h1>"]
    if len(rows) > 1 and len(rows[0]) > 1:
        head = rows[0]
        body_parts.append("<table><thead><tr>" + "".join(f"<th>{html_lib.escape(str(c))}</th>" for c in head) + "</tr></thead><tbody>")
        for row in rows[1:200]:
            body_parts.append("<tr>" + "".join(f"<td>{html_lib.escape(str(c))}</td>" for c in row) + "</tr>")
        body_parts.append("</tbody></table>")
    else:
        in_list = False
        for line in clean.splitlines():
            stripped = line.strip()
            if not stripped:
                if in_list:
                    body_parts.append("</ul>")
                    in_list = False
                continue
            heading = re.match(r"^(#{1,4})\s+(.+)$", stripped)
            bullet = re.match(r"^[-*]\s+(.+)$", stripped)
            if heading:
                if in_list:
                    body_parts.append("</ul>")
                    in_list = False
                level = min(4, len(heading.group(1)) + 1)
                body_parts.append(f"<h{level}>{html_lib.escape(heading.group(2))}</h{level}>")
            elif bullet:
                if not in_list:
                    body_parts.append("<ul>")
                    in_list = True
                body_parts.append(f"<li>{html_lib.escape(bullet.group(1))}</li>")
            else:
                if in_list:
                    body_parts.append("</ul>")
                    in_list = False
                body_parts.append(f"<p>{html_lib.escape(stripped)}</p>")
        if in_list:
            body_parts.append("</ul>")
    return """<!doctype html>
<html lang="en"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>{title}</title><style>
body{{margin:0;background:#07070f;color:#f4f4ff;font:16px/1.6 Inter,system-ui,Segoe UI,sans-serif;padding:48px;}}
main{{max-width:980px;margin:auto;}} h1{{font-size:36px;letter-spacing:0;margin:0 0 24px;color:#fff;}}
h2,h3,h4{{color:#c7c2ff;margin-top:28px}} p,li,td,th{{color:#deddf4}} table{{border-collapse:collapse;width:100%;margin:22px 0;background:#111124}}
td,th{{border:1px solid #302f55;padding:10px 12px;text-align:left}} th{{background:#1d1a3f;color:#fff}}
</style></head><body><main>{body}</main></body></html>""".format(
        title=html_lib.escape(title), body="\n".join(body_parts)
    )

def _json_from_generated_text(prompt: str, title: str, content: str) -> str:
    clean = _clean_generated_content(content)
    candidate = clean
    fenced = re.search(r"```json\s*([\s\S]*?)```", clean, re.I)
    if fenced:
        candidate = fenced.group(1).strip()
    if candidate.lstrip().startswith(("{", "[")):
        try:
            parsed = json.loads(candidate)
            return json.dumps(parsed, ensure_ascii=False, indent=2)
        except Exception:
            pass
    rows = _rows_from_generated_text(clean) if ("|" in clean or "," in clean) else []
    if len(rows) > 1 and len(rows[0]) > 1:
        headers = [str(x) for x in rows[0]]
        items = []
        for row in rows[1:500]:
            items.append({headers[i] if i < len(headers) else f"column_{i+1}": row[i] for i in range(len(row))})
        payload: Any = {"title": title, "items": items}
    else:
        payload = {
            "title": title,
            "prompt": prompt,
            "content": clean,
            "lines": [ln.strip() for ln in clean.splitlines() if ln.strip()][:500],
        }
    return json.dumps(payload, ensure_ascii=False, indent=2)

def _write_generated_docx(path: Path, title: str, content: str) -> None:
    clean = _clean_generated_content(content)
    if HAS_DOCX:
        doc = DocxDocument()
        doc.add_heading(title, 0)
        for line in clean.splitlines():
            stripped = line.strip()
            if not stripped:
                continue
            heading = re.match(r"^(#{1,4})\s+(.+)$", stripped)
            bullet = re.match(r"^[-*]\s+(.+)$", stripped)
            if heading:
                doc.add_heading(heading.group(2), min(4, len(heading.group(1))))
            elif bullet:
                doc.add_paragraph(bullet.group(1), style="List Bullet")
            else:
                doc.add_paragraph(stripped)
        doc.save(str(path))
        return
    paragraphs = [f"<w:p><w:r><w:t>{html_lib.escape(title)}</w:t></w:r></w:p>"]
    for line in clean.splitlines()[:800]:
        if line.strip():
            paragraphs.append(f"<w:p><w:r><w:t xml:space=\"preserve\">{html_lib.escape(line.strip())}</w:t></w:r></w:p>")
    document_xml = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"><w:body>{''.join(paragraphs)}<w:sectPr/></w:body></w:document>"""
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/><Default Extension="xml" ContentType="application/xml"/><Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/></Types>""")
        z.writestr("_rels/.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/></Relationships>""")
        z.writestr("word/document.xml", document_xml)

def _slide_sections_from_generated_text(title: str, content: str) -> List[Tuple[str, List[str]]]:
    lines = [ln.strip() for ln in _clean_generated_content(content).splitlines() if ln.strip()]
    sections: List[Tuple[str, List[str]]] = []
    current_title = title
    current_items: List[str] = []
    for line in lines:
        m = re.match(r"^(?:#{1,3}\s+|slide\s*\d+[:.)-]\s*)(.+)$", line, re.I)
        if m:
            if current_items:
                sections.append((current_title, current_items[:7]))
            current_title = m.group(1).strip()
            current_items = []
            continue
        item = re.sub(r"^[-*\d.)\s]+", "", line).strip()
        if item:
            current_items.append(item)
        if len(current_items) >= 7:
            sections.append((current_title, current_items[:7]))
            current_title = title
            current_items = []
    if current_items:
        sections.append((current_title, current_items[:7]))
    if not sections:
        sections = [(title, ["Generated by JAZZ"])]
    return sections[:12]

def _write_generated_pptx(path: Path, title: str, content: str) -> None:
    sections = _slide_sections_from_generated_text(title, content)
    if HAS_PPTX and Presentation is not None:
        prs = Presentation()
        for idx, (slide_title, items) in enumerate(sections):
            layout = prs.slide_layouts[1] if idx else prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_title
            if idx and len(slide.placeholders) > 1:
                body = slide.placeholders[1].text_frame
                body.clear()
                for item in items:
                    p = body.add_paragraph()
                    p.text = item
                    p.level = 0
        prs.save(str(path))
        return
    def slide_xml(slide_title: str, items: List[str]) -> str:
        bullets = "".join(
            f'<a:p><a:pPr lvl="0"><a:buChar char="&#8226;"/></a:pPr><a:r><a:rPr lang="en-US" sz="2200"><a:solidFill><a:srgbClr val="F4F4FF"/></a:solidFill></a:rPr><a:t>{html_lib.escape(item)}</a:t></a:r></a:p>'
            for item in items[:7]
        )
        return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cSld><p:bg><p:bgPr><a:solidFill><a:srgbClr val="07070F"/></a:solidFill><a:effectLst/></p:bgPr></p:bg><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr><p:sp><p:nvSpPr><p:cNvPr id="2" name="Title"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr><a:xfrm><a:off x="685800" y="457200"/><a:ext cx="10515600" cy="914400"/></a:xfrm></p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr lang="en-US" sz="4000" b="1"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:rPr><a:t>{html_lib.escape(slide_title)}</a:t></a:r></a:p></p:txBody></p:sp><p:sp><p:nvSpPr><p:cNvPr id="3" name="Content"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr><a:xfrm><a:off x="914400" y="1524000"/><a:ext cx="10058400" cy="4420000"/></a:xfrm></p:spPr><p:txBody><a:bodyPr/><a:lstStyle/>{bullets}</p:txBody></p:sp></p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:sld>'''
    content_types = ['<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/><Default Extension="xml" ContentType="application/xml"/><Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/><Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/><Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/><Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>']
    content_types += [f'<Override PartName="/ppt/slides/slide{i}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>' for i in range(1, len(sections)+1)]
    content_types.append("</Types>")
    presentation_rels = ['<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>']
    presentation_rels += [f'<Relationship Id="rId{i+1}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide{i}.xml"/>' for i in range(1, len(sections)+1)]
    presentation_rels.append("</Relationships>")
    sld_ids = "".join(f'<p:sldId id="{256+i}" r:id="rId{i+1}"/>' for i in range(1, len(sections)+1))
    presentation_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId1"/></p:sldMasterIdLst><p:sldIdLst>{sld_ids}</p:sldIdLst><p:sldSz cx="12192000" cy="6858000" type="screen16x9"/><p:notesSz cx="6858000" cy="9144000"/></p:presentation>'''
    blank_tree = '<p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr></p:spTree></p:cSld>'
    slide_layout = f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?><p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" type="blank" preserve="1">{blank_tree}<p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:sldLayout>'
    slide_master = f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?><p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">{blank_tree}<p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/><p:sldLayoutIdLst><p:sldLayoutId id="2147483649" r:id="rId1"/></p:sldLayoutIdLst></p:sldMaster>'
    theme = '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="JAZZ"><a:themeElements><a:clrScheme name="JAZZ"><a:dk1><a:srgbClr val="07070F"/></a:dk1><a:lt1><a:srgbClr val="FFFFFF"/></a:lt1><a:dk2><a:srgbClr val="111124"/></a:dk2><a:lt2><a:srgbClr val="F4F4FF"/></a:lt2><a:accent1><a:srgbClr val="7C6BFF"/></a:accent1><a:accent2><a:srgbClr val="22D3A6"/></a:accent2><a:accent3><a:srgbClr val="60A5FA"/></a:accent3><a:accent4><a:srgbClr val="F59E0B"/></a:accent4><a:accent5><a:srgbClr val="F472B6"/></a:accent5><a:accent6><a:srgbClr val="A78BFA"/></a:accent6><a:hlink><a:srgbClr val="60A5FA"/></a:hlink><a:folHlink><a:srgbClr val="A78BFA"/></a:folHlink></a:clrScheme><a:fontScheme name="JAZZ"><a:majorFont><a:latin typeface="Aptos Display"/></a:majorFont><a:minorFont><a:latin typeface="Aptos"/></a:minorFont></a:fontScheme><a:fmtScheme name="JAZZ"><a:fillStyleLst><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:fillStyleLst><a:lnStyleLst><a:ln w="6350"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln></a:lnStyleLst><a:effectStyleLst><a:effectStyle><a:effectLst/></a:effectStyle></a:effectStyleLst><a:bgFillStyleLst><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:bgFillStyleLst></a:fmtScheme></a:themeElements></a:theme>'
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", "".join(content_types))
        z.writestr("_rels/.rels", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/></Relationships>')
        z.writestr("ppt/presentation.xml", presentation_xml)
        z.writestr("ppt/_rels/presentation.xml.rels", "".join(presentation_rels))
        z.writestr("ppt/slideMasters/slideMaster1.xml", slide_master)
        z.writestr("ppt/slideMasters/_rels/slideMaster1.xml.rels", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/><Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/></Relationships>')
        z.writestr("ppt/slideLayouts/slideLayout1.xml", slide_layout)
        z.writestr("ppt/slideLayouts/_rels/slideLayout1.xml.rels", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="../slideMasters/slideMaster1.xml"/></Relationships>')
        z.writestr("ppt/theme/theme1.xml", theme)
        for i, (slide_title, items) in enumerate(sections, start=1):
            z.writestr(f"ppt/slides/slide{i}.xml", slide_xml(slide_title, items))
            z.writestr(f"ppt/slides/_rels/slide{i}.xml.rels", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/></Relationships>')

def _write_generated_zip(path: Path, title: str, prompt: str, content: str) -> None:
    clean = _clean_generated_content(content)
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("README.txt", clean)
        z.writestr("content.md", clean)
        z.writestr("content.html", _html_from_generated_text(title, clean))
        rows = _rows_from_generated_text(clean)
        csv_buf = io.StringIO()
        writer = csv.writer(csv_buf)
        writer.writerows(rows)
        z.writestr("data.csv", csv_buf.getvalue())
        z.writestr("metadata.json", json.dumps({"title": title, "prompt": prompt, "files": ["README.txt", "content.md", "content.html", "data.csv"]}, indent=2))

def _write_generated_pdf(path: Path, title: str, content: str) -> None:
    if HAS_FITZ:
        doc = fitz.open()
        margin = 54
        page = doc.new_page()
        rect = fitz.Rect(margin, margin, page.rect.width - margin, page.rect.height - margin)
        story = f"{title}\n\n{content}"
        # Simple manual pagination keeps this dependency-light and stable.
        chunks = [story[i:i+2400] for i in range(0, len(story), 2400)] or [""]
        doc.delete_page(0)
        for chunk in chunks:
            page = doc.new_page()
            page.insert_textbox(rect, chunk, fontsize=10, fontname="helv", color=(0, 0, 0), align=0)
        doc.save(str(path))
        doc.close()
        return
    escaped = content.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    lines = escaped.splitlines() or [escaped]
    rendered = []
    y = 760
    for line in lines[:42]:
        rendered.append(f"BT /F1 10 Tf 50 {y} Td ({line[:110]}) Tj ET")
        y -= 16
    stream = "\n".join(rendered).encode("latin-1", "replace")
    objects = [
        b"1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj\n",
        b"2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj\n",
        b"3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >> endobj\n",
        b"4 0 obj << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> endobj\n",
        b"5 0 obj << /Length " + str(len(stream)).encode() + b" >> stream\n" + stream + b"\nendstream endobj\n",
    ]
    out = io.BytesIO(); out.write(b"%PDF-1.4\n")
    offsets = []
    for obj in objects:
        offsets.append(out.tell()); out.write(obj)
    xref = out.tell()
    out.write(f"xref\n0 {len(objects)+1}\n0000000000 65535 f \n".encode())
    for off in offsets:
        out.write(f"{off:010d} 00000 n \n".encode())
    out.write(f"trailer << /Size {len(objects)+1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF".encode())
    path.write_bytes(out.getvalue())

def _build_generated_file(req: FileGenerateReq, user_id: str) -> Dict[str, Any]:
    ext = _generated_ext(req.file_type)
    mime = _GENERATED_MIME[ext]
    original = _infer_generated_filename(req.prompt, req.content, ext) if not req.filename else _safe_generated_filename(req.filename, ext)
    stored = f"{_new_id()}_{original}"
    path = UPLOADS_DIR / stored
    title = Path(original).stem.replace("_", " ").strip() or "JAZZ Generated File"
    clean_content = _clean_generated_content(req.content)
    extra: Dict[str, Any] = {}
    if ext == "txt":
        path.write_text(clean_content, encoding="utf-8")
    elif ext == "md":
        path.write_text(clean_content, encoding="utf-8")
    elif ext == "tex":
        clean_content = _normalise_latex_document(clean_content) if _looks_like_latex(clean_content) else clean_content
        path.write_text(clean_content, encoding="utf-8")
        extra = {"analysis": _latex_analysis(clean_content) if _looks_like_latex(clean_content) else {}}
    elif ext == "html":
        path.write_text(_html_from_generated_text(title, clean_content), encoding="utf-8")
    elif ext == "json":
        path.write_text(_json_from_generated_text(req.prompt, title, clean_content), encoding="utf-8")
    elif ext == "csv":
        rows = _rows_from_generated_text(clean_content)
        with path.open("w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerows(rows)
    elif ext == "xlsx":
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Generated"
        for row in _rows_from_generated_text(clean_content):
            ws.append(row)
        wb.save(str(path))
    elif ext == "pdf":
        if _looks_like_latex(clean_content) or ("latex" in (req.prompt or "").lower()):
            extra = _write_latex_pdf(path, title, clean_content)
        else:
            _write_generated_pdf(path, title, clean_content)
    elif ext == "docx":
        _write_generated_docx(path, title, clean_content)
    elif ext == "pptx":
        _write_generated_pptx(path, title, clean_content)
    elif ext == "zip":
        _write_generated_zip(path, title, req.prompt, clean_content)
    text_preview, _ = _extract_text_sync(path, mime)
    return {
        "filename": stored,
        "original": original,
        "size": path.stat().st_size,
        "type": mime,
        "text_preview": (text_preview or clean_content)[:12000],
        "text_chars": len(text_preview or clean_content),
        "generated": True,
        "download_url": f"/files/download/{urllib.parse.quote(stored)}",
        **extra,
    }

def _build_latex_artifact(req: LatexCompileReq, user_id: str) -> Dict[str, Any]:
    source = _normalise_latex_document(req.source)
    analysis = _latex_analysis(source)
    if req.output == "analysis":
        return {"ok": True, "analysis": analysis, "text_preview": _latex_plain_preview(source)[:12000]}
    ext = "tex" if req.output == "tex" else "pdf"
    mime = _GENERATED_MIME[ext]
    base_name = req.filename or analysis.get("title") or "jazz-latex"
    original = _safe_generated_filename(base_name, ext)
    stored = f"{_new_id()}_{original}"
    path = UPLOADS_DIR / stored
    extra: Dict[str, Any] = {"analysis": analysis}
    if ext == "tex":
        path.write_text(source, encoding="utf-8")
    else:
        extra.update(_write_latex_pdf(path, Path(original).stem.replace("_", " "), source))
    text_preview, _ = _extract_text_sync(path, mime)
    return {
        "filename": stored,
        "original": original,
        "size": path.stat().st_size,
        "type": mime,
        "text_preview": (text_preview or _latex_plain_preview(source))[:12000],
        "text_chars": len(text_preview or source),
        "generated": True,
        "download_url": f"/files/download/{urllib.parse.quote(stored)}",
        **extra,
    }

def _latex_compile_request_from_message(message: str) -> Optional[LatexCompileReq]:
    text = message or ""
    low = text.lower()
    if not _looks_like_latex(text):
        return None
    if not re.search(r"\b(convert|compile|compiler|render|export|download|save|generate|create|make)\b", low):
        return None
    if not re.search(r"\b(pdf|latex|tex|overleaf)\b", low):
        return None
    output = "tex" if re.search(r"\b(tex|latex source|source file)\b", low) and "pdf" not in low else "pdf"
    filename = None
    m = re.search(r"(?:called|named|filename|file name|as)\s+[`\"']?([A-Za-z0-9][A-Za-z0-9._ -]{0,80}?\.(?:pdf|tex))", text, re.I)
    if m:
        filename = m.group(1).strip()
    return LatexCompileReq(source=_extract_latex_source(text), filename=filename, output=output)

@app.post("/latex/compile")
async def compile_latex(req: LatexCompileReq, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub", "")
    result = await asyncio.get_running_loop().run_in_executor(_executor, _build_latex_artifact, req, uid)
    return result

@app.post("/files/generate")
async def generate_file(req: FileGenerateReq, user: Dict = Depends(_get_current_user)):
    uid = user.get("id") or user.get("sub", "")
    result = await asyncio.get_running_loop().run_in_executor(_executor, _build_generated_file, req, uid)
    return result

@app.get("/files/download/{stored_name}")
async def download_generated_or_uploaded_file(stored_name: str, user: Dict = Depends(_get_current_user)):
    safe = Path(stored_name).name
    if not safe or safe != stored_name or not re.fullmatch(r"[A-Za-z0-9._ -]+", safe):
        raise HTTPException(400, "Invalid file name")
    path = (UPLOADS_DIR / safe).resolve()
    if UPLOADS_DIR.resolve() not in path.parents or not path.exists() or not path.is_file():
        raise HTTPException(404, "File not found")
    download_name = re.sub(r"^[a-z0-9_\\-]{8,40}_", "", safe, flags=re.I) or safe
    return FileResponse(str(path), media_type="application/octet-stream", filename=download_name)

# ══════════════════════════════════════════════════════════════════════════════
# §39b  PLATFORM CONNECTOR MANAGEMENT (ADMIN)
# ══════════════════════════════════════════════════════════════════════════════

@app.get("/admin/platform-connectors")
async def admin_list_platform_connectors(admin: Dict = Depends(_require_admin)):
    """Return all connectors with enabled status, setup status, and env key presence."""
    rows = await db_fetchall(
        "SELECT id, connector_type, is_enabled, display_name, category, icon,"
        " setup_status, requires_oauth, env_keys, admin_notes, updated_at"
        " FROM platform_connectors ORDER BY category, display_name")
    # Refresh setup status live (env may have changed without restart)
    result = []
    for r in rows:
        env_keys = json.loads(r.get("env_keys") or "[]")
        status   = _calc_setup_status(env_keys)
        env_detail = {k: bool(os.getenv(k, "").strip()) for k in env_keys}
        result.append({**dict(r), "setup_status": status, "env_detail": env_detail})
    return {"connectors": result}

@app.patch("/admin/platform-connectors/{connector_type}")
async def admin_toggle_platform_connector(
    connector_type: str,
    body: dict,
    admin: Dict = Depends(_require_admin)
):
    """Enable/disable a connector platform-wide. Optionally update admin_notes."""
    row = await db_fetchone(
        "SELECT id, setup_status, env_keys FROM platform_connectors WHERE connector_type=?",
        (connector_type,))
    if not row:
        raise HTTPException(404, f"Connector '{connector_type}' not found")

    is_enabled = body.get("is_enabled")
    admin_notes = body.get("admin_notes")

    # Prevent enabling when env keys are missing
    if is_enabled:
        env_keys = json.loads(row.get("env_keys") or "[]")
        status   = _calc_setup_status(env_keys)
        if status == "not_configured" and env_keys:
            raise HTTPException(400,
                f"Cannot enable: required environment variables not set: {', '.join(env_keys)}")

    updates, vals = [], []
    if is_enabled is not None:
        updates.append("is_enabled=?"); vals.append(int(bool(is_enabled)))
    if admin_notes is not None:
        updates.append("admin_notes=?"); vals.append(admin_notes)
    if not updates:
        return {"ok": True}

    updates.append("updated_at=?"); vals.append(_utcnow()); vals.append(connector_type)
    await db_execute(
        f"UPDATE platform_connectors SET {', '.join(updates)} WHERE connector_type=?",
        tuple(vals))

    action = "enable_connector" if is_enabled else "disable_connector"
    await db_execute(
        "INSERT INTO audit_log(id,actor_id,target_id,action,created_at) VALUES(?,?,?,?,?)",
        (_new_id(), admin.get("id") or admin.get("sub",""), connector_type, action, _utcnow()))
    return {"ok": True, "connector_type": connector_type, "is_enabled": bool(is_enabled)}

@app.post("/admin/platform-connectors/bulk")
async def admin_bulk_toggle_connectors(body: dict, admin: Dict = Depends(_require_admin)):
    """Bulk enable/disable multiple connectors at once."""
    connector_types = body.get("connector_types", [])
    is_enabled = body.get("is_enabled", False)
    if not connector_types:
        raise HTTPException(400, "connector_types list required")
    now = _utcnow()
    updated = 0
    for ctype in connector_types:
        row = await db_fetchone(
            "SELECT env_keys FROM platform_connectors WHERE connector_type=?", (ctype,))
        if not row:
            continue
        if is_enabled:
            env_keys = json.loads(row.get("env_keys") or "[]")
            if _calc_setup_status(env_keys) == "not_configured" and env_keys:
                continue  # skip unconfigured ones silently on bulk
        await db_execute(
            "UPDATE platform_connectors SET is_enabled=?,updated_at=? WHERE connector_type=?",
            (int(bool(is_enabled)), now, ctype))
        updated += 1
    return {"ok": True, "updated": updated}

@app.get("/platform-connectors")
async def get_enabled_platform_connectors(user: Dict = Depends(_get_current_user)):
    """Returns only admin-enabled connector types — used by the client UI."""
    rows = await db_fetchall(
        "SELECT connector_type, display_name, category, icon, setup_status"
        " FROM platform_connectors WHERE is_enabled=1 ORDER BY category, display_name")
    return {"enabled": [r["connector_type"] for r in rows], "connectors": rows}

# ══════════════════════════════════════════════════════════════════════════════
# §40  ENTRY POINT
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    import uvicorn
    port = int(os.getenv("PORT", 8000))
    logger.info("🎵 Starting JAZZ AI v14.0 on port %d", port)
    uvicorn.run(app, host="0.0.0.0", port=port,
                reload=False, log_level="info", timeout_keep_alive=75)
